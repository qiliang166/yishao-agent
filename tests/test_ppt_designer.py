import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'backend'))
from pptx import Presentation
from pptx.util import Inches, Pt
from services.ppt_designer import (
    DesignSystem, extract_design, build_slide,
    build_cover, build_toc, build_technique, build_content, build_summary
)


def make_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


# ── Task 1: DesignSystem + extract_design ──

def test_extract_design_fallback():
    """Design system falls back to defaults when rules is empty."""
    ds = extract_design({}, None)
    assert ds.primary_color == (0xC0, 0x2E, 0x2E)
    assert ds.title_size_pt == 36
    assert ds.body_size_pt == 18
    assert ds.font_name == "Microsoft YaHei"


def test_extract_design_from_rules():
    """Design system reads colors and fonts from rules."""
    rules = {
        "design_rules": {
            "colors": {"primary": "#1a73e8", "accent": "#ff6d01", "background": "#ffffff", "text": "#333333"},
            "fonts": {"title_size": 40, "body_size": 16, "font_name": "SimHei"}
        }
    }
    ds = extract_design(rules, None)
    assert ds.primary_color == (0x1a, 0x73, 0xe8)
    assert ds.title_size_pt == 40
    assert ds.body_size_pt == 16
    assert ds.font_name == "SimHei"


def test_extract_design_from_typography():
    """Body/title sizes come from typography_profile when rules don't specify."""
    profile = {"body_font_size_pt": 14.0, "title_font_size_pt": 32.0, "line_height_ratio": 1.3}
    ds = extract_design({}, profile)
    assert ds.title_size_pt == 32
    assert ds.body_size_pt == 14


# ── Task 2: Builders ──

def test_build_cover_creates_slide():
    prs = make_prs()
    ds = DesignSystem()
    zones = {"title": "测试菜品", "subtitle": "制作工艺", "date": "2024-01-01"}
    build_cover(prs, zones, ds)
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    text_shapes = [sh for sh in slide.shapes if sh.has_text_frame]
    assert len(text_shapes) >= 2


def test_build_slide_dispatches():
    prs = make_prs()
    ds = DesignSystem()
    build_slide(prs, "cover", {"title": "T", "subtitle": "S"}, ds)
    assert len(prs.slides) == 1


def test_build_toc_with_items():
    prs = make_prs()
    ds = DesignSystem()
    zones = {"heading": "目录", "items": "1. 选材\n2. 备料\n3. 烹饪\n4. 装盘"}
    build_toc(prs, zones, ds)
    assert len(prs.slides) == 1


def test_build_technique_with_params():
    prs = make_prs()
    ds = DesignSystem()
    zones = {"heading": "焯水技法", "operation": "冷水下锅，大火烧开，撇去浮沫",
             "principle": "去腥除血水", "params": "水温: 100°C | 时间: 3分钟"}
    build_technique(prs, zones, ds)
    assert len(prs.slides) == 1


def test_build_content_long_text():
    prs = make_prs()
    ds = DesignSystem()
    zones = {"heading": "详细说明", "body": "这是一段很长的正文内容。" * 50}
    build_content(prs, zones, ds)
    assert len(prs.slides) == 1


def test_build_summary():
    prs = make_prs()
    ds = DesignSystem()
    zones = {"heading": "总结", "points": "要点一\n要点二\n要点三"}
    build_summary(prs, zones, ds)
    assert len(prs.slides) == 1
