import sys, os, json, tempfile
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'backend'))
from pptx import Presentation
from services.ppt_service import generate_ppt
from services.ppt_designer import extract_design, build_slide, DesignSystem

SAMPLE_RULES = {
    "layout_types": [
        {"id": "cover", "name": "Cover", "zones": ["title", "subtitle", "date"]},
        {"id": "toc", "name": "TOC", "zones": ["heading", "items"]},
        {"id": "technique", "name": "Technique", "zones": ["heading", "operation", "principle", "params"]},
        {"id": "content", "name": "Content", "zones": ["heading", "body"]},
        {"id": "summary", "name": "Summary", "zones": ["heading", "points"]},
    ],
    "design_rules": {
        "colors": {"primary": "#C02E2E", "accent": "#FF6D01", "background": "#FFFFFF", "text": "#333333"},
        "fonts": {"title_size": 36, "body_size": 18, "font_name": "Microsoft YaHei"}
    }
}

SAMPLE_PLAN = [
    {"type": "cover", "zones": {"title": "Red Braised Pork Belly", "subtitle": "Traditional Technique & Modern Standard", "date": "2024-01-01"}},
    {"type": "toc", "zones": {"heading": "Contents", "items": "01 Ingredient Selection\n02 Preparation\n03 Cooking Technique\n04 Plating"}},
    {"type": "technique", "zones": {"heading": "Blanching", "operation": "Cut pork belly into 3cm cubes. Place in cold water with ginger slices and cooking wine. Bring to boil, skim off foam, drain.", "principle": "High temperature coagulates protein. Blood and impurities rise as foam. Skimming removes odor and enhances freshness.", "params": "Water temp: 100C | Time: 3-5 min | Cube: 3cm"}},
    {"type": "technique", "zones": {"heading": "Caramel Coloring", "operation": "Put rock sugar in cold wok. Heat on low until amber color with bubbles. Quickly add meat and stir-fry evenly.", "principle": "Caramelization decomposes sugar into reddish-brown compounds that adhere to meat surface.", "params": "Rock sugar: 30g | Temp: 160-180C | Time: 2-3 min"}},
    {"type": "content", "zones": {"heading": "Sauce Reduction Control", "body": "Sauce reduction is the final step before serving. Turn heat to high and stir continuously.\nKey indicators:\n1. Sauce thickness - clings to meat without dripping\n2. Color - bright red, glossy, no darkening\n3. Salt level - sauce concentrates during reduction; season lightly earlier"}},
    {"type": "summary", "zones": {"heading": "Key Points Summary", "points": "Selection: 3-layer pork belly, 3:7 fat ratio\nBlanching: Cold water start, boil without rolling\nCaramel: Amber color with bubbles = optimal\nReduction: Sauce clings to spatula = done"}},
]

def test_design_engine_full_flow():
    """Build a complete PPT from sample plan using design engine."""
    prs = Presentation()
    ds = extract_design(SAMPLE_RULES, None)

    for sd in SAMPLE_PLAN:
        build_slide(prs, sd["type"], sd.get("zones", {}), ds)

    assert len(prs.slides) == 6, f"Expected 6 slides, got {len(prs.slides)}"

    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
        tmp_path = f.name
    prs.save(tmp_path)
    saved_size = os.path.getsize(tmp_path)
    assert saved_size > 10000, f"PPTX too small: {saved_size} bytes"
    print(f"Saved PPTX: {tmp_path} ({saved_size} bytes)")
    try:
        os.unlink(tmp_path)
    except Exception:
        pass

    # Verify slide content
    cover = prs.slides[0]
    texts = []
    for sh in cover.shapes:
        if sh.has_text_frame:
            texts.append(sh.text_frame.text)
    assert "Red Braised Pork Belly" in ''.join(texts), f"Cover text mismatch: {texts}"

    # Verify technique slide has operation text
    tech = prs.slides[2]
    tech_texts = []
    for sh in tech.shapes:
        if sh.has_text_frame:
            tech_texts.append(sh.text_frame.text)
    combined = ''.join(tech_texts)
    assert "Cut pork belly" in combined, f"Technique text mismatch: {tech_texts[:3]}"


def test_generate_ppt_with_plan():
    """generate_ppt accepts slide_plan and produces valid PPTX."""
    with tempfile.TemporaryDirectory() as tmpdir:
        path = generate_ppt(
            content="test", template_id=None, branding=None,
            output_dir=tmpdir, provider_id="", model="",
            slide_plan=SAMPLE_PLAN
        )
        assert os.path.exists(path), f"PPTX not found: {path}"
        assert os.path.getsize(path) > 10000, f"PPTX too small: {os.path.getsize(path)} bytes"
        prs = Presentation(path)
        assert len(prs.slides) == 6, f"Expected 6 slides, got {len(prs.slides)}"
        print(f"PPTX: {path} ({os.path.getsize(path)} bytes, {len(prs.slides)} slides)")


def test_api_plan_endpoint():
    """Verify /api/ppt/plan endpoint returns valid structure."""
    import urllib.request
    data = json.dumps({
        "content": "test",
        "template_id": "default-dao"
    }).encode('utf-8')
    req = urllib.request.Request('http://localhost:8765/api/ppt/plan',
        data=data, headers={'Content-Type': 'application/json'})
    resp = urllib.request.urlopen(req)
    result = json.loads(resp.read())
    assert "slide_plan" in result, f"Missing slide_plan in response: {list(result.keys())}"
    print(f"Plan endpoint: {len(result['slide_plan']) if result['slide_plan'] else 0} slides")


def test_api_generate_with_plan():
    """Verify /api/ppt/generate with slide_plan produces a valid PPTX."""
    import urllib.request
    data = json.dumps({
        "content": "test",
        "template_id": "default-dao",
        "slide_plan": [
            {"type": "cover", "zones": {"title": "API Test", "subtitle": "Verify", "date": "2024"}},
            {"type": "summary", "zones": {"heading": "Done", "points": "Point A\nPoint B"}},
        ]
    }).encode('utf-8')
    req = urllib.request.Request('http://localhost:8765/api/ppt/generate',
        data=data, headers={'Content-Type': 'application/json'})
    resp = urllib.request.urlopen(req)
    result = json.loads(resp.read())
    assert "filename" in result, f"Missing filename: {list(result.keys())}"
    assert result.get("slide_plan") is not None, "slide_plan not returned"
    filepath = os.path.join(os.path.dirname(__file__), '..', 'backend', 'data', 'exports', result["filename"])
    assert os.path.exists(filepath), f"PPTX not found at {filepath}"
    prs = Presentation(filepath)
    assert len(prs.slides) == 2, f"Expected 2 slides, got {len(prs.slides)}"
    print(f"Generate API: {result['filename']} ({len(prs.slides)} slides)")
