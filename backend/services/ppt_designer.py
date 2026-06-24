"""Design system engine for PPT generation.

Creates slides from scratch using python-pptx, driven by a DesignSystem
extracted from template rules and typography profiles."""
from dataclasses import dataclass, field
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


@dataclass
class DesignSystem:
    """Colors, fonts, and spacing for slide creation."""
    # Colors as RGB tuples
    primary_color: tuple = (0xC0, 0x2E, 0x2E)       # red
    accent_color: tuple = (0xFF, 0x6D, 0x01)         # orange
    background_color: tuple = (0xFF, 0xFF, 0xFF)      # white
    text_color: tuple = (0x33, 0x33, 0x33)            # dark gray
    light_text_color: tuple = (0xFF, 0xFF, 0xFF)       # white
    # Fonts
    font_name: str = "Microsoft YaHei"
    title_size_pt: int = 36
    subtitle_size_pt: int = 20
    body_size_pt: int = 18
    small_size_pt: int = 12
    # Spacing in EMU
    slide_width: int = int(13.333 * 914400)
    slide_height: int = int(7.5 * 914400)
    margin_left: int = int(1.0 * 914400)
    margin_right: int = int(1.0 * 914400)
    margin_top: int = int(0.8 * 914400)
    margin_bottom: int = int(0.6 * 914400)
    line_spacing_ratio: float = 1.3


def _hex_to_rgb(hex_str: str) -> tuple:
    """Convert '#1a73e8' to (0x1a, 0x73, 0xe8)."""
    h = hex_str.lstrip('#')
    if len(h) == 3:
        h = ''.join(c * 2 for c in h)
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def extract_design(rules: dict, typography_profile: Optional[dict] = None) -> DesignSystem:
    """Build a DesignSystem from template rules and typography profile.

    Priority: rules.design_rules > typography_profile > hardcoded defaults.
    """
    ds = DesignSystem()

    design_rules = rules.get("design_rules", {}) if rules else {}

    # Colors
    colors = design_rules.get("colors", {})
    if "primary" in colors:
        ds.primary_color = _hex_to_rgb(colors["primary"])
    if "accent" in colors:
        ds.accent_color = _hex_to_rgb(colors["accent"])
    if "background" in colors:
        ds.background_color = _hex_to_rgb(colors["background"])
    if "text" in colors:
        ds.text_color = _hex_to_rgb(colors["text"])
    if "light_text" in colors:
        ds.light_text_color = _hex_to_rgb(colors["light_text"])

    # Fonts
    fonts = design_rules.get("fonts", {})
    if "font_name" in fonts:
        ds.font_name = fonts["font_name"]
    elif "title" in fonts:
        ds.font_name = fonts["title"]
    # Theme font reference (e.g. +mj-ea) — fall back
    if ds.font_name.startswith('+'):
        ds.font_name = "Microsoft YaHei"
    if "title_size" in fonts:
        val = fonts["title_size"]
        if isinstance(val, str):
            val = val.replace("pt", "").strip()
        ds.title_size_pt = int(val)
    if "body_size" in fonts:
        val = fonts["body_size"]
        if isinstance(val, str):
            val = val.replace("pt", "").strip()
        ds.body_size_pt = int(val)

    # Typography profile overrides
    if typography_profile:
        if "title_font_size_pt" in typography_profile:
            ds.title_size_pt = int(typography_profile["title_font_size_pt"])
        if "body_font_size_pt" in typography_profile:
            ds.body_size_pt = int(typography_profile["body_font_size_pt"])
        if "line_height_ratio" in typography_profile:
            ds.line_spacing_ratio = typography_profile["line_height_ratio"]

    return ds


# ── Internal helpers ──

def _rgb(clr: tuple) -> RGBColor:
    return RGBColor(*clr)


def _add_textbox(slide, left, top, width, height, text: str = "",
                 font_size_pt: int = 18, color=None, bold: bool = False,
                 alignment=PP_ALIGN.LEFT, font_name: str = "Microsoft YaHei",
                 word_wrap: bool = True):
    """Add a textbox with one paragraph, returning (shape, text_frame)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size_pt)
    p.font.bold = bold
    if color:
        p.font.color.rgb = _rgb(color)
    p.font.name = font_name
    p.alignment = alignment
    return txBox, tf


def _add_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()  # no border
    return shape


# ── Builder registry ──

BUILDERS = {}


def _register(*types: str):
    def deco(fn):
        for t in types:
            BUILDERS[t] = fn
        return fn
    return deco


def build_slide(prs: Presentation, slide_type: str, zones: dict,
                design: DesignSystem):
    """Create a slide using the appropriate builder for slide_type.

    Falls back to build_content for unknown types.
    """
    builder = BUILDERS.get(slide_type, build_content)
    builder(prs, zones, design)


@_register("cover")
def build_cover(prs: Presentation, zones: dict, design: DesignSystem):
    """Full-bleed primary background, centered title + subtitle + date."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(design.primary_color)

    sw = design.slide_width
    sh = design.slide_height

    title = zones.get("title", "").strip()
    subtitle = zones.get("subtitle", "").strip()
    date = zones.get("date", "").strip()

    # Title — centered
    if title:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(2.0 * 914400)),
                     Emu(design.slide_width - design.margin_left * 2), Emu(int(1.6 * 914400)),
                     title, font_size_pt=design.title_size_pt + 8,
                     color=design.light_text_color, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    # Subtitle
    if subtitle:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(3.8 * 914400)),
                     Emu(design.slide_width - design.margin_left * 2), Emu(int(0.8 * 914400)),
                     subtitle, font_size_pt=design.subtitle_size_pt,
                     color=design.light_text_color, bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    # Date — bottom right
    if date:
        _add_textbox(slide, Emu(int(8.0 * 914400)), Emu(int(6.5 * 914400)),
                     Emu(int(4.3 * 914400)), Emu(int(0.5 * 914400)),
                     date, font_size_pt=design.small_size_pt,
                     color=design.light_text_color, bold=False,
                     alignment=PP_ALIGN.RIGHT, font_name=design.font_name)


@_register("toc")
def build_toc(prs: Presentation, zones: dict, design: DesignSystem):
    """Left heading + numbered items with accent bars."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "目录").strip()
    items_text = zones.get("items", "").strip()

    # Heading
    _add_textbox(slide, Emu(design.margin_left), Emu(design.margin_top),
                 Emu(int(3.0 * 914400)), Emu(int(0.8 * 914400)), heading,
                 font_size_pt=design.title_size_pt,
                 color=design.primary_color, bold=True,
                 font_name=design.font_name)

    # Accent line under heading
    _add_rect(slide, Emu(design.margin_left), Emu(int(1.8 * 914400)),
              Emu(int(1.5 * 914400)), Emu(int(0.05 * 914400)),
              design.accent_color)

    # Items
    if items_text:
        items = [it.strip() for it in items_text.split('\n') if it.strip()]
        y = int(2.2 * 914400)
        for idx, item in enumerate(items):
            # Number badge
            num_text = f"{idx + 1:02d}"
            _add_textbox(slide, Emu(design.margin_left), Emu(y),
                         Emu(int(0.6 * 914400)), Emu(int(0.5 * 914400)),
                         num_text, font_size_pt=design.body_size_pt,
                         color=design.accent_color, bold=True,
                         font_name=design.font_name)
            # Item text
            _add_textbox(slide, Emu(int(2.0 * 914400)), Emu(y),
                         Emu(int(9.5 * 914400)), Emu(int(0.5 * 914400)),
                         item, font_size_pt=design.body_size_pt,
                         color=design.text_color, bold=False,
                         font_name=design.font_name)
            y += int(0.7 * 914400)


@_register("technique")
def build_technique(prs: Presentation, zones: dict, design: DesignSystem):
    """Title bar + operation steps + principle box + params row."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "").strip()
    operation = zones.get("operation", "").strip()
    principle = zones.get("principle", "").strip()
    params = zones.get("params", "").strip()

    # Title bar — primary color background
    _add_rect(slide, 0, 0, design.slide_width, int(1.2 * 914400), design.primary_color)
    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(0.15 * 914400)),
                     Emu(int(11.3 * 914400)), Emu(int(0.9 * 914400)),
                     heading, font_size_pt=design.title_size_pt,
                     color=design.light_text_color, bold=True,
                     font_name=design.font_name)

    y = int(1.5 * 914400)

    # Operation steps
    if operation:
        _add_textbox(slide, Emu(design.margin_left), Emu(y),
                     Emu(int(2.0 * 914400)), Emu(int(0.4 * 914400)),
                     "▎操作步骤", font_size_pt=design.subtitle_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)
        y += int(0.5 * 914400)
        line_count = max(1, len(operation) // 40 + operation.count('\n') + 1)
        _, tf = _add_textbox(slide, Emu(design.margin_left), Emu(y),
                             Emu(int(11.3 * 914400)), Emu(int(2.0 * 914400)),
                             operation, font_size_pt=design.body_size_pt,
                             color=design.text_color, bold=False,
                             font_name=design.font_name)
        y += int(line_count * 0.45 * 914400)

    # Principle box
    if principle:
        y += int(0.3 * 914400)
        _add_rect(slide, Emu(design.margin_left), Emu(y),
                  Emu(int(11.3 * 914400)), Emu(int(0.05 * 914400)),
                  design.accent_color)
        y += int(0.2 * 914400)
        _add_textbox(slide, Emu(design.margin_left), Emu(y),
                     Emu(int(11.3 * 914400)), Emu(int(0.8 * 914400)),
                     f"原理: {principle}", font_size_pt=design.body_size_pt,
                     color=design.accent_color, bold=False,
                     font_name=design.font_name)
        y += int(0.5 * 914400)

    # Params
    if params:
        y += int(0.2 * 914400)
        _add_textbox(slide, Emu(design.margin_left), Emu(y),
                     Emu(int(11.3 * 914400)), Emu(int(0.5 * 914400)),
                     f"参数: {params}", font_size_pt=design.small_size_pt,
                     color=design.text_color, bold=False,
                     font_name=design.font_name)


@_register("content")
def build_content(prs: Presentation, zones: dict, design: DesignSystem):
    """Generic content slide: heading + body text."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "").strip()
    body = zones.get("body", "").strip()

    # Heading
    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(design.margin_top),
                     Emu(int(11.3 * 914400)), Emu(int(0.8 * 914400)), heading,
                     font_size_pt=design.title_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)
        # Underline
        _add_rect(slide, Emu(design.margin_left), Emu(int(1.2 * 914400)),
                  Emu(int(11.3 * 914400)), Emu(int(0.03 * 914400)),
                  design.primary_color)

    # Body
    if body:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(1.5 * 914400)),
                     Emu(int(11.3 * 914400)), Emu(int(5.0 * 914400)),
                     body, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     font_name=design.font_name)


@_register("summary")
def build_summary(prs: Presentation, zones: dict, design: DesignSystem):
    """Summary slide with accent left bar and key points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "总结").strip()
    points = zones.get("points", "").strip()

    # Left accent bar
    _add_rect(slide, Emu(int(0.5 * 914400)), Emu(int(1.5 * 914400)),
              Emu(int(0.08 * 914400)), Emu(int(4.5 * 914400)),
              design.accent_color)

    # Heading
    _add_textbox(slide, Emu(int(1.0 * 914400)), Emu(int(1.5 * 914400)),
                 Emu(int(11.3 * 914400)), Emu(int(0.8 * 914400)),
                 heading, font_size_pt=design.title_size_pt,
                 color=design.primary_color, bold=True,
                 font_name=design.font_name)

    # Points
    if points:
        point_list = [p.strip() for p in points.split('\n') if p.strip()]
        y = int(2.5 * 914400)
        for pt_text in point_list:
            _add_textbox(slide, Emu(int(1.0 * 914400)), Emu(y),
                         Emu(int(11.3 * 914400)), Emu(int(0.5 * 914400)),
                         f"✦ {pt_text}", font_size_pt=design.body_size_pt,
                         color=design.text_color, bold=False,
                         font_name=design.font_name)
            y += int(0.55 * 914400)


@_register("section")
def build_section(prs: Presentation, zones: dict, design: DesignSystem):
    """Section divider: primary background + heading + subtitle."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(design.primary_color)

    heading = zones.get("heading", "").strip()
    subtitle = zones.get("subtitle", "").strip()

    if heading:
        _add_textbox(slide, Emu(int(1.5 * 914400)), Emu(int(2.5 * 914400)),
                     Emu(int(10.3 * 914400)), Emu(int(1.2 * 914400)),
                     heading, font_size_pt=design.title_size_pt + 4,
                     color=design.light_text_color, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    if subtitle:
        _add_textbox(slide, Emu(int(1.5 * 914400)), Emu(int(3.8 * 914400)),
                     Emu(int(10.3 * 914400)), Emu(int(0.8 * 914400)),
                     subtitle, font_size_pt=design.subtitle_size_pt,
                     color=design.light_text_color, bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)


# ── Swiss / Extended builders ──

@_register("chapter")
def build_chapter(prs: Presentation, zones: dict, design: DesignSystem):
    """Chapter divider: large number + heading on accent background."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(design.accent_color)

    number = zones.get("number", "").strip()
    heading = zones.get("heading", "").strip()

    if number:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(1.5 * 914400)),
                     Emu(int(11.3 * 914400)), Emu(int(2.0 * 914400)),
                     number, font_size_pt=design.title_size_pt + 24,
                     color=design.light_text_color, bold=True,
                     alignment=PP_ALIGN.LEFT, font_name=design.font_name)
    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(3.8 * 914400)),
                     Emu(int(11.3 * 914400)), Emu(int(1.2 * 914400)),
                     heading, font_size_pt=design.title_size_pt,
                     color=design.light_text_color, bold=True,
                     font_name=design.font_name)
    # Accent line
    _add_rect(slide, Emu(design.margin_left), Emu(int(5.1 * 914400)),
              Emu(int(2.0 * 914400)), Emu(int(0.05 * 914400)),
              design.light_text_color)


@_register("closing")
def build_closing(prs: Presentation, zones: dict, design: DesignSystem):
    """Closing slide: quote + signature, centered."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(design.primary_color)

    quote = zones.get("quote", "").strip()
    signature = zones.get("signature", "").strip()
    summary = zones.get("summary", "").strip()

    y = int(2.0 * 914400)
    display_text = quote or summary
    if display_text:
        _, tf = _add_textbox(slide, Emu(int(1.5 * 914400)), Emu(y),
                             Emu(int(10.3 * 914400)), Emu(int(2.5 * 914400)),
                             display_text, font_size_pt=design.title_size_pt - 4,
                             color=design.light_text_color, bold=False,
                             alignment=PP_ALIGN.CENTER, font_name=design.font_name)
        y += int(3.0 * 914400)
    if signature:
        _add_textbox(slide, Emu(int(1.5 * 914400)), Emu(y),
                     Emu(int(10.3 * 914400)), Emu(int(0.6 * 914400)),
                     f"— {signature}", font_size_pt=design.body_size_pt,
                     color=design.light_text_color, bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)


@_register("data_hero", "key_params")
def build_data_hero(prs: Presentation, zones: dict, design: DesignSystem):
    """Big number KPI with label and context."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "").strip()
    big_number = zones.get("big_number", "").strip()
    label = zones.get("label", "").strip()
    context = zones.get("context", zones.get("notes", "")).strip()

    # Heading
    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(design.margin_top),
                     Emu(int(11.3 * 914400)), Emu(int(0.8 * 914400)), heading,
                     font_size_pt=design.title_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)
        _add_rect(slide, Emu(design.margin_left), Emu(int(1.2 * 914400)),
                  Emu(int(11.3 * 914400)), Emu(int(0.03 * 914400)),
                  design.primary_color)

    # Big number — centered, huge
    if big_number:
        _add_textbox(slide, Emu(int(2.0 * 914400)), Emu(int(1.8 * 914400)),
                     Emu(int(9.3 * 914400)), Emu(int(2.5 * 914400)),
                     big_number, font_size_pt=design.title_size_pt + 24,
                     color=design.accent_color, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    # Label under the number
    if label:
        _add_textbox(slide, Emu(int(2.0 * 914400)), Emu(int(4.3 * 914400)),
                     Emu(int(9.3 * 914400)), Emu(int(0.6 * 914400)),
                     label, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    # Context / notes
    if context:
        _add_textbox(slide, Emu(int(2.0 * 914400)), Emu(int(5.2 * 914400)),
                     Emu(int(9.3 * 914400)), Emu(int(1.5 * 914400)),
                     context, font_size_pt=design.small_size_pt,
                     color=design.text_color, bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)


@_register("duo_compare", "comparison")
def build_duo_compare(prs: Presentation, zones: dict, design: DesignSystem):
    """Two-column comparison: left vs right with divider."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "").strip()
    left = zones.get("left", "").strip()
    right = zones.get("right", "").strip()
    left_label = zones.get("left_label", "传统").strip()
    right_label = zones.get("right_label", "创新").strip()

    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(design.margin_top),
                     Emu(int(11.3 * 914400)), Emu(int(0.8 * 914400)), heading,
                     font_size_pt=design.title_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)

    # Vertical divider
    mid_x = int(6.5 * 914400)
    _add_rect(slide, Emu(mid_x), Emu(int(1.2 * 914400)),
              Emu(int(0.03 * 914400)), Emu(int(5.8 * 914400)),
              design.accent_color)

    col_w = int(4.8 * 914400)
    y = int(1.5 * 914400)

    # Left
    _add_textbox(slide, Emu(design.margin_left), Emu(y),
                 Emu(col_w), Emu(int(0.5 * 914400)),
                 left_label, font_size_pt=design.subtitle_size_pt,
                 color=design.primary_color, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=design.font_name)
    if left:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(2.2 * 914400)),
                     Emu(col_w), Emu(int(4.5 * 914400)),
                     left, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    # Right
    right_x = int(7.2 * 914400)
    _add_textbox(slide, Emu(right_x), Emu(y),
                 Emu(col_w), Emu(int(0.5 * 914400)),
                 right_label, font_size_pt=design.subtitle_size_pt,
                 color=design.accent_color, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=design.font_name)
    if right:
        _add_textbox(slide, Emu(right_x), Emu(int(2.2 * 914400)),
                     Emu(col_w), Emu(int(4.5 * 914400)),
                     right, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)


@_register("process_flow")
def build_process_flow(prs: Presentation, zones: dict, design: DesignSystem):
    """Horizontal process flow: step circles connected by arrows."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "").strip()
    steps_text = zones.get("steps", "").strip()

    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(design.margin_top),
                     Emu(int(11.3 * 914400)), Emu(int(0.7 * 914400)), heading,
                     font_size_pt=design.title_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)

    if steps_text:
        steps = [s.strip() for s in steps_text.replace('\n', '|').split('|') if s.strip()]
        if not steps:
            steps = [steps_text]
        n = len(steps)
        total_w = int(10.0 * 914400)
        start_x = int(1.6 * 914400)
        gap = total_w // max(n, 1) if n > 0 else total_w
        y = int(3.0 * 914400)
        circle_r = int(0.35 * 914400)

        for i, step in enumerate(steps):
            cx = start_x + gap * i + gap // 2
            # Circle
            shape = slide.shapes.add_shape(
                9,  # MSO_SHAPE.OVAL
                Emu(cx - circle_r), Emu(y - circle_r),
                Emu(circle_r * 2), Emu(circle_r * 2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(design.accent_color)
            shape.line.fill.background()
            # Number in circle
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(design.body_size_pt)
            p.font.bold = True
            p.font.color.rgb = _rgb(design.light_text_color)
            p.font.name = design.font_name
            p.alignment = PP_ALIGN.CENTER
            # Step label
            _add_textbox(slide, Emu(cx - int(1.2 * 914400)), Emu(y + circle_r + int(0.15 * 914400)),
                         Emu(int(2.4 * 914400)), Emu(int(0.6 * 914400)),
                         step, font_size_pt=design.small_size_pt,
                         color=design.text_color, bold=False,
                         alignment=PP_ALIGN.CENTER, font_name=design.font_name)

        # Arrow connector line
        if n > 1:
            line_y = y
            line_start = start_x + gap // 2 + circle_r
            line_end = start_x + gap * (n - 1) + gap // 2 - circle_r
            _add_rect(slide, Emu(line_start), Emu(line_y),
                      Emu(line_end - line_start), Emu(int(0.04 * 914400)),
                      design.accent_color)


@_register("grid_cards")
def build_grid_cards(prs: Presentation, zones: dict, design: DesignSystem):
    """2x2 or 3x2 card grid with accent top border."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "").strip()
    cards_text = zones.get("cards", zones.get("items", "")).strip()

    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(design.margin_top),
                     Emu(int(11.3 * 914400)), Emu(int(0.7 * 914400)), heading,
                     font_size_pt=design.title_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)
        _add_rect(slide, Emu(design.margin_left), Emu(int(1.1 * 914400)),
                  Emu(int(11.3 * 914400)), Emu(int(0.03 * 914400)),
                  design.primary_color)

    if cards_text:
        cards = [c.strip() for c in cards_text.replace('\n', '|').split('|') if c.strip()]
        if not cards:
            cards = [cards_text]
        n = len(cards)
        cols = 3 if n >= 6 else (2 if n >= 4 else n)
        rows = (n + cols - 1) // cols
        card_w = int(10.0 / cols * 914400)
        card_h = int(4.5 / max(rows, 1) * 914400)
        x0 = int(1.5 * 914400)
        y0 = int(1.4 * 914400)

        for i, card in enumerate(cards):
            col = i % cols
            row = i // cols
            cx = x0 + col * (card_w + int(0.3 * 914400))
            cy = y0 + row * (card_h + int(0.2 * 914400))

            # Card background
            _add_rect(slide, Emu(cx), Emu(cy), Emu(card_w), Emu(card_h),
                      design.background_color if design.background_color != (0xFF, 0xFF, 0xFF)
                      else (0xF5, 0xF5, 0xF5))
            # Accent top bar
            _add_rect(slide, Emu(cx), Emu(cy), Emu(card_w), Emu(int(0.06 * 914400)),
                      design.accent_color)
            # Card text
            _add_textbox(slide, Emu(cx + int(0.3 * 914400)),
                         Emu(cy + int(0.2 * 914400)),
                         Emu(card_w - int(0.6 * 914400)),
                         Emu(card_h - int(0.4 * 914400)),
                         card, font_size_pt=design.body_size_pt - 2,
                         color=design.text_color, bold=False,
                         font_name=design.font_name)


@_register("food_archive")
def build_food_archive(prs: Presentation, zones: dict, design: DesignSystem):
    """Ingredient archive card: food name + params + mechanism + substitutes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    food_name = zones.get("food_name", "").strip()
    params = zones.get("params", "").strip()
    mechanism = zones.get("mechanism", "").strip()
    substitutes = zones.get("substitutes", "").strip()

    # Title bar
    _add_rect(slide, 0, 0, design.slide_width, int(1.0 * 914400), design.primary_color)
    if food_name:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(0.15 * 914400)),
                     Emu(int(11.3 * 914400)), Emu(int(0.7 * 914400)),
                     food_name, font_size_pt=design.title_size_pt,
                     color=design.light_text_color, bold=True,
                     font_name=design.font_name)

    y = int(1.3 * 914400)
    sections = [
        ("黄金参数", params, design.accent_color),
        ("作用机理", mechanism, design.primary_color),
        ("替代与风险", substitutes, design.text_color),
    ]

    for label, content, label_color in sections:
        if not content:
            continue
        _add_textbox(slide, Emu(design.margin_left), Emu(y),
                     Emu(int(2.5 * 914400)), Emu(int(0.4 * 914400)),
                     f"▎{label}", font_size_pt=design.subtitle_size_pt,
                     color=label_color, bold=True,
                     font_name=design.font_name)
        y += int(0.5 * 914400)
        line_count = max(1, len(content) // 50 + content.count('\n') + 1)
        _add_textbox(slide, Emu(int(1.0 * 914400)), Emu(y),
                     Emu(int(11.3 * 914400)), Emu(int(line_count * 0.45 * 914400)),
                     content, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     font_name=design.font_name)
        y += int(line_count * 0.5 * 914400)


@_register("skill_card")
def build_skill_card(prs: Presentation, zones: dict, design: DesignSystem):
    """Technique skill card: name + description + flowchart ref + migration."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    skill_name = zones.get("skill_name", "").strip()
    description = zones.get("description", "").strip()
    flowchart = zones.get("flowchart", "").strip()
    migration = zones.get("migration", "").strip()

    # Header — accent background
    _add_rect(slide, 0, 0, design.slide_width, int(1.0 * 914400), design.accent_color)
    if skill_name:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(0.15 * 914400)),
                     Emu(int(11.3 * 914400)), Emu(int(0.7 * 914400)),
                     skill_name, font_size_pt=design.title_size_pt,
                     color=design.light_text_color, bold=True,
                     font_name=design.font_name)

    y = int(1.3 * 914400)
    sections = [
        ("技法描述", description),
        ("可迁移至", migration),
    ]

    for label, content in sections:
        if not content:
            continue
        _add_textbox(slide, Emu(design.margin_left), Emu(y),
                     Emu(int(3.0 * 914400)), Emu(int(0.4 * 914400)),
                     f"▎{label}", font_size_pt=design.subtitle_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)
        y += int(0.5 * 914400)
        line_count = max(1, len(content) // 50 + content.count('\n') + 1)
        _add_textbox(slide, Emu(int(1.0 * 914400)), Emu(y),
                     Emu(int(11.3 * 914400)), Emu(int(line_count * 0.45 * 914400)),
                     content, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     font_name=design.font_name)
        y += int(line_count * 0.5 * 914400)

    if flowchart:
        y += int(0.2 * 914400)
        _add_textbox(slide, Emu(design.margin_left), Emu(y),
                     Emu(int(11.3 * 914400)), Emu(int(0.5 * 914400)),
                     f"→ 流程: {flowchart}", font_size_pt=design.small_size_pt,
                     color=design.accent_color, bold=False,
                     font_name=design.font_name)


@_register("troubleshoot")
def build_troubleshoot(prs: Presentation, zones: dict, design: DesignSystem):
    """Troubleshooting diagnosis page: problem → cause → solution → prevention."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "故障诊断").strip()
    problem = zones.get("problem", "").strip()
    cause = zones.get("cause", "").strip()
    solution = zones.get("solution", "").strip()
    prevention = zones.get("prevention", "").strip()

    # Header
    _add_rect(slide, 0, 0, design.slide_width, int(0.9 * 914400), design.primary_color)
    _add_textbox(slide, Emu(design.margin_left), Emu(int(0.1 * 914400)),
                 Emu(int(11.3 * 914400)), Emu(int(0.7 * 914400)),
                 f"⚠️ {heading}", font_size_pt=design.title_size_pt,
                 color=design.light_text_color, bold=True,
                 font_name=design.font_name)

    y = int(1.2 * 914400)
    items = [
        ("问题现象", problem, (0xD4, 0x3E, 0x2E)),
        ("原因分析", cause, design.accent_color),
        ("解决方案", solution, (0x1A, 0x8D, 0x3F)),
        ("预防措施", prevention, design.primary_color),
    ]

    for label, content, item_color in items:
        if not content:
            continue
        _add_textbox(slide, Emu(design.margin_left), Emu(y),
                     Emu(int(3.0 * 914400)), Emu(int(0.4 * 914400)),
                     f"▎{label}", font_size_pt=design.subtitle_size_pt,
                     color=item_color, bold=True,
                     font_name=design.font_name)
        y += int(0.45 * 914400)
        line_count = max(1, len(content) // 55 + content.count('\n') + 1)
        _add_textbox(slide, Emu(int(1.0 * 914400)), Emu(y),
                     Emu(int(11.3 * 914400)), Emu(int(line_count * 0.4 * 914400)),
                     content, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     font_name=design.font_name)
        y += int(line_count * 0.45 * 914400)


@_register("quote")
def build_quote(prs: Presentation, zones: dict, design: DesignSystem):
    """Large quote slide with accent left bar and attribution."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    quote_text = zones.get("quote", zones.get("body", "")).strip()
    author = zones.get("author", zones.get("signature", "")).strip()
    context = zones.get("context", "").strip()

    # Accent left bar
    _add_rect(slide, Emu(int(1.0 * 914400)), Emu(int(1.5 * 914400)),
              Emu(int(0.1 * 914400)), Emu(int(4.5 * 914400)),
              design.accent_color)

    # Quote text
    if quote_text:
        _add_textbox(slide, Emu(int(1.6 * 914400)), Emu(int(2.0 * 914400)),
                     Emu(int(9.7 * 914400)), Emu(int(3.0 * 914400)),
                     quote_text, font_size_pt=design.title_size_pt - 2,
                     color=design.text_color, bold=False,
                     font_name=design.font_name)

    # Author
    if author:
        _add_textbox(slide, Emu(int(1.6 * 914400)), Emu(int(5.5 * 914400)),
                     Emu(int(9.7 * 914400)), Emu(int(0.5 * 914400)),
                     f"— {author}", font_size_pt=design.body_size_pt,
                     color=design.accent_color, bold=False,
                     alignment=PP_ALIGN.RIGHT, font_name=design.font_name)

    # Context
    if context:
        _add_textbox(slide, Emu(int(1.6 * 914400)), Emu(int(6.0 * 914400)),
                     Emu(int(9.7 * 914400)), Emu(int(0.4 * 914400)),
                     context, font_size_pt=design.small_size_pt,
                     color=design.text_color, bold=False,
                     alignment=PP_ALIGN.RIGHT, font_name=design.font_name)


@_register("timeline")
def build_timeline(prs: Presentation, zones: dict, design: DesignSystem):
    """Vertical timeline: events along an accent line."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    heading = zones.get("heading", "").strip()
    events_texts = zones.get("events", zones.get("items", "")).strip()

    if heading:
        _add_textbox(slide, Emu(design.margin_left), Emu(design.margin_top),
                     Emu(int(11.3 * 914400)), Emu(int(0.7 * 914400)), heading,
                     font_size_pt=design.title_size_pt,
                     color=design.primary_color, bold=True,
                     font_name=design.font_name)

    # Vertical line
    _add_rect(slide, Emu(int(2.5 * 914400)), Emu(int(1.5 * 914400)),
              Emu(int(0.04 * 914400)), Emu(int(5.2 * 914400)),
              design.accent_color)

    if events_texts:
        events = [e.strip() for e in events_texts.replace('\n', '|').split('|') if e.strip()]
        if not events:
            events = [events_texts]
        y = int(1.8 * 914400)
        for i, event in enumerate(events):
            # Dot on timeline
            dot = slide.shapes.add_shape(
                9, Emu(int(2.32 * 914400)), Emu(y + int(0.1 * 914400)),
                Emu(int(0.2 * 914400)), Emu(int(0.2 * 914400)))
            dot.fill.solid()
            dot.fill.fore_color.rgb = _rgb(design.accent_color)
            dot.line.fill.background()
            # Event text
            _add_textbox(slide, Emu(int(3.0 * 914400)), Emu(y),
                         Emu(int(9.0 * 914400)), Emu(int(0.6 * 914400)),
                         event, font_size_pt=design.body_size_pt,
                         color=design.text_color, bold=False,
                         font_name=design.font_name)
            y += int(0.8 * 914400)


@_register("image_hero")
def build_image_hero(prs: Presentation, zones: dict, design: DesignSystem):
    """Full-width image placeholder with caption."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    caption = zones.get("caption", zones.get("heading", "")).strip()
    image_zone = zones.get("image", "").strip()

    # Image placeholder (gray block)
    _add_rect(slide, Emu(int(0.5 * 914400)), Emu(int(0.5 * 914400)),
              Emu(int(12.3 * 914400)), Emu(int(5.5 * 914400)),
              (0xE0, 0xE0, 0xE0))
    # Placeholder icon text
    _add_textbox(slide, Emu(int(3.0 * 914400)), Emu(int(2.5 * 914400)),
                 Emu(int(7.3 * 914400)), Emu(int(1.0 * 914400)),
                 "📷 图片区域", font_size_pt=design.title_size_pt,
                 color=(0x99, 0x99, 0x99), bold=False,
                 alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    if image_zone:
        _add_textbox(slide, Emu(int(3.0 * 914400)), Emu(int(3.5 * 914400)),
                     Emu(int(7.3 * 914400)), Emu(int(0.6 * 914400)),
                     image_zone, font_size_pt=design.small_size_pt,
                     color=(0x99, 0x99, 0x99), bold=False,
                     alignment=PP_ALIGN.CENTER, font_name=design.font_name)

    # Caption at bottom
    if caption:
        _add_textbox(slide, Emu(design.margin_left), Emu(int(6.3 * 914400)),
                     Emu(int(11.3 * 914400)), Emu(int(0.6 * 914400)),
                     caption, font_size_pt=design.body_size_pt,
                     color=design.text_color, bold=False,
                     font_name=design.font_name)


def _extract_dominant_colors(pptx_path: str) -> dict:
    """Mechanically extract dominant colors from a PPTX template.

    Returns {primary, accent, background, text, light_text} as hex strings.
    """
    from pptx import Presentation
    from collections import Counter

    prs = Presentation(pptx_path)

    bg_colors = []
    fill_colors = []
    text_colors = []

    for slide in prs.slides:
        try:
            bg = slide.background.fill
            if bg.type is not None:
                try:
                    rgb = bg.fore_color.rgb
                    bg_colors.append(str(rgb))
                except Exception:
                    pass
        except Exception:
            pass

        for shape in slide.shapes:
            try:
                if hasattr(shape, 'fill') and shape.fill is not None:
                    try:
                        rgb = shape.fill.fore_color.rgb
                        fill_colors.append(str(rgb))
                    except Exception:
                        pass
            except Exception:
                pass

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                text_colors.append(str(run.font.color.rgb))
                        except Exception:
                            pass

    def most_common(colors, fallback):
        if not colors:
            return fallback
        c = Counter(colors)
        for clr, _ in c.most_common():
            h = clr.upper()
            if h not in ('FFFFFF', '000000'):
                return f"#{h}"
        return f"#{c.most_common(1)[0][0].upper()}"

    non_white = [c for c in fill_colors if c.upper() not in ('FFFFFF', '000000')]
    primary = most_common(non_white, "C02E2E")

    white_bg = [c for c in bg_colors if c.upper() == 'FFFFFF']
    background = "FFFFFF" if white_bg else most_common(bg_colors, "FFFFFF")

    dark_text = [c for c in text_colors if _is_dark(c)]
    text = most_common(dark_text, "333333") if dark_text else "333333"

    non_white_fills = [c for c in non_white if c.upper() != primary.lstrip('#').upper()]
    accent = most_common(non_white_fills, "FF6D01")

    light_text_colors = [c for c in text_colors if c.upper() == 'FFFFFF']
    light_text = "FFFFFF" if light_text_colors else "FFFFFF"

    return {
        "primary": primary if primary.startswith('#') else f"#{primary}",
        "accent": accent if accent.startswith('#') else f"#{accent}",
        "background": background if background.startswith('#') else f"#{background}",
        "text": text if text.startswith('#') else f"#{text}",
        "light_text": light_text if light_text.startswith('#') else f"#{light_text}",
    }


def _is_dark(hex_str: str) -> bool:
    """True if hex color is perceptually dark."""
    h = hex_str.lstrip('#').upper()
    if len(h) != 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < 128
