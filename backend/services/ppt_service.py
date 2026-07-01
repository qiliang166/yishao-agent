"""PPT generation service — SVG output with PPTX fallback.

Uses PPT-Agent Bento Grid methodology: style YAML → AI outline → Bento Grid layout → SVG output.
"""
import os
import re
import json
import asyncio
import logging
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from database import get_db

_logger = logging.getLogger("uvicorn")


def _extract_typography(prs) -> dict:
    """Extract typography info from a PPTX Presentation object."""
    fonts = set()
    sizes = set()
    bold_count = 0
    total_runs = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        total_runs += 1
                        f = run.font
                        if f.name:
                            fonts.add(f.name)
                        if f.size:
                            sizes.add(round(f.size / 12700, 1))
                        if f.bold:
                            bold_count += 1
    return {
        "fonts": sorted(list(fonts)),
        "sizes_pt": sorted(list(sizes)),
        "bold_ratio": round(bold_count / total_runs, 2) if total_runs > 0 else 0,
        "total_text_runs": total_runs
    }


def _clean_json_response(response: str) -> str:
    """Strip markdown code fences and extract JSON from AI response."""
    import re
    text = response.strip()
    # Remove markdown code block wrappers
    text = re.sub(r'^```(?:json)?\s*\n?', '', text, flags=re.IGNORECASE)
    text = re.sub(r'\n?```\s*$', '', text)
    # Remove PPT-Agent [PPT_OUTLINE] markers
    text = re.sub(r'\[PPT_OUTLINE\]\s*', '', text)
    text = re.sub(r'\[/PPT_OUTLINE\]\s*', '', text)
    return text.strip()


def _safe_run_async(coro):
    """Run an async coroutine from any thread — works in FastAPI thread pools.

    On Windows, asyncio.run() can leak "Event loop is closed" errors from
    httpx/AsyncOpenAI cleanup callbacks. We always spawn a dedicated daemon
    thread to fully isolate the event loop lifecycle.
    """
    import concurrent.futures
    import threading

    result = [None]
    error = [None]

    def _runner():
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                result[0] = loop.run_until_complete(coro)
            finally:
                try:
                    loop.call_soon(lambda: None)
                    loop.run_until_complete(asyncio.sleep(0))
                    loop.close()
                except RuntimeError:
                    pass  # Event loop already closed by httpx cleanup on Windows
        except Exception as e:
            error[0] = e

    t = threading.Thread(target=_runner, daemon=True)
    t.start()
    t.join(timeout=300)
    if t.is_alive():
        raise TimeoutError("LLM API call timed out after 300s")
    if error[0]:
        raise error[0]
    return result[0]

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _set_font_all(run, font_name: str):
    """Set both Latin and East Asian font on a run.

    run.font.name only sets the Latin font. For Chinese characters,
    the East Asian (ea) font must also be set via XML.
    """
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(f"{{{_NS}}}ea")
    if ea is None:
        ea = etree.SubElement(rPr, f"{{{_NS}}}ea")
    ea.set("typeface", font_name)

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
WORKSPACE_ROOT = os.path.dirname(BASE_DIR)  # d:\YISHAOAGENT
EXPORT_DIR = os.path.join(BASE_DIR, "data", "exports")
os.makedirs(EXPORT_DIR, exist_ok=True)


def _load_branding() -> tuple:
    """Load copyright and signature from DB settings. Returns (copyright_str, signature_str)."""
    try:
        db = get_db()
        rows = db.execute(
            "SELECT key, value FROM settings WHERE key IN ('branding_copyright', 'branding_signature')"
        ).fetchall()
        db.close()
        copyright_str = ""
        sig_str = ""
        for r in rows:
            if r["key"] == "branding_copyright" and r["value"]:
                copyright_str = r["value"].strip()
            elif r["key"] == "branding_signature" and r["value"]:
                sig_str = r["value"].strip()
        if not copyright_str and not sig_str:
            return ("&copy; 2024 美食研究所 &middot; 保密文档", "商务部监制")
        return (copyright_str, sig_str)
    except Exception:
        return ("&copy; 2024 美食研究所 &middot; 保密文档", "商务部监制")


def _load_style_from_template(template_id: str = None) -> str:
    """Extract style_id from template rules. Defaults to 'business'."""
    if not template_id:
        return "business"
    db = get_db()
    try:
        row = db.execute(
            "SELECT rules FROM templates WHERE id = ?", (template_id,)
        ).fetchone()
        if row and row["rules"]:
            rules = json.loads(row["rules"])
            style_id = rules.get("style_id", "")
            if style_id:
                return style_id
    except Exception:
        pass
    finally:
        db.close()
    return "business"


_COLUMN_ID_RE = re.compile(r'^[a-zA-Z0-9_-]+$')

def _validate_column_id(column_id: str) -> str:
    """Validate column_id for use in file paths and DB queries.

    Returns the sanitized column_id if valid (alphanumeric + underscore + hyphen),
    or empty string if the input is unsafe. Prevents path traversal via '../' etc.
    """
    if column_id and _COLUMN_ID_RE.match(column_id):
        return column_id
    return ""



# In-memory status for PPT generation progress polling
_ppt_status: dict = {}

def get_ppt_status(project_id: str) -> dict | None:
    """Return current PPT generation status for a project, or None if idle."""
    return _ppt_status.get(project_id)


# In-memory log for PPT generation progress display
_ppt_log: dict[str, list[dict]] = {}

def get_ppt_log(project_id: str) -> list:
    """Return current PPT generation log for a project."""
    return _ppt_log.get(project_id, [])

def _append_log(project_id: str, message: str):
    """Append a timestamped log entry for a project."""
    if not project_id:
        return
    import datetime as _dt
    entry = {"time": _dt.datetime.now().strftime("%H:%M:%S"), "message": message}
    if project_id not in _ppt_log:
        _ppt_log[project_id] = []
    _ppt_log[project_id].append(entry)


def generate_ppt(content: str, template_id: str = None, branding: dict = None,
                 output_dir: str = None, provider_id: str = "",
                 model: str = "", slide_plan: list = None, format: str = "svg",
                 project_name: str = "", column_id: str = "",
                 color_scheme: str = "deep-blue", temperature: float = 0.3,
                 project_id: str = "",
                 temp_keyword: float = 0, temp_research: float = 0,
                 temp_outline: float = 0, temp_fill: float = 0,
                 temp_cards: float = 0, temp_html: float = 0,
                 temp_svg_batch: float = 0, temp_svg_single: float = 0,
                 temp_review: float = 0, temp_fix: float = 0,
                 temp_holistic: float = 0, temp_holistic_fix: float = 0,
                 temp_stage_outline: float = 0, temp_stage_generation: float = 0,
                 temp_stage_review: float = 0) -> str:
    """Generate presentation from content. Default: SVG (PPT-Agent Bento Grid).

    All 12 stage temperatures follow the pattern: 0 = use `temperature` as fallback."""

    def _t(val, default=0.3):
        return val if val > 0 else temperature

    st = dict(
        keyword=_t(temp_keyword, 0.3),
        research=_t(temp_research, 0.7),
        outline=_t(temp_outline, 1.0),
        fill=_t(temp_fill, 1.0),
        cards=_t(temp_cards, 0.7),
        html=_t(temp_html, 0.8),
        svg_batch=_t(temp_svg_batch, 0.7),
        svg_single=_t(temp_svg_single, 0.7),
        review=_t(temp_review, 0.3),
        fix=_t(temp_fix, 0.7),
        holistic=_t(temp_holistic, 0.3),
        holistic_fix=_t(temp_holistic_fix, 0.7),
    )

    # Stage-level overrides: when set, override all per-step temps in that stage
    if temp_stage_outline > 0:
        st.update(keyword=temp_stage_outline, research=temp_stage_outline,
                  outline=temp_stage_outline, fill=temp_stage_outline)
    if temp_stage_generation > 0:
        st.update(cards=temp_stage_generation, html=temp_stage_generation,
                  svg_batch=temp_stage_generation, svg_single=temp_stage_generation)
    if temp_stage_review > 0:
        st.update(review=temp_stage_review, fix=temp_stage_review,
                  holistic=temp_stage_review, holistic_fix=temp_stage_review)

    # Track status for progress polling
    if project_id:
        _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在生成大纲...", "message": "AI 分析内容中"}
        _ppt_log[project_id] = []  # clear stale logs from previous run
        _append_log(project_id, "开始生成 PPT，正在进行内容分析...")

    # Build human-readable dir name: "{项目名}_{类型}" (e.g. "测试1_道术PPT")
    safe_proj = "".join(c for c in project_name if c.isalnum() or c in "._- ()（）").strip() if project_name else ""
    dir_name = f"{safe_proj}_{column_id}" if safe_proj and column_id else ""
    prs = None
    slide_data = None
    rules = {}
    typography_profile = None
    style_id = "business"

    if slide_plan is not None:
        slide_data = slide_plan

    if template_id:
        db = get_db()
        try:
            row = db.execute(
                "SELECT file_path, rules, typography_profile, branding_config FROM templates WHERE id = ?",
                (template_id,)).fetchone()
            if row:
                if row["file_path"]:
                    fp = row["file_path"]
                    if not os.path.isabs(fp):
                        fp_ws = os.path.join(WORKSPACE_ROOT, fp)
                        fp_be = os.path.join(BASE_DIR, fp)
                        if os.path.exists(fp_ws):
                            fp = fp_ws
                        elif os.path.exists(fp_be):
                            fp = fp_be
                    if os.path.exists(fp) and fp.endswith('.pptx'):
                        prs = Presentation(fp)
                if not branding and row["branding_config"]:
                    try:
                        branding = json.loads(row["branding_config"])
                    except Exception:
                        pass
                if row["typography_profile"]:
                    try:
                        typography_profile = json.loads(row["typography_profile"])
                    except Exception:
                        pass
                if row["rules"]:
                    try:
                        rules = json.loads(row["rules"])
                    except Exception:
                        pass

                # Extract style_id from template rules (for SVG mode)
                style_id = rules.get("style_id", "business")

                print(f"[PPT-DBG] AI check: pid={bool(provider_id)} model={bool(model)} "
                      f"style={style_id} "
                      f"rules_empty={not rules} "
                      f"content_len={len(content.strip()) if content else 0}", flush=True)

                if slide_data is None and provider_id and model and content.strip():
                    print("[PPT-DBG] Using AI staged generation", flush=True)
                    col_prompt = ""
                    col_skill = ""
                    try:
                        cfg2 = None
                        if column_id:
                            cfg2 = db.execute(
                                "SELECT prompt, skill, rules FROM column_configs WHERE column_id = ?",
                                (column_id,)
                            ).fetchone()
                        if cfg2:
                            col_prompt = cfg2["prompt"] or ""
                            col_skill = cfg2["skill"] or ""
                            # Merge column_configs rules into template rules (template wins)
                            try:
                                col_rules = json.loads(cfg2["rules"] or "{}")
                                if col_rules and isinstance(col_rules, dict):
                                    merged = dict(col_rules)
                                    merged.update(rules)
                                    rules = merged
                            except Exception:
                                pass
                    except Exception:
                        pass
                    slide_data = _generate_slides_staged(provider_id, model, rules, content,
                                                         col_prompt, col_skill,
                                                         temperature=temperature,
                                                         st=st, column_id=column_id,
                                                         color_scheme=color_scheme,
                                                         project_id=project_id)
                    print(f"[PPT-DBG] AI result: {bool(slide_data)}, "
                          f"slides={len(slide_data) if slide_data else 0}", flush=True)
                    if not slide_data:
                        print("[PPT-DBG] AI generation failed — no fallback, returning error", flush=True)
                else:
                    if not slide_data:
                        print("[PPT-DBG] No AI provider/model and no saved plan — cannot generate", flush=True)
        finally:
            db.close()

    # ── HTML path (new two-phase pipeline: slides have 'html' field) ──
    if format != "pptx" and slide_data and isinstance(slide_data, list) and len(slide_data) > 0:
        first = slide_data[0]
        if isinstance(first, dict) and "html" not in first and "heading" in first and provider_id and model:
            # Outline format — need to run Phase 5a (Structure) + Phase 5b (HTML)
            print(f"[PPT-DBG] Outline input detected: {len(slide_data)} slides, running Structure+HTML", flush=True)
            if project_id:
                _append_log(project_id, f"大纲已有 {len(slide_data)} 页，跳过分析，直接规划布局...")
                _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在规划布局...", "message": f"大纲已有 {len(slide_data)} 页，进入布局规划"}
            from services.llm_service import generate as llm_generate
            structure = _stage2_structure(provider_id, model, llm_generate,
                                           slide_data, style_id=style_id,
                                           temperature=st['cards'],
                                           column_id=column_id,
                                           color_scheme=color_scheme)
            if structure:
                if project_id:
                    _append_log(project_id, f"布局规划完成，共 {len(structure)} 页，开始并行生成 HTML")
                    _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在生成页面...", "message": f"{len(structure)} 页布局已规划，并行生成 HTML 中", "slides_done": 0, "slides_total": len(structure)}
                html_slides = _stage2_html_per_slide(provider_id, model, llm_generate,
                                                      structure, style_id=style_id,
                                                      temperature=st['html'],
                                                      column_id=column_id,
                                                      color_scheme=color_scheme,
                                                      project_id=project_id)
                if html_slides:
                    slide_data = html_slides
                    first = slide_data[0] if slide_data else {}
                    if project_id:
                        _append_log(project_id, f"PPT 生成完成，共 {len(html_slides)} 页")
        if isinstance(first, dict) and "html" in first:
            # New HTML pipeline — slides already have complete HTML, just assemble
            print(f"[PPT-DBG] HTML pipeline: {len(slide_data)} slides with inline HTML", flush=True)
            title = first.get("heading", "") or "Presentation"

            # Build output directory
            html_dir = output_dir if output_dir else os.path.join(EXPORT_DIR, "html_decks")
            if dir_name:
                html_dir = os.path.join(html_dir, dir_name)
            os.makedirs(html_dir, exist_ok=True)

            scheme_data = _load_scheme_data(style_id, color_scheme)
            gen_canvas_w, gen_canvas_h = _get_canvas_dimensions(column_id)
            deck_html = _assemble_html_deck(slide_data, title, style_id, scheme_data, canvas_w=gen_canvas_w, canvas_h=gen_canvas_h)
            if scheme_data:
                deck_html = _resolve_color_vars(deck_html, scheme_data, css_vars=True)
            html_path = os.path.join(html_dir, "index.html")
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(deck_html)
            _logger.info(f"HTML deck written: {html_path}")

            # Save unresolved variable version for future color scheme switching
            vars_slide_data = [{**s, "html": s.get("html_vars", s.get("html", ""))} for s in slide_data]
            deck_vars = _assemble_html_deck(vars_slide_data, title, style_id, scheme_data, canvas_w=gen_canvas_w, canvas_h=gen_canvas_h)
            vars_path = os.path.join(html_dir, "index_vars.html")
            with open(vars_path, "w", encoding="utf-8") as f:
                f.write(deck_vars)
            _logger.info(f"Variable-version deck written: {vars_path}")

            # Save individual slide files for future granular edits
            _save_slide_files(html_dir, slide_data)
            _logger.info(f"Individual slide files saved to {os.path.join(html_dir, SLIDES_DIR)}")

            # Speaker notes
            notes_path = os.path.join(html_dir, "speaker-notes.md")
            try:
                notes_md = _generate_speaker_notes(slide_data, title)
                with open(notes_path, "w", encoding="utf-8") as f:
                    f.write(notes_md)
            except Exception as e:
                _logger.warning(f"Speaker notes failed (non-critical): {e}")

            if project_id:
                _ppt_status.pop(project_id, None)
            return html_path, slide_data

        # ── Old SVG path (legacy: slides have 'zones' field) ──
        print(f"[PPT-DBG] Rendering SVG deck: {len(slide_data)} slides, style={style_id}", flush=True)
        from services.svg_designer import DeckDesigner
        from services.llm_service import generate as llm_generate

        designer = DeckDesigner(style_name=style_id, branding=branding)
        title = "Presentation"
        if isinstance(slide_data[0], dict):
            z = slide_data[0].get("zones", {})
            title = z.get("heading", "") or z.get("title", "") or "Presentation"

        try:
            # ── Try AI-SVG direct generation (PPT-Agent slide-core quality) ──
            ai_svg_slides = None
            if provider_id and model:
                print(f"[PPT-DBG] Attempting AI-SVG direct generation...", flush=True)
                try:
                    ai_svg_slides = _stage3_svg(provider_id, model, llm_generate,
                                                 slide_data, style_id=style_id,
                                                 rules=rules, temperature=temperature,
                                                 st=st, column_id=column_id,
                                                 color_scheme=color_scheme)
                except Exception as e:
                    print(f"[PPT-DBG] AI-SVG generation failed: {e}", flush=True)

            if ai_svg_slides and len(ai_svg_slides) > 0:
                slides_info, preview_html, svg_dir, _run_id = designer.render_deck_from_ai_svg(
                    ai_svg_slides, title, output_dir, dir_name
                )
                print(f"[PPT-DBG] AI-SVG deck: {len(slides_info)} slides → {svg_dir}", flush=True)
            else:
                slides_info, preview_html, svg_dir, _run_id = designer.render_deck(
                    slide_data, title, output_dir, dir_name
                )
                print(f"[PPT-DBG] Code-rendered SVG deck: {len(slides_info)} slides → {svg_dir}", flush=True)

            html_path = os.path.join(svg_dir, "index.html")

            # ── Speaker notes (PPT-Agent Phase 7 deliverable) ──
            notes_path = os.path.join(svg_dir, "speaker-notes.md")
            try:
                notes_md = _generate_speaker_notes(slide_data, title)
                with open(notes_path, "w", encoding="utf-8") as f:
                    f.write(notes_md)
                _logger.info(f"Speaker notes written: {notes_path}")
            except Exception as e:
                _logger.warning(f"Speaker notes generation failed (non-critical): {e}")

            if project_id:
                _ppt_status.pop(project_id, None)
            return html_path, slide_data
        except Exception as e:
            print(f"[PPT-DBG] SVG rendering failed ({e}), falling back to PPTX", flush=True)
            format = "pptx"  # fall through to PPTX path

    # ── PPTX fallback ──
    if format == "pptx":
        if prs is None:
            prs = Presentation()

        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        if slide_data and isinstance(slide_data, list) and len(slide_data) > 0:
            from services.ppt_designer import extract_design, build_slide
            _remove_all_slides(prs)
            design = extract_design(rules, typography_profile)
            for sd in slide_data:
                if not isinstance(sd, dict):
                    continue
                slide_type = sd.get("type", "content")
                zones = sd.get("zones", {})
                if not isinstance(zones, dict):
                    zones = {}
                build_slide(prs, slide_type, zones, design)
        else:
            _remove_all_slides(prs)
            _mechanical_fill(prs, content)

        if branding:
            for slide in prs.slides:
                left = Inches(0.5)
                top = Inches(7.0)
                width = Inches(12.3)
                height = Inches(0.4)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                if branding.get("copyright"):
                    tf.text = branding["copyright"]
                if branding.get("signature"):
                    p = tf.add_paragraph()
                    p.text = branding["signature"]
                    p.alignment = PP_ALIGN.RIGHT

        if typography_profile:
            spacing = _compute_spacing_from_profile(typography_profile)
        else:
            db = get_db()
            try:
                spacing = _load_spacing_rules(db)
            finally:
                db.close()
        _normalize_formatting(prs, spacing)

        filename = f"ppt_{os.urandom(4).hex()}.pptx"
        target_dir = output_dir if output_dir else EXPORT_DIR
        os.makedirs(target_dir, exist_ok=True)
        filepath = os.path.join(target_dir, filename)
        prs.save(filepath)
        if project_id:
            _ppt_status.pop(project_id, None)
        return filepath, slide_data

    # No viable path
    if project_id:
        _ppt_status.pop(project_id, None)
    return None, None


# ── Staged AI slide generation ──

def _generate_speaker_notes(slide_data: list, title: str) -> str:
    """Generate speaker notes markdown from slide data — PPT-Agent Phase 7 deliverable.

    Extracts talking points, transitions, and timing from outline data.
    Output format matches PPT-Agent's speaker-notes.md structure.
    """
    import re

    def _extract_from_html(html: str, field: str) -> str:
        """Extract heading or body content from inline HTML when JSON fields are empty."""
        if not html:
            return ""
        if field == "heading":
            # Try h1/h2/h3 with font-size > 30px
            m = re.search(r'<(?:h[123]|div)\s[^>]*?font-size:\s*(\d+)px[^>]*?>\s*(.+?)\s*</(?:h[123]|div)>', html)
            if m and int(m.group(1)) >= 30:
                return m.group(2).strip()
            # Fallback: any heading-like div
            m = re.search(r'font-size:\s*(3[6-9]|[4-9]\d)px[^>]*?>\s*(.{1,80}?)\s*</div>', html)
            if m:
                return m.group(2).strip()
            return ""
        if field == "body":
            # Extract longest text block
            texts = re.findall(r'font-size:\s*1[6-8]px[^>]*?>\s*(.{30,500}?)\s*</(?:p|div)>', html)
            if texts:
                return max(texts, key=len).strip()
            return ""
        return ""

    lines = [f"# Speaker Notes: {title}", "", f"**Total slides**: {len(slide_data)}", ""]
    total_minutes = 0

    for i, sd in enumerate(slide_data):
        if not isinstance(sd, dict):
            continue
        seq = sd.get("seq", i + 1)
        zones = sd.get("zones", {}) if isinstance(sd.get("zones"), dict) else {}
        html = sd.get("html", "")
        heading = zones.get("heading", "") or sd.get("heading", "") or _extract_from_html(html, "heading") or f"Slide {seq}"
        body = zones.get("body", "") or sd.get("body", "") or _extract_from_html(html, "body")
        notes = zones.get("notes", "") or sd.get("notes", "")
        key_points = sd.get("key_points", []) or zones.get("key_points", [])
        stype = sd.get("type", "content")

        # Estimate timing: cover 1min, content 2-3min, quote/section 0.5min
        time_map = {"cover": 1, "toc": 1.5, "quote": 0.5, "section": 0.5,
                     "summary": 1.5, "closing": 1, "copyright": 0.5}
        minutes = time_map.get(stype, 2)
        total_minutes += minutes

        lines.append(f"---")
        lines.append(f"## Slide {seq:02d}: {heading}")
        lines.append(f"**Type**: {stype} | **Time**: ~{minutes} min")
        lines.append("")

        if key_points:
            lines.append("**Talking Points:**")
            for kp in key_points[:5]:
                if isinstance(kp, str) and kp.strip():
                    lines.append(f"- {kp.strip()}")
            lines.append("")

        if notes and notes.strip():
            lines.append(f"**Notes**: {notes.strip()}")
            lines.append("")

        if body and body.strip():
            # Take first 2 sentences of body as context
            sentences = body.strip().split("。")
            context = "。".join(sentences[:2])
            if len(sentences) > 1:
                context += "。"
            lines.append(f"**Context**: {context.strip()}")
            lines.append("")

        # Transition hint
        if i < len(slide_data) - 1:
            next_slide = slide_data[i + 1]
            if isinstance(next_slide, dict):
                next_zones = next_slide.get("zones", {}) if isinstance(next_slide.get("zones"), dict) else {}
                next_html = next_slide.get("html", "")
                next_heading = next_zones.get("heading", "") or next_slide.get("heading", "") or _extract_from_html(next_html, "heading")
                if next_heading:
                    lines.append(f'**Transition**: "Now moving on to {next_heading}..."')
                    lines.append("")

    lines.append("---")
    lines.append(f"## Estimated Total Time: ~{total_minutes} minutes")
    lines.append(f"*Generated by YISHAOAGENT PPT-Agent pipeline*")

    return "\n".join(lines)


def _resolve_review_models(provider_id: str, model: str) -> tuple[str, str, str]:
    """Try to find a different model for review (cross-model validation).

    PPT-Agent advantage: Claude generates, Gemini reviews — different models
    catch different errors. We replicate this: if multiple providers/models are
    configured, pick a different one for review. If only one model exists,
    return the same model but flag for self-review mode.

    Returns (review_provider_id, review_model, mode):
      mode = "cross_model" | "self_review"
    """
    try:
        db = get_db()
        rows = db.execute(
            "SELECT id, model FROM llm_providers WHERE is_enabled = 1"
        ).fetchall()
        db.close()

        candidates = []
        for row in rows:
            pid = row["id"]
            # Get configured model from provider record
            prov_model = row.get("model", "")
            candidates.append((pid, prov_model))

        # Find a different provider or different model
        for pid, prov_model in candidates:
            if pid != provider_id:
                # Different provider — best case
                _logger.info(f"Cross-model review: gen={provider_id}/{model}, review={pid}/{prov_model or model}")
                return pid, prov_model or model, "cross_model"
            elif prov_model and prov_model != model:
                # Same provider, different model
                _logger.info(f"Cross-model review: gen={provider_id}/{model}, review={pid}/{prov_model}")
                return pid, prov_model, "cross_model"

        # If we get here, only one model available
        _logger.info(f"Self-review mode: only {provider_id}/{model} available")
        return provider_id, model, "self_review"
    except Exception as e:
        _logger.info(f"Model resolution failed ({e}), falling back to self-review")
        return provider_id, model, "self_review"


def _phase2_research(provider_id: str, model: str, llm_generate, sop_content: str,
                     temperature: float = 0.3,
                     temp_keyword: float = 0, temp_research: float = 0) -> str:
    """Phase 2 (Research): Direct deep analysis of SOP using LLM trained knowledge.

    No web search — the LLM already has sufficient domain knowledge.
    One LLM call produces the full structured research context.
    """
    if not sop_content or not sop_content.strip():
        return ""

    research_system = _load_research_prompt()

    research_user = f"""## SOP 文档（唯一分析对象）
{sop_content}

## 任务
深度分析上述 SOP 文档，输出结构化的研究上下文。这将用于后续的 PPT 大纲规划和内容提取。"""

    try:
        response = _safe_run_async(llm_generate(provider_id, model,
            research_system, research_user, temperature=temp_research or temperature))
        _logger.info(f"Phase 2 research: {len(response)} chars of analysis")
        return response
    except Exception as e:
        _logger.warning(f"Phase 2 research failed: {e}")
        return ""


def _generate_slides_staged(provider_id: str, model: str, rules: dict, sop_content: str,
                            system_prompt: str = "", skill_template: str = "",
                            temperature: float = 0.3,
                            st: dict = None, column_id: str = "",
                            color_scheme: str = "deep-blue",
                            project_id: str = "") -> list | None:
    """Two-phase HTML slide generation pipeline.

    Phase 2 (Research): Deep SOP analysis → research context.
    Phase 4 (Outline): Extract & organize content → outline with body text.
    Phase 5a (Structure): AI plans slide types, layouts, card allocations — lightweight.
    Phase 5b (HTML): Per-slide parallel HTML generation with design-system.md as guide.

    AI generates complete inline-styled HTML. No JSON cards, no SVG code rendering."""
    from services.llm_service import generate as llm_generate
    if st is None:
        st = {}

    if project_id:
        _append_log(project_id, "进入分阶段生成，开始深度分析内容...")

    # ── Phase 2: Research (SOP deep analysis → research-context) ──
    research_context = ""
    try:
        research_context = _phase2_research(provider_id, model, llm_generate, sop_content,
                                            temp_keyword=st.get('keyword', temperature),
                                            temp_research=st.get('research', temperature))
        if research_context:
            _logger.info(f"Phase 2 research done: {len(research_context)} chars")
            if project_id:
                _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在分析内容...", "message": "研究完成，进入大纲规划"}
                _append_log(project_id, "内容分析完成，进入大纲规划")
    except Exception as e:
        _logger.warning(f"Phase 2 research failed (proceeding without): {e}")

    # ── Phase 4: Content extraction (outline + body per slide) ──
    stage1 = _stage1_content(provider_id, model, llm_generate, rules, sop_content,
                             system_prompt, skill_template, research_context=research_context,
                             temp_outline=st.get('outline', temperature),
                             temp_fill=st.get('fill', temperature),
                             column_id=column_id)
    if not stage1:
        return None
    _logger.info(f"Phase 4 outline: {len(stage1)} slides extracted")
    if project_id:
        _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在规划布局...", "message": f"已提取 {len(stage1)} 页大纲"}
        _append_log(project_id, f"大纲提取完成，共 {len(stage1)} 页")

    # ── Two-Phase HTML Pipeline ──
    style_id = rules.get("style_id", "business")

    cw, ch = _get_canvas_dimensions(column_id)
    is_a4 = ch > cw

    # Phase 1: Structure planning (lightweight, one LLM call for all slides)
    # ── col3/A4: skip structure planning — template defines fixed 6-part structure ──
    if is_a4:
        structure = stage1  # Use stage1 directly, no layout/card planning needed
        # Normalize stage1 format for _stage2_html_per_slide compatibility
        for s in structure:
            if 'type' not in s:
                s['type'] = s.get('page_type', 'content')
            s.setdefault('layout', '')
            s.setdefault('cards', [])
            s.setdefault('has_chart', False)
            s.setdefault('chart_hint', '')
        _logger.info(f"A4 document mode: {len(structure)} slides, skipping structure planning")
    else:
        structure = _stage2_structure(provider_id, model, llm_generate,
                                      stage1, style_id=style_id,
                                      temperature=st.get('cards', temperature),
                                      column_id=column_id,
                                      color_scheme=color_scheme)
        if not structure:
            return None
        _logger.info(f"Phase 1 structure: {len(structure)} slides planned")
    if project_id:
        if is_a4:
            _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在生成文档...", "message": f"{len(structure)} 页文档，并行生成 HTML 中", "slides_done": 0, "slides_total": len(structure)}
            _append_log(project_id, f"A4 文档模式，共 {len(structure)} 页，开始并行生成 HTML")
        else:
            _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在生成页面...", "message": f"{len(structure)} 页布局已规划，并行生成 HTML 中", "slides_done": 0, "slides_total": len(structure)}
            _append_log(project_id, f"布局规划完成，共 {len(structure)} 页，开始并行生成 HTML")

    # Phase 2: Per-slide HTML generation (parallel, design-system.md guided)
    html_slides = _stage2_html_per_slide(provider_id, model, llm_generate,
                                         structure, style_id=style_id,
                                         temperature=st.get('html', temperature),
                                         column_id=column_id,
                                         color_scheme=color_scheme,
                                         project_id=project_id)
    if not html_slides:
        _logger.warning("Phase 2 HTML generation failed, using fallback")
        html_slides = _fallback_stage1_to_html_slides(stage1, style_id)

    _logger.info(f"Phase 2 HTML: {len(html_slides)} slides generated")
    if project_id:
        _append_log(project_id, f"{'A4 文档' if is_a4 else 'PPT'} 生成完成，共 {len(html_slides)} 页")
    return html_slides


def _generate_outline_only(provider_id, model, rules, sop_content,
                           system_prompt="", skill_template="",
                           temperature: float = 0.3,
                           st: dict = None, project_id: str = "",
                           column_id: str = "") -> tuple:
    """Run only Phase 2 (Research) + Phase 4 (Outline+Content Fill).

    Returns (outline_json, outline_text) — outline_json is the structured
    JSON data, outline_text is natural-language text for human reading/editing.
    """
    from services.llm_service import generate as llm_generate
    if st is None:
        st = {}

    if project_id:
        _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在生成大纲...", "message": "AI 分析内容中"}
        _ppt_log[project_id] = []  # clear stale logs from previous run
        _append_log(project_id, "开始生成大纲，正在进行内容分析...")

    research_context = ""
    try:
        research_context = _phase2_research(provider_id, model, llm_generate, sop_content,
                                            temp_keyword=st.get('keyword', temperature),
                                            temp_research=st.get('research', temperature))
        if research_context:
            _logger.info(f"Outline-only research done: {len(research_context)} chars")
            if project_id:
                _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在生成大纲...", "message": "AI 提取大纲中"}
                _append_log(project_id, "内容分析完成，进入大纲提取")
    except Exception as e:
        _logger.warning(f"Outline-only research failed (proceeding without): {e}")

    stage1 = _stage1_content(provider_id, model, llm_generate, rules, sop_content,
                              system_prompt, skill_template, research_context=research_context,
                              temp_outline=st.get('outline', temperature),
                              temp_fill=st.get('fill', temperature),
                              column_id=column_id)
    if not stage1:
        if project_id:
            _ppt_status.pop(project_id, None)
        return None, ""

    outline_text = _slides_to_human_text(stage1)
    if project_id:
        _ppt_status[project_id] = {"phase": "done", "phase_label": "大纲生成完成", "message": f"共 {len(stage1)} 页"}
        _append_log(project_id, f"大纲生成完成，共 {len(stage1)} 页")
    return stage1, outline_text


def _slides_to_human_text(slides: list) -> str:
    """Convert structured slide outline into natural-language text.

    No field labels, no machine tokens — just readable paragraphs that
    anyone can understand and edit.
    """
    type_labels = {
        "cover": "封面", "toc": "目录", "section": "章节分隔", "chapter": "章节页",
        "content": "内容页", "data": "数据页", "data_hero": "数据突出",
        "technique": "技法页", "principle": "原则页", "process_flow": "流程图",
        "process_timeline": "流程时间线", "timeline": "时间线",
        "comparison": "对比页", "duo_compare": "双栏对比",
        "table": "表格页", "grid_cards": "卡片组", "image_grid": "图片网格",
        "quote": "引言页", "image_hero": "图片突出",
        "troubleshoot": "问题排查",
        "appendix": "附录", "copyright": "版权页",
        "closing": "结尾页", "summary": "总结页",
        "document": "A4文档",
    }

    parts = []
    for s in slides:
        seq = s.get("seq", len(parts) + 1)
        heading = s.get("heading", "")
        page_type = s.get("page_type", "content")
        type_cn = type_labels.get(page_type, page_type)
        lead = s.get("lead", "")
        body = s.get("body", "")
        key_points = s.get("key_points", [])
        notes = s.get("notes", "")
        kicker = s.get("kicker", "")

        block = f"第{seq}页 — {type_cn}"
        if kicker:
            block += f" · {kicker}"
        block += "\n"

        if heading:
            block += f"\n{heading}\n"
        if lead:
            block += f"\n{lead}\n"
        if body:
            block += f"\n{body}\n"
        if key_points:
            kp_text = " · ".join(str(kp) for kp in key_points if isinstance(kp, str) and kp.strip())
            if kp_text:
                block += f"\n关键要点：{kp_text}\n"
        if notes:
            block += f"\n（{notes}）\n"

        parts.append(block.strip())

    return "\n\n———\n\n".join(parts)


def _human_text_to_json(provider_id, model, human_text: str, original_json: list) -> list | None:
    """Use LLM to convert edited natural-language text back to structured JSON.

    Low temperature (0.1) ensures near-deterministic output. The LLM is
    instructed to preserve all content — no additions, deletions, or rewrites.
    Falls back to original_json on any failure.
    """
    from services.llm_service import generate as llm_generate

    if not provider_id or not model:
        _logger.warning("No LLM configured for human_text_to_json — returning original JSON")
        return original_json

    # Provide original JSON as schema reference so LLM knows the exact structure
    original_schema = json.dumps(original_json[:2], ensure_ascii=False, indent=2) if len(original_json) > 1 else json.dumps(original_json, ensure_ascii=False, indent=2)

    # Extract unique page_types from original JSON as the allowed set
    existing_types = set()
    if isinstance(original_json, list):
        for s in original_json:
            if isinstance(s, dict) and s.get("page_type"):
                existing_types.add(s["page_type"])
    type_list = ", ".join(sorted(existing_types)) if existing_types else "cover, toc, content, summary"

    t2j_template = _load_text_to_json_prompt()
    if not t2j_template:
        t2j_template = """你是数据整理专家。你的唯一任务是将人类编辑的自然文本转回结构化 JSON。

严格规则：
1. 从原文中提取每一页的序号(seq)、类型(page_type)、标题(heading)、副标题(lead)、正文(body)、关键要点(key_points)、备注(notes)、章节标签(kicker)
2. page_type 只使用: {type_list}
3. 不新增任何内容，不删除任何内容，不改写原文
4. 原文中没提到的字段不要编造
5. key_points 是字符串数组
6. 输出纯 JSON 数组，不要用 markdown 包裹"""
    system = t2j_template.format(type_list=type_list)

    user = f"""## 原始 JSON 结构参考（字段名和类型以此为基准）
```json
{original_schema}
```

## 用户编辑后的文本（唯一内容来源）
{human_text}

## 任务
将上述文本的每一页转为 JSON，保持原始结构的所有字段名和类型。输出纯 JSON 数组。"""

    for attempt in range(2):
        try:
            response = _safe_run_async(llm_generate(provider_id, model, system, user, temperature=0.1))
            response = _clean_json_response(response)
            data = json.loads(response)
            if isinstance(data, list) and len(data) > 0:
                # Validate: every slide must have at least heading or body
                valid = [s for s in data if isinstance(s, dict) and (s.get("heading") or s.get("body"))]
                if len(valid) > 0:
                    _logger.info(f"Human text → JSON: {len(valid)} slides converted (attempt {attempt+1})")
                    return valid
            _logger.warning(f"Human text → JSON attempt {attempt+1}: invalid output structure")
        except Exception as e:
            _logger.warning(f"Human text → JSON attempt {attempt+1} failed: {e}")

    _logger.error("Human text → JSON: all attempts failed, returning original JSON as fallback")
    return original_json


def _stage1_content(provider_id, model, llm_generate, rules, sop_content,
                    system_prompt: str = "", skill_template: str = "",
                    research_context: str = "", temperature: float = 0.3,
                    temp_outline: float = 0, temp_fill: float = 0,
                    column_id: str = "") -> list | None:
    """Stage 1: Two-phase — outline first, then fill content per slide in batches.

    research_context from Phase 2 (research-core) provides pre-analyzed SOP structure,
    key insights, data points, and suggested focus areas.
    """
    # Load outline-architect prompt (UI-configurable via column_configs rules, disk fallback)
    outline_spec = ""
    if rules.get("outline_architect_prompt"):
        outline_spec = rules["outline_architect_prompt"]
    if not outline_spec:
        outline_spec = _load_outline_spec(column_id)

    # Load cognitive design principles (referenced by outline-architect.md)
    cognitive_spec_stage1 = rules.get("cognitive_design_principles", "") if rules else ""
    if not cognitive_spec_stage1:
        cognitive_spec_stage1 = _load_cognitive_spec(column_id)

    # Always combine column role prompt + VI design rules (never discard either)
    role_block = system_prompt if system_prompt else ""
    vi_block = f"""{outline_spec}

## 认知设计原则（必须遵守）
{cognitive_spec_stage1}
"""
    pyramid_rules = _load_outline_rules()
    if not pyramid_rules:
        pyramid_rules = """重要补充：从提供的 SOP 文章中提取内容，严格按金字塔原理组织大纲。
核心纪律：大纲中的每个「技法」/「步骤」/「章节」必须独占一页，绝不合并。
即使 SOP 很短，封面页和总结页也不可省略。
每页一主题，内容聚焦不堆砌。输出纯 JSON，不要用 markdown 包裹。"""

    base_system = "\n\n".join(b for b in [role_block, vi_block, pyramid_rules] if b.strip())

    skill_block = ""
    if skill_template and not skill_template.strip().startswith('|'):
        skill_block = f"""## 幻灯片结构模板（必须严格遵循的页面结构和字段）
{skill_template}

"""

    # ── Research context block (from Phase 2) ──
    research_block = ""
    if research_context:
        research_block = f"""## 研究上下文（SOP 深度分析结果）
{research_context[:4000]}

---

"""

    # ── Phase 1: Generate outline (headings + page types + layout hints, lightweight) ──
    style_id = rules.get("style_id", "business")
    page_type_prompt = _build_page_type_prompt(style_id, column_id)

    # Column-aware output requirements
    cw_stage1, ch_stage1 = _get_canvas_dimensions(column_id) if column_id else (1280, 720)
    is_a4_stage1 = ch_stage1 > cw_stage1
    if is_a4_stage1:
        output_reqs = """## 输出要求
- 严格遵循上方「文档结构模板」的栏目章节结构和 JSON 格式
- 栏目结构、构建块类型、硬约束均以模板为准，不得自行增删章节
- 仅输出 JSON，不输出其他文字"""
    else:
        output_reqs = """## 输出要求
- 严格遵循上方「幻灯片结构模板」的栏目章节结构和 JSON 格式
- 栏目结构、页面类型、硬约束均以模板为准，不得自行增删章节
- layout_hint 从以下选择: single_focus, two_column, two_column_asymmetric, three_column, hero_grid, mixed_grid, dashboard, timeline, horizontal_split, full_bleed
- visual_weight 从以下选择: low, medium, high
- 仅输出 JSON，不输出其他文字"""

    outline_user = f"""{research_block}{skill_block}## SOP 文章（唯一内容来源）
{sop_content}

{page_type_prompt}

{output_reqs}"""

    outline = None
    for attempt in range(2):
        try:
            response = _safe_run_async(llm_generate(provider_id, model,
                base_system, outline_user, temperature=temp_outline or temperature))
            response = _clean_json_response(response)
            data = json.loads(response)
            slides = data.get("slides", data) if isinstance(data, dict) else data
            if isinstance(slides, list) and len(slides) > 0:
                outline = slides
                break
        except Exception as e:
            _logger.warning(f"Stage 1 outline attempt {attempt+1} failed: {e}")

    if not outline:
        _logger.error("Stage 1 outline generation failed after retries")
        return None
    _logger.info(f"Stage 1 outline: {len(outline)} slides")

    # ── Phase 2: Fill content per slide in parallel batches ──
    BATCH_SIZE = 4
    fill_system = _load_fill_content_prompt()
    if not fill_system:
        fill_system = (
            "你是内容编辑专家。根据 SOP 文章和金字塔原理为指定幻灯片填充正文内容。"
            "遵循结论先行、以上统下、归类分组(MECE)、逻辑递进四大原则。"
            "从 SOP 中提取归纳对应部分的内容，不编造。输出纯 JSON，不要用 markdown 包裹。"
            "正文必须分段落：核心结论单独一段，支撑细节分段展开。用 \\n\\n 分隔段落，"
            "每段不超过 180 字。并列要点用编号列表（1. 2. 3. 换行分隔）。"
            "禁止把全部内容塞进一个不换行的长段落。"
        )

    from concurrent.futures import ThreadPoolExecutor, as_completed

    def _fill_batch(batch: list) -> list:
        """Fill one batch of slides. Returns filled slide dicts, or [] on failure."""
        batch_json = json.dumps(batch, ensure_ascii=False, indent=2)
        fill_user = f"""## SOP 文章（唯一内容来源）
{sop_content}

## 需要填充的幻灯片（只输出这些幻灯片的 body 和 notes，不要更改其他字段）
{batch_json}

## 输出格式
```json
{{"slides": [{{"seq":1,"heading":"原样保留","body":"从此 SOP 提取归纳的正文（结论先行，分段落，用 \\n\\n 分隔段落，每段不超过180字，并列要点用编号列表）","notes":"备注或反面后果(可选)","layout_hint":"原样保留","visual_weight":"原样保留","key_points":["原文保留"]}}, ...]}}
```
铁律：
- 只输出以上幻灯片，不增减
- heading, layout_hint, visual_weight, key_points 保持原样
- body 从此 SOP 对应部分提取归纳，每页至少 80 字，结论先行，必须分段落
- 仅输出 JSON"""

        for attempt in range(2):
            try:
                response = _safe_run_async(llm_generate(provider_id, model,
                    fill_system, fill_user, temperature=temp_fill or temperature))
                response = _clean_json_response(response)
                data = json.loads(response)
                filled = data.get("slides", data) if isinstance(data, dict) else data
                if isinstance(filled, list):
                    return filled
            except Exception as e:
                if attempt == 1:
                    _logger.warning(f"Stage 1 fill batch failed: {e}")
        return []

    # Create batches
    batches = [outline[i:i + BATCH_SIZE] for i in range(0, len(outline), BATCH_SIZE)]

    # Parallel fill all batches
    with ThreadPoolExecutor(max_workers=len(batches)) as ex:
        futures = {ex.submit(_fill_batch, b): b for b in batches}
        for future in as_completed(futures):
            try:
                filled = future.result()
                for f in filled:
                    seq = f.get("seq", 0)
                    idx = seq - 1
                    if 0 <= idx < len(outline):
                        outline[idx]["body"] = f.get("body", outline[idx].get("body", ""))
                        outline[idx]["notes"] = f.get("notes", outline[idx].get("notes", ""))
                        if f.get("layout_hint"):
                            outline[idx]["layout_hint"] = f["layout_hint"]
                        if f.get("visual_weight"):
                            outline[idx]["visual_weight"] = f["visual_weight"]
            except Exception as e:
                _logger.warning(f"Stage 1 fill batch failed: {e}")

    return outline


# ── Scenario prompt files ──
SCENARIOS_DIR = os.path.join(BASE_DIR, "resources", "scenarios")
SCENARIO_FILES = [
    "design-system.md",
    "outline-architect.md",
    "cognitive-design-principles.md",
    "reviewer.md",
    "svg-generator.md",
    "bento-grid-layout.md",
]


def _load_scenario_file(filename: str, column_id: str = "", canvas_w: int = 0, canvas_h: int = 0) -> str:
    """Load a scenario prompt file with fallback chain.

    Priority: scenarios/{column_id}/ → scenarios/_default/ → prompts/ (legacy)
    If canvas_w/canvas_h are provided, substitute {{canvas_w}} and {{canvas_h}} placeholders.
    """
    content = ""
    # 1) Per-column custom file
    if column_id:
        p = os.path.join(SCENARIOS_DIR, column_id, filename)
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                content = f.read()
    # 2) Default scenario template
    if not content:
        p = os.path.join(SCENARIOS_DIR, "_default", filename)
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                content = f.read()
    # 3) Legacy prompts directory
    if not content:
        p = os.path.join(BASE_DIR, "resources", "prompts", filename)
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                content = f.read()
    # Substitute canvas dimensions if provided
    if canvas_w and canvas_h:
        content = content.replace("{{canvas_w}}", str(canvas_w))
        content = content.replace("{{canvas_h}}", str(canvas_h))
    return content


def _load_svg_prompt_specs(column_id: str = "") -> tuple[str, str]:
    """Load svg-generator.md and bento-grid-layout.md from scenario files."""
    cw, ch = _get_canvas_dimensions(column_id)
    return (
        _load_scenario_file("svg-generator.md", column_id, cw, ch),
        _load_scenario_file("bento-grid-layout.md", column_id, cw, ch),
    )


def _load_outline_spec(column_id: str = "") -> str:
    """Load outline-architect.md from scenario files."""
    return _load_scenario_file("outline-architect.md", column_id)


def _get_canvas_dimensions(column_id: str = "") -> tuple[int, int]:
    """Get canvas dimensions from column_configs.rules JSON.

    Looks for {"canvas": {"width": W, "height": H}} in the column's rules.
    Falls back to 1280x720 (16:9 presentation default).
    """
    default_w, default_h = 1280, 720
    column_id = _validate_column_id(column_id)
    if not column_id:
        return (default_w, default_h)
    try:
        db = get_db()
        row = db.execute(
            "SELECT rules FROM column_configs WHERE column_id = ?",
            (column_id,)
        ).fetchone()
        db.close()
        if row and row["rules"]:
            rules = json.loads(row["rules"])
            canvas = rules.get("canvas", {})
            if isinstance(canvas, dict):
                w = canvas.get("width")
                h = canvas.get("height")
                if w and h:
                    return (int(w), int(h))
    except Exception:
        pass
    return (default_w, default_h)


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color like '#1a365d' or '#1A365D' to (r, g, b) tuple."""
    h = hex_color.lstrip("#")
    if len(h) == 6:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    if len(h) == 3:
        return (int(h[0] * 2, 16), int(h[1] * 2, 16), int(h[2] * 2, 16))
    return (0, 0, 0)


def _auto_fix_hardcoded_hex(html: str, scheme: dict, slide_seq: int) -> str:
    """Scan LLM-generated HTML for hardcoded hex values and replace with {{placeholder}} vars.

    THREE-PASS strategy:
    1. Exact match: hex in active scheme → its {{placeholder}}
    2. Semantic match: hex NOT in scheme → map by hue/saturation to the semantically
       correct scheme variable (red→semantic_negative, green→semantic_positive, etc.)
    3. Fallback: if no semantic category matches, use RGB Euclidean distance
    #ffffff is intentionally excluded — white text on dark backgrounds is correct
    per iron law #5 (封面对比度铁律).

    Returns the fixed HTML with ALL hardcoded hex replaced by placeholders.
    """
    if not scheme or not html:
        return html

    import re as _re_hex
    import math

    # ── Build hex → (placeholder_name, (r,g,b,h,s,l)) lookup for all scheme colors ──
    all_scheme_colors: list[tuple[str, str, tuple[int, int, int], tuple[float, float, float]]] = []

    def _rgb_to_hsl(r: int, g: int, b: int) -> tuple[float, float, float]:
        """Convert RGB (0-255) to HSL (h:0-360, s:0-1, l:0-1)."""
        rn, gn, bn = r / 255.0, g / 255.0, b / 255.0
        mx, mn = max(rn, gn, bn), min(rn, gn, bn)
        l = (mx + mn) / 2.0
        if mx == mn:
            return (0.0, 0.0, l)
        d = mx - mn
        s = d / (2.0 - mx - mn) if l > 0.5 else d / (mx + mn)
        if mx == rn:
            h = ((gn - bn) / d) % 6.0
        elif mx == gn:
            h = (bn - rn) / d + 2.0
        else:
            h = (rn - gn) / d + 4.0
        return (h * 60.0, s, l)

    def _add_scheme_entry(key: str, hex_val: str):
        if hex_val and hex_val.startswith("#") and len(hex_val) == 7:
            r, g, b = int(hex_val[1:3], 16), int(hex_val[3:5], 16), int(hex_val[5:7], 16)
            h, s, l = _rgb_to_hsl(r, g, b)
            all_scheme_colors.append((hex_val.lower(), key, (r, g, b), (h, s, l)))

    base_keys = ["primary", "secondary", "accent", "background", "text", "card_bg"]
    for key in base_keys:
        _add_scheme_entry(key, scheme.get(key, ""))

    chart_colors = scheme.get("chart_colors", [])
    if isinstance(chart_colors, list):
        for i, c in enumerate(chart_colors):
            _add_scheme_entry(f"chart_{i}", c)

    semantic = scheme.get("semantic", {})
    if isinstance(semantic, dict):
        for k, v in semantic.items():
            _add_scheme_entry(f"semantic_{k}", v)

    if not all_scheme_colors:
        return html

    # Build reverse lookup for exact match (exclude #ffffff)
    reverse: dict[str, str] = {}
    scheme_rgb: dict[str, tuple[int, int, int]] = {}
    scheme_hsl: dict[str, tuple[float, float, float]] = {}
    for hex_lower, var_name, rgb, hsl in all_scheme_colors:
        if hex_lower != "#ffffff":
            if hex_lower not in reverse:
                reverse[hex_lower] = var_name
            scheme_rgb[var_name] = rgb
            scheme_hsl[var_name] = hsl

    # ── Build semantic categories from scheme colors ──
    def _hue_category(h: float) -> str:
        """Classify hue angle (0-360) into semantic category."""
        if h < 20 or h >= 340:
            return "red"
        if h < 45:
            return "orange"
        if h < 65:
            return "yellow"
        if h < 170:
            return "green"
        if h < 270:
            return "blue"
        return "purple"

    # Map scheme vars to semantic roles
    semantic_red_vars: list[str] = []      # primary, semantic_negative, red chart colors
    semantic_warm_vars: list[str] = []     # accent, orange/yellow chart colors
    semantic_green_vars: list[str] = []    # semantic_positive, green chart colors
    semantic_cool_vars: list[str] = []     # blue chart colors
    semantic_purple_vars: list[str] = []   # purple chart colors, secondary
    semantic_neutral_dark: list[str] = []  # text
    semantic_neutral_light: list[str] = [] # background, card_bg

    for hex_lower, var_name, rgb, (h, s, l) in all_scheme_colors:
        if hex_lower == "#ffffff":
            continue
        cat = _hue_category(h)
        # Low saturation → neutral
        if s < 0.12:
            if l < 0.5:
                if var_name not in semantic_neutral_dark:
                    semantic_neutral_dark.append(var_name)
            else:
                if var_name not in semantic_neutral_light:
                    semantic_neutral_light.append(var_name)
            continue
        # Semantic color roles
        if var_name == "semantic_negative" or (var_name == "primary" and cat == "red"):
            if var_name not in semantic_red_vars:
                semantic_red_vars.append(var_name)
        elif var_name == "semantic_positive":
            if var_name not in semantic_green_vars:
                semantic_green_vars.append(var_name)
        # Categorize by hue
        if cat in ("red",):
            if var_name not in semantic_red_vars:
                semantic_red_vars.append(var_name)
        elif cat in ("orange", "yellow"):
            if var_name not in semantic_warm_vars:
                semantic_warm_vars.append(var_name)
        elif cat in ("green",):
            if var_name not in semantic_green_vars:
                semantic_green_vars.append(var_name)
        elif cat in ("blue",):
            if var_name not in semantic_cool_vars:
                semantic_cool_vars.append(var_name)
        elif cat in ("purple",):
            if var_name not in semantic_purple_vars:
                semantic_purple_vars.append(var_name)

    # Fallback: if a category is empty, use all vars
    all_var_names = [n for _, n, _, _ in all_scheme_colors if _ != "#ffffff"]

    def _best_in_category(hex_expanded: str, candidates: list[str]) -> str | None:
        """Pick the closest color by RGB distance within the given candidate list."""
        if not candidates:
            return None
        r = int(hex_expanded[1:3], 16)
        g = int(hex_expanded[3:5], 16)
        b = int(hex_expanded[5:7], 16)
        best_name = None
        best_dist = float('inf')
        for var_name in candidates:
            if var_name not in scheme_rgb:
                continue
            sr, sg, sb = scheme_rgb[var_name]
            dist = (r - sr) ** 2 + (g - sg) ** 2 + (b - sb) ** 2
            if dist < best_dist:
                best_dist = dist
                best_name = var_name
        return best_name

    # ── Find all unique hex values in HTML (excluding #ffffff/#fff) ──
    def _expand_hex(h: str) -> str:
        """Expand 3-digit hex to 6-digit: #abc → #aabbcc"""
        h = h.lower()
        if len(h) == 4 and h[0] == '#':  # #abc
            return '#' + h[1]*2 + h[2]*2 + h[3]*2
        return h

    all_hex_raw = set(h.lower() for h in _re_hex.findall(r'#[0-9a-fA-F]{3}(?:[0-9a-fA-F]{3})?', html))
    all_hex_expanded = {_expand_hex(h) for h in all_hex_raw}
    all_hex_expanded.discard("#ffffff")

    if not all_hex_expanded:
        return html

    # ── Build a map from original hex → expanded hex (for 3-digit lookup) ──
    hex_originals: dict[str, str] = {}  # expanded → first original form seen
    for raw_h in all_hex_raw:
        exp = _expand_hex(raw_h)
        if exp not in hex_originals:
            hex_originals[exp] = raw_h

    # ── For each hex, find its replacement ──
    replacement_map: dict[str, str] = {}
    exact_count = 0
    semantic_count = 0
    fallback_count = 0
    semantic_details: list[tuple[str, str, str]] = []  # (hex, matched_var, reason)

    for hex_expanded in all_hex_expanded:
        original_form = hex_originals.get(hex_expanded, hex_expanded)
        if hex_expanded in reverse:
            replacement_map[original_form] = reverse[hex_expanded]
            exact_count += 1
            continue

        # ── Semantic-aware matching ──
        r = int(hex_expanded[1:3], 16)
        g = int(hex_expanded[3:5], 16)
        b = int(hex_expanded[5:7], 16)
        h, s, l = _rgb_to_hsl(r, g, b)
        cat = _hue_category(h)

        best_name = None
        reason = ""

        # 1) Neutral (low saturation): use lightness to pick text vs background
        if s < 0.12:
            if l < 0.5:
                best_name = _best_in_category(hex_expanded, semantic_neutral_dark or ["text"])
                reason = f"neutral-dark (sat={s:.2f} l={l:.2f})"
            else:
                best_name = _best_in_category(hex_expanded, semantic_neutral_light or ["background"])
                reason = f"neutral-light (sat={s:.2f} l={l:.2f})"
        # 2) Warm red → semantic_negative, primary, secondary
        elif cat == "red":
            best_name = _best_in_category(hex_expanded, semantic_red_vars or all_var_names)
            reason = f"red-hue ({h:.0f}deg)"
        # 3) Orange/yellow → accent, warm chart colors
        elif cat in ("orange", "yellow"):
            best_name = _best_in_category(hex_expanded, semantic_warm_vars or all_var_names)
            reason = f"{cat}-hue ({h:.0f}deg)"
        # 4) Green → semantic_positive, green chart colors
        elif cat == "green":
            best_name = _best_in_category(hex_expanded, semantic_green_vars or all_var_names)
            reason = f"green-hue ({h:.0f}deg)"
        # 5) Blue → cool chart colors
        elif cat == "blue":
            best_name = _best_in_category(hex_expanded, semantic_cool_vars or all_var_names)
            reason = f"blue-hue ({h:.0f}deg)"
        # 6) Purple → purple chart colors
        elif cat == "purple":
            best_name = _best_in_category(hex_expanded, semantic_purple_vars or all_var_names)
            reason = f"purple-hue ({h:.0f}deg)"

        if best_name:
            replacement_map[original_form] = best_name
            semantic_count += 1
            if len(semantic_details) < 8:
                semantic_details.append((original_form, best_name, reason))
        else:
            # ═══ Pure RGB fallback (should rarely be reached) ═══
            best_name = _best_in_category(hex_expanded, all_var_names)
            if best_name:
                replacement_map[original_form] = best_name
                fallback_count += 1

    # ── Replace all hex values with {{placeholder}} vars ──
    for hex_original, var_name in replacement_map.items():
        hex_pattern = '(?i)' + hex_original
        placeholder = '{{' + var_name + '}}'
        html = _re_hex.sub(hex_pattern, placeholder, html)

    if exact_count > 0:
        _logger.info(
            f"[HEX-FIX] Slide {slide_seq}: exact-matched {exact_count} hex → {{{{placeholder}}}} vars"
        )
    if semantic_count > 0:
        _logger.info(
            f"[HEX-FIX] Slide {slide_seq}: semantic-matched {semantic_count} hex: {semantic_details}"
        )
    if fallback_count > 0:
        _logger.info(
            f"[HEX-FIX] Slide {slide_seq}: fallback RGB-matched {fallback_count} hex"
        )

    return html


def _fix_malformed_hex(html: str, slide_seq: int) -> str:
    """Fix LLM-garbled hex values like #fffffffff → #ffffff.

    The LLM sometimes corrupts hex colors by repeating characters beyond the
    standard 6. This normalizes # followed by 7+ identical hex chars → 6 chars.
    Valid 8-digit hex (RRGGBBAA) is NOT touched — only all-same-char sequences.
    Must run BEFORE _auto_fix_hardcoded_hex.
    """
    import re as _re_mh

    def _normalize(m):
        full = m.group(0)
        chars = full[1:]
        if len(chars) > 6 and len(set(chars)) == 1:
            fixed = full[:7]
            _logger.info(
                f"[MALFORMED-HEX] Slide {slide_seq}: {full} → {fixed} "
                f"(garbled repeat, normalized to 6 chars)"
            )
            return fixed
        return full

    return _re_mh.sub(r'#[0-9a-fA-F]{7,}', _normalize, html)


def _hex_luminance(hex_color: str) -> float | None:
    """Compute relative luminance (0-255) from a hex color string."""
    if not hex_color or not hex_color.startswith("#") or len(hex_color) != 7:
        return None
    try:
        r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
        return 0.2126 * r + 0.7152 * g + 0.0722 * b
    except ValueError:
        return None


def _resolve_placeholder_value(value: str, scheme: dict) -> str:
    """Resolve a single {{placeholder}} against scheme. Returns hex if resolved, else original."""
    m = re.match(r'\{\{(\w+)\}\}', str(value).strip())
    if m:
        key = m.group(1)
        resolved = scheme.get(key, "")
        if resolved and resolved.startswith("#"):
            return resolved
    return str(value)


def _get_effective_page_bg_luminance(style_id: str, page_type: str, scheme: dict) -> float | None:
    """Compute the effective background luminance for a specific page type.

    Checks tokens.yaml slide_type_overrides[page_type] for card_bg or background,
    resolves {{placeholder}} references against scheme, and computes relative luminance.
    Falls back to scheme.background if no override is found.

    This is the single source of truth for "is this page dark or light?" —
    no hardcoded lists, no assumptions about which styles have dark pages.
    """
    import yaml

    # Prevent path traversal: style_id must be a simple directory name
    if not style_id or not re.match(r'^[a-zA-Z0-9_-]+$', style_id):
        return _hex_luminance(scheme.get("background", ""))

    vi_base = os.path.join(BASE_DIR, "resources", "vi")
    tokens_path = os.path.join(vi_base, style_id, "tokens.yaml")
    # Verify resolved path stays within the expected base
    real_base = os.path.realpath(vi_base)
    if os.path.exists(tokens_path):
        if not os.path.realpath(tokens_path).startswith(real_base + os.sep):
            return _hex_luminance(scheme.get("background", ""))
        try:
            with open(tokens_path, "r", encoding="utf-8") as f:
                tokens = yaml.safe_load(f.read())
            overrides = tokens.get("slide_type_overrides", {})
            page_override = overrides.get(page_type, {})
            bg_ref = page_override.get("card_bg") or page_override.get("background")
            if bg_ref:
                bg_hex = _resolve_placeholder_value(str(bg_ref), scheme)
                lum = _hex_luminance(bg_hex)
                if lum is not None:
                    return lum
        except Exception:
            pass

    # Fallback: global scheme background
    return _hex_luminance(scheme.get("background", ""))


def _is_light_background(scheme: dict) -> bool:
    """Return True if the scheme's global background color is light (luminance > 128)."""
    lum = _hex_luminance(scheme.get("background", ""))
    return lum is not None and lum > 128


def _auto_fix_white_on_light(html: str, scheme: dict, slide_seq: int,
                             style_id: str = None, page_type: str = None) -> str:
    """Fix LLM hardcoded white text on light backgrounds.

    Uses page-type-specific background from tokens.yaml slide_type_overrides
    when available (via _get_effective_page_bg_luminance), falling back to the
    global scheme background. This ensures pages with dark overrides in tokens.yaml
    (e.g. business cover/section/summary/quote, vintage section/summary) correctly
    keep their white text — no hardcoded exclusion list needed.

    Must run AFTER _auto_fix_hardcoded_hex (which intentionally skips #ffffff)
    and BEFORE _resolve_color_vars (which converts {{placeholder}} → hex).
    """
    if not scheme or not html:
        return html

    # Determine effective background luminance for THIS page type
    if style_id and page_type:
        bg_luminance = _get_effective_page_bg_luminance(style_id, page_type, scheme)
    else:
        bg_luminance = _hex_luminance(scheme.get("background", ""))

    # If the effective background is dark, white text is correct — skip
    if bg_luminance is None or bg_luminance <= 128:
        return html

    text_hex = scheme.get("text", "")
    if not text_hex or not text_hex.startswith("#"):
        return html

    import re as _re_wol

    white_patterns = ['#ffffff', '#FFFFFF', '#fff', '#FFF']
    white_hex_re = r'(?i)#(?:fff|ffffff)\b'

    # Count for logging
    fix_count = 0

    # ── Layer 1: Solid white text → {{text}} placeholder ──
    for white in white_patterns:
        # Match white as a color value in CSS
        patterns = [
            (f'color:{white}', f'color:{{{{text}}}}'),
            (f'color: {white}', f'color: {{{{text}}}}'),
        ]
        for pat, repl in patterns:
            count = html.count(pat)
            if count > 0:
                html = html.replace(pat, repl)
                fix_count += count

    # ── Layer 2: rgba(255,255,255,N) → rgba(var(--text-rgb),N) ──
    rgba_patterns = [
        (r'rgba\(\s*255\s*,\s*255\s*,\s*255\s*,', 'rgba(var(--text-rgb),'),
        (r'rgba\(255,255,255,', 'rgba(var(--text-rgb),'),
    ]
    for pat, repl in rgba_patterns:
        matches = list(_re_wol.finditer(pat, html))
        if matches:
            html = _re_wol.sub(pat, repl, html)
            fix_count += len(matches)

    # ── Layer 3: white borders on light backgrounds → use accent or text ──
    # Only fix border-color (not border shorthand which is complex)
    accent_hex = scheme.get("accent", "")
    if accent_hex and accent_hex.startswith("#"):
        for white in white_patterns:
            border_pat = f'border-color:{white}'
            if border_pat in html:
                html = html.replace(border_pat, f'border-color:{accent_hex}')
                fix_count += 1

    if fix_count > 0:
        _logger.info(
            f"[WHITE-FIX] Slide {slide_seq}: fixed {fix_count} white-on-light "
            f"occurrences → {{{{text}}}} (bg luminance > 128)"
        )

    return html


# Core theme variables that must NEVER be overridden on individual slides.
# The LLM sometimes invents local reassignments like --primary: var(--card_bg)
# which inverts the color theme and causes invisible text.
_PROTECTED_CSS_VARS = [
    "primary", "primary_rgb", "primary_r", "primary_g", "primary_b",
    "secondary", "secondary_rgb", "secondary_r", "secondary_g", "secondary_b",
    "text", "text_rgb", "text_r", "text_g", "text_b",
    "background", "background_rgb", "background_r", "background_g", "background_b",
    "card_bg", "card_bg_rgb", "card_bg_r", "card_bg_g", "card_bg_b",
    "accent", "accent_rgb", "accent_r", "accent_g", "accent_b",
]


def _strip_local_var_overrides(html: str, slide_seq: int) -> str:
    """Remove LLM-invented CSS variable reassignments from individual slides.

    The LLM sometimes adds local style overrides like:
      style="... --primary: var(--card_bg); --text: var(--card_bg); --card_bg: var(--chart-1); ..."

    This inverts the color theme — on light-background styles, it turns text
    light-on-light (invisible). These core theme variables must only be defined
    once at :root level and inherited, never reassigned per-slide.
    """
    if not html:
        return html

    import re as _re_strip

    fix_count = 0

    for var_name in _PROTECTED_CSS_VARS:
        # Match: --varname: <anything up to ; or >
        # Pattern: --varname followed by optional whitespace, colon, then any non-empty value until ; or end of style attr
        pattern = rf'--{var_name}\s*:\s*[^;"]+(?:;\s*)?'
        matches = list(_re_strip.finditer(pattern, html))
        if matches:
            html = _re_strip.sub(pattern, '', html)
            fix_count += len(matches)

    if fix_count > 0:
        _logger.info(
            f"[STRIP-VARS] Slide {slide_seq}: removed {fix_count} local CSS "
            f"variable overrides (core theme vars must be inherited from :root)"
        )

    return html


def _auto_fix_font_size(html: str, slide_seq: int, is_a4: bool = False) -> str:
    """Enforce VI typography minimums per spec.

    Valid sizes: 14px (captions), 16-18px (body), 22-24px (card titles),
    36-44px (page titles), 58-65px (cover titles).
    10-13px: below minimum → bump to 14px (PPT) or keep 10px/12px (A4).
    15px: dead zone → bump to 16px.
    """
    import re as _re_fs
    total_fixes = 0
    if is_a4:
        # A4: 10px = header/footer (keep), 11px → 12px, 12px = body (keep), 13px → 14px
        for sz in ['11', '13']:
            before = len(_re_fs.findall(rf'font-size:\s*{sz}px', html))
            html = html.replace(f'font-size:{sz}px', f'font-size:{"12" if sz == "11" else "14"}px')
            html = html.replace(f'font-size: {sz}px', f'font-size: {"12" if sz == "11" else "14"}px')
            total_fixes += before
    else:
        # PPT: 10-13px → 14px (below absolute minimum)
        for sz in ['10', '11', '12', '13']:
            before = len(_re_fs.findall(rf'font-size:\s*{sz}px', html))
            html = html.replace(f'font-size:{sz}px', 'font-size:14px')
            html = html.replace(f'font-size: {sz}px', 'font-size: 14px')
            total_fixes += before
    # 15px → 16px (invalid gap between caption and body)
    before15 = len(_re_fs.findall(r'font-size:\s*15px', html))
    html = html.replace('font-size:15px', 'font-size:16px')
    html = html.replace('font-size: 15px', 'font-size: 16px')
    total_fixes += before15
    if total_fixes > 0:
        _logger.info(f"[FONT-FIX] Slide {slide_seq}: {total_fixes} instances fixed")
    return html


def _resolve_color_vars(text: str, scheme: dict, css_vars: bool = False) -> str:
    """Replace {{primary}}, {{accent_rgb}}, {{chart_0}}, etc. with values.

    When css_vars=False (default): resolves to hex/rgb values for LLM prompts.
    When css_vars=True: resolves to CSS var() references for final HTML output.

    Supported variable forms for each color key:
      {{key}}          → hex (#1a365d) or var(--key)
      {{key_r}}        → red int (26) or var(--key-r)
      {{key_g}}        → green int (54) or var(--key-g)
      {{key_b}}        → blue int (93) or var(--key-b)
      {{key_rgb}}      → rgb string (26, 54, 93) or var(--key-rgb)
    """
    import re as _re

    vars_map: dict[str, str] = {}

    # Base color keys
    color_keys = ["primary", "secondary", "accent", "background", "text", "card_bg"]
    for key in color_keys:
        val = scheme.get(key, "")
        if val and val.startswith("#"):
            if css_vars:
                vars_map[key] = f"var(--{key})"
                vars_map[f"{key}_r"] = f"var(--{key}-r)"
                vars_map[f"{key}_g"] = f"var(--{key}-g)"
                vars_map[f"{key}_b"] = f"var(--{key}-b)"
                vars_map[f"{key}_rgb"] = f"var(--{key}-rgb)"
            else:
                vars_map[key] = val
                r, g, b = _hex_to_rgb(val)
                vars_map[f"{key}_r"] = str(r)
                vars_map[f"{key}_g"] = str(g)
                vars_map[f"{key}_b"] = str(b)
                vars_map[f"{key}_rgb"] = f"{r}, {g}, {b}"

    # Chart colors
    chart_colors = scheme.get("chart_colors", [])
    if isinstance(chart_colors, list):
        for i, c in enumerate(chart_colors):
            if c and c.startswith("#"):
                if css_vars:
                    vars_map[f"chart_{i}"] = f"var(--chart-{i})"
                    vars_map[f"chart_{i}_rgb"] = f"var(--chart-{i}-rgb)"
                else:
                    vars_map[f"chart_{i}"] = c
                    r, g, b = _hex_to_rgb(c)
                    vars_map[f"chart_{i}_rgb"] = f"{r}, {g}, {b}"

    # Semantic colors
    semantic = scheme.get("semantic", {})
    if isinstance(semantic, dict):
        for k, v in semantic.items():
            if v and v.startswith("#"):
                if css_vars:
                    vars_map[f"semantic_{k}"] = f"var(--semantic-{k})"
                    vars_map[f"semantic_{k}_rgb"] = f"var(--semantic-{k}-rgb)"
                else:
                    vars_map[f"semantic_{k}"] = v
                    r, g, b = _hex_to_rgb(v)
                    vars_map[f"semantic_{k}_rgb"] = f"{r}, {g}, {b}"

    # Replace all {{var}} placeholders
    def _replacer(m):
        var_name = m.group(1)
        return vars_map.get(var_name, m.group(0))

    return _re.sub(r"\{\{(\w+)\}\}", _replacer, text)


def _load_style_yaml_text(style_id: str, color_scheme: str = "deep-blue", resolve_vars: bool = True) -> str:
    """Load YAML structured data for a style, with optional color scheme override.

    Priority: 1) {style_id}/tokens.yaml (directory structure),
    2) standalone {style_id}.yaml file in data/styles/.

    When color_scheme is specified and tokens.yaml has a color_schemes block,
    the selected scheme's values are injected as the flat color_scheme block.

    When resolve_vars=False, {{primary}} placeholders are kept as-is (for LLM prompts)
    and the color_scheme hex block is stripped from the YAML output.
    """
    import yaml
    styles_dir = os.path.join(BASE_DIR, "data", "styles")
    vi_dir = os.path.join(BASE_DIR, "resources", "vi")

    # 1) Directory structure: data/vi/{style_id}/tokens.yaml
    tokens_path = os.path.join(vi_dir, style_id, "tokens.yaml")
    if os.path.exists(tokens_path):
        with open(tokens_path, "r", encoding="utf-8") as f:
            raw = f.read()
        try:
            parsed = yaml.safe_load(raw)
            if parsed and isinstance(parsed, dict) and "color_schemes" in parsed:
                schemes = parsed.pop("color_schemes", {})
                parsed.pop("color_scheme", None)  # remove legacy flat key
                active_scheme: dict = {}
                if color_scheme in schemes:
                    active_scheme = schemes[color_scheme]
                elif schemes:
                    first = next(iter(schemes.keys()))
                    active_scheme = schemes[first]
                import io as _io
                out = _io.StringIO()
                if resolve_vars:
                    parsed["color_scheme"] = active_scheme
                yaml.dump(parsed, out, allow_unicode=True, default_flow_style=False, sort_keys=False)
                yaml_text = out.getvalue()
                # Resolve color variables in non-scheme sections (gradients, decoration, etc.)
                if resolve_vars and active_scheme:
                    yaml_text = _resolve_color_vars(yaml_text, active_scheme)
                return yaml_text
        except Exception:
            pass
        return raw

    # 2) Fall back to standalone YAML file
    p = os.path.join(styles_dir, f"{style_id}.yaml")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def _load_style_vi(style_id: str, color_scheme: str = "deep-blue", resolve_vars: bool = True) -> str:
    """Load the Visual Identity System document from data/vi/{style_id}/vi.md.

    Resolves {{color}} variables against the active color_scheme.
    When resolve_vars=False, keeps {{primary}} etc. as placeholders (for LLM prompts).
    """
    vi_md_path = os.path.join(BASE_DIR, "resources", "vi", style_id, "vi.md")
    if os.path.exists(vi_md_path):
        with open(vi_md_path, "r", encoding="utf-8") as f:
            result = f.read()
        if resolve_vars:
            scheme = _load_scheme_data(style_id, color_scheme)
            if scheme:
                result = _resolve_color_vars(result, scheme)
        return result
    return ""


def _load_style_prompt(style_id: str, color_scheme: str = "deep-blue") -> str:
    """Load the style-specific LLM persona + task prompt.

    Injects the color scheme persona_hint from tokens.yaml if available,
    replacing {{PERSONA_HINT}} in the prompt template.
    """
    import yaml
    prompt_path = os.path.join(BASE_DIR, "resources", "vi", style_id, "prompt.md")
    if not os.path.exists(prompt_path):
        return ""

    with open(prompt_path, "r", encoding="utf-8") as f:
        prompt = f.read()

    # Inject persona_hint from tokens.yaml if available
    tokens_path = os.path.join(BASE_DIR, "resources", "vi", style_id, "tokens.yaml")
    if os.path.exists(tokens_path):
        try:
            with open(tokens_path, "r", encoding="utf-8") as f:
                tokens = yaml.safe_load(f.read())
            schemes = tokens.get("color_schemes", {}) if isinstance(tokens, dict) else {}
            scheme = schemes.get(color_scheme, {})
            if isinstance(scheme, dict):
                hint = scheme.get("persona_hint", "")
                if hint:
                    prompt = prompt.replace("{{PERSONA_HINT}}", hint)
        except Exception:
            pass

    return prompt


def _list_color_schemes(style_id: str) -> list:
    """List available color schemes for a style from tokens.yaml.

    Returns list of {id, label, primary, accent, background, text, card_bg}
    for use in UI pickers.
    """
    import yaml
    tokens_path = os.path.join(BASE_DIR, "resources", "vi", style_id, "tokens.yaml")
    if not os.path.exists(tokens_path):
        return []
    try:
        with open(tokens_path, "r", encoding="utf-8") as f:
            tokens = yaml.safe_load(f.read())
        if not isinstance(tokens, dict):
            return []
        schemes = tokens.get("color_schemes", {})
        result = []
        for sid, scheme in schemes.items():
            if isinstance(scheme, dict):
                result.append({
                    "id": sid,
                    "label": scheme.get("label", sid),
                    "primary": scheme.get("primary", ""),
                    "accent": scheme.get("accent", ""),
                    "background": scheme.get("background", ""),
                    "text": scheme.get("text", ""),
                    "card_bg": scheme.get("card_bg", ""),
                })
        return result
    except Exception:
        return []


def _recolor_slide_html(style_id: str, slide_html: str, new_color_scheme: str, source_scheme: str) -> str:
    """Direct color replacement: map old scheme colors → new scheme colors.

    1. Use source_scheme as the authoritative current scheme (caller reads from result.json)
    2. Build a role-for-role mapping from old scheme to new scheme
    3. Apply all replacements in a single pass
    No AI involved — pure deterministic string replacement.
    """
    import yaml, re
    tokens_path = os.path.join(BASE_DIR, "resources", "vi", style_id, "tokens.yaml")
    if not os.path.exists(tokens_path):
        return slide_html

    with open(tokens_path, "r", encoding="utf-8") as f:
        tokens = yaml.safe_load(f.read())
    if not isinstance(tokens, dict):
        return slide_html

    schemes = tokens.get("color_schemes", {})
    if not schemes or new_color_scheme not in schemes:
        return slide_html

    target = schemes[new_color_scheme]
    base_keys = ["primary", "secondary", "accent", "background", "text", "card_bg"]

    # source_scheme is authoritative — caller reads it from result.json
    if source_scheme not in schemes:
        return slide_html

    best_scheme = source_scheme
    old = schemes[best_scheme]
    if best_scheme == new_color_scheme:
        return slide_html  # Already using this scheme

    # Step 2: Build old_scheme → new_scheme replacement map
    # IMPORTANT: Use insertion order to resolve conflicts — "first role wins".
    # When a scheme uses the same hex for multiple roles (e.g. dark-gold
    # primary=text=#1a1a2e), the first role in base_keys order keeps the mapping,
    # preventing later roles (text) from overwriting earlier ones (primary).
    replacements: dict[str, str] = {}

    def _add(old_val: str, new_val: str):
        if old_val and new_val and len(old_val) >= 4 and len(new_val) >= 4:
            old_lower = old_val.lower()
            if old_lower == new_val.lower():
                return
            if old_lower not in replacements:
                replacements[old_lower] = new_val.lower()

    # Base keys (by role priority: primary > secondary > accent > background > text > card_bg)
    for k in base_keys:
        _add(old.get(k, ""), target.get(k, ""))

    # Chart colors (positional)
    old_charts = old.get("chart_colors") or []
    new_charts = target.get("chart_colors") or []
    for i in range(min(len(old_charts), len(new_charts))):
        _add(old_charts[i], new_charts[i])

    # Semantic colors (by key)
    old_sem = old.get("semantic") or {}
    new_sem = target.get("semantic") or {}
    for k in old_sem:
        if k in new_sem:
            _add(old_sem[k], new_sem[k])

    if not replacements:
        return slide_html

    # Step 3: Apply replacements (hex + rgba + rgb)
    result = slide_html
    for old_c, new_c in replacements.items():
        # Hex literal replacement
        result = re.sub(re.escape(old_c), new_c, result, flags=re.IGNORECASE)
        # rgba / rgb variants: convert hex→rgb components and build regex patterns
        if old_c.startswith("#") and new_c.startswith("#"):
            r_old, g_old, b_old = _hex_to_rgb(old_c)
            r_new, g_new, b_new = _hex_to_rgb(new_c)
            # rgba(R, G, B, opacity) → preserve opacity
            result = re.sub(
                rf"rgba\(\s*{r_old}\s*,\s*{g_old}\s*,\s*{b_old}\s*,\s*([\d.]+)\s*\)",
                rf"rgba({r_new}, {g_new}, {b_new}, \1)",
                result,
                flags=re.IGNORECASE,
            )
            # rgb(R, G, B) without alpha
            result = re.sub(
                rf"rgb\(\s*{r_old}\s*,\s*{g_old}\s*,\s*{b_old}\s*\)",
                rf"rgb({r_new}, {g_new}, {b_new})",
                result,
                flags=re.IGNORECASE,
            )

    return result


def _load_scheme_data(style_id: str, color_scheme: str = "deep-blue") -> dict:
    """Load a single color scheme dict from tokens.yaml by name. Returns empty dict if not found."""
    import yaml
    tokens_path = os.path.join(BASE_DIR, "resources", "vi", style_id, "tokens.yaml")
    if not os.path.exists(tokens_path):
        return {}
    try:
        with open(tokens_path, "r", encoding="utf-8") as f:
            tokens = yaml.safe_load(f.read())
        if not isinstance(tokens, dict):
            return {}
        schemes = tokens.get("color_schemes", {})
        if color_scheme in schemes:
            return schemes[color_scheme]
        if schemes:
            return next(iter(schemes.values()))
    except Exception:
        pass
    return {}


def _load_style_vi_section(style_id: str, section: str, color_scheme: str = "deep-blue", resolve_vars: bool = True, column_id: str = "") -> str:
    """Load a specific VI sub-file from data/vi/{style_id}/.

    section: 'cover', 'content', 'data', or 'summary'.
    Loads both vi.md (general) and the section-specific file, concatenated.
    Resolves {{color}} variables against the active color_scheme.
    When resolve_vars=False, keeps {{primary}} etc. as placeholders (for LLM prompts).
    Checks vi/{style_id}/{column_id}/ first for column-specific VI overrides,
    falls back to vi/{style_id}/.
    """
    column_id = _validate_column_id(column_id)
    vi_dir = os.path.join(BASE_DIR, "resources", "vi", style_id)
    parts = []

    # Section-specific file FIRST — guaranteed included before truncation
    if section != "toc":
        col_section_file = None
        if column_id:
            col_section_file = os.path.join(vi_dir, column_id, f"{section}.md")
        section_file = os.path.join(vi_dir, f"{section}.md")
        # Also check blocks/ and templates/ subdirectories (document blocks/templates)
        blocks_file = os.path.join(vi_dir, "blocks", f"{section}.md")
        templates_file = os.path.join(vi_dir, "templates", f"{section}.md")
        # Priority depends on column type:
        # - A4/portrait (col3): column > blocks/ > templates/ ONLY
        #   (NEVER load PPT page type files like content.md/data.md for A4)
        # - PPT/landscape (col4/col5): column > style top-level > blocks/ > templates/
        cw_prio, ch_prio = _get_canvas_dimensions(column_id) if column_id else (1280, 720)
        is_a4_prio = ch_prio > cw_prio
        if is_a4_prio:
            candidates = (col_section_file, blocks_file, templates_file)
        else:
            candidates = (col_section_file, section_file, blocks_file, templates_file)
        chosen = None
        for candidate in candidates:
            if candidate and os.path.exists(candidate):
                chosen = candidate
                break
        if chosen:
            with open(chosen, "r", encoding="utf-8") as f:
                parts.append(f.read())

    # General vi.md second — fills remaining space after section rules
    vi_md_path = os.path.join(vi_dir, "vi.md")
    if column_id:
        col_vi = os.path.join(vi_dir, column_id, "vi.md")
        if os.path.exists(col_vi):
            vi_md_path = col_vi
    if os.path.exists(vi_md_path):
        with open(vi_md_path, "r", encoding="utf-8") as f:
            parts.append(f.read())

    if not parts:
        return _load_style_vi(style_id, color_scheme, resolve_vars=resolve_vars)

    result = "\n\n---\n\n".join(parts)

    if resolve_vars:
        scheme = _load_scheme_data(style_id, color_scheme)
        if scheme:
            result = _resolve_color_vars(result, scheme)

    return result
# Page type order — lazily built from index.md (single source of truth)
_page_type_order_cache: dict[str, list[str]] = {}

def _get_page_type_order(style_id: str = "business") -> list[str]:
    """Get page type names in canonical order from index.md. Cached per style."""
    if style_id not in _page_type_order_cache:
        types = _scan_vi_page_types(style_id)
        _page_type_order_cache[style_id] = [t["type"] for t in types]
    return _page_type_order_cache[style_id]

def _page_type_sort_key(ptype: str) -> int:
    """Sort key for VI file listing — mirrors frontend sectionSortKey().

    Order: 总纲 < 设计原则 < 设计元素 < 页面类型 < 文档构建块 < 文档模板 < 列专属覆写
    """
    # 总纲
    if ptype in ('vi', 'prompt', 'tokens', 'index'):
        return {'vi': -4, 'prompt': -3, 'tokens': -2, 'index': -1}[ptype]
    # 设计原则
    if ptype in ('principles', 'consistency', 'richness', 'checklist', 'images', 'data_rules', 'decorations'):
        return ['principles', 'consistency', 'richness', 'checklist', 'images', 'data_rules', 'decorations'].index(ptype)
    # 设计元素
    if ptype in ('colors', 'typography', 'card_styles', 'charts', 'layouts', 'card_roles', 'chart_decision', 'icons'):
        return ['colors', 'typography', 'card_styles', 'charts', 'layouts', 'card_roles', 'chart_decision', 'icons'].index(ptype) + 10
    # 页面类型
    try:
        return _get_page_type_order().index(ptype) + 100
    except ValueError:
        pass
    # 文档构建块
    if ptype.startswith('blocks/'):
        return 200
    # 文档模板
    if ptype.startswith('templates/'):
        return 300
    # 列专属覆写
    if ptype.startswith('col3/') or ptype.startswith('col4/') or ptype.startswith('col5/'):
        return 400
    return 999


def _scan_vi_page_types(style_id: str) -> list[dict]:
    """Read VI index.md to get the canonical page type list.

    Parses the "页面类型" markdown table from index.md.
    Only returns page types (excludes design principles and elements).
    """
    vi_dir = os.path.join(BASE_DIR, "resources", "vi")
    index_path = os.path.join(vi_dir, style_id, "index.md")

    if not os.path.exists(index_path):
        return _fallback_page_types()

    try:
        with open(index_path, "r", encoding="utf-8") as f:
            content = f.read()
    except Exception:
        return _fallback_page_types()

    # Parse the "页面类型" markdown table
    types = []
    in_page_type_section = False
    has_number_col = False  # True if table has 编号 column (new format)
    for line in content.split("\n"):
        line = line.strip()
        if line.startswith("## 页面类型"):
            in_page_type_section = True
            continue
        if in_page_type_section and line.startswith("## "):
            break  # next section starts
        if not in_page_type_section or not line.startswith("|") or "---" in line:
            continue
        if "类型名" in line:
            has_number_col = "编号" in line
            continue
        # Parse: | P01 | cover | 封面 | 开场页... | cover.md |  (new)
        #    or: | cover | 封面 | 开场页... | cover.md |         (old)
        cells = [c.strip() for c in line.split("|")[1:-1]]
        if has_number_col and len(cells) >= 4:
            types.append({
                "type": cells[1],   # skip 编号 column
                "label": cells[2],
                "purpose": cells[3],
            })
        elif not has_number_col and len(cells) >= 3:
            types.append({
                "type": cells[0],
                "label": cells[1],
                "purpose": cells[2],
            })

    if not types:
        return _fallback_page_types()
    return types


def _fallback_page_types() -> list[dict]:
    """Minimal page type set when VI directory doesn't exist."""
    return [
        {"type": "cover", "label": "封面", "purpose": "开场页，全屏视觉冲击"},
        {"type": "toc", "label": "目录", "purpose": "内容导航与章节概览"},
        {"type": "content", "label": "内容页", "purpose": "通用内容展示"},
        {"type": "summary", "label": "总结页", "purpose": "结尾收束与核心要点回顾"},
    ]


def _scan_vi_document_blocks(style_id: str) -> list[dict]:
    """Read VI index.md to get document building blocks + templates.

    Parses "文档构建块" (B01-B08) and "文档模板" (T01) tables from index.md.
    Used for A4 document columns (col3) instead of PPT page types.
    """
    vi_dir = os.path.join(BASE_DIR, "resources", "vi")
    index_path = os.path.join(vi_dir, style_id, "index.md")

    if not os.path.exists(index_path):
        return _fallback_document_blocks()

    try:
        with open(index_path, "r", encoding="utf-8") as f:
            content = f.read()
    except Exception:
        return _fallback_document_blocks()

    blocks = []
    current_section = None

    for line in content.split("\n"):
        line = line.strip()
        if line.startswith("## 文档构建块"):
            current_section = "blocks"
            continue
        if line.startswith("## 文档模板"):
            current_section = "templates"
            continue
        if current_section and line.startswith("## "):
            break  # next major section
        if not current_section or not line.startswith("|") or "---" in line:
            continue
        cells = [c.strip() for c in line.split("|")[1:-1]]
        if len(cells) >= 2 and cells[1] in ("块名", "模板名"):
            continue  # header row

        if len(cells) >= 4:
            blocks.append({
                "type": cells[1],   # block name (header, title, etc.)
                "label": cells[2],  # 中文名
                "purpose": cells[3],  # 用途
                "section": current_section,  # "blocks" or "templates"
            })

    if not blocks:
        return _fallback_document_blocks()
    return blocks


def _fallback_document_blocks() -> list[dict]:
    """Minimal document blocks when VI directory doesn't exist."""
    return [
        {"type": "header", "label": "页头", "purpose": "文档页头区域"},
        {"type": "title", "label": "标题", "purpose": "主标题行"},
        {"type": "info_block", "label": "基本信息块", "purpose": "标签-值对+图片占位"},
        {"type": "table_block", "label": "表格块", "purpose": "多列数据网格"},
        {"type": "text_block", "label": "文字块", "purpose": "自由段落文字"},
        {"type": "list_block", "label": "列表块", "purpose": "编号或项目符号列表"},
        {"type": "closing", "label": "结尾块", "purpose": "文档结尾区"},
        {"type": "footer", "label": "页脚", "purpose": "文档页脚"},
    ]


def _build_page_type_prompt(style_id: str, column_id: str = "") -> str:
    """Generate the page_type selection prompt block from VI index.md.

    For A4/portrait columns (col3): returns document building blocks (B01-B08)
    and document templates (T01). PPT page types are explicitly forbidden.
    For PPT/landscape columns (col4, col5): returns PPT page types (P01-P26).
    """
    cw, ch = _get_canvas_dimensions(column_id) if column_id else (1280, 720)
    is_a4 = ch > cw

    if is_a4:
        blocks = _scan_vi_document_blocks(style_id)
        lines = [
            "## 可用文档构建块（来源：VI 索引 index.md）",
            f"当前模板: {style_id}，A4 文档模式，共 {len(blocks)} 种构建块：",
            "",
        ]
        for b in blocks:
            section_label = "模板" if b.get("section") == "templates" else "构建块"
            lines.append(f"- **{b['type']}**（{b['label']}）[{section_label}]：{b['purpose']}")
        lines.extend([
            "",
            "## 设计规范来源",
            "所有设计规范均从 VI 索引文件（index.md）获取：",
            "- **文档构建块详细规范**：读取对应的 blocks/ 目录下的 .md 文件（如 blocks/header.md、blocks/table_block.md 等）",
            "- **文档模板**：读取 templates/ 目录下的 .md 文件",
            "- **设计原则**（7项）：principles.md / consistency.md / richness.md / checklist.md / images.md / data_rules.md / decorations.md",
            "- **设计参数**（8项）：colors.md / typography.md / card_styles.md / charts.md / layouts.md / card_roles.md / chart_decision.md / icons.md",
            "",
            "## 规则",
            f"1. 只使用上述 {len(blocks)} 种构建块类型，不要编造新类型。",
            "2. 为每页选择最匹配内容特征的构建块类型。",
            "3. 需要某类型的详细规范时，读取该类型对应的 .md 文件。",
            "4. 需要设计参数（颜色/字号/圆角等）时，读取设计元素对应的 .md 文件。",
            "5. ⛔ 禁止使用 PPT 页面类型（cover/toc/section/content/data/summary/closing 等），A4 文档只能使用上述文档构建块。",
        ])
    else:
        types = _scan_vi_page_types(style_id)
        # Exclude P27 "document" (A4文档) for PPT columns
        types = [t for t in types if t.get("type") != "document"]
        lines = [
            "## 可用页面类型（来源：VI 索引 index.md）",
            f"当前模板: {style_id}，共 {len(types)} 种页面类型：",
            "",
        ]
        for t in types:
            lines.append(f"- **{t['type']}**（{t['label']}）：{t['purpose']}")
        lines.extend([
            "",
            "## 设计规范来源",
            "所有设计规范均从 VI 索引文件（index.md）获取：",
            "- **页面类型详细规范**：读取对应的页面类型 .md 文件（如 cover.md、content.md 等）",
            "- **设计原则**（7项）：principles.md / consistency.md / richness.md / checklist.md / images.md / data_rules.md / decorations.md",
            "- **设计参数**（8项）：colors.md / typography.md / card_styles.md / charts.md / layouts.md / card_roles.md / chart_decision.md / icons.md",
            "",
            "## 规则",
            f"1. 只使用上述 {len(types)} 种类型，不要编造新类型。",
            "2. 为每页选择最匹配内容特征的 page_type。",
            "3. 需要某类型的详细布局规范时，读取该类型对应的 .md 文件。",
            "4. 需要设计参数（颜色/字号/圆角等）时，读取设计元素对应的 .md 文件。",
        ])
    return "\n".join(lines)

def _load_cognitive_spec(column_id: str = "") -> str:
    """Load cognitive-design-principles.md from scenario files."""
    return _load_scenario_file("cognitive-design-principles.md", column_id)


def _load_reviewer_spec(column_id: str = "") -> str:
    """Load reviewer.md from scenario files."""
    cw, ch = _get_canvas_dimensions(column_id)
    return _load_scenario_file("reviewer.md", column_id, cw, ch)


def _load_design_system(column_id: str = "") -> str:
    """Load design-system.md — the HTML slide design guide for LLM."""
    cw, ch = _get_canvas_dimensions(column_id)
    return _load_scenario_file("design-system.md", column_id, cw, ch)


def _resolve_font_range(style_yaml_text: str) -> dict:
    """Extract font size thresholds from style YAML for the LLM prompt."""
    import yaml
    try:
        data = yaml.safe_load(style_yaml_text)
        typography = data.get("typography", {})
        return {
            "h1": typography.get("cover_title", {}).get("size", "58-65px"),
            "h2": typography.get("page_title", {}).get("size", "36-44px"),
            "h3": typography.get("card_title", {}).get("size", "22-24px"),
            "body": typography.get("body", {}).get("size", "16-18px"),
            "caption": typography.get("caption", {}).get("size", "12-14px"),
        }
    except Exception:
        return {}


def _validate_cards(slides: list) -> list[str]:
    """Structural validation of AI-generated card data.

    Checks: cards ≤5 per slide (Miller's Law), each card has role, layout valid,
    type present, card has title or body, cover ≤3 info units.
    Returns list of error strings. Empty list = all valid.
    """
    VALID_LAYOUTS = {
        "single_focus", "two_column", "two_column_asymmetric", "three_column",
        "hero_grid", "mixed_grid", "dashboard", "timeline", "horizontal_split", "full_bleed"
    }
    VALID_TYPES = {
        "cover", "toc", "section", "chapter",
        "content", "data", "data_hero",
        "technique", "principle", "process_flow", "process_timeline", "timeline",
        "comparison", "duo_compare", "table", "grid_cards", "image_grid",
        "quote", "image_hero",
        "troubleshoot", "appendix", "copyright",
        "closing", "summary",
    }
    issues = []
    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            issues.append(f"Slide {i}: not a dict, got {type(slide).__name__}")
            continue
        sid = slide.get("seq", i + 1)
        stype = slide.get("type", "")
        layout = slide.get("layout", "")
        zones = slide.get("zones", {})

        if not stype:
            issues.append(f"Slide {sid}: missing 'type' field")
        elif stype not in VALID_TYPES:
            issues.append(f"Slide {sid}: unknown type '{stype}'")

        if layout and layout not in VALID_LAYOUTS:
            issues.append(f"Slide {sid}: invalid layout '{layout}'")

        if not isinstance(zones, dict):
            issues.append(f"Slide {sid}: 'zones' must be a dict")
            continue

        cards = zones.get("cards", [])
        if not isinstance(cards, list):
            issues.append(f"Slide {sid}: zones.cards must be a list")
            continue

        num_cards = len(cards)
        if num_cards == 0:
            heading = zones.get("heading", "") or zones.get("title", "") or zones.get("kicker", "")
            body = zones.get("body", "") or zones.get("lead", "") or zones.get("body_text", "")
            if not heading and not body:
                issues.append(f"Slide {sid}: no cards and no heading/body in zones")
        elif num_cards > 5:
            issues.append(f"Slide {sid}: {num_cards} cards exceeds Miller's Law limit (max 5)")

        for ci, card in enumerate(cards):
            if not isinstance(card, dict):
                issues.append(f"Slide {sid} card {ci}: not a dict")
                continue
            if "role" not in card:
                issues.append(f"Slide {sid} card {ci}: missing 'role'")
            has_content = card.get("title") or card.get("body") or card.get("chart")
            if not has_content:
                issues.append(f"Slide {sid} card {ci}: no title, body, or chart")

        if stype == "cover" and num_cards > 3:
            issues.append(f"Slide {sid}: cover has {num_cards} cards (max 3 info units)")

    return issues


def _stage2_cards(provider_id, model, llm_generate, rules, stage1_slides,
                   style_id: str = "business", temperature: float = 0.3,
                   column_id: str = "", color_scheme: str = "deep-blue") -> list | None:
    """Phase 5+6: AI selects Bento Grid layout + fills card content per slide.

    Input: stage1_slides [{seq, heading, page_type, layout_hint, body, ...}, ...]
    Output: [{type, layout, zones: {kicker, heading, lead, cards: [{role, title, body, chart}]}}, ...]

    Includes validation loop: generate → validate → if issues → fix (max 2 rounds).
    On generation failure, converts raw zones as fallback.
    """
    design_system = _load_design_system(column_id)
    style_yaml = _load_style_yaml_text(style_id, color_scheme)
    cognitive_spec = _load_cognitive_spec(column_id)

    spec_version = rules.get("spec_version", "2.2.1")

    # ── System prompt: AI as designer with strict structural rules ──
    cards_system_core = _load_cards_system_prompt()
    if not cards_system_core:
        cards_system_core = """你是一位演示文稿设计师。你必须严格按照设计系统为每页幻灯片生成结构化的卡片数据。

## 你的任务

根据大纲内容，为每页幻灯片做出设计决策：
1. **选择布局** — 根据内容语义从 10 种布局中选择（参考第十节决策树），封面必须用 full_bleed
2. **确定卡片** — 按布局→卡片映射表确定 role 和数量（参考第十一节卡片目录），每页 ≤5 张
3. **数据可视化** — 识别大纲中的数字并转化为 chart（参考第十二节 chart 决策树），有数字必有图表
4. **色彩分配** — 遵循色彩角色分工：accent=页面框架装饰，chart_colors=卡片色条轮换，primary=标题
5. **文案精炼** — 将大纲 body 文字转化为精炼的卡片 title（≤48字）+ body

## 硬性规则

- 封面页 layout=full_bleed，cards ≤3 个，禁止 hero 卡带 chart_colors 色块背景
- 每卡必有 role + (title 或 body 或 chart)
- 卡片色条颜色按 chart_colors[0]→[1]→[2]→[3]→[4] 轮换，禁止所有卡片同一颜色
- 数据页（含 %/数字/占比）→ 优先 dashboard 或 mixed_grid 布局
- 对比内容（优劣/A vs B）→ two_column 布局
- 流程/步骤 → timeline 布局

输出纯 JSON，不要用 markdown 包裹。"""

    cards_system = f"""{cards_system_core}

{design_system}

## 风格 Token（颜色/字体/阴影/圆角/渐变）
```yaml
{style_yaml}
```

## 认知设计原则
{cognitive_spec}"""

    # ── Batch generation with validation loop ──
    BATCH_SIZE = 3
    MAX_FIX_ROUNDS = 2
    result = []

    for batch_start in range(0, len(stage1_slides), BATCH_SIZE):
        batch = stage1_slides[batch_start:batch_start + BATCH_SIZE]
        batch_json = json.dumps(batch, ensure_ascii=False, indent=2)

        cards_user = f"""## 大纲（content-core 输出，每页已填充 body 正文）
{batch_json}

## 任务
为以上幻灯片选择布局并填充卡片内容。每页输出格式：
```json
{{"slides": [
  {{"seq": 1, "type": "cover", "layout": "full_bleed",
    "zones": {{"heading": "标题", "kicker": "标签", "lead": "副标题",
      "cards": [{{"role": "hero", "title": "主标题", "body": "正文内容"}}]
    }}
  }},
  {{"seq": 2, "type": "content", "layout": "hero_grid",
    "zones": {{"heading": "页标题", "kicker": "章节标签",
      "cards": [
        {{"role": "hero", "title": "核心观点", "body": "详细内容"}},
        {{"role": "metric", "title": "数据1", "chart": {{"type": "big_number", "value": 85, "label": "%"}}}}
      ]
    }}
  }}
]}}
```

铁律：
- 只输出此批次的幻灯片，按 seq 顺序
{_build_page_type_prompt(style_id, column_id)}
- layout 从以下选择: single_focus, two_column, two_column_asymmetric, three_column, hero_grid, mixed_grid, dashboard, timeline, horizontal_split, full_bleed
- ⛔ layout=single_focus 仅限 cover/quote/section 页使用。summary/closing/content/data/comparison/process/timeline 等所有其余类型一律禁止 single_focus
- 每卡必有 role (hero/metric/card_0/card_1/left/right/cell_0_0 等)
- 每卡必有 title 或 body 或 chart
- 封面页 layout=full_bleed，cards 最多 3 个
- 图表用 chart 字段: {{type: "big_number"|"donut"|"bar"|"progress_bar"|"timeline"|"sparkline", ...}}
- {spec_version} 规范：识别并转化数据为图表
- 仅输出 JSON"""

        generated = None
        fix_round = 0
        while fix_round <= MAX_FIX_ROUNDS and not generated:
            try:
                response = _safe_run_async(llm_generate(provider_id, model,
                    cards_system, cards_user, temperature=temperature))
                response = _clean_json_response(response)
                data = json.loads(response)
                slides = data.get("slides", data) if isinstance(data, dict) else data

                if not isinstance(slides, list) or len(slides) == 0:
                    raise ValueError("AI returned empty or non-list slides")

                # ── Validate ──
                errors = _validate_cards(slides)
                if errors:
                    err_text = "\n".join(f"- {e}" for e in errors)
                    _logger.warning(f"Batch {batch_start//BATCH_SIZE+1} round {fix_round}: {len(errors)} validation errors")

                    if fix_round < MAX_FIX_ROUNDS:
                        # Feed errors back into fix prompt
                        cards_user = f"""## 上一轮生成的验证错误（请逐一修复）
{err_text}

## 原始大纲（必须在修复后保留内容完整性）
{batch_json}

## 任务
修复上述验证错误后重新生成。每页输出格式不变。仅输出 JSON。"""
                        fix_round += 1
                        continue
                    else:
                        _logger.warning(f"Batch {batch_start//BATCH_SIZE+1}: max fix rounds reached, accepting with {len(errors)} errors")

                # ── Augment each slide with missing fields from stage1 outline ──
                for slide in slides:
                    if isinstance(slide, dict):
                        seq = slide.get("seq", 0)
                        s1 = next((s for s in batch if s.get("seq") == seq), None)
                        if s1:
                            zones = slide.setdefault("zones", {})
                            if not zones.get("heading") and s1.get("heading"):
                                zones["heading"] = s1["heading"]
                            if not slide.get("layout") and s1.get("layout_hint"):
                                slide["layout"] = s1["layout_hint"]
                            if not slide.get("type") and s1.get("page_type"):
                                slide["type"] = s1["page_type"]
                            if not zones.get("body") and not zones.get("cards") and s1.get("body"):
                                zones["body"] = s1["body"]
                            if not zones.get("kicker") and s1.get("kicker"):
                                zones["kicker"] = s1["kicker"]
                            if not zones.get("lead") and s1.get("lead"):
                                zones["lead"] = s1["lead"]

                generated = slides

            except Exception as e:
                err_msg = str(e)[:200]
                _logger.warning(f"Batch {batch_start//BATCH_SIZE+1} round {fix_round} failed: {err_msg}")
                if fix_round < MAX_FIX_ROUNDS:
                    cards_user = f"""## 上一轮生成失败: {err_msg}，请重试

## 原始大纲
{batch_json}

## 任务
重新为以上幻灯片选择布局并填充卡片内容。仅输出 JSON。"""
                    fix_round += 1
                else:
                    break

        if not generated:
            # ── Fallback: convert raw stage1 zones to card format ──
            _logger.warning(f"Batch {batch_start//BATCH_SIZE+1}: generation failed, using raw zone fallback")
            generated = _fallback_zones_to_cards(batch)

        result.extend(generated)

    return result if result else None


def _fallback_zones_to_cards(slides: list) -> list:
    """Convert raw stage1 outline slides to minimal card-structured format.

    Fallback when AI card generation fails — ensures basic structural validity
    so render_deck() can still produce output.
    """
    result = []
    for s in slides:
        if not isinstance(s, dict):
            continue
        seq = s.get("seq", 0)
        heading = s.get("heading", "")
        ptype = s.get("page_type", "content")
        layout = s.get("layout_hint", "hero_grid")
        body = s.get("body", "")
        notes = s.get("notes", "")
        kicker = s.get("kicker", "")
        key_points = s.get("key_points", [])

        body_text = body
        if key_points and not body_text:
            body_text = "\n".join(f"• {kp}" for kp in key_points if isinstance(kp, str))

        zones = {
            "heading": heading,
            "kicker": kicker,
            "body": body_text,
            "cards": [{
                "role": "hero",
                "title": heading,
                "body": body_text or notes,
                "kicker": kicker,
            }]
        }

        result.append({
            "seq": seq,
            "type": ptype,
            "layout": layout,
            "zones": zones,
        })

    return result


# ═══════════════════════════════════════════════════════════════════════════════
# Two-Phase HTML Pipeline (replaces old JSON card → SVG flow)
# Phase 1: Structure planning (lightweight, one call)
# Phase 2: Per-slide HTML generation (parallel, design-system.md guided)
# ═══════════════════════════════════════════════════════════════════════════════

def _fallback_stage1_structure(stage1_slides: list) -> list:
    """Build slide structure from stage1 outline deterministically — no LLM needed."""
    fallback_layouts = [
        "full_bleed", "three_column", "two_column", "hero_grid",
        "mixed_grid", "dashboard", "two_column_asymmetric", "full_bleed"
    ]
    result = []
    for s in stage1_slides:
        if not isinstance(s, dict):
            continue
        seq = s.get("seq", 0)
        layout_idx = (seq - 1) % len(fallback_layouts)
        ptype = s.get("page_type", "content")
        heading = s.get("heading", "")
        cards = [{"role": "hero", "content_hint": heading}]
        has_chart = any(kw in (heading + s.get("body", ""))
                       for kw in ["%", "数据", "指标", "占比", "率", "值", "量"])
        result.append({
            "seq": seq,
            "type": ptype,
            "layout": s.get("layout_hint") or fallback_layouts[layout_idx],
            "heading": heading,
            "body": s.get("body", ""),
            "key_points": s.get("key_points", []),
            "kicker": s.get("kicker", ""),
            "lead": s.get("lead", ""),
            "notes": s.get("notes", ""),
            "cards": cards,
            "has_chart": has_chart,
            "chart_hint": "big_number" if has_chart else "",
        })
    return result


def _stage2_structure(provider_id, model, llm_generate, stage1_slides,
                      style_id: str = "business", temperature: float = 0.3,
                      column_id: str = "", color_scheme: str = "deep-blue") -> list | None:
    """Phase 1 of two-phase HTML pipeline: Lightweight structure planning.

    One LLM call for ALL slides. AI decides: slide types, layouts, card count/roles,
    whether each slide needs a chart. Output is a structural blueprint that Phase 2
    uses to generate per-slide HTML.

    This is intentionally lightweight — no JSON card filling, just structure decisions.
    """
    design_system = _load_design_system(column_id)
    style_yaml = _load_style_yaml_text(style_id, color_scheme, resolve_vars=False)
    style_vi = _load_style_vi(style_id, color_scheme, resolve_vars=False)
    style_prompt = _load_style_prompt(style_id, color_scheme)

    if style_prompt:
        persona_block = style_prompt
    else:
        persona_block = "你是演示文稿结构规划师。你必须严格按照以下设计系统为每页幻灯片做出结构设计决策。"

    vi_block = ""
    if style_vi:
        vi_block = f"""

## 视觉识别系统 (VIS) — 本风格的权威视觉规范
{style_vi}
"""

    struct_output = _load_structure_output_prompt()
    if not struct_output:
        struct_output = """2. layout: 从第十节布局库中选择
3. cards: 每张卡指定 role 和 content_hint
4. has_chart: 有数字/百分比时 = true
5. chart_hint: big_number/donut/bar/progress_bar/timeline/sparkline

页面数量严格等于大纲给出的页数，不增不减。

输出纯 JSON，不要用 markdown 包裹。"""

    system = f"""{persona_block}

{design_system}

## 风格 Token
```yaml
{style_yaml}
```
{vi_block}
## 输出格式要求

为每页输出结构决策：
{_build_page_type_prompt(style_id, column_id)}
{struct_output}"""

    # ── DEBUG: monitor structure prompt for hex ──
    import re as _re_mon3
    struct_hex = len(_re_mon3.findall(r'#[0-9a-fA-F]{6}', system))
    _logger.info(f"[MONITOR] Structure system prompt: {len(system)} chars, hex_count={struct_hex}")
    debug_dir = os.path.join(BASE_DIR, "data", "debug")
    os.makedirs(debug_dir, exist_ok=True)
    with open(os.path.join(debug_dir, "last_structure_prompt.txt"), "w", encoding="utf-8") as _df3:
        _df3.write(system)

    outline_json = json.dumps(stage1_slides, ensure_ascii=False, indent=2)
    user = f"""## 大纲（{len(stage1_slides)} 页）
{outline_json}

## 任务
为每页做出结构决策，输出格式：
```json
{{"slides": [
  {{"seq":1,"type":"cover","layout":"full_bleed","cards":[{{"role":"hero","content_hint":"主标题+副标题"}}],"has_chart":false,"chart_hint":""}},
  {{"seq":2,"type":"content","layout":"hero_grid","cards":[{{"role":"hero","content_hint":"核心观点"}},{{"role":"metric","content_hint":"关键指标"}}],"has_chart":true,"chart_hint":"big_number"}},
  ...
]}}
```

严格遵守设计系统中的所有规则。仅输出 JSON。"""

    for attempt in range(2):
        try:
            response = _safe_run_async(llm_generate(provider_id, model,
                system, user, temperature=temperature))
            response = _clean_json_response(response)
            data = json.loads(response)
            slides = data.get("slides", data) if isinstance(data, dict) else data
            if isinstance(slides, list) and len(slides) > 0:
                # Merge with stage1 content
                for s in slides:
                    seq = s.get("seq", 0)
                    s1 = next((x for x in stage1_slides if x.get("seq") == seq), None)
                    if s1:
                        s["heading"] = s1.get("heading", "")
                        s["body"] = s1.get("body", "")
                        s["key_points"] = s1.get("key_points", [])
                        s["kicker"] = s1.get("kicker", "")
                        s["lead"] = s1.get("lead", "")
                        s["notes"] = s1.get("notes", "")
                _logger.info(f"Stage 2 structure: {len(slides)} slides planned "
                           f"({sum(1 for s in slides if s.get('has_chart'))} with charts)")
                return slides
        except Exception as e:
            _logger.warning(f"Stage 2 structure attempt {attempt+1} failed: {e}")

    # Fallback: deterministic structure from stage1
    _logger.info("Stage 2 structure: using deterministic fallback")
    return _fallback_stage1_structure(stage1_slides)


def _fallback_single_slide_html(slide: dict, style_id: str,
                                 canvas_w: int = 1280, canvas_h: int = 720) -> dict:
    """Generate minimal HTML for a single slide when AI generation fails."""
    seq = slide.get("seq", 0)
    stype = slide.get("type", "content")
    layout = slide.get("layout", "hero_grid")
    heading = slide.get("heading", "Untitled")
    body = slide.get("body", "")
    kicker = slide.get("kicker", "")
    lead = slide.get("lead", "")
    key_points = slide.get("key_points", [])

    body_html = ""
    if body:
        body_html = f'<p style="font-size:18px;line-height:1.7;color:{{{{text}}}};margin:0 0 16px">{body[:500]}</p>'
    if key_points:
        pts = "".join(f'<li style="font-size:16px;line-height:1.6;margin-bottom:6px">{kp}</li>'
                      for kp in key_points[:5] if isinstance(kp, str))
        body_html += f'<ul style="padding-left:20px;margin:0">{pts}</ul>'

    kicker_html = f'<div style="font-size:14px;color:{{{{text}}}};opacity:0.55;letter-spacing:2px;text-transform:uppercase;margin-bottom:8px">{kicker}</div>' if kicker else ""
    lead_html = f'<div style="font-size:20px;color:{{{{text}}}};opacity:0.7;margin-bottom:24px">{lead}</div>' if lead else ""

    html = f"""<section style="width:{canvas_w}px;height:{canvas_h}px;background:{{{{background}}}};position:relative;overflow:hidden;font-family:system-ui,-apple-system,sans-serif">
<div style="position:absolute;top:0;left:0;right:0;height:4px;background:{{{{accent}}}}"></div>
<div style="padding:60px 80px 40px">
{kicker_html}
<h1 style="font-size:38px;font-weight:700;color:{{{{primary}}}};margin:0 0 8px;letter-spacing:-0.3px">{heading}</h1>
<div style="width:40px;height:3px;background:{{{{accent}}}};border-radius:2px;margin-bottom:32px"></div>
{lead_html}
{body_html}
</div>
<div style="position:absolute;bottom:28px;right:48px;display:flex;align-items:center;gap:8px;font-size:13px;color:{{{{text}}}};opacity:0.35">
<span style="width:6px;height:6px;border-radius:50%;background:{{{{accent}}}}"></span>
{seq}
</div>
</section>"""

    return {**slide, "html": html, "html_vars": html}


def _stage2_html_per_slide(provider_id, model, llm_generate, structure_slides,
                           style_id: str = "business", parallel: int = 3,
                           temperature: float = 0.3, column_id: str = "",
                           color_scheme: str = "deep-blue",
                           project_id: str = "") -> list | None:
    """Phase 2 of two-phase HTML pipeline: Per-slide parallel HTML generation.

    Each slide gets its own LLM call with the full design-system.md as the design guide.
    Calls run in parallel (default 3 at a time) via ThreadPoolExecutor.

    LLM generates complete inline-styled HTML for a slide section.
    Canvas dimensions are determined by the column's canvas config (default 1280×720).
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed
    import re

    canvas_w, canvas_h = _get_canvas_dimensions(column_id)

    # resolve_vars=False: LLM gets {{primary}} placeholders, not hex values
    style_yaml = _load_style_yaml_text(style_id, color_scheme, resolve_vars=False)
    active_scheme = _load_scheme_data(style_id, color_scheme)
    font_range = _resolve_font_range(style_yaml)
    style_prompt = _load_style_prompt(style_id, color_scheme)

    # Verify modular core files exist
    core_identity = os.path.join(BASE_DIR, "resources", "prompts", "core", "always", "identity.md")
    if not os.path.exists(core_identity):
        _logger.error("core/always/identity.md not found — cannot build per-slide prompts")
        return None

    font_info = ""
    if font_range:
        font_info = "\n".join(f"- {k}: {v}" for k, v in font_range.items())
        font_info = f"\n## 字号参考（从 style YAML 提取）\n{font_info}"

    if canvas_h > canvas_w:
        persona_block = "你是文档排版师。你必须严格按照《A4 文档排版系统》为每一页生成完整的 HTML。使用 8 列统一表格网格，禁止卡片布局和绝对定位。"
    elif style_prompt:
        persona_block = style_prompt
    else:
        persona_block = "你是演示文稿设计艺术总监。你必须严格按照《幻灯片 HTML 设计系统 v2》为每一页生成完整的 HTML。"

    total = len(structure_slides)
    is_a4 = canvas_h > canvas_w

    # For A4 documents: load design system explicitly (normally loaded in _stage2_structure,
    # but that stage is skipped for col3)
    doc_design_system = ""
    if is_a4:
        doc_design_system = _load_design_system(column_id)
        if doc_design_system:
            doc_design_system = f"\n{doc_design_system}\n"

    # Load format specification (right/wrong examples for every syntax category)
    # Included for BOTH PPT and A4 pipelines — prevents LLM from generating
    # malformed hex, missing CSS units, wrong variable syntax, etc.
    format_spec_path = os.path.join(BASE_DIR, "resources", "prompts", "core", "always", "format-spec.md")
    format_spec = ""
    if os.path.exists(format_spec_path):
        with open(format_spec_path, encoding="utf-8") as _fs:
            format_spec = _fs.read()

    import threading as _threading
    _done_lock = _threading.Lock()
    _done_count = [0]

    def _gen_one(slide, idx):
        seq = slide.get("seq", idx + 1)
        stype = slide.get("type", "content")
        layout = slide.get("layout", "hero_grid")
        heading = slide.get("heading", "")
        body = slide.get("body", "")
        key_points = slide.get("key_points", [])
        kicker = slide.get("kicker", "")
        lead = slide.get("lead", "")
        cards = slide.get("cards", [])
        has_chart = slide.get("has_chart", False)
        chart_hint = slide.get("chart_hint", "")
        notes = slide.get("notes", "")

        # ── Per-slide lean system prompt ──
        # Build a tailored system prompt: core rules + slide-type-specific sections
        if is_a4:
            # A4 document: core system comes from design-system.md (loaded via _load_design_system)
            # Skip build_slide_prompt — PPT card/layout rules don't apply to documents
            core_system = ""
        else:
            core_system = build_slide_prompt(stype, layout, has_chart)
        # Append VI section for this slide type (design instruction, not user content)
        vi_section = _load_style_vi_section(style_id, stype, color_scheme, resolve_vars=False, column_id=column_id)
        vi_append = ""
        if vi_section:
            vi_append = f"\n\n## {stype} 类型专属视觉规范\n{vi_section}"

        html_output_inst = _load_html_output_prompt(column_id)
        if html_output_inst:
            # Substitute template variables — {{canvas_w}} and {{canvas_h}} are
            # placeholders in scenarios/_default/html-output.md. Column-specific
            # overrides use the same variable names so one substitution covers all.
            html_output_inst = html_output_inst.replace("{{canvas_w}}", str(canvas_w))
            html_output_inst = html_output_inst.replace("{{canvas_h}}", str(canvas_h))
        else:
            html_output_inst = f"""输出一个 HTML 代码块，用 ```html ... ``` 包裹。仅输出 HTML。

样式要求：
- 尺寸: width:{canvas_w}px; height:{canvas_h}px, position:relative, overflow:hidden
- 全部内联 inline style，禁止 class/id/<style>/@import/@font-face
- 全部 CSS 属性必须带单位（px）
- 禁止输出 <!DOCTYPE html>/<html>/<head>/<body>/<title>/<meta>/<link>
- 输出只包含一个 div 容器（width:{canvas_w}px;height:{canvas_h}px）及其子元素"""

        # ── Cover slide: inject quantified color rules from tokens ──
        cover_color_rules = ""
        if is_a4 and stype == "cover" and active_scheme:
            cover_color_rules = _build_cover_color_rules(style_id, color_scheme, active_scheme)
            if cover_color_rules:
                cover_color_rules = f"\n{cover_color_rules}\n"

        tailored_system = f"""{persona_block}

{core_system}
{doc_design_system}
## 风格 Token（仅排版，无颜色 hex）
```yaml
{style_yaml}
```
{font_info}
{vi_append}
{cover_color_rules}
{html_output_inst}
{format_spec}"""

        content_parts = [
            f"页码: {seq}/{total}",
            f"类型: {stype}",
        ]
        if not is_a4:
            content_parts.append(f"布局: {layout}")
        content_parts.append(f"标题: {heading}")
        if kicker:
            content_parts.append(f"章节标签: {kicker}")
        if lead:
            content_parts.append(f"副标题: {lead}")
        if body:
            content_parts.append(f"正文内容: {body[:1000]}")
        if key_points:
            content_parts.append(f"关键点: {'; '.join(str(kp) for kp in key_points[:5])}")
        if not is_a4 and cards:
            cards_desc = "; ".join(
                f"[{c.get('role','card')}] {c.get('content_hint','')}" for c in cards)
            content_parts.append(f"卡片分配: {cards_desc}")
        if notes:
            content_parts.append(f"备注: {notes[:300]}")
        if not is_a4 and has_chart:
            content_parts.append(
                f"需要数据图表: {chart_hint or '根据内容中的数据选择合适的图表类型(big_number/donut/bar/progress_bar)'}")

        user = "\n".join(content_parts)

        # ── DEBUG: monitor prompt for color hex leakage ──
        if idx == 0:
            import re as _re_monitor
            hex_count = len(_re_monitor.findall(r'#[0-9a-fA-F]{6}', tailored_system))
            has_colors_md = "四、色彩系统" in tailored_system
            has_semantics_md = "十四、色彩语义" in tailored_system
            has_var_law = "颜色变量铁律" in tailored_system
            _logger.info(f"[MONITOR] Slide {seq} system prompt: {len(tailored_system)} chars, "
                        f"hex_count={hex_count}, colors_md={has_colors_md}, "
                        f"color_semantics_md={has_semantics_md}, var_iron_law={has_var_law}")
            # Dump to file for inspection
            debug_dir = os.path.join(BASE_DIR, "data", "debug")
            os.makedirs(debug_dir, exist_ok=True)
            with open(os.path.join(debug_dir, "last_system_prompt.txt"), "w", encoding="utf-8") as _df:
                _df.write(tailored_system)

        for attempt in range(2):
            try:
                response = _safe_run_async(llm_generate(provider_id, model,
                    tailored_system, user, temperature=temperature))
                # Extract HTML
                m = re.search(r'```html\s*\n(.*?)\n```', response, re.DOTALL)
                if m:
                    html = m.group(1).strip()
                else:
                    m = re.search(r'(<section[\s\S]*?</section>)', response, re.IGNORECASE)
                    if m:
                        html = m.group(1).strip()
                    else:
                        html = response.strip()
                        if html.startswith("```"):
                            html = re.sub(r'^```\w*\n?', '', html)
                            html = re.sub(r'\n?```$', '', html)

                if html and len(html) > 300:
                    # ── DEBUG: monitor LLM output for hex leakage ──
                    if idx == 0:
                        import re as _re_monitor2
                        html_hex_count = len(_re_monitor2.findall(r'#[0-9a-fA-F]{6}', html))
                        html_var_count = len(_re_monitor2.findall(r'\{\{[a-z_]+\}\}', html))
                        _logger.info(f"[MONITOR] Slide {seq} LLM output: {len(html)} chars, "
                                    f"hex_count={html_hex_count}, var_count={html_var_count}")
                        with open(os.path.join(debug_dir, "last_llm_output.txt"), "w", encoding="utf-8") as _df2:
                            _df2.write(html)

                    # Post-process: fix common LLM HTML errors
                    html = _fix_llm_html_errors(html, is_a4=is_a4)
                    # Post-process: enforce cover rules on UNRESOLVED HTML
                    # Must run BEFORE _auto_fix_hardcoded_hex (which skips #ffffff)
                    # and BEFORE _resolve_color_vars (which destroys CSS variable info).
                    # After resolution, is_variable_bg check fails because var(--card_bg)
                    # has been replaced with its hex value.
                    if is_a4 and stype == "cover" and active_scheme:
                        html = _enforce_cover_rules(html, active_scheme, style_id)
                    # Post-process: scan & replace hardcoded hex with {{placeholder}} vars
                    html = _fix_malformed_hex(html, seq)
                    html = _auto_fix_hardcoded_hex(html, active_scheme, seq)
                    # Post-process: fix white text on light backgrounds (LLM habit from dark styles)
                    # Must run AFTER _auto_fix_hardcoded_hex (which intentionally skips #ffffff)
                    # and BEFORE _resolve_color_vars (which converts {{placeholder}} → hex)
                    # IMPORTANT: Skip slides whose tokens.yaml overrides specify
                    # dark backgrounds with white text (cover, section, summary, quote).
                    # These correctly use white-on-dark — overriding would break them.
                    if active_scheme:
                        html = _auto_fix_white_on_light(html, active_scheme, seq,
                                                        style_id=style_id, page_type=stype)
                    # Post-process: strip LLM-invented CSS variable reassignments
                    # (e.g. --primary: var(--card_bg) inverts the theme → invisible text)
                    html = _strip_local_var_overrides(html, seq)
                    # Post-process: enforce VI typography minimums (15px → 16px)
                    html = _auto_fix_font_size(html, seq, is_a4=is_a4)
                    # Save unresolved HTML for later recolor (before hex resolution)
                    html_vars = html
                    # Resolve {{primary}} etc. → actual hex from active color scheme
                    if active_scheme:
                        html = _resolve_color_vars(html, active_scheme, css_vars=True)
                    # Check for truncation: HTML must end with a properly closed tag
                    stripped = html.rstrip()
                    if not stripped.endswith('>'):
                        _logger.warning(f"Slide {seq} attempt {attempt+1}: "
                                        f"truncated (ends with '{stripped[-50:]}')")
                        continue  # retry

                    # Check div balance: unbalanced divs collapse browser rendering.
                    # Both missing AND extra </div> are dangerous — extra </div> prematurely
                    # closes the slide-wrapper, corrupting the NEXT slide's DOM nesting.
                    open_divs = len(re.findall(r'<div\b', html))
                    close_divs = len(re.findall(r'</div>', html))
                    if open_divs != close_divs:
                        delta = open_divs - close_divs
                        if attempt < 1:
                            _logger.warning(f"Slide {seq} attempt {attempt+1}: "
                                            f"unbalanced divs ({open_divs} open vs {close_divs} close, "
                                            f"delta={delta}), retrying")
                            continue
                        elif delta > 0:
                            _logger.warning(f"Slide {seq}: unbalanced divs ({open_divs} open vs "
                                            f"{close_divs} close), auto-fixing by appending {delta} </div>")
                            html += '</div>' * delta
                            html_vars += '</div>' * delta
                        elif delta < 0:
                            _logger.warning(f"Slide {seq}: extra </div> ({open_divs} open vs "
                                            f"{close_divs} close, surplus {abs(delta)}), "
                                            f"stripping excess closing tags")
                            for _ in range(abs(delta)):
                                last_close = html.rfind('</div>')
                                if last_close > 0:
                                    html = html[:last_close] + html[last_close + 6:]
                                    vc = html_vars.rfind('</div>')
                                    if vc > 0:
                                        html_vars = html_vars[:vc] + html_vars[vc + 6:]

                    # Check SVG balance: unclosed <svg> corrupts DOM parsing just like unclosed <div>
                    open_svgs = len(re.findall(r'<svg\b', html))
                    close_svgs = len(re.findall(r'</svg>', html))
                    if open_svgs != close_svgs:
                        delta = open_svgs - close_svgs
                        if attempt < 1:
                            _logger.warning(f"Slide {seq} attempt {attempt+1}: "
                                            f"unbalanced SVG tags ({open_svgs} open vs {close_svgs} close, "
                                            f"delta={delta}), retrying")
                            continue
                        elif delta > 0:
                            _logger.warning(f"Slide {seq}: {delta} unclosed SVG(s), "
                                            f"appending </svg> at end of slide HTML")
                            html += '</svg>' * delta
                            html_vars += '</svg>' * delta
                        elif delta < 0:
                            _logger.warning(f"Slide {seq}: extra </svg> ({open_svgs} open vs "
                                            f"{close_svgs} close, surplus {abs(delta)}), stripping")
                            for _ in range(abs(delta)):
                                last_svg = html.rfind('</svg>')
                                if last_svg > 0:
                                    html = html[:last_svg] + html[last_svg + 6:]
                                    vs = html_vars.rfind('</svg>')
                                    if vs > 0:
                                        html_vars = html_vars[:vs] + html_vars[vs + 6:]

                    # --- Code-based structural checks (not LLM self-review) ---
                    retry_msg = ""

                    # Check 1: container violations (forbidden centered-card patterns)
                    violation = _detect_container_violation(html, stype)
                    if violation:
                        retry_msg = f"⛔ 容器违规({violation})，已废弃。重新生成必须：内容区用 left:60px; right:60px 基准容器，禁止 left:180px 和 transform:translateX(-50%)。"

                    # Check 2: content overflow (too many content blocks per card)
                    if not retry_msg:
                        overflow = _detect_content_overflow(html, stype, layout)
                        if overflow:
                            cards_str = ", ".join(
                                f"卡片{i}有{c}个内容块" for i, c in overflow.items()
                            )
                            retry_msg = (
                                f"⛔ 内容溢出({cards_str})，已废弃。"
                                f"重新生成必须：将内容拆分到更多卡片中。"
                                f"每张卡片最多容纳 8 个信息单元（标题/指标/步骤/段落），"
                                f"超出必须使用多卡片布局或拆分到多页。"
                                f"当前检测到 {overflow} 张卡片内容超载。"
                            )

                    # Check 3: full-screen opaque mask
                    if not retry_msg:
                        mask = _detect_fullscreen_mask(html)
                        if mask:
                            retry_msg = (
                                f"⛔ 全屏遮罩({mask})，已废弃。"
                                f"重新生成必须：禁止使用覆盖整个画布的不透明 div。"
                                f"装饰元素必须使用半透明背景或限制在局部区域。"
                            )

                    if retry_msg and attempt < 1:
                        _logger.warning(f"Slide {seq} attempt {attempt+1}: {retry_msg[:120]}")
                        user += f"\n\n{retry_msg}"
                        continue
                    elif retry_msg:
                        _logger.warning(f"Slide {seq}: {retry_msg[:120]} "
                                        f"after retries, accepting (reviewer should catch)")

                    _logger.info(f"Slide {seq}: HTML {len(html)} chars")
                    # Update progress for status polling
                    if project_id:
                        with _done_lock:
                            _done_count[0] += 1
                            _ppt_status[project_id] = {"phase": "generating", "phase_label": "正在生成页面...", "message": f"已完成 {_done_count[0]}/{total} 页", "slides_done": _done_count[0], "slides_total": total}
                            if _done_count[0] % 3 == 0 or _done_count[0] == total:
                                _append_log(project_id, f"HTML 生成进度: {_done_count[0]}/{total} 页")
                    return {**slide, "html": html, "html_vars": html_vars}
                else:
                    _logger.warning(f"Slide {seq} attempt {attempt+1}: HTML too short ({len(html)} chars)")
            except Exception as e:
                import traceback as _tb
                _logger.warning(f"Slide {seq} HTML attempt {attempt+1} failed: {e}")
                _logger.warning(f"Slide {seq} traceback: {_tb.format_exc()[-500:]}")

        _logger.warning(f"Slide {seq}: all attempts failed, using fallback")
        return _fallback_single_slide_html(slide, style_id, canvas_w, canvas_h)

    _logger.info(f"Stage 2 HTML: dispatching {total} slides in parallel (workers={parallel})")
    t0 = __import__("time").time()

    result = []
    with ThreadPoolExecutor(max_workers=parallel) as ex:
        futures = {ex.submit(_gen_one, s, i): i for i, s in enumerate(structure_slides)}
        for f in as_completed(futures):
            try:
                r = f.result()
                if r:
                    result.append(r)
            except Exception as e:
                _logger.error(f"Slide HTML generation crashed: {e}")

    elapsed = __import__("time").time() - t0
    result.sort(key=lambda s: s.get("seq", 0))
    _logger.info(f"Stage 2 HTML: {len(result)}/{total} slides generated in {elapsed:.1f}s")
    return result if result else None


def _fallback_stage1_to_html_slides(stage1_slides: list, style_id: str = "business") -> list:
    """Convert stage1 outline directly to HTML slides — fallback when AI HTML fails."""
    structure = _fallback_stage1_structure(stage1_slides)
    result = []
    for s in structure:
        result.append(_fallback_single_slide_html(s, style_id))
    return result


def _detect_container_violation(html: str, page_type: str = "") -> str | None:
    """Scan HTML for forbidden container patterns.

    Returns violation description string if found, None if clean.
    Cover, quote, and section pages are exempt from container width checks.
    """
    import re

    exempt_types = {"cover", "quote", "section", "summary", "closing"}
    if page_type in exempt_types:
        return None

    # Check 1: Forbidden left:180px on content area (not decorative)
    if re.search(r'left\s*:\s*180px', html):
        return "left:180px (forbidden margin, use left:60px; right:60px)"

    # Check 2: Centered positioning with transform
    if re.search(r'transform\s*:\s*translateX\s*\(\s*-50%\s*\)', html):
        return "transform:translateX(-50%) (centered card, use base container)"

    # Check 3: Non-standard left margins on content cards
    # Only flag if width is also specified (indicating a centered card, not a full-width container)
    narrow_card = re.search(r'left\s*:\s*(30|36|40|48|80|88)px.*?width\s*:\s*\d+px', html)
    if narrow_card:
        return f"left:{narrow_card.group(1)}px + fixed width (non-standard margin)"

    # Check 4: Fixed-width card with no right edge constraint (orphan card)
    # Pattern: a content div with width:920px or width:960px and no right:60px sibling
    if re.search(r'(?:width\s*:\s*920px|width\s*:\s*960px)', html):
        # Make sure it's not a decorative element
        if not re.search(r'right\s*:\s*60px', html):
            return "narrow card (920/960px) without right-edge anchor"

    return None


def _parse_style(style_str: str) -> dict:
    """Parse inline CSS style string into a dict of property:value pairs."""
    if not style_str:
        return {}
    result = {}
    for part in style_str.split(";"):
        part = part.strip()
        if ":" in part:
            k, v = part.split(":", 1)
            result[k.strip()] = v.strip()
    return result


def _detect_fullscreen_mask(html: str) -> str | None:
    """Detect full-screen opaque divs that cover all slide content.

    Pattern: a div with position:absolute covering the full 1280x720 canvas
    with an opaque background (no transparency).
    """
    from bs4 import BeautifulSoup
    import re

    soup = BeautifulSoup(html, "html.parser")

    for div in soup.find_all("div"):
        s = _parse_style(div.get("style", ""))
        pos = s.get("position", "")
        if pos not in ("absolute", "fixed"):
            continue
        top = s.get("top", "").strip()
        left = s.get("left", "").strip()
        width = s.get("width", "")
        height = s.get("height", "")
        # Check full-viewport coverage
        covers_top_left = (top in ("0", "0px") and left in ("0", "0px"))
        covers_size = (
            width in ("100%", "100vw", "1280px")
            and height in ("100%", "100vh", "720px")
        )
        if not (covers_top_left and covers_size):
            continue
        # Skip empty background layers — an empty full-screen div is a
        # legitimate slide background, not a content-hiding mask.
        if not div.get_text(strip=True) and not div.find():
            continue
        # Check if background is opaque
        bg = s.get("background", "")
        bg_color = s.get("background-color", "")
        combined = f"{bg} {bg_color}"
        # Look for solid colors (not rgba with alpha < 1, not transparent)
        if re.search(r'#[0-9a-fA-F]{6}', combined):
            return f"opaque full-screen mask (bg={bg or bg_color})"
        if re.search(
            r'rgba?\s*\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*\)', combined
        ):
            return f"opaque full-screen mask (bg={bg or bg_color})"
        if re.search(
            r'rgba?\s*\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*1\s*\)', combined
        ):
            return f"opaque full-screen mask (bg={bg or bg_color})"
        # Named colors (white, black, etc.) without transparency
        named_opaques = [
            "white", "#fff", "#ffffff", "black", "#000", "#000000",
            "red", "blue", "green", "gray", "grey",
        ]
        for c in named_opaques:
            if c in combined.lower() and "transparent" not in combined.lower():
                # Check it's not rgba with alpha < 1
                if "rgba" not in combined.lower():
                    return f"opaque full-screen mask (bg={bg or bg_color})"

    return None


def _detect_content_overflow(html: str, page_type: str = "", layout: str = "") -> dict | None:
    """Detect cards whose content exceeds their physical capacity.

    Parses the HTML DOM tree, finds card containers, and counts distinct
    content blocks within each.  Returns {card_index: block_count} for any
    card exceeding the threshold, or None if all cards are within capacity.

    This is purely mechanical — no LLM involved.  The physical constraint:
    a card inside a 720px slide (130px top + 50px bottom = 540px available)
    with 28px padding can hold at most 8 content blocks before overflow.
    """
    from bs4 import BeautifulSoup
    import re
    import math

    soup = BeautifulSoup(html, "html.parser")

    # ── Find the base content container (4th layer) ──
    # Pattern: position:absolute + left:60px + right:60px + display:flex
    base = None
    for div in soup.find_all("div"):
        s = _parse_style(div.get("style", ""))
        if (
            s.get("position") == "absolute"
            and s.get("left") == "60px"
            and s.get("right") == "60px"
            and "flex" in s.get("display", "")
        ):
            base = div
            break

    if not base:
        return None  # not a standard content page layout, skip

    # ── Calculate available card height ──
    top_px = 130
    bottom_px = 50
    top_str = _parse_style(base.get("style", "")).get("top", "")
    bot_str = _parse_style(base.get("style", "")).get("bottom", "")
    try:
        top_px = int(top_str.replace("px", ""))
    except (ValueError, AttributeError):
        pass
    try:
        bottom_px = int(bot_str.replace("px", ""))
    except (ValueError, AttributeError):
        pass
    available_height = 720 - top_px - bottom_px  # typically 540px

    # ── Find card elements (direct flex children of the base container) ──
    cards = []
    for child in base.find_all("div", recursive=False):
        child_style = child.get("style", "")
        if "flex:" in child_style or "flex " in child_style:
            cards.append(child)

    if not cards:
        return None  # no cards found, skip

    # ── Per-card thresholds ──
    exempt_types = {"cover", "quote", "section", "summary", "closing"}
    if page_type in exempt_types:
        return None

    # A card can physically hold ~8 content blocks before overflow
    # (540px - 56px padding = 484px; ~55px per block → 8.8 blocks)
    max_blocks_per_card = 8

    overflow: dict = {}
    for i, card in enumerate(cards):
        count = _count_content_blocks(card)
        if count > max_blocks_per_card:
            overflow[i] = count

    return overflow if overflow else None


def _count_content_blocks(card_element) -> int:
    """Count distinct content blocks within a card element.

    A content block is any visible structural unit that consumes vertical space:
    icon+text row, step item, metric display, text paragraph, or note box.

    Walks all descendant elements (div, span, p) and counts those that match
    content-block patterns.  Uses a set of element ids to avoid double-counting
    a parent container AND its children.
    """
    import re

    count = 0
    counted_ids: set = set()

    for elem in card_element.descendants:
        if not hasattr(elem, "name"):
            continue  # NavigableString, skip
        if elem.name not in ("div", "span", "p"):
            continue

        style = elem.get("style", "") if hasattr(elem, "get") else ""
        eid = id(elem)

        # --- Exclusions ---
        # Absolute-positioned (accent bars, overlays)
        if "position:absolute" in style or "position: absolute" in style:
            continue
        # Accent bar: thin strip
        if re.search(r'width\s*:\s*4px', style) and re.search(r'(?:top|bottom)\s*:\s*0', style):
            continue
        # Invisible or near-invisible
        if re.search(r'opacity\s*:\s*0(?:\.0)?\b', style):
            continue
        # Pure structural dividers (1px height, no text)
        own = elem.find(string=True, recursive=False)
        own_text = own.strip() if own else ""
        if re.search(r'height\s*:\s*1px', style) and len(own_text) == 0:
            continue

        # --- Content block patterns ---
        has_svg = elem.find("svg") is not None if hasattr(elem, "find") else False

        # Pattern A: Big number metric (font-size >= 38px)
        if re.search(r'font-size\s*:\s*(?:3[89]|[4-9]\d)px', style):
            if eid not in counted_ids:
                count += 1
                counted_ids.add(eid)
            continue

        # Pattern B: Step row — parent div containing a numbered circle
        # (width:18px; height:18px; border-radius:9px) AND text
        has_numbered_circle = False
        for span in elem.find_all("span", recursive=False):
            s = span.get("style", "") if hasattr(span, "get") else ""
            if re.search(
                r'width\s*:\s*1[8-9]px.*?height\s*:\s*1[8-9]px.*?border-radius\s*:\s*(?:9|10|50)px', s
            ):
                has_numbered_circle = True
                break
        if has_numbered_circle:
            if eid not in counted_ids:
                count += 1
                counted_ids.add(eid)
            continue

        # Pattern C: Icon+text flex row (SVG icon + span text, display:flex, align-items)
        if has_svg and elem.find("span") is not None:
            if "display:flex" in style and "align-items" in style:
                if eid not in counted_ids:
                    count += 1
                    counted_ids.add(eid)
                continue

        # Pattern D: Standalone paragraph with substantial text (>30 chars)
        if elem.name == "p" and len(elem.get_text(strip=True)) > 30:
            if eid not in counted_ids:
                count += 1
                counted_ids.add(eid)
            continue

        # Pattern E: Warning/note box — colored background, padding, substantial text
        if (
            "padding" in style
            and "border-radius" in style
            and re.search(r'background\s*:\s*#[0-9a-fA-F]{6}', style)
            and len(elem.get_text(strip=True)) > 25
        ):
            if eid not in counted_ids:
                count += 1
                counted_ids.add(eid)

    return count


def build_slide_prompt(page_type: str, layout: str, has_chart: bool) -> str:
    """Build tailored system prompt from modular core/ files — no regex, no full-document parse.

    File structure:
        core/always/     → loaded for every slide
        core/by_type/    → loaded based on page_type
        core/by_feature/ → loaded conditionally (charts)
        core/by_layout/  → per-layout file

    A typical content slide gets ~6-8K chars instead of the full 14K design-system.md.
    """
    import re

    _CORE = os.path.join(BASE_DIR, "resources", "prompts", "core")

    def _read(*parts: str) -> str | None:
        p = os.path.join(_CORE, *parts)
        if os.path.exists(p):
            with open(p, encoding="utf-8") as f:
                return f.read()
        return None

    parts: list[str] = []

    # ── Always (7 files, ~3K chars total) ──
    for name in ("identity", "iron-laws", "colors", "color-semantics", "structure",
                  "richness", "typography", "consistency", "checklist",
                  "format-spec"):
        text = _read("always", f"{name}.md")
        if text:
            parts.append(text)

    # ── By type: card design depth (skip for bookend pages) ──
    no_card_types = {"cover", "quote", "section", "summary", "closing"}
    if page_type not in no_card_types:
        text = _read("by_type", "cards.md")
        if text:
            parts.append(text)

    # ── By type: card role directory (multi-card layouts only) ──
    multi_card = {"two_column", "two_column_asymmetric", "three_column",
                  "hero_grid", "mixed_grid", "dashboard", "horizontal_split"}
    if layout in multi_card:
        text = _read("by_type", "card-roles.md")
        if text:
            parts.append(text)

    # ── By type: decoration (extract relevant sub-section) ──
    deco = _read("by_type", "decoration.md")
    if deco:
        is_bookend = page_type in {"cover", "summary", "closing", "section"}
        deco_kw = "封面/总结页装饰" if is_bookend else "内容页装饰"
        deco_match = re.search(
            rf'(### {re.escape(deco_kw)}.*?)(?=### |\Z)',
            deco, re.DOTALL
        )
        if deco_match:
            parts.append(f"## 五、装饰系统\n{deco_match.group(1)}")
        else:
            parts.append(deco)

    # ── By type: illustration ──
    text = _read("by_type", "illustration.md")
    if text:
        parts.append(text)

    # ── By feature: charts ──
    if has_chart:
        text = _read("by_feature", "charts.md")
        if text:
            parts.append(text)

    # ── By layout: single per-layout file ──
    text = _read("by_layout", f"{layout}.md")
    if text:
        parts.append(f"## 当前布局\n{text}")

    return "\n\n".join(parts)


# ── Core prompt file loaders ──

def _load_research_prompt() -> str:
    """Load Phase 2 research system prompt from core/research.md."""
    p = os.path.join(BASE_DIR, "resources", "prompts", "core", "research.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return "你是专业的内容研究员。对提供的 SOP 文档进行深度结构化分析。"


def _load_outline_rules() -> str:
    """Load outline generation rules from core/outline-rules.md."""
    p = os.path.join(BASE_DIR, "resources", "prompts", "core", "outline-rules.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def _load_fill_content_prompt() -> str:
    """Load Phase 2 fill content system prompt from core/fill-content.md."""
    p = os.path.join(BASE_DIR, "resources", "prompts", "core", "fill-content.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return "你是内容编辑专家。根据 SOP 文章和金字塔原理为指定幻灯片填充正文内容。"


def _load_text_to_json_prompt() -> str:
    """Load text-to-JSON conversion prompt template from core/text-to-json.md."""
    p = os.path.join(BASE_DIR, "resources", "prompts", "core", "text-to-json.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return "你是数据整理专家。你的唯一任务是将人类编辑的自然文本转回结构化 JSON。"


def _load_cards_system_prompt() -> str:
    """Load cards generation system prompt from core/cards-system.md."""
    p = os.path.join(BASE_DIR, "resources", "prompts", "core", "cards-system.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def _load_structure_output_prompt() -> str:
    """Load structure planning output format from core/structure-output.md."""
    p = os.path.join(BASE_DIR, "resources", "prompts", "core", "structure-output.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def _load_html_output_prompt(column_id: str = "") -> str:
    """Load HTML output format instructions — per-column override supported.

    Resolution order (same pattern as _load_scenario_file):
    1) scenarios/{column_id}/html-output.md  (per-column override)
    2) scenarios/_default/html-output.md     (shared default)
    3) prompts/core/html-output.md           (legacy fallback)
    """
    # 1) Per-column custom file
    if column_id:
        col_file = os.path.join(SCENARIOS_DIR, column_id, "html-output.md")
        if os.path.exists(col_file):
            with open(col_file, "r", encoding="utf-8") as f:
                return f.read()
    # 2) Default scenario template
    default_file = os.path.join(SCENARIOS_DIR, "_default", "html-output.md")
    if os.path.exists(default_file):
        with open(default_file, "r", encoding="utf-8") as f:
            return f.read()
    # 3) Legacy core prompt (fallback — keep for backward compatibility)
    core_file = os.path.join(BASE_DIR, "resources", "prompts", "core", "html-output.md")
    if os.path.exists(core_file):
        with open(core_file, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def _load_edit_agent_prompt() -> str:
    """Load edit agent system prompt template from core/edit-agent.md."""
    p = os.path.join(BASE_DIR, "resources", "prompts", "core", "edit-agent.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return f.read()
    return ""


def _build_edit_system_prompt(style_id: str, color_scheme: str = "deep-blue") -> str:
    """Build edit-slide system prompt — agent with knowledge lookup tool.

    Architecture:
      1. Agent identity — who it is
      2. Knowledge index — topics available for lookup
      3. Working object — the HTML slide
      4. Workflow — discuss → lookup knowledge → confirm → execute
    """
    persona = _load_style_prompt(style_id, color_scheme)
    if not persona:
        persona = "你是演示文稿设计艺术总监。"

    agent_tmpl = _load_edit_agent_prompt()
    if not agent_tmpl:
        return f"""## 角色

{persona}

你的工作对象是用户提供的单页幻灯片 HTML。你的能力：
- 分析当前设计是否符合设计系统规范
- 使用工具查阅具体的设计规范
- 与用户讨论改进方案
- 在用户确认后输出修改后的 HTML

## 可用工具

你有以下工具可以使用：
1. **list_knowledge_topics** — 列出所有可查询的设计知识主题
2. **lookup_knowledge** — 按主题名查询具体的设计规范内容

## 工作流

1. **接收用户消息** → 判断意图：讨论 / 简单明确修改 / 模糊修改需求
2. **简单明确修改**（如"标题改42px""删第三段"）→ 直接执行，不需要查知识库
3. **模糊需求/涉及设计规范**（如"颜色感觉不对""这里布局合理吗"）→ 先用工具查相关规范，再分析回复
4. **用户确认**（如"改""做吧""就按这个"）→ 输出修改后的完整 HTML

## 输出格式

- 讨论/分析 → 纯文本
- 执行修改 → 完整 slide HTML，用 ```html ... ``` 包裹

## 约束

- 画布: 1280×720px，全部内联样式
- 内容区基准边距: left:60px; right:60px，禁止修改
- 顶部 accent 色条、标题短线、页码标记属于页面框架，禁止删除
- 卡片容器必须包含 overflow:hidden
- 修改范围不超过用户要求，风格匹配已有元素"""

    return agent_tmpl.format(persona=persona)


async def _run_edit_agent(provider_id: str, model: str,
                           system_prompt: str, user_message: str,
                           temperature: float = 0.3) -> str:
    """Run the edit agent with knowledge-lookup tools.

    The agent can call list_knowledge_topics and lookup_knowledge to
    consult design rules before responding.  Supports Anthropic native
    and OpenAI-compatible APIs.
    """
    from services.llm_service import _is_anthropic, _mk_anthropic
    from openai import AsyncOpenAI
    from database import get_db

    db = get_db()
    provider = None
    try:
        row = db.execute(
            "SELECT * FROM llm_providers WHERE id=? AND is_enabled=1",
            (provider_id,)
        ).fetchone()
        if row:
            provider = dict(row)
    finally:
        db.close()

    if not provider:
        raise ValueError(f"Provider {provider_id} not found or disabled")

    # ── Tool definitions (Anthropic format) ──
    tools = [
        {
            "name": "list_knowledge_topics",
            "description": "列出所有可用的设计知识主题",
            "input_schema": {
                "type": "object",
                "properties": {},
                "required": []
            }
        },
        {
            "name": "lookup_knowledge",
            "description": "按主题名查询具体的设计规范内容",
            "input_schema": {
                "type": "object",
                "properties": {
                    "topic": {
                        "type": "string",
                        "description": "知识主题名称，如 cards, colors, typography, structure"
                    }
                },
                "required": ["topic"]
            }
        }
    ]

    # OpenAI format tools
    openai_tools = [{
        "type": "function",
        "function": {
            "name": t["name"],
            "description": t["description"],
            "parameters": t["input_schema"]
        }
    } for t in tools]

    def _exec_tool(name: str, args: dict) -> str:
        if name == "list_knowledge_topics":
            return _list_knowledge_topics()
        elif name == "lookup_knowledge":
            topic = args.get("topic", "")
            result = _lookup_knowledge(topic)
            return result if result else f"(未找到知识主题: {topic})"
        return "(未知工具)"

    max_rounds = 5

    if _is_anthropic(provider):
        client = _mk_anthropic(provider)
        messages = [{"role": "user", "content": user_message}]

        for _ in range(max_rounds):
            response = await client.messages.create(
                model=model,
                max_tokens=16384,
                system=system_prompt,
                messages=messages,
                tools=tools,
                temperature=temperature,
            )

            # Collect text and tool_use blocks
            text_blocks = []
            tool_blocks = []
            for block in response.content:
                if block.type == "text":
                    text_blocks.append(block.text)
                elif block.type == "tool_use":
                    tool_blocks.append(block)

            if tool_blocks:
                # Add assistant response to messages
                messages.append({"role": "assistant", "content": [b.to_dict() for b in response.content]})
                # Execute tools and add results
                tool_results = []
                for tb in tool_blocks:
                    result = _exec_tool(tb.name, tb.input)
                    tool_results.append({
                        "type": "tool_result",
                        "tool_use_id": tb.id,
                        "content": result,
                    })
                messages.append({"role": "user", "content": tool_results})
            else:
                return "\n".join(text_blocks).strip()

        return "\n".join(text_blocks).strip() if text_blocks else ""

    else:
        # ── OpenAI-compatible path ──
        client = AsyncOpenAI(
            api_key=provider["api_key"],
            base_url=provider["base_url"],
            timeout=120.0,
        )
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_message},
        ]

        for _ in range(max_rounds):
            response = await client.chat.completions.create(
                model=model,
                messages=messages,
                tools=openai_tools,
                temperature=temperature,
                max_tokens=16384,
            )
            msg = response.choices[0].message

            if msg.tool_calls:
                tool_calls = msg.tool_calls
                messages.append({"role": "assistant", "content": msg.content or "", "tool_calls": [tc.to_dict() for tc in tool_calls]})
                for tc in tool_calls:
                    args = json.loads(tc.function.arguments)
                    result = _exec_tool(tc.function.name, args)
                    messages.append({"role": "tool", "tool_call_id": tc.id, "content": result})
                    # Some providers (e.g. DeepSeek) require the role name "tool"
            else:
                return msg.content.strip() if msg.content else ""

        return msg.content.strip() if msg.content else ""


def _list_knowledge_topics() -> str:
    """Scan the core/ directory and return available knowledge topics."""
    _CORE = os.path.join(BASE_DIR, "resources", "prompts", "core")
    if not os.path.isdir(_CORE):
        return "(知识库目录不存在)"
    lines: list[str] = []
    for root, dirs, files in os.walk(_CORE):
        for f in sorted(files):
            if f.endswith(".md"):
                rel = os.path.relpath(os.path.join(root, f), _CORE)
                path = os.path.join(root, f)
                try:
                    with open(path, encoding="utf-8") as fh:
                        first = fh.readline().strip()
                        if first.startswith("## "):
                            desc = first[3:]
                        elif first.startswith("# "):
                            desc = first[2:]
                        else:
                            desc = ""
                except Exception:
                    desc = ""
                lines.append(f"- `{rel.replace(os.sep, '/').replace('.md', '')}`: {desc}")
    return "\n".join(lines)


def _lookup_knowledge(topic: str) -> str | None:
    """Read a knowledge file by topic name or path fragment."""
    _CORE = os.path.join(BASE_DIR, "resources", "prompts", "core")
    if not os.path.isdir(_CORE):
        return None
    q = topic.lower().strip().replace(" ", "-").replace("_", "-")
    for root, dirs, files in os.walk(_CORE):
        for f in files:
            if f.endswith(".md"):
                name = f.replace(".md", "").lower().replace("_", "-")
                rel = os.path.relpath(os.path.join(root, f), _CORE).replace(os.sep, "/").lower().replace(".md", "").replace("_", "-")
                if q == name or q in name or q in rel or name in q:
                    path = os.path.join(root, f)
                    with open(path, encoding="utf-8") as fh:
                        return fh.read()
    return None


def _fix_llm_html_errors(html: str, is_a4: bool = False) -> str:
    """Fix common LLM HTML generation mistakes.

    1. Missing px units on CSS position/size properties (e.g., top:28 → top:28px)
    2. <!DOCTYPE html> or <html>/<head>/<body> tags inside slide content
    3. Truncated content ending mid-tag
    """
    import re

    # ── 1. Add missing px units ──
    # CSS properties that require length units for non-zero values.
    # Careful: don't add px to z-index, opacity, font-weight, flex, order, etc.
    dim_props = r'(?:top|right|bottom|left|width|height|min-width|max-width|min-height|max-height|margin|margin-top|margin-right|margin-bottom|margin-left|padding|padding-top|padding-right|padding-bottom|padding-left|font-size|border-radius|gap|row-gap|column-gap|letter-spacing|border-width|outline-width|text-indent|word-spacing)'
    # Match: prop: NUMBER (not followed by a unit or another digit)
    html = re.sub(
        rf'({dim_props})\s*:\s*(\d+)\s*(?=[;!"\'>\s])',
        r'\1: \2px',
        html
    )

    # ── 2. Strip DOCTYPE and outer HTML structure if LLM generated a full page ──
    html = re.sub(r'<!DOCTYPE\s+html[^>]*>', '', html, flags=re.IGNORECASE)
    html = re.sub(r'</?html[^>]*>', '', html, flags=re.IGNORECASE)
    html = re.sub(r'</?head[^>]*>', '', html, flags=re.IGNORECASE)
    html = re.sub(r'</?body[^>]*>', '', html, flags=re.IGNORECASE)
    # Remove <title>, <meta>, <link> tags the LLM might generate
    html = re.sub(r'<title[^>]*>.*?</title>', '', html, flags=re.IGNORECASE | re.DOTALL)
    html = re.sub(r'<meta[^>]*>', '', html, flags=re.IGNORECASE)
    html = re.sub(r'<link[^>]*>', '', html, flags=re.IGNORECASE)

    # ── 3. Fix truncated ending (content ending mid-tag) ──
    # If the HTML ends with an unclosed tag or mid-attribute, remove the fragment
    last_gt = html.rfind('>')
    last_lt = html.rfind('<')
    if last_lt > last_gt:
        # Content ends inside a tag — truncate to last complete tag
        html = html[:last_gt + 1]

    # ── 4. Detect truncated style attributes (e.g., style="font-size</div>) ──
    # The last style=" in the HTML should have a closing " before the next >
    last_style = html.rfind('style="')
    if last_style > 0:
        after = html[last_style + 7:]
        close_quote = after.find('"')
        next_gt = after.find('>')
        # If the quote never closes, or > comes before ", the style is broken
        if close_quote < 0:
            html = html[:last_style]
        elif next_gt > 0 and next_gt < close_quote:
            # "font-size</div> — style ends without value or closing quote
            html = html[:last_style]

    # ── 5. Fix accent bar border-radius mismatch ──
    # A 4px-wide left-edge strip with its own border-radius can never
    # geometrically match the parent card's corner curves.
    html = re.sub(
        r'border-radius:\s*\d+px\s+0\s+0\s+\d+px\s*;?',
        '',
        html
    )
    # Fix card containers: position:relative is required so absolute-positioned
    # accent bars inside the card use the card as their containing block.
    # overflow:hidden clips them to the card's rounded corners.
    def _fix_card_container(m):
        tag = m.group(0)
        if 'position:' not in tag and 'background:#1e293b;' in tag:
            tag = tag.replace(
                'background:#1e293b;',
                'position:relative;background:#1e293b;'
            )
        if 'overflow:' not in tag:
            tag = tag.replace('border-radius:', 'overflow:hidden;border-radius:')
        return tag
    html = re.sub(
        r'<div[^>]*border-radius:\s*\d+px[^>]*>',
        _fix_card_container,
        html
    )

    # ── 6. Enforce minimum font-size ──
    # PPT (1280x720): minimum 14px — anything smaller is unreadable
    # A4 (794x1123): minimum 12px — VI template uses 12px for body, 10px for header/footer
    _min_fs = 12 if is_a4 else 14
    def _bump_font_size(m):
        sz = int(m.group(1))
        if sz < _min_fs:
            return f'font-size:{_min_fs}px'
        return m.group(0)
    html = re.sub(r'font-size:(\d+)px', _bump_font_size, html)

    # ── 7. Fix dangerous line-height: 1px on visible elements ──
    # AI sometimes generates line-height:1px on decorative text, causing overlap.
    # Replace with a safe minimum. Only fix when the same element has visible font-size.
    def _fix_line_height(m):
        tag = m.group(0)
        # Check if the same style attr has a font-size that's visible (>20px)
        fs_match = re.search(r'font-size:\s*(\d+)px', tag)
        if fs_match and int(fs_match.group(1)) > 20:
            # Large decorative text: line-height:1 is safe (prevents overlap)
            return re.sub(r'line-height:\s*1px', 'line-height:1', tag)
        else:
            # Small text or unknown: use safe 1.2
            return re.sub(r'line-height:\s*1px', 'line-height:1.2', tag)
    html = re.sub(
        r'<[^>]*line-height:\s*1px[^>]*>',
        _fix_line_height,
        html
    )

    # ── 8. Enforce minimum margin-top (PPT only: 1280x720 canvas needs larger spacing) ──
    if not is_a4:
        html = re.sub(r'margin-top:\s*4px', 'margin-top:12px', html)
        html = re.sub(r'margin-top:\s*6px', 'margin-top:12px', html)

    # ── 9. Fix SVG icon margin-bottom too small ──
    # Icons before numbers need ≥20px gap
    def _fix_icon_margin(m):
        tag = m.group(0)
        mb_match = re.search(r'margin-bottom:\s*(\d+)px', tag)
        if mb_match and int(mb_match.group(1)) < 20:
            tag = tag.replace(mb_match.group(0), 'margin-bottom:20px')
        return tag
    html = re.sub(
        r'<svg[^>]*margin-bottom:\s*\d+px[^>]*>',
        _fix_icon_margin,
        html
    )

    return html


# ═══════════════════════════════════════════════════════════════════
# Per-slide file storage — eliminates fragile HTML string splicing
# ═══════════════════════════════════════════════════════════════════

SLIDES_DIR = "slides"


def _save_slide_files(run_dir: str, slides: list):
    """Save each slide's resolved and unresolved HTML as individual files.

    Directory layout:
      {run_dir}/slides/slide_01.html        (resolved: var(--primary))
      {run_dir}/slides/slide_01_vars.html   (unresolved: {{primary}})
      ...
    """
    import os as _os
    slides_dir = _os.path.join(run_dir, SLIDES_DIR)
    _os.makedirs(slides_dir, exist_ok=True)
    for s in slides:
        seq = s.get("seq", 0)
        html = s.get("html", "")
        html_vars = s.get("html_vars", html)
        if html:
            path = _os.path.join(slides_dir, f"slide_{seq:02d}.html")
            with open(path, "w", encoding="utf-8") as f:
                f.write(html)
        if html_vars:
            path = _os.path.join(slides_dir, f"slide_{seq:02d}_vars.html")
            with open(path, "w", encoding="utf-8") as f:
                f.write(html_vars)


def _load_slide_files(run_dir: str) -> list[dict]:
    """Load all individual slide HTML files from a run directory.

    Returns a list of dicts with keys: seq, html, html_vars.
    Slides are sorted by seq.
    """
    import os as _os
    import re as _re
    slides_dir = _os.path.join(run_dir, SLIDES_DIR)
    if not _os.path.isdir(slides_dir):
        return []
    slides: dict[int, dict] = {}
    for fname in _os.listdir(slides_dir):
        m = _re.match(r"slide_(\d+)\.html$", fname)
        if not m:
            continue
        seq = int(m.group(1))
        if seq not in slides:
            slides[seq] = {"seq": seq}
        path = _os.path.join(slides_dir, fname)
        with open(path, "r", encoding="utf-8") as f:
            slides[seq]["html"] = f.read()
    # Also load vars versions
    for fname in _os.listdir(slides_dir):
        m = _re.match(r"slide_(\d+)_vars\.html$", fname)
        if not m:
            continue
        seq = int(m.group(1))
        if seq not in slides:
            slides[seq] = {"seq": seq}
        path = _os.path.join(slides_dir, fname)
        with open(path, "r", encoding="utf-8") as f:
            slides[seq]["html_vars"] = f.read()
    return [slides[k] for k in sorted(slides.keys())]


def _extract_slides_from_html(full_html: str) -> list[dict]:
    """Extract individual slide HTML from a full assembled HTML document.

    Handles both slide-wrapper (div) and section-based formats.
    Returns list of dicts with keys: seq, html.
    """
    import re as _re
    slides = []

    # Try slide-wrapper format
    wrapper_pat = _re.compile(
        r'<div\s+class="slide-wrapper"\s+data-seq="(\d+)"\s*>(.*?)</div>\s*(?=<div\s+class="slide-wrapper"|</body>|</html>|$)',
        _re.DOTALL
    )
    for m in wrapper_pat.finditer(full_html):
        seq = int(m.group(1))
        inner = m.group(2).strip()
        # Inner content is the actual slide HTML (may be wrapped in another div)
        slides.append({"seq": seq, "html": inner})

    if not slides:
        # Try section-based format
        section_pat = _re.compile(
            r'<section\s+class="slide"\s+data-seq="(\d+)"[^>]*>(.*?)</section>',
            _re.DOTALL
        )
        for m in section_pat.finditer(full_html):
            seq = int(m.group(1))
            slides.append({"seq": seq, "html": m.group(2).strip()})

    slides.sort(key=lambda s: s["seq"])
    return slides


def _build_root_vars(scheme_data: dict) -> str:
    """Build :root CSS block defining all color variables with actual hex/rgb values."""
    lines = [":root {"]

    color_keys = ["primary", "secondary", "accent", "background", "text", "card_bg"]
    for key in color_keys:
        val = scheme_data.get(key, "")
        if val and val.startswith("#"):
            lines.append(f"  --{key}: {val};")
            r, g, b = _hex_to_rgb(val)
            lines.append(f"  --{key}-r: {r};")
            lines.append(f"  --{key}-g: {g};")
            lines.append(f"  --{key}-b: {b};")
            lines.append(f"  --{key}-rgb: {r}, {g}, {b};")

    chart_colors = scheme_data.get("chart_colors", [])
    if isinstance(chart_colors, list):
        for i, c in enumerate(chart_colors):
            if c and c.startswith("#"):
                lines.append(f"  --chart-{i}: {c};")
                r, g, b = _hex_to_rgb(c)
                lines.append(f"  --chart-{i}-rgb: {r}, {g}, {b};")

    semantic = scheme_data.get("semantic", {})
    if isinstance(semantic, dict):
        for k, v in semantic.items():
            if v and v.startswith("#"):
                lines.append(f"  --semantic-{k}: {v};")
                r, g, b = _hex_to_rgb(v)
                lines.append(f"  --semantic-{k}-rgb: {r}, {g}, {b};")

    lines.append("}")
    return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════════
# A4 overflow pagination engine
# ═══════════════════════════════════════════════════════════════════


def _estimate_block_height(block_html: str) -> int:
    """Estimate vertical pixel height of an HTML content block."""
    import re as _re
    h = 0

    # Explicit heights
    for hm in _re.findall(r'height:\s*(\d+)px', block_html):
        h += int(hm)

    # Table rows
    rows = len(_re.findall(r'<tr\b', block_html))
    if rows > 0:
        pad = 14
        pm = _re.search(r'padding(?:-top|-bottom)?:\s*(\d+)px', block_html)
        if pm:
            pad = int(pm.group(1))
        h = max(h, rows * (pad * 2 + 16))

    # Text content (rough)
    text = _re.sub(r'<[^>]+>', '', block_html).strip()
    if text:
        line_h = 24
        lines = max(1, len(text) / 40)
        h = max(h, int(lines * line_h))

    # Margins
    for m in _re.findall(r'margin-(?:top|bottom):\s*(\d+)px', block_html):
        h += int(m)
    for p in _re.findall(r'padding(?:-top|-bottom)?:\s*(\d+)px', block_html):
        h += int(p)

    return max(h, 40)


def _extract_content_blocks(inner_html: str) -> list[str]:
    """Extract top-level content blocks (tables, divs, etc.) from HTML.

    Uses unified depth-aware parsing: finds the next opening tag, determines
    its tag name, then traces <tagname / </tagname> depth until the matching
    closing tag is found. Handles nested elements of the same type correctly.
    """
    blocks = []
    pos = 0
    n = len(inner_html)

    # Void/self-closing elements that never have children
    _void_tags = {'br', 'hr', 'img', 'input', 'meta', 'link', 'area',
                  'base', 'col', 'embed', 'source', 'track', 'wbr'}

    while pos < n:
        # Skip whitespace
        while pos < n and inner_html[pos] in ' \t\n\r':
            pos += 1
        if pos >= n:
            break

        if inner_html[pos] != '<':
            # Text node — collect until next <
            end = inner_html.find('<', pos)
            if end < 0:
                blocks.append(inner_html[pos:])
                break
            blocks.append(inner_html[pos:end])
            pos = end
            continue

        # Extract tag name: letters/digits/hyphens after <
        j = pos + 1
        while j < n and (inner_html[j].isalnum() or inner_html[j] == '-'):
            j += 1
        if j == pos + 1:
            # Not a regular tag (e.g. <!--, <!doctype)
            pos += 1
            continue
        tag_name = inner_html[pos + 1:j].lower()

        # Void elements — extract self-closing block
        if tag_name in _void_tags:
            gt = inner_html.find('>', pos)
            if gt > 0:
                blocks.append(inner_html[pos:gt + 1])
                pos = gt + 1
            else:
                pos += 1
            continue

        # Container element — depth-aware extraction
        start = pos
        depth = 0
        open_pat = f'<{tag_name}'
        close_pat = f'</{tag_name}>'

        while pos < n:
            next_open = inner_html.find(open_pat, pos)
            next_close = inner_html.find(close_pat, pos)

            if next_close == -1:
                pos = n
                break

            if next_open != -1 and next_open < next_close:
                depth += 1
                pos = next_open + len(tag_name) + 1
            else:
                depth -= 1
                if depth == 0:
                    pos = next_close + len(close_pat)
                    blocks.append(inner_html[start:pos])
                    break
                pos = next_close + len(close_pat)

        if depth != 0 and pos >= n:
            # Unclosed container — include rest as one block
            blocks.append(inner_html[start:])
            break

    return blocks


def _split_a4_html_content(html: str, content_max_h: int) -> list[str]:
    """Split an A4 page that overflows its content area into multiple pages.

    Finds the main content container (position:absolute div with overflow-y),
    splits its child blocks across pages, and rebuilds the HTML structure
    for each page. Returns list of complete page HTML strings.
    """
    import re as _re

    # Remove overflow scrolling
    html = _re.sub(r'overflow(?:-y)?:\s*(?:auto|scroll)\s*;?', '', html)

    # Find the content container — a positioned div that holds the bulk of content
    # Pattern: position:absolute with top: + bottom: or overflow-y that we just removed
    content_re = _re.compile(
        r'(<div[^>]*position:\s*absolute[^>]*'
        r'(?:top:\s*\d+px[^>]*(?:bottom|right)[^>]*|'
        r'bottom:\s*\d+px[^>]*top:\s*\d+px[^>]*)'
        r'[^>]*>)'
    )
    matches = list(content_re.finditer(html))

    if not matches:
        return [html]

    # Pick the content container with the most inner content, not the
    # last regex match (which might be a decorative accent bar).
    content_match = matches[0]
    content_match_end = -1
    best_size = -1
    for m in matches:
        m_inner_start = m.end()
        m_depth = 0
        m_pos = m.start()
        m_end = -1
        while m_pos < len(html):
            no = html.find('<div', m_pos)
            nc = html.find('</div>', m_pos)
            if nc == -1:
                break
            if no != -1 and no < nc:
                m_depth += 1
                m_pos = no + 4
            else:
                m_depth -= 1
                if m_depth == 0:
                    m_end = nc
                    break
                m_pos = nc + 6
        if m_end > 0:
            inner_size = m_end - m_inner_start
            if inner_size > best_size:
                best_size = inner_size
                content_match = m
                content_match_end = m_end

    content_tag_start = content_match.start()
    content_inner_start = content_match.end()

    # Use the pre-computed end from the best-match scan above, or find it now
    if content_match_end > 0:
        content_end = content_match_end
    else:
        depth = 0
        pos = content_match.start()
        content_end = -1
        while pos < len(html):
            no = html.find('<div', pos)
            nc = html.find('</div>', pos)
            if nc == -1:
                break
            if no != -1 and no < nc:
                depth += 1
                pos = no + 4
            else:
                depth -= 1
                if depth == 0:
                    content_end = nc
                    break
                pos = nc + 6

    if content_end < 0:
        return [html]

    inner = html[content_inner_start:content_end]
    blocks = _extract_content_blocks(inner)

    # Detect single oversized block (common: large <table> with many rows)
    _table_open = ""
    if len(blocks) <= 1:
        if blocks and blocks[0].strip().startswith('<table'):
            solo = blocks[0]
            rows = _re.findall(r'<tr[^>]*>.*?</tr>', solo, _re.DOTALL)
            if len(rows) > 1:
                table_open_m = _re.match(r'<table[^>]*>', solo)
                _table_open = table_open_m.group(0) if table_open_m else '<table>'
                blocks = rows  # Split into row-level blocks
            else:
                return [html]
        else:
            return [html]

    # Estimate heights and group into pages
    heights = [_estimate_block_height(b) for b in blocks]
    total_est = sum(heights)
    if total_est <= content_max_h * 1.1:
        return [html]

    pages = []
    current_blocks = []
    current_h = 0

    for block, h in zip(blocks, heights):
        if current_h + h > content_max_h and current_blocks:
            pages.append(current_blocks)
            current_blocks = [block]
            current_h = h
        else:
            current_blocks.append(block)
            current_h += h

    if current_blocks:
        pages.append(current_blocks)

    if len(pages) <= 1:
        return [html]

    before = html[:content_inner_start]
    after = html[content_end:]

    result = []
    for page_blocks in pages:
        page_inner = '\n'.join(page_blocks)
        if _table_open:
            page_inner = _table_open + '\n' + page_inner + '\n</table>'
        result.append(before + page_inner + after)

    return result


def _preprocess_a4_slides(slides: list, canvas_w: int, canvas_h: int) -> list:
    """Pre-process slides for A4/portrait mode: split overflow pages.

    Scans non-cover slides for content overflow, splits overflowing pages
    at block boundaries, and re-numbers all slides. Returns a new list
    with overflow continuations inserted.

    This runs BEFORE header/footer injection — only content splitting.
    """
    if canvas_h <= canvas_w:
        return slides

    content_max_h = canvas_h - 80  # header + footer

    result = []
    for s in slides:
        html = s.get("html", "")
        if not html:
            result.append(s)
            continue

        seq = s.get("seq", 1)
        if seq == 1:
            result.append(s)
            continue

        has_overflow = bool(re.search(r'overflow(?:-y)?:\s*(?:auto|scroll)', html, re.IGNORECASE))
        if not has_overflow:
            result.append(s)
            continue

        pages = _split_a4_html_content(html, content_max_h)
        for j, page_html in enumerate(pages):
            result.append({**s, "html": page_html, "_overflow_split": j > 0})

    # Re-number sequentially
    for i, s in enumerate(result):
        s["seq"] = i + 1

    return result


def _extract_outermost_div(html: str) -> str:
    """Extract the outermost <div> container from LLM-generated HTML.

    The LLM is instructed to output exactly one outer <div> container
    (width:...;height:...). We find the first <div and trace element
    depth to locate its matching </div>, returning a self-contained
    div tree that cannot leak into or out of the slide wrapper.

    This replaces fragile regex-based div counting that could strip
    closing tags from the wrong position when divs are structurally
    (not numerically) imbalanced.
    """
    first_open = html.find('<div')
    if first_open < 0:
        return html

    depth = 0
    pos = first_open
    while pos < len(html):
        next_open = html.find('<div', pos)
        next_close = html.find('</div>', pos)

        if next_close == -1:
            break

        if next_open != -1 and next_open < next_close:
            depth += 1
            pos = next_open + 4
        else:
            depth -= 1
            if depth == 0:
                return html[first_open:next_close + 6]
            pos = next_close + 6

    # Depth never returned to 0 — outermost <div> is unclosed.
    # Close it ourselves so the slide wrapper stays intact.
    return html[first_open:] + '</div>'


def _wcag_relative_luminance(hex_color: str) -> float:
    """WCAG 2.1 relative luminance of a hex color."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return 0.0
    r, g, b = int(h[0:2], 16) / 255.0, int(h[2:4], 16) / 255.0, int(h[4:6], 16) / 255.0
    rs = r / 12.92 if r <= 0.04045 else ((r + 0.055) / 1.055) ** 2.4
    gs = g / 12.92 if g <= 0.04045 else ((g + 0.055) / 1.055) ** 2.4
    bs = b / 12.92 if b <= 0.04045 else ((b + 0.055) / 1.055) ** 2.4
    return 0.2126 * rs + 0.7152 * gs + 0.0722 * bs


def _wcag_contrast_ratio(hex1: str, hex2: str) -> float:
    """WCAG 2.1 contrast ratio between two hex colors."""
    l1 = _wcag_relative_luminance(hex1)
    l2 = _wcag_relative_luminance(hex2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _load_cover_overrides(style_id: str) -> dict:
    """Load cover slide_type_overrides from tokens.yaml.

    Returns dict with 'card_bg' and 'text' keys (raw values from YAML,
    may be {{placeholder}} or hex). Returns empty dict on failure.
    """
    import yaml
    tokens_path = os.path.join(BASE_DIR, "resources", "vi", style_id, "tokens.yaml")
    if not os.path.exists(tokens_path):
        return {}
    try:
        with open(tokens_path, "r", encoding="utf-8") as f:
            tokens = yaml.safe_load(f.read())
        return tokens.get("slide_type_overrides", {}).get("cover", {})
    except Exception:
        return {}


def _placeholder_to_css_var(placeholder: str) -> str:
    """Resolve {{name}} → var(--name). Hex values pass through unchanged."""
    if not placeholder:
        return ""
    m = re.match(r'\{\{(\w+)\}\}', placeholder)
    if m:
        return f"var(--{m.group(1).replace('_', '-')})"
    return placeholder


def _build_cover_color_rules(style_id: str, color_scheme: str, scheme: dict) -> str:
    """Generate quantified cover color rules with WCAG contrast data for LLM prompt.

    Reads slide_type_overrides.cover from tokens.yaml, resolves colors,
    calculates WCAG AA contrast ratios, and returns a markdown rule block
    that tells the LLM exactly which CSS variables to use and WHY.
    """
    cover = _load_cover_overrides(style_id)
    if not cover:
        return ""

    bg_placeholder = cover.get("card_bg", "")
    text_placeholder = cover.get("text", "")
    if not bg_placeholder or not text_placeholder:
        return ""

    bg_var = _placeholder_to_css_var(bg_placeholder)
    text_var = _placeholder_to_css_var(text_placeholder)

    # Resolve placeholders to actual hex values from the active scheme
    def _resolve(p):
        m = re.match(r'\{\{(\w+)\}\}', p)
        if m:
            return scheme.get(m.group(1), p)
        return p

    bg_hex = _resolve(bg_placeholder)
    text_hex = _resolve(text_placeholder)

    # Calculate WCAG contrast
    cr = 0.0
    white_cr = 0.0
    if bg_hex.startswith("#") and text_hex.startswith("#"):
        cr = _wcag_contrast_ratio(bg_hex, text_hex)
        white_cr = _wcag_contrast_ratio(bg_hex, "#ffffff")

    cr_status = "PASS (>= 4.5:1)" if cr >= 4.5 else "FAIL (< 4.5:1)"
    is_light_bg = _wcag_relative_luminance(bg_hex) > 0.5 if bg_hex.startswith("#") else False

    # Build quantified rule block
    rules = f"""## 封面颜色规则（从活跃色系自动派生 — 量化标准，必须严格执行）

| 属性 | CSS 变量 | 实际色值 | WCAG 说明 |
|------|---------|---------|----------|
| 封面背景 | `background: {bg_var}` | {bg_hex} | — |
| 封面标题文字 | `color: {text_var}` | {text_hex} | 对比度 {cr:.1f}:1 → {cr_status} |
| #ffffff 在封面背景上 | — | — | 对比度仅 {white_cr:.1f}:1 {"→ 不可读，严禁使用" if white_cr < 3.0 else "→ 不满足 AA，禁止用于正文"} |

**硬约束（违反 = 废稿）：**
- 封面标题 `color` 必须是 `{text_var}` —— 严禁使用 hex 值（包括 #ffffff、{bg_hex} 等）
- 封面背景 `background` 必须是 `{bg_var}` —— 严禁改为其他 CSS 变量或 hex 值
{"- 封面背景为浅色（L > 0.5），文字必须使用深色 CSS 变量 —— 严禁 #ffffff" if is_light_bg else "- 封面背景为深色（L < 0.5），文字色由 tokens 定义"}
- 装饰文字/辅助信息使用 `rgba(var(--text-rgb), N)` 控制透明度，N 取值 0.08-0.65"""

    return rules


def _enforce_cover_rules(html: str, scheme_data: dict, style_id: str = "business") -> str:
    """Enforce VI cover rules — fully tokens-driven, single unified path for all styles.

    Zero hardcoded assumptions:
    - No branching on background type (variable vs. non-variable)
    - No forced gradient (templates control their own backgrounds)
    - No hardcoded CSS variable name lists

    Tokens say X → code enforces X. That's it.
    """
    import re as _re

    cover = _load_cover_overrides(style_id)
    if not cover:
        return html

    bg_var = _placeholder_to_css_var(cover.get("card_bg", ""))
    text_val = _placeholder_to_css_var(cover.get("text", ""))
    if not bg_var or not text_val:
        return html

    # ── Layer 1: Fix hardcoded white text → correct value from tokens ──
    # Covers all forms: 6-digit, 3-digit, lowercase, uppercase, with/without space.
    # If tokens say text should be white, these are no-ops (correct by definition).
    for hardcoded in ('#ffffff', '#FFFFFF', '#fff', '#FFF'):
        html = html.replace(f'color:{hardcoded}', f'color:{text_val}')
        html = html.replace(f'color: {hardcoded}', f'color: {text_val}')
        html = html.replace(f'color:{hardcoded};', f'color:{text_val};')
        html = html.replace(f'color: {hardcoded};', f'color: {text_val};')

    # ── Layer 2: Fix hardcoded rgba(R,G,B, → rgba(var(--text-rgb), ──
    # LLM sometimes copies the resolved hex RGB into rgba() instead of using
    # the CSS variable form. Detect from scheme_data and fix.
    text_hex = scheme_data.get("text", "")
    if text_hex.startswith("#") and len(text_hex) == 7:
        r, g, b = int(text_hex[1:3], 16), int(text_hex[3:5], 16), int(text_hex[5:7], 16)
        html = html.replace(f'rgba({r},{g},{b},', 'rgba(var(--text-rgb),')
        html = html.replace(f'rgba({r}, {g}, {b},', 'rgba(var(--text-rgb),')

    # ── Layer 3: Fix hardcoded primary RGB in rgba form ──
    # Same issue: LLM may use resolved primary RGB instead of var(--primary-rgb).
    primary_hex = scheme_data.get("primary", "")
    if primary_hex.startswith("#") and len(primary_hex) == 7:
        r, g, b = int(primary_hex[1:3], 16), int(primary_hex[3:5], 16), int(primary_hex[5:7], 16)
        html = html.replace(f'rgba({r},{g},{b},', 'rgba(var(--primary-rgb),')
        html = html.replace(f'rgba({r}, {g}, {b},', 'rgba(var(--primary-rgb),')

    # ── Layer 4: Fix var(--text) → correct text variable from tokens ──
    # Template may use var(--text) but tokens override to e.g. var(--primary) or #ffffff.
    if text_val != 'var(--text)':
        new = text_val
        html = html.replace('color:var(--text)', f'color:{new}')
        html = html.replace('color: var(--text)', f'color: {new}')
        html = html.replace('color:var(--text);', f'color:{new};')
        html = html.replace('color: var(--text);', f'color: {new};')

    # ── Layer 5: Fix background hardcoded hex → correct CSS variable ──
    # If the cover background token resolves to a scheme color, replace any
    # hardcoded-hex version of it with the CSS variable form.
    m = _re.match(r'\{\{(\w+)\}\}', cover.get("card_bg", ""))
    if m:
        bg_hex = scheme_data.get(m.group(1), "")
        if bg_hex.startswith("#"):
            html = html.replace(f'background:{bg_hex}', f'background:{bg_var}')
            html = html.replace(f'background: {bg_hex}', f'background: {bg_var}')

    return html


def _strip_page_numbers(html: str) -> str:
    """Strip AI-generated page number divs using depth-balanced matching.

    Replaces the fragile regex approach (which could match wrong divs and
    create div imbalance by only stripping one </div> from a nested tree).
    """
    # Find divs that look like page numbers: bottom+right positioning + NN / MM text
    import re as _re_pn
    pn_pattern = _re_pn.compile(
        r'<div\b[^>]*position:\s*absolute[^>]*bottom:\s*\d+px[^>]*right:\s*\d+px[^>]*>'
    )
    page_num_text = _re_pn.compile(r'\d{1,2}\s*/\s*\d{1,2}')

    search_from = 0
    while True:
        m = pn_pattern.search(html, search_from)
        if not m:
            break
        open_start = m.start()
        # Check if this div contains page number text within reasonable range
        look_ahead = html[open_start:open_start + 600]
        if not page_num_text.search(look_ahead):
            search_from = open_start + 1
            continue
        # Depth-balanced removal: trace from opening <div to matching </div>
        depth = 1
        pos = open_start + len('<div')
        while pos < len(html) and depth > 0:
            nxt = html.find('<div', pos)
            end = html.find('</div>', pos)
            if end == -1:
                break
            if nxt != -1 and nxt < end:
                depth += 1
                pos = nxt + 4
            else:
                depth -= 1
                if depth == 0:
                    # Remove the entire balanced div (including trailing newline)
                    while open_start > 0 and html[open_start - 1] in ('\n', '\r'):
                        open_start -= 1
                    html = html[:open_start] + html[end + 6:]
                    search_from = open_start  # restart search from removal point
                    break
                pos = end + 6
        else:
            search_from = open_start + 1
    return html


def _assemble_html_deck(slides: list, title: str = "Presentation",
                        style_id: str = "business", scheme_data: dict = None,
                        total_slides: int = None,
                        canvas_w: int = 1280, canvas_h: int = 720) -> str:
    """Wrap individual slide HTML sections into a complete HTML document.

    Also injects unified page numbers and strips AI-generated ones.
    """
    import re

    # Pre-process A4/portrait slides: split overflow pages before assembly
    slides = _preprocess_a4_slides(slides, canvas_w, canvas_h)

    valid_slides = [s for s in slides if s.get("html")]
    total = total_slides or max((s.get("seq", 0) for s in valid_slides), default=len(valid_slides))

    wrapped_parts = []
    for i, s in enumerate(slides):
        html = s.get("html", "")
        if not html:
            continue
        slide_num = s.get("seq", i + 1)

        # Strip AI-generated page number divs (depth-balanced — no regex div leakage)
        html = _strip_page_numbers(html)

        # ═══ Extract outermost div: prevents cross-slide DOM corruption ═══
        # The LLM is told to output exactly one outer <div>. We extract it
        # by tracing depth from the first <div> to its matching </div>.
        # This is structural — a slide's internal divs cannot leak into or
        # close the wrapper, even if the LLM outputs extra or missing tags.
        html = _extract_outermost_div(html)

        # Inject unified page number (skip cover slide; only for landscape/PPT:
        # portrait/A4 documents use header/footer page numbers defined in VI)
        if slide_num > 1 and canvas_w >= canvas_h:
            pn_tag = (
                f'<div style="position:absolute;bottom:16px;right:48px;'
                f'background:{{{{primary}}}};padding:5px 14px;'
                f'border-radius:16px;z-index:50;box-shadow:0 0 0 2px rgba(0,0,0,0.3);">'
                f'<span style="font-size:14px;color:{{{{background}}}};'
                f'font-weight:500;letter-spacing:0.3px;">'
                f'{slide_num:02d} / {total:02d}</span>'
                f'</div>'
            )
            # Insert before the outermost closing tag (prefer </section>)
            last_section = html.rfind('</section>')
            if last_section > 0:
                html = html[:last_section] + pn_tag + '\n' + html[last_section:]
            else:
                # Trace depth from first <div to find TRUE outermost </div>
                first_open = html.find('<div')
                if first_open >= 0:
                    depth = 0
                    pos = first_open
                    outer_close = -1
                    while pos < len(html):
                        no = html.find('<div', pos)
                        nc = html.find('</div>', pos)
                        if nc == -1:
                            break
                        if no != -1 and no < nc:
                            depth += 1
                            pos = no + 4
                        else:
                            depth -= 1
                            if depth == 0:
                                outer_close = nc
                                break
                            pos = nc + 6
                    if outer_close > 0:
                        html = html[:outer_close] + pn_tag + '\n' + html[outer_close:]
                    else:
                        # Outer <div> not properly closed — close it, insert before
                        html = html.rstrip() + '\n' + pn_tag + '\n</div>'
                else:
                    html = html + '\n' + pn_tag

        # ── Branding variable substitution (ALL slides, code-enforced) ──
        # VI templates use {{BRAND_COPYRIGHT}} / {{BRAND_SIGNATURE}} placeholders.
        # Code replaces them with DB values — no structural injection needed.
        _copyright_str, _sig_str = _load_branding()
        import html as _html_escape
        html = html.replace('{{BRAND_COPYRIGHT}}', _html_escape.escape(_copyright_str))
        html = html.replace('{{BRAND_SIGNATURE}}', _html_escape.escape(_sig_str))

        # Cover slide: code-enforced design rules (VI §cover)
        # LLM cannot be trusted to follow cover rules — code enforces them.
        if slide_num == 1 and scheme_data:
            html = _enforce_cover_rules(html, scheme_data, style_id)
            # Cover font-size: upscale subtitle/metadata 12px→14px in content area
            # (above bottom branding div)
            _btm = html.rfind("position:absolute;bottom:")
            if _btm > 0:
                _content_area = html[:_btm]
                _branding_area = html[_btm:]
                import re as _re_cv
                _content_area = _re_cv.sub(
                    r'<(?:td|div)\b[^>]*font-size\s*:\s*12px[^>]*>',
                    lambda m: m.group(0).replace('font-size:12px', 'font-size:14px')
                              .replace('font-size: 12px', 'font-size: 14px'),
                    _content_area
                )
                html = _content_area + _branding_area

        wrapped_parts.append(f'<div class="slide-wrapper" data-seq="{slide_num}">{html}</div>')

    wrapped = "\n".join(wrapped_parts)

    # Global div balance safety check — catches cross-slide leakage
    total_open = len(re.findall(r'<div\b', wrapped))
    total_close = len(re.findall(r'</div>', wrapped))
    if total_open != total_close:
        delta = total_open - total_close
        _logger.warning(f"_assemble_html_deck: global div imbalance ({total_open} open vs "
                        f"{total_close} close, delta={delta}), applying emergency fix")
        if delta > 0:
            wrapped += '</div>' * delta
        elif delta < 0:
            for _ in range(abs(delta)):
                last_close = wrapped.rfind('</div>')
                if last_close > 0:
                    wrapped = wrapped[:last_close] + wrapped[last_close + 6:]

    root_vars = _build_root_vars(scheme_data) if scheme_data else ""

    is_portrait = canvas_h > canvas_w
    body_bg = "#ffffff" if is_portrait else "{{{{secondary}}}}"

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
{root_vars}
  body {{
    background: {body_bg};
    display: flex;
    flex-direction: column;
    align-items: center;
    gap: 24px;
    padding: 24px;
    font-family: system-ui, -apple-system, sans-serif;
  }}
  .slide-wrapper {{
    width: {canvas_w}px;
    height: {canvas_h}px;
    overflow: hidden;
    border-radius: 4px;
    box-shadow: 0 4px 24px rgba(0,0,0,0.5);
    flex-shrink: 0;
  }}
</style>
</head>
<body>
{wrapped}
</body>
</html>"""


def _extract_svg(response: str) -> str | None:
    """Extract SVG XML from AI response — handles markdown-wrapped and raw SVG."""
    import re

    # Try ```svg ... ``` first
    m = re.search(r'```svg\s*\n(.*?)\n```', response, re.DOTALL | re.IGNORECASE)
    if m:
        return m.group(1).strip()

    # Try ```xml ... ```
    m = re.search(r'```xml\s*\n(.*?)\n```', response, re.DOTALL | re.IGNORECASE)
    if m:
        return m.group(1).strip()

    # Try <svg ... </svg> directly
    m = re.search(r'(<svg[\s\S]*?</svg>)', response, re.IGNORECASE)
    if m:
        return m.group(1).strip()

    return None


def _check_svg(svg: str, idx: int) -> list[str]:
    """Validate AI-generated SVG: XML well-formedness, viewBox, font sizes.

    Returns list of error strings. Empty list = all valid.
    """
    import re
    import xml.etree.ElementTree as ET

    issues = []

    if not re.search(r'<svg[\s>]', svg, re.IGNORECASE):
        issues.append(f"Slide {idx}: no <svg> tag found")
        return issues

    m = re.search(r'viewBox\s*=\s*["\']([^"\']+)["\']', svg)
    if m:
        vb = m.group(1)
        parts = vb.split()
        if len(parts) == 4:
            try:
                w, h = float(parts[2]), float(parts[3])
                if abs(w - 1280) > 2 or abs(h - 720) > 2:
                    issues.append(f"Slide {idx}: viewBox {w}x{h}, expected 1280x720")
            except ValueError:
                issues.append(f"Slide {idx}: unparseable viewBox '{vb}'")
    else:
        issues.append(f"Slide {idx}: missing viewBox")

    if svg.count('<svg') != svg.count('</svg>'):
        issues.append(f"Slide {idx}: mismatched <svg>/</svg> tags")

    try:
        ET.fromstring(svg)
    except ET.ParseError as e:
        issues.append(f"Slide {idx}: XML parse error: {str(e)[:100]}")

    font_sizes = re.findall(r'font-size\s*=\s*["\'](\d+(?:\.\d+)?)\s*(?:px)?["\']', svg)
    small_fonts = [float(s) for s in font_sizes if float(s) < 10]
    if small_fonts:
        issues.append(f"Slide {idx}: font sizes below 10px: {small_fonts}")
    tiny_fonts = [float(s) for s in font_sizes if float(s) < 12]
    if tiny_fonts:
        issues.append(f"Slide {idx}: font sizes below 12px (WARN): {tiny_fonts[:5]}")

    # Safe-area boundary check (PPT-Agent slide-core validation #4)
    # Check for content placed outside safe area. Exclude:
    # - x=0,y=0 (background rect), x=60 (safe left edge), y=40-60 (header zone)
    # - y=690-710 (footer zone with page numbers/branding)
    safe_violations = []
    x_vals = re.findall(r'(?:x|translate)\s*\(\s*(-?\d+(?:\.\d+)?)\s*[,\)]', svg)
    y_vals = re.findall(r'(?:y|translate)\s*\(\s*\d+\s*,\s*(-?\d+(?:\.\d+)?)\s*\)', svg)
    direct_x = re.findall(r'\sx\s*=\s*["\'](-?\d+(?:\.\d+)?)\s*(?:px)?["\']', svg)
    direct_y = re.findall(r'\sy\s*=\s*["\'](-?\d+(?:\.\d+)?)\s*(?:px)?["\']', svg)
    # Exclude x=0 (background rect), x=60 (safe left margin)
    all_x_low = [float(v) for v in x_vals + direct_x if float(v) < 40]
    # Exclude y=0 (background), y >= 28 (header zone)
    all_y_low = [float(v) for v in y_vals + direct_y if float(v) < 28]
    # Exclude y >= 685 but < 715 (footer zone is expected)
    all_y_high = [float(v) for v in y_vals + direct_y if float(v) > 710]
    if all_x_low:
        safe_violations.append(f"x < 40: {all_x_low[:3]}")
    if all_y_low:
        safe_violations.append(f"y < 28: {all_y_low[:3]}")
    if all_y_high:
        safe_violations.append(f"y > 710: {all_y_high[:3]}")
    if safe_violations:
        issues.append(f"Slide {idx}: safe-area boundary issues: {'; '.join(safe_violations)}")

    return issues


def _stage3_svg(provider_id, model, llm_generate, slide_data,
                 style_id: str = "business", rules: dict = None,
                 batch_size: int = 2, temperature: float = 0.3,
                 st: dict = None, column_id: str = "",
                 color_scheme: str = "deep-blue") -> list | None:
    """Phase 6 (AI-SVG): AI generates complete SVG XML per slide — PPT-Agent quality.

    Unlike the code-rendered path (AI→JSON→Code→SVG), this has AI write SVG directly,
    giving it full creative control: decorative elements, custom layouts, gradients,
    accent borders, semi-transparent overlays, info cards, etc.

    Args:
        provider_id: LLM provider ID
        model: model name
        llm_generate: async generate function
        slide_data: [{type, layout, zones: {heading, kicker, lead, cards: [{role, title, body, chart}]}}]
        style_id: style YAML name
        rules: optional rules dict
        batch_size: slides per LLM call (1-2, lower = better quality)
        st: stage temperatures dict with keys: svg_batch, svg_single, review, fix, holistic, holistic_fix

    Returns:
        [{seq, file, svg_content, type, label}] or None on failure
    """
    import re
    if st is None:
        st = {}

    gt_batch = st.get('svg_batch', temperature)
    gt_single = st.get('svg_single', temperature)
    rt_review = st.get('review', temperature)
    rt_fix = st.get('fix', temperature)
    rt_holistic = st.get('holistic', temperature)
    rt_holistic_fix = st.get('holistic_fix', temperature)

    svg_spec, bento_spec = _load_svg_prompt_specs(column_id)
    cognitive_spec = _load_cognitive_spec(column_id)
    reviewer_spec = _load_reviewer_spec(column_id)
    style_yaml = _load_style_yaml_text(style_id, color_scheme)

    if not svg_spec:
        _logger.warning("svg-generator.md not found, AI-SVG generation disabled")
        return None
    if not style_yaml:
        _logger.warning(f"Style YAML '{style_id}' not found, AI-SVG generation disabled")
        return None

    svg_system = f"""{svg_spec}

{bento_spec}

## 认知设计原则
{cognitive_spec}

## 设计评审标准（用于自检）
{reviewer_spec}

## 风格 Token（颜色/字体/阴影/圆角/渐变/发光）
{style_yaml}

---
严格按照以上 Style YAML 中的 color_scheme、typography、card_style、decoration、gradients 和 elevation 规范生成 SVG。所有颜色和字体必须从 YAML 精确提取，不得使用样例替代。"""

    result = []
    total = len([s for s in slide_data if isinstance(s, dict)])

    for batch_start in range(0, len(slide_data), batch_size):
        batch = slide_data[batch_start:batch_start + batch_size]
        valid_batch = [s for s in batch if isinstance(s, dict)]
        if not valid_batch:
            continue

        slide_descriptions = []
        for slide in valid_batch:
            seq = slide.get("seq", 0)
            slide_type = slide.get("type", "content")
            zones = slide.get("zones", {})
            heading = zones.get("heading", "")
            kicker = zones.get("kicker", "")
            lead = zones.get("lead", "")
            cards = zones.get("cards", [])
            body = zones.get("body", "")

            desc = f"--- 幻灯片 {seq}/{total} ---\n"
            desc += f"页码: {seq:02d} / {total:02d}\n"
            desc += f"类型: {slide_type}\n"
            desc += f"标题: {heading}\n"
            if kicker:
                desc += f"章节标签: {kicker}\n"
            if lead:
                desc += f"副标题: {lead}\n"
            if body:
                desc += f"正文: {body[:600]}\n"

            if cards:
                desc += f"卡片 ({len(cards)} 张):\n"
                for ci, card in enumerate(cards):
                    role = card.get("role", "unknown")
                    ctitle = card.get("title", "")
                    cbody = card.get("body", "")
                    chart = card.get("chart")
                    desc += f"  [{role}] {ctitle}"
                    if cbody:
                        desc += f" — {cbody[:300]}"
                    if chart:
                        desc += f" [图表: {json.dumps(chart, ensure_ascii=False)}]"
                    desc += "\n"

            slide_descriptions.append(desc)

        all_descs = "\n".join(slide_descriptions)

        svg_user = f"""为以下 {len(valid_batch)} 页幻灯片生成完整的 SVG XML。

{all_descs}

## 输出要求
为每一页输出一个完整的 SVG，按顺序用 ```svg ... ``` 包裹。每个 SVG 之间用空行分隔。
仅输出 SVG 代码块，不输出其他文字。"""

        generated_for_batch = []

        for attempt in range(2):
            try:
                response = _safe_run_async(llm_generate(provider_id, model,
                    svg_system, svg_user, temperature=gt_batch))

                svg_matches = re.findall(r'```svg\s*\n(.*?)\n```', response, re.DOTALL | re.IGNORECASE)
                if not svg_matches:
                    svg_matches = re.findall(r'(<svg[\s\S]*?</svg>)', response, re.IGNORECASE)

                if svg_matches and len(svg_matches) >= len(valid_batch):
                    for j, slide in enumerate(valid_batch):
                        seq = slide.get("seq", 0)
                        svg_content = svg_matches[j].strip()
                        errors = _check_svg(svg_content, seq)
                        if errors:
                            _logger.warning(f"Slide {seq} SVG check ({len(errors)} issues): {'; '.join(errors[:3])}")

                        generated_for_batch.append({
                            "seq": seq,
                            "file": f"slide-{seq:02d}.svg",
                            "svg_content": svg_content,
                            "type": slide.get("type", "content"),
                            "label": slide.get("zones", {}).get("heading", f"Slide {seq}"),
                        })
                    break
                else:
                    found = len(svg_matches) if svg_matches else 0
                    _logger.warning(
                        f"Batch {batch_start//batch_size+1} attempt {attempt+1}: "
                        f"found {found} SVGs, need {len(valid_batch)}"
                    )
                    if attempt == 0 and found > 0:
                        svg_user = f"上一轮只输出了 {found}/{len(valid_batch)} 个 SVG。请为所有 {len(valid_batch)} 页重新生成。\n\n{all_descs}"

            except Exception as e:
                _logger.warning(f"Batch {batch_start//batch_size+1} attempt {attempt+1} failed: {e}")
                if attempt == 0:
                    svg_user = f"生成失败: {str(e)[:200]}。请重试。\n\n{all_descs}"

        if len(generated_for_batch) == len(valid_batch):
            result.extend(generated_for_batch)
            _logger.info(f"Batch {batch_start//batch_size+1}: {len(generated_for_batch)} SVGs generated")
        else:
            _logger.warning(
                f"Batch {batch_start//batch_size+1}: only {len(generated_for_batch)}/{len(valid_batch)} "
                f"SVGs — retrying individually"
            )
            result.extend(generated_for_batch)
            # Retry each missing slide individually
            for slide in valid_batch:
                seq = slide.get("seq", 0)
                already_got = any(g.get("seq") == seq for g in generated_for_batch)
                if already_got:
                    continue
                # Build single-slide prompt
                zones = slide.get("zones", {})
                heading = zones.get("heading", "")
                kicker = zones.get("kicker", "")
                lead = zones.get("lead", "")
                cards = zones.get("cards", [])
                body = zones.get("body", "")
                single_desc = f"页码: {seq:02d} / {total:02d}\n类型: {slide.get('type', 'content')}\n标题: {heading}\n"
                if kicker:
                    single_desc += f"章节标签: {kicker}\n"
                if lead:
                    single_desc += f"副标题: {lead}\n"
                if body:
                    single_desc += f"正文: {body[:600]}\n"
                if cards:
                    single_desc += f"卡片 ({len(cards)} 张):\n"
                    for ci, card in enumerate(cards):
                        single_desc += f"  [{card.get('role', 'unknown')}] {card.get('title', '')}"
                        if card.get('body'):
                            single_desc += f" — {card['body'][:300]}"
                        if card.get('chart'):
                            single_desc += f" [图表]"
                        single_desc += "\n"

                single_user = f"""为以下 1 页幻灯片生成完整的 SVG XML。

{single_desc}

输出一个完整的 SVG，用 ```svg ... ``` 包裹。仅输出 SVG 代码块。"""

                single_ok = False
                for sa in range(3):
                    try:
                        resp = _safe_run_async(llm_generate(provider_id, model,
                            svg_system, single_user, temperature=gt_single))
                        svg_content = _extract_svg(resp)
                        if svg_content and "<svg" in svg_content and "</svg>" in svg_content:
                            errors = _check_svg(svg_content, seq)
                            if errors:
                                _logger.warning(f"Slide {seq} individual retry {sa+1} SVG check ({len(errors)} issues)")
                            result.append({
                                "seq": seq,
                                "file": f"slide-{seq:02d}.svg",
                                "svg_content": svg_content,
                                "type": slide.get("type", "content"),
                                "label": heading or f"Slide {seq}",
                            })
                            _logger.info(f"Slide {seq}: individual retry {sa+1} SUCCESS")
                            single_ok = True
                            break
                        else:
                            _logger.warning(f"Slide {seq} individual retry {sa+1}: no valid SVG extracted")
                    except Exception as e:
                        _logger.warning(f"Slide {seq} individual retry {sa+1} failed: {e}")
                if not single_ok:
                    _logger.error(f"Slide {seq}: all individual retries exhausted — slide will be missing")

    if not result:
        return None

    review_pid, review_m, review_mode = provider_id, model, "self_review"
    if provider_id and model:
        review_pid, review_m, review_mode = _resolve_review_models(provider_id, model)

    # ── Phase 6b: Per-slide review loop (PPT-Agent review-core equivalent) ──
    if provider_id and model:
        result = _review_and_fix_slides(review_pid, review_m, llm_generate, result,
                                        style_yaml, style_id, svg_system,
                                        review_mode=review_mode,
                                        temp_review=rt_review, temp_fix=rt_fix,
                                        column_id=column_id)

    # ── Phase 6c: Holistic deck review (cross-slide consistency) ──
    if provider_id and model and result and len(result) >= 3:
        result = _holistic_review(review_pid, review_m, llm_generate, result,
                                  slide_data, style_yaml, style_id, svg_system,
                                  review_mode=review_mode,
                                  temp_holistic=rt_holistic,
                                  temp_holistic_fix=rt_holistic_fix,
                                  column_id=column_id)

    return result


def _review_and_fix_slides(provider_id, model, llm_generate, slides, style_yaml,
                           style_id, svg_system, max_rounds=2,
                           review_mode: str = "self_review",
                           temperature: float = 0.3,
                           temp_review: float = 0, temp_fix: float = 0,
                           column_id: str = ""):
    """PPT-Agent review-core equivalent: review each slide, fix if score < 7.

    Uses the same LLM (DeepSeek/Moonshot) with PPT-Agent's reviewer.md prompt.
    Returns updated slides list with fixes applied.
    """
    import re, json

    reviewer_spec = _load_reviewer_spec(column_id)
    if not reviewer_spec:
        _logger.info("Reviewer spec not available, skipping review loop")
        return slides

    # ── Self-critique injection for single-model review ──
    self_critique_block = ""
    if review_mode == "self_review":
        self_critique_block = (
            "\n\n## CRITICAL: Self-Review Mode\n"
            "You are reviewing SVGs generated by YOUR OWN model. This means you share\n"
            "the same blind spots. To compensate:\n"
            "1. Be EXTRA skeptical — assume every color, font size, and layout choice\n"
            "   might be wrong. Verify each against the Style YAML above.\n"
            "2. Check for patterns YOU tend to overuse (same layout, same accent placement,\n"
            "   same card structure across slides — variety is required).\n"
            "3. Look for MISSING elements: shadows you forgot, gradients you skipped,\n"
            "   decorative elements the style requires but you omitted.\n"
            "4. Reduce inflated scores: if your first instinct is 8+, re-examine and\n"
            "   look harder for issues. Self-review tends to overrate by 1-2 points.\n"
        )

    review_system = f"""{reviewer_spec}

## Style YAML (reference for this review)
```yaml
{style_yaml}
```
{self_critique_block}

You are reviewing slides for style "{style_id}". Score each SVG on all 5 criteria.
Output MUST include the Suggestions JSON block with typed, actionable suggestions.
If the slide passes all quality gates, output an empty Suggestions JSON array []."""

    for round_num in range(max_rounds):
        fixes_needed = False
        for i, slide in enumerate(slides):
            svg_content = slide.get("svg_content", "")
            if not svg_content:
                continue

            seq = slide.get("seq", i + 1)
            slide_type = slide.get("type", "content")

            review_user = f"""## Slide {seq}: {slide.get('label', 'Untitled')} (type: {slide_type})

```svg
{svg_content[:10000]}
```

Review this slide and provide your full Quality Gate scores + Suggestions JSON."""

            try:
                response = _safe_run_async(llm_generate(provider_id, model,
                    review_system, review_user, temperature=temp_review or temperature))
            except Exception as e:
                _logger.warning(f"Review round {round_num+1} slide {seq} failed: {e}")
                continue

            # Parse score from response
            score_match = re.search(r'overall_score\s*\|\s*(\d+)', response)
            score = int(score_match.group(1)) if score_match else 0
            pass_match = re.search(r'pass\s*\|\s*(true|false)', response)
            passed = pass_match.group(1) == "true" if pass_match else False

            # Parse Suggestions JSON
            fixes_json = []
            json_match = re.search(r'```json\s*\n(.*?)\n```', response, re.DOTALL)
            if json_match:
                try:
                    fixes_json = json.loads(json_match.group(1))
                except json.JSONDecodeError:
                    pass

            _logger.info(
                f"Slide {seq} review round {round_num+1}: "
                f"score={score}/10 pass={passed} suggestions={len(fixes_json)}"
            )

            if score >= 7 and passed:
                continue  # This slide passes

            fixes_needed = True

            # Determine fix strategy based on suggestion types
            suggestion_types = {s.get("type", "") for s in fixes_json}

            if "full_rethink" in suggestion_types:
                # Regenerate from scratch with guidance
                guidance = fixes_json[0].get("details", {}).get("guidance", "")
                fix_prompt = (
                    f"\n\n## 重新设计要求（审查反馈）\n"
                    f"上一版评分 {score}/10。请完全重新设计本页：\n{guidance}\n"
                    f"严格遵循上方所有 Style YAML 和 SVG 规范。"
                )
            elif "layout_restructure" in suggestion_types:
                # Regenerate with layout constraint
                constraint = fixes_json[0].get("details", {}).get("constraint", "")
                suggested = fixes_json[0].get("details", {}).get("suggested_layout", "")
                fix_prompt = (
                    f"\n\n## 布局调整要求（审查反馈）\n"
                    f"上一版评分 {score}/10。使用布局 {suggested}：{constraint}\n"
                )
            elif "content_reduction" in suggestion_types:
                target = fixes_json[0].get("details", {}).get("target_info_units", 4)
                what = fixes_json[0].get("details", {}).get("what_to_remove", "")
                fix_prompt = (
                    f"\n\n## 内容精简要求（审查反馈）\n"
                    f"上一版评分 {score}/10。目标信息单元 {target} 个。\n移除：{what}\n"
                )
            elif "attribute_change" in suggestion_types:
                # Patch specific attributes
                patches = "; ".join(
                    f"{s.get('details',{}).get('attribute','?')}: "
                    f"{s.get('details',{}).get('current','?')} → "
                    f"{s.get('details',{}).get('target','?')}"
                    for s in fixes_json if s.get("type") == "attribute_change"
                )
                fix_prompt = (
                    f"\n\n## 精确属性修正（审查反馈）\n"
                    f"上一版评分 {score}/10。请修正以下属性：{patches}\n"
                )
            else:
                fix_prompt = (
                    f"\n\n## 质量改进（审查反馈）\n"
                    f"上一版评分 {score}/10。请改进整体质量。\n"
                )

            # Regenerate with fix prompt
            single_desc = f"页码: {seq:02d}\n类型: {slide_type}\n标题: {slide.get('label', '')}\n"
            fix_user = f"""为以下 1 页幻灯片生成改进版 SVG。

{single_desc}
{fix_prompt}

输出一个完整的 SVG，用 ```svg ... ``` 包裹。仅输出 SVG 代码块。"""

            for fa in range(2):
                try:
                    resp = _safe_run_async(llm_generate(provider_id, model,
                        svg_system, fix_user, temperature=temp_fix or temperature))
                    new_svg = _extract_svg(resp)
                    if new_svg and "<svg" in new_svg and "</svg>" in new_svg:
                        slide["svg_content"] = new_svg
                        _logger.info(f"Slide {seq} fix round {round_num+1} applied ({len(new_svg)} bytes)")
                        break
                except Exception as e:
                    _logger.warning(f"Slide {seq} fix attempt {fa+1} failed: {e}")

        if not fixes_needed:
            _logger.info(f"Review round {round_num+1}: all slides pass quality gate")
            break

    return slides


def _holistic_review(provider_id, model, llm_generate, slides, slide_data,
                     style_yaml, style_id, svg_system,
                     review_mode: str = "self_review",
                     temperature: float = 0.3,
                     temp_holistic: float = 0, temp_holistic_fix: float = 0,
                     column_id: str = ""):
    """Phase 6c: Holistic deck review — cross-slide consistency evaluation.

    PPT-Agent review-core holistic mode (reviewer.md:286-331): reads ALL slide
    SVGs and evaluates 5-Dimension cross-slide consistency. Unlike the previous
    metadata-only approach, this sends structural fingerprints with actual
    color values, font sizes, shadow defs, and stripped SVG skeletons so the
    LLM can detect real inconsistencies across slides.
    """
    import re, json, yaml

    reviewer_spec = _load_reviewer_spec(column_id)
    if not reviewer_spec:
        _logger.info("Reviewer spec not available, skipping holistic review")
        return slides

    if len(slides) < 3:
        _logger.info("Deck too small (<3 slides), skipping holistic review")
        return slides

    # ── Build structural fingerprint + stripped SVG per slide ──
    fingerprint_parts = []

    for i, s in enumerate(slides):
        seq = s.get("seq", i + 1)
        svg = s.get("svg_content", "")
        stype = s.get("type", "content")
        label = s.get("label", f"Slide {seq}")

        # Extract structural attributes from SVG
        fills = sorted(set(re.findall(r'fill\s*=\s*"([^"]+)"', svg)))
        strokes = sorted(set(re.findall(r'stroke\s*=\s*"([^"]+)"', svg)))
        hex_colors = sorted(set(
            c for c in fills + strokes
            if re.match(r'^#[0-9a-fA-F]{3,8}$', c)
        ))
        font_sizes = sorted(set(re.findall(r'font-size\s*=\s*"(\d+(?:\.\d+)?)', svg)))
        font_families = sorted(set(re.findall(r'font-family\s*=\s*"([^"]+)"', svg)))
        rx_vals = sorted(set(re.findall(r'rx\s*=\s*"(\d+(?:\.\d+)?)"', svg)))
        shadow_ids = re.findall(r'<filter\s+id="([^"]*shadow[^"]*)"', svg)
        gradient_ids = re.findall(r'<linearGradient\s+id="([^"]+)"', svg)

        # Layout detection
        layout = s.get("layout", "")
        if not layout:
            if re.search(r'three.column|3.col', svg):
                layout = "three_column"
            elif re.search(r'two.column|2.col', svg):
                layout = "two_column"
            elif re.search(r'single.focus|full.bleed', svg):
                layout = "single_focus"
            elif re.search(r'hero.grid', svg):
                layout = "hero_grid"
            else:
                layout = "mixed"

        # Visual weight
        low_types = {"quote", "image", "image_hero", "section", "copyright", "appendix"}
        high_types = {"data_hero", "comparison", "duo_compare", "table"}
        weight = "low" if stype in low_types else ("high" if stype in high_types else "medium")

        # Stripped SVG: remove text content but keep structure
        # Strip <text>...</text> body content, keep all structural tags + defs
        stripped = re.sub(r'(<text[^>]*>).*?(</text>)', r'\1...\2', svg)
        # Keep only first 6000 chars of stripped SVG
        stripped = stripped[:6000]
        if len(svg) > 6000:
            stripped += "\n<!-- ... truncated ... -->"

        fp = (
            f"## Slide {seq:02d}: {label}\n"
            f"- Type: {stype} | Layout: {layout} | Weight: {weight}\n"
            f"- Colors in use: {', '.join(hex_colors[:20])}\n"
            f"- Font sizes: {', '.join(font_sizes[:15])}\n"
            f"- Font families: {', '.join(font_families[:5])}\n"
            f"- Border radii: {', '.join(rx_vals[:10])}\n"
            f"- Shadow filters: {', '.join(shadow_ids) if shadow_ids else 'none'}\n"
            f"- Gradients: {', '.join(gradient_ids) if gradient_ids else 'none'}\n"
            f"\n```svg\n{stripped}\n```\n"
        )
        fingerprint_parts.append(fp)

    all_fingerprints = "\n---\n".join(fingerprint_parts)

    # Summary table for quick overview
    summary_rows = []
    for i, s in enumerate(slides):
        seq = s.get("seq", i + 1)
        svg = s.get("svg_content", "")
        stype = s.get("type", "content")
        label = s.get("label", f"Slide {seq}")
        low_types = {"quote", "image", "image_hero", "section", "copyright", "appendix"}
        high_types = {"data_hero", "comparison", "duo_compare", "table"}
        weight = "low" if stype in low_types else ("high" if stype in high_types else "medium")
        layout = s.get("layout", "mixed")
        card_count = len(re.findall(r'<g[^>]*transform', svg))
        summary_rows.append(
            f"| {seq:02d} | {stype} | {layout} | {weight} | {card_count} | {label} |"
        )

    summary_table = "\n".join([
        "| # | Type | Layout | Weight | Cards | Title |",
        "|---|------|--------|--------|-------|-------|",
    ] + summary_rows)

    # Self-critique injection for single-model holistic review
    holistic_critique = ""
    if review_mode == "self_review":
        holistic_critique = (
            "\n\n## CRITICAL: Self-Review Mode\n"
            "You are reviewing a deck generated by YOUR OWN model. Cross-slide issues\n"
            "are YOUR blind spots — be extra vigilant:\n"
            "1. Compare STRUCTURAL ATTRIBUTES numerically: are all slide titles really\n"
            "   the same font-size? Check the fingerprint data, don't assume.\n"
            "2. SCAN for monotony: your model tends to repeat the same layout pattern.\n"
            "   Flag any 3+ consecutive slides with identical layout.\n"
            "3. COLOR DRIFT: check if hex colors are consistent across slides.\n"
            "   Your model often uses slightly different hex codes for 'the same' color.\n"
            "4. Reduce scores by 1 point: self-review systematically overrates by 1-2 points.\n"
        )

    holistic_system = f"""{reviewer_spec}

## Style YAML (reference tokens)
```yaml
{style_yaml}
```
{holistic_critique}

You are performing HOLISTIC DECK REVIEW for style "{style_id}". You have access to each slide's structural fingerprint (exact colors, font sizes, shadows, gradients) AND stripped SVG skeleton. Use these to detect REAL cross-slide inconsistencies, not guess from metadata."""

    holistic_user = f"""## Deck Summary: {len(slides)} slides

{summary_table}

## Per-Slide Structural Fingerprints + Stripped SVGs

{all_fingerprints}

## Task

Evaluate cross-slide consistency using the 5-Dimension Framework. Use the actual color values, font sizes, shadow defs, and border radii from the fingerprints above to detect inconsistencies — do not guess from metadata alone.

1. **Visual Rhythm** (25%): Do layouts alternate? 3+ consecutive same layout = trigger.
2. **Color Story** (20%): Compare hex_colors across slides. Accent on >60% of slides = trigger. Missing accent on climax slide = trigger.
3. **Narrative Arc** (20%): 3+ consecutive high-weight slides without low-weight breathing slide = trigger.
4. **Style Consistency** (20%): Compare font_sizes, font_families, rx_vals, shadow_ids across slides. >30% variance on any attribute = trigger.
5. **Pacing** (15%): 4+ consecutive content/data/comparison slides with no breathing slide = trigger.

Output:
- Holistic Scoring table (each dimension 0-10, weighted total)
- deck_coordination Suggestions JSON with affected_slides and concrete fix recommendations
- If coherence >= 7 AND no P1 issues, output empty JSON array []"""

    try:
        response = _safe_run_async(llm_generate(provider_id, model,
            holistic_system, holistic_user, temperature=temp_holistic or temperature))
    except Exception as e:
        _logger.warning(f"Holistic review failed: {e}")
        return slides

    # Parse coherence score
    coherence_match = re.search(r'Overall Coherence.*?\|\s*\**(\d+(?:\.\d+)?)\**', response)
    coherence = float(coherence_match.group(1)) if coherence_match else 0

    # Parse deck_coordination suggestions
    fixes_json = []
    json_match = re.search(r'```json\s*\n(.*?)\n```', response, re.DOTALL)
    if json_match:
        try:
            fixes_json = json.loads(json_match.group(1))
        except json.JSONDecodeError:
            pass

    _logger.info(
        f"Holistic review: coherence={coherence}/10, "
        f"{len(fixes_json)} deck_coordination suggestions"
    )

    if coherence >= 7 and not fixes_json:
        _logger.info("Holistic review: PASSED")
        return slides

    # ── Apply priority-1 deck_coordination fixes ──
    p1_fixes = [f for f in fixes_json if f.get("priority") == 1]
    if not p1_fixes:
        _logger.info(f"Holistic review: score={coherence}, no P1 fixes — advisory only")
        return slides

    _logger.info(f"Holistic review: applying {len(p1_fixes)} priority-1 fixes")
    for fix in p1_fixes:
        affected = fix.get("details", {}).get("affected_slides", [])
        suggestion = fix.get("details", {}).get("suggestion", fix.get("description", ""))
        issue_type = fix.get("details", {}).get("issue_type", "")

        for slide_idx in affected:
            if isinstance(slide_idx, int) and 1 <= slide_idx <= len(slides):
                slide = slides[slide_idx - 1]
            else:
                continue

            seq = slide.get("seq", slide_idx)
            fix_prompt = (
                f"\n\n## 整体协调修正（Holistic Review P1）\n"
                f"问题类型: {issue_type}\n修正说明: {suggestion}\n"
                f"严格遵循上方所有 Style YAML 和 SVG 规范。"
            )

            single_desc = f"页码: {seq:02d}\n类型: {slide.get('type', 'content')}\n标题: {slide.get('label', '')}\n"
            fix_user = f"""为以下 1 页幻灯片生成改进版 SVG。

{single_desc}
{fix_prompt}

输出一个完整的 SVG，用 ```svg ... ``` 包裹。仅输出 SVG 代码块。"""

            for fa in range(2):
                try:
                    resp = _safe_run_async(llm_generate(provider_id, model,
                        svg_system, fix_user, temperature=temp_holistic_fix or temperature))
                    new_svg = _extract_svg(resp)
                    if new_svg and "<svg" in new_svg and "</svg>" in new_svg:
                        slide["svg_content"] = new_svg
                        _logger.info(f"Holistic fix applied to slide {seq} ({len(new_svg)} bytes)")
                        break
                except Exception as e:
                    _logger.warning(f"Holistic fix slide {seq} attempt {fa+1} failed: {e}")

    return slides