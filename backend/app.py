import os
import shutil
import uuid
import hashlib
from fastapi import FastAPI, HTTPException, UploadFile, File, Body, Request
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from database import init_db, get_db
from models import (ProjectCreate, ProjectUpdate, StepResultSave, LLMGenerateRequest,
    LLMRefineRequest, SynthesizeRequest, PPTGenerateRequest, PPTPlanRequest, PPTEditSlideRequest,
    PPTSlideSourceRequest, PPTRegenerateSlideRequest,
    ImageGenerateRequest, SourceMaterialCreate, SourceMaterialUpdate,
    ProjectItemCreate, ProjectItemUpdate, ProjectItemResultSave)
from typing import Optional
import json
import logging
from services.llm_service import test_connection, generate, generate_stream, refine, get_provider
from services.prompt_service import (
    list_prompts, get_prompt, create_prompt, update_prompt, delete_prompt,
    rollback_version, diff_versions, set_default, export_prompts, import_prompts,
)

DEFAULT_SITE_NAME = "Yishao Agent"

init_db()

app = FastAPI(title=DEFAULT_SITE_NAME)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
WORKSPACE_ROOT = os.path.dirname(BASE_DIR)  # d:\YISHAOAGENT
AUDIO_DIR = os.path.join(BASE_DIR, "data", "audio")
os.makedirs(AUDIO_DIR, exist_ok=True)
EXPORT_DIR = os.path.join(BASE_DIR, "data", "exports")
os.makedirs(EXPORT_DIR, exist_ok=True)
LOGO_DIR = os.path.join(BASE_DIR, "data", "logos")
os.makedirs(LOGO_DIR, exist_ok=True)
THUMBNAIL_DIR = os.path.join(BASE_DIR, "data", "thumbnails")
os.makedirs(THUMBNAIL_DIR, exist_ok=True)

# Run-id → actual directory mapping for SVG preview serving
# Persisted to data/run_dirs.json so it survives restarts
_RUN_DIRS_FILE = os.path.join(BASE_DIR, "data", "run_dirs.json")

def _load_run_dirs() -> dict[str, str]:
    try:
        with open(_RUN_DIRS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def _save_run_dirs(mapping: dict[str, str]):
    try:
        os.makedirs(os.path.dirname(_RUN_DIRS_FILE), exist_ok=True)
        with open(_RUN_DIRS_FILE, "w", encoding="utf-8") as f:
            json.dump(mapping, f)
    except Exception:
        pass

_run_dirs: dict[str, str] = _load_run_dirs()

def _scan_output_bases() -> list[str]:
    """Scan for possible output directories (project storage dirs)."""
    bases = []
    for output_root in [
        os.path.join(WORKSPACE_ROOT, "output"),
        os.path.join(WORKSPACE_ROOT, "data", "output"),
    ]:
        if os.path.isdir(output_root):
            for name in os.listdir(output_root):
                d = os.path.join(output_root, name)
                if os.path.isdir(d):
                    bases.append(d)
    return bases

@app.get("/api/exports/{run_id}/{filename:path}")
def api_serve_export_file(run_id: str, filename: str):
    import starlette.responses as _sr
    run_dir = _run_dirs.get(run_id)
    if not run_dir:
        # Fallback: look in EXPORT_DIR and scan subdirs
        candidate = os.path.join(EXPORT_DIR, run_id)
        if os.path.isdir(candidate):
            run_dir = candidate
    if not run_dir:
        # Last resort: scan known output dirs for run_id
        for base in [EXPORT_DIR] + _scan_output_bases():
            candidate = os.path.join(base, run_id)
            if os.path.isdir(candidate):
                run_dir = candidate
                _run_dirs[run_id] = run_dir
                _save_run_dirs(_run_dirs)
                break
    if not run_dir:
        run_dir = os.path.join(EXPORT_DIR, run_id)
    filepath = os.path.normpath(os.path.join(run_dir, filename))
    if not filepath.startswith(os.path.normpath(run_dir)):
        raise HTTPException(status_code=403, detail="Path traversal denied")
    if not os.path.isfile(filepath):
        raise HTTPException(status_code=404, detail="Not Found")
    headers = {}
    if filepath.endswith('.html'):
        headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    return _sr.FileResponse(filepath, headers=headers)

import re


def _get_global_save_path() -> str:
    """Read global save_path from settings, or return default."""
    db = get_db()
    try:
        row = db.execute("SELECT value FROM settings WHERE key = 'save_path'").fetchone()
        if row and row["value"]:
            return row["value"]
    finally:
        db.close()
    return os.path.join(BASE_DIR, "data", "output")


def _sanitize_folder_name(name: str) -> str:
    """Remove characters invalid for Windows folder names."""
    return re.sub(r'[<>:"/\\|?*]', '_', name).strip().rstrip('.') or "unnamed"


def _get_setting(key: str) -> str:
    """Read a single setting value from the database."""
    db = get_db()
    try:
        row = db.execute("SELECT value FROM settings WHERE key = ?", (key,)).fetchone()
        return row["value"] if row else ""
    finally:
        db.close()


def _get_site_name() -> str:
    """Read site name from settings (brand_name), with fallback."""
    name = _get_setting("brand_name")
    return name if name else DEFAULT_SITE_NAME


SALT = "yishao-agent-salt-2026"


def _hash_password(password: str) -> str:
    return hashlib.sha256((SALT + password).encode("utf-8")).hexdigest()


def _get_project(project_id: str):
    """Return project record dict or None."""
    db = get_db()
    try:
        row = db.execute("SELECT * FROM projects WHERE id = ?", (project_id,)).fetchone()
        return dict(row) if row else None
    finally:
        db.close()


def resolve_project_storage(project_id: str, auto_create: bool = True) -> str:
    """Return the effective storage directory for a project."""
    db = get_db()
    try:
        proj = db.execute("SELECT * FROM projects WHERE id = ?", (project_id,)).fetchone()
        if not proj:
            raise FileNotFoundError(f"Project {project_id} not found")

        if proj["storage_path"]:
            path = os.path.normpath(proj["storage_path"])
            if auto_create:
                os.makedirs(path, exist_ok=True)
            return path

        base = _get_global_save_path()
        folder = _sanitize_folder_name(proj["name"])
        path = os.path.normpath(os.path.join(base, folder))

        if auto_create:
            os.makedirs(path, exist_ok=True)
            db.execute(
                "UPDATE projects SET storage_path = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?",
                (path, project_id))
            db.commit()

        return path
    finally:
        db.close()

# ── Health check ──

@app.get("/api/health")
def health():
    return {"status": "ok", "app": _get_site_name()}


# ── Projects ──

@app.get("/api/projects")
def list_projects(page: int = 1, page_size: int = 20):
    db = get_db()
    try:
        total = db.execute("SELECT COUNT(*) FROM projects").fetchone()[0]
        offset = (page - 1) * page_size
        rows = db.execute(
            "SELECT * FROM projects ORDER BY updated_at DESC LIMIT ? OFFSET ?",
            (page_size, offset)
        ).fetchall()
        return {"projects": [dict(r) for r in rows], "total": total, "page": page, "page_size": page_size}
    finally:
        db.close()


@app.post("/api/projects")
def create_project(req: ProjectCreate):
    pid = uuid.uuid4().hex[:12]
    db = get_db()
    try:
        # Compute default storage path
        base = _get_global_save_path()
        folder = _sanitize_folder_name(req.name)
        storage_path = os.path.normpath(req.storage_path or os.path.join(base, folder))
        os.makedirs(storage_path, exist_ok=True)

        # Generate project code KH{YYMMDD}-{seq} (seq resets daily)
        from datetime import date
        today = date.today().strftime("%y%m%d")  # "260627"
        today_prefix = f"KH{today}-%"
        today_count = db.execute(
            "SELECT COUNT(*) FROM projects WHERE project_code LIKE ?", (today_prefix,)
        ).fetchone()[0]
        project_code = f"KH{today}-{today_count + 1:04d}"

        db.execute(
            "INSERT INTO projects (id, name, source_type, storage_path, project_code) VALUES (?, ?, ?, ?, ?)",
            (pid, req.name, req.source_type, storage_path, project_code))
        db.commit()
        row = db.execute("SELECT * FROM projects WHERE id = ?", (pid,)).fetchone()
        return dict(row)
    finally:
        db.close()


@app.get("/api/projects/{project_id}/videos")
def api_project_videos(project_id: str):
    """List video files in the project's storage folder."""
    path = resolve_project_storage(project_id, auto_create=False)
    videos = []
    if os.path.exists(path):
        for f in os.listdir(path):
            if f.lower().endswith(('.mp4', '.mkv', '.webm', '.avi', '.mov', '.flv')):
                full = os.path.join(path, f)
                videos.append({"filename": f, "path": full, "size": os.path.getsize(full)})
    return {"videos": videos, "storage_path": path}


@app.get("/api/projects/{project_id}/files")
def api_project_files(project_id: str):
    """List all generated output files in the project's storage folder."""
    path = resolve_project_storage(project_id, auto_create=False)
    files = []
    if os.path.exists(path):
        for f in sorted(os.listdir(path), key=lambda x: os.path.getmtime(os.path.join(path, x)), reverse=True):
            full = os.path.join(path, f)
            if os.path.isfile(full):
                ext = os.path.splitext(f)[1].lower()
                type_map = {'.pptx': 'PPT', '.docx': 'Word', '.txt': 'Text',
                           '.mp3': 'MP3', '.wav': 'Audio', '.mp4': 'Video'}
                file_type = type_map.get(ext, 'Other')
                files.append({
                    "filename": f,
                    "type": file_type,
                    "size": os.path.getsize(full),
                    "modified": os.path.getmtime(full),
                    "download_url": f"/api/download/{f}?project_id={project_id}"
                })
    return {"files": files, "storage_path": path}


@app.get("/api/projects/{project_id}")
def get_project(project_id: str):
    db = get_db()
    try:
        row = db.execute("SELECT * FROM projects WHERE id = ?", (project_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Project not found")
        return dict(row)
    finally:
        db.close()


@app.put("/api/projects/{project_id}")
def update_project(project_id: str, req: ProjectUpdate):
    db = get_db()
    try:
        existing = db.execute("SELECT id FROM projects WHERE id = ?", (project_id,)).fetchone()
        if not existing:
            raise HTTPException(status_code=404, detail="Project not found")

        if req.name is not None:
            db.execute("UPDATE projects SET name = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?", (req.name, project_id))
        if req.status is not None:
            db.execute("UPDATE projects SET status = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?", (req.status, project_id))
        if req.storage_path is not None:
            if req.storage_path.strip():
                os.makedirs(req.storage_path, exist_ok=True)
            db.execute("UPDATE projects SET storage_path = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?", (os.path.normpath(req.storage_path) if req.storage_path.strip() else req.storage_path, project_id))
        if req.is_locked is not None:
            db.execute("UPDATE projects SET is_locked = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?", (req.is_locked, project_id))
        db.commit()
        row = db.execute("SELECT * FROM projects WHERE id = ?", (project_id,)).fetchone()
        return dict(row)
    finally:
        db.close()


@app.get("/api/fs/dirs")
def api_list_fs_dirs(path: str = ""):
    """List subdirectories at any filesystem path. Empty path returns drive letters on Windows."""
    import platform
    # Empty path on Windows: return drive letters
    if not path:
        if platform.system() == "Windows":
            import string
            drives = []
            for letter in string.ascii_uppercase:
                p = f"{letter}:\\"
                if os.path.exists(p):
                    drives.append(p)
            return {"ok": True, "dirs": sorted(drives), "path": "", "parent": None}
        else:
            path = "/"
    target = os.path.normpath(path)
    if not os.path.exists(target):
        return {"ok": True, "dirs": [], "path": path, "parent": None}
    if not os.path.isdir(target):
        raise HTTPException(400, "Path is not a directory")
    try:
        entries = os.listdir(target)
    except PermissionError:
        return {"ok": True, "dirs": [], "path": path, "parent": os.path.normpath(os.path.join(target, '..')) if path else None}
    dirs = sorted([
        e for e in entries
        if os.path.isdir(os.path.join(target, e)) and not e.startswith('.')
    ])
    # Parent: full path of parent, or null if at root
    parent_dir = os.path.dirname(target)
    if parent_dir == target:
        parent = None  # At filesystem root
    else:
        parent = parent_dir
    return {"ok": True, "dirs": dirs, "path": path, "parent": parent}


@app.post("/api/fs/mkdir")
async def api_create_fs_dir(request: Request):
    """Create a new directory at the given parent path."""
    req = await request.json()
    parent = req.get("parent", "")
    name = req.get("name", "").strip()
    if not parent or not name:
        raise HTTPException(400, "parent and name are required")
    # Reject names that contain path separators or traversal
    if '/' in name or '\\' in name or name in ('.', '..'):
        raise HTTPException(400, "Invalid directory name")
    if not os.path.exists(parent) or not os.path.isdir(parent):
        raise HTTPException(400, "Parent directory does not exist")
    target = os.path.join(parent, name)
    if os.path.exists(target):
        raise HTTPException(409, "Directory already exists")
    os.makedirs(target)
    return {"ok": True, "path": target}


@app.get("/api/projects/{project_id}/directories")
def api_list_project_directories(project_id: str, subdir: str = ""):
    """List subdirectories under the project storage path (or a subdirectory thereof)."""
    base = resolve_project_storage(project_id, auto_create=False)
    if not os.path.exists(base):
        return {"ok": True, "dirs": [], "base": base, "subdir": subdir}
    target = os.path.normpath(os.path.join(base, subdir))
    # Prevent path traversal
    if os.path.commonpath([os.path.abspath(target), os.path.abspath(base)]) != os.path.abspath(base):
        raise HTTPException(403, "Path traversal denied")
    try:
        entries = os.listdir(target)
    except FileNotFoundError:
        return {"ok": True, "dirs": [], "base": base, "subdir": subdir}
    dirs = sorted([
        e for e in entries
        if os.path.isdir(os.path.join(target, e)) and not e.startswith('.')
    ])
    parent = None
    if subdir:
        parent = os.path.normpath(os.path.join(subdir, '..'))
        if parent == '.':
            parent = ''
    return {"ok": True, "dirs": dirs, "base": base, "subdir": subdir, "parent": parent}


@app.post("/api/projects/{project_id}/save-file")
def api_save_file_to_project(project_id: str, req: dict):
    """Save content to a file. If target_dir is provided (absolute path), use it directly.
    Otherwise resolve relative to the project's storage directory."""
    filename = req.get("filename", "document.txt")
    content = req.get("content", "")
    encoding = req.get("encoding", "text")
    target_dir = req.get("target_dir", "")
    subdir = req.get("subdir", "")

    if target_dir and os.path.isabs(target_dir):
        # Use the absolute path directly (from filesystem browser)
        if not os.path.exists(target_dir):
            raise HTTPException(400, f"Target directory does not exist: {target_dir}")
        if not os.path.isdir(target_dir):
            raise HTTPException(400, f"Target path is not a directory: {target_dir}")
        os.makedirs(target_dir, exist_ok=True)
    elif subdir:
        path = resolve_project_storage(project_id)
        os.makedirs(path, exist_ok=True)
        target_dir = os.path.normpath(os.path.join(path, subdir))
        if os.path.commonpath([os.path.abspath(target_dir), os.path.abspath(path)]) != os.path.abspath(path):
            raise HTTPException(403, "Path traversal denied")
        os.makedirs(target_dir, exist_ok=True)
    else:
        target_dir = resolve_project_storage(project_id)
        os.makedirs(target_dir, exist_ok=True)

    filepath = os.path.join(target_dir, filename)
    # Avoid overwriting: append (1), (2), etc. if file exists
    if os.path.exists(filepath):
        base_name, ext = os.path.splitext(filename)
        n = 1
        while os.path.exists(os.path.join(target_dir, f"{base_name}({n}){ext}")):
            n += 1
        filename = f"{base_name}({n}){ext}"
        filepath = os.path.join(target_dir, filename)
    import base64
    if encoding == "base64":
        with open(filepath, "wb") as f:
            f.write(base64.b64decode(content))
    else:
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(content)
    return {"ok": True, "path": filepath, "filename": filename, "target_dir": target_dir}


def _delete_project_files(project_id: str, db):
    """Remove project storage directory and all related data."""
    db.execute("DELETE FROM source_materials WHERE project_id = ?", (project_id,))
    db.execute("DELETE FROM project_item_results WHERE project_item_id IN (SELECT id FROM project_items WHERE project_id = ?)", (project_id,))
    db.execute("DELETE FROM project_items WHERE project_id = ?", (project_id,))
    db.execute("DELETE FROM step_results WHERE project_id = ?", (project_id,))
    try:
        path = resolve_project_storage(project_id, auto_create=False)
        if os.path.exists(path):
            shutil.rmtree(path)
    except Exception:
        pass


@app.delete("/api/projects/{project_id}")
def delete_project(project_id: str):
    db = get_db()
    try:
        row = db.execute("SELECT is_locked FROM projects WHERE id = ?", (project_id,)).fetchone()
        if not row:
            raise HTTPException(404, "Project not found")
        if row["is_locked"]:
            raise HTTPException(403, "项目已锁定，无法删除")
        _delete_project_files(project_id, db)
        db.execute("DELETE FROM projects WHERE id = ?", (project_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/projects/batch-delete")
def batch_delete_projects(req: dict):
    ids = req.get("ids", [])
    if not ids:
        raise HTTPException(400, "ids required")
    db = get_db()
    try:
        placeholders = ",".join(["?"] * len(ids))
        locked_rows = db.execute(
            f"SELECT id, name FROM projects WHERE id IN ({placeholders}) AND is_locked = 1", ids
        ).fetchall()
        locked_ids = {r["id"] for r in locked_rows}
        unlocked = [pid for pid in ids if pid not in locked_ids]
        if not unlocked:
            names = ", ".join(r["name"] for r in locked_rows)
            raise HTTPException(403, f"所选项目均已锁定，无法删除: {names}")
        for pid in unlocked:
            _delete_project_files(pid, db)
        db.execute(
            f"DELETE FROM projects WHERE id IN ({','.join(['?'] * len(unlocked))})", unlocked
        )
        db.commit()
        skipped = len(ids) - len(unlocked)
        msg = f"已删除 {len(unlocked)} 个项目"
        if skipped > 0:
            names = ", ".join(r["name"] for r in locked_rows)
            msg += f"，{skipped} 个已锁定跳过: {names}"
        return {"ok": True, "deleted": len(unlocked), "skipped": skipped, "message": msg}
    finally:
        db.close()


# ── Step Results ──

@app.get("/api/projects/{project_id}/steps")
def get_steps(project_id: str):
    db = get_db()
    try:
        rows = db.execute(
            "SELECT * FROM step_results WHERE project_id = ? ORDER BY step_name", (project_id,)
        ).fetchall()
        return {"steps": [dict(r) for r in rows]}
    finally:
        db.close()


@app.put("/api/projects/{project_id}/steps/{step_name}")
def save_step(project_id: str, step_name: str, req: StepResultSave):
    db = get_db()
    try:
        existing = db.execute(
            "SELECT id FROM step_results WHERE project_id = ? AND step_name = ?",
            (project_id, step_name)).fetchone()
        if existing:
            db.execute(
                "UPDATE step_results SET content = ?, content_type = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?",
                (req.content, req.content_type, existing["id"]))
        else:
            db.execute(
                "INSERT INTO step_results (project_id, step_name, content, content_type) VALUES (?, ?, ?, ?)",
                (project_id, step_name, req.content, req.content_type))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


def save_step_meta(project_id: str, step_name: str, content: str):
    """Save result metadata to step_results (fire-and-forget, no fail on error)."""
    try:
        db = get_db()
        try:
            existing = db.execute(
                "SELECT id FROM step_results WHERE project_id = ? AND step_name = ?",
                (project_id, step_name)).fetchone()
            if existing:
                db.execute(
                    "UPDATE step_results SET content = ?, content_type = 'json', updated_at = CURRENT_TIMESTAMP WHERE id = ?",
                    (content, existing["id"]))
            else:
                db.execute(
                    "INSERT INTO step_results (project_id, step_name, content, content_type) VALUES (?, ?, ?, 'json')",
                    (project_id, step_name, content))
            db.commit()
        finally:
            db.close()
    except Exception:
        pass  # non-critical, don't fail the request


# ── Source Materials (multi-format input) ──

@app.get("/api/projects/{project_id}/materials")
def list_materials(project_id: str):
    db = get_db()
    try:
        rows = db.execute(
            "SELECT * FROM source_materials WHERE project_id = ? ORDER BY created_at DESC",
            (project_id,)).fetchall()
        return {"materials": [dict(r) for r in rows]}
    finally:
        db.close()


@app.post("/api/projects/{project_id}/materials")
def add_material(project_id: str, req: SourceMaterialCreate):
    db = get_db()
    try:
        mat_id = f"sm-{project_id}-{uuid.uuid4().hex[:8]}"
        db.execute(
            "INSERT INTO source_materials (id, project_id, source_type, source_name, "
            "raw_content, processed_content, status) VALUES (?, ?, ?, ?, ?, ?, ?)",
            (mat_id, project_id, req.source_type, req.source_name,
             req.raw_content, req.processed_content or req.raw_content, req.status))
        db.commit()
        row = db.execute("SELECT * FROM source_materials WHERE id = ?", (mat_id,)).fetchone()
        return dict(row)
    finally:
        db.close()


@app.post("/api/projects/{project_id}/materials/upload")
async def upload_material(project_id: str, file: UploadFile = File(...)):
    from services.file_parser import parse_bytes

    data = await file.read()
    result = parse_bytes(data, file.filename or "unknown")

    db = get_db()
    try:
        mat_id = f"sm-{project_id}-{uuid.uuid4().hex[:8]}"
        source_type = os.path.splitext(file.filename or "")[1].lower().lstrip(".")
        db.execute(
            "INSERT INTO source_materials (id, project_id, source_type, source_name, "
            "raw_content, processed_content, status) VALUES (?, ?, ?, ?, ?, ?, ?)",
            (mat_id, project_id, source_type, file.filename or "",
             result.get("text", ""), result.get("text", ""),
             "processed" if result.get("status") == "ok" else "error"))
        db.commit()
        row = db.execute("SELECT * FROM source_materials WHERE id = ?", (mat_id,)).fetchone()
        return {
            "material": dict(row),
            "parse_result": result,
        }
    finally:
        db.close()


@app.delete("/api/projects/{project_id}/materials/{material_id}")
def delete_material(project_id: str, material_id: str):
    db = get_db()
    try:
        db.execute("DELETE FROM source_materials WHERE id = ? AND project_id = ?",
                   (material_id, project_id))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.put("/api/projects/{project_id}/materials/{material_id}")
def update_material(project_id: str, material_id: str, req: SourceMaterialUpdate):
    db = get_db()
    try:
        existing = db.execute("SELECT id FROM source_materials WHERE id = ? AND project_id = ?",
                              (material_id, project_id)).fetchone()
        if not existing:
            raise HTTPException(404, "Material not found")
        updates = {}
        for k in ("source_name", "raw_content", "processed_content", "status"):
            v = getattr(req, k, None)
            if v is not None:
                updates[k] = v
        if updates:
            cols = ", ".join(f"{k} = ?" for k in updates)
            vals = list(updates.values())
            db.execute(f"UPDATE source_materials SET {cols} WHERE id = ?", vals + [material_id])
            db.commit()
        row = db.execute("SELECT * FROM source_materials WHERE id = ?", (material_id,)).fetchone()
        return dict(row)
    finally:
        db.close()


# ── Project Items (dynamic output steps) ──

@app.get("/api/projects/{project_id}/items")
def list_project_items(project_id: str):
    db = get_db()
    try:
        rows = db.execute(
            "SELECT pi.*, "
            "  (SELECT COUNT(*) FROM project_item_results WHERE project_item_id = pi.id) as result_count "
            "FROM project_items pi WHERE pi.project_id = ? ORDER BY pi.sort_order",
            (project_id,)).fetchall()
        return {"items": [dict(r) for r in rows]}
    finally:
        db.close()


@app.post("/api/projects/{project_id}/items")
def create_project_item(project_id: str, req: ProjectItemCreate):
    db = get_db()
    try:
        item_id = f"pi-{project_id}-{uuid.uuid4().hex[:8]}"
        db.execute(
            "INSERT INTO project_items (id, project_id, name, prompt, skill, "
            "output_mode, config_json, source_item_id, sort_order) "
            "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)",
            (item_id, project_id, req.name, req.prompt, req.skill,
             req.output_mode, req.config_json, req.source_item_id, req.sort_order))
        db.commit()
        row = db.execute("SELECT * FROM project_items WHERE id = ?", (item_id,)).fetchone()
        return dict(row)
    finally:
        db.close()


@app.put("/api/projects/{project_id}/items/{item_id}")
def update_project_item(project_id: str, item_id: str, req: ProjectItemUpdate):
    db = get_db()
    try:
        existing = db.execute("SELECT id FROM project_items WHERE id = ? AND project_id = ?",
                              (item_id, project_id)).fetchone()
        if not existing:
            raise HTTPException(404, "Item not found")
        updates = {}
        for k in ("name", "prompt", "skill", "output_mode", "config_json",
                   "source_item_id", "sort_order", "status"):
            v = getattr(req, k, None)
            if v is not None:
                updates[k] = v
        if updates:
            updates["updated_at"] = "CURRENT_TIMESTAMP"
            cols = ", ".join(f"{k} = ?" for k in updates if k != "updated_at")
            vals = [updates[k] for k in updates if k != "updated_at"]
            db.execute(f"UPDATE project_items SET {cols}, updated_at = CURRENT_TIMESTAMP WHERE id = ?",
                       vals + [item_id])
            db.commit()
        row = db.execute("SELECT * FROM project_items WHERE id = ?", (item_id,)).fetchone()
        return dict(row)
    finally:
        db.close()


@app.delete("/api/projects/{project_id}/items/{item_id}")
def delete_project_item(project_id: str, item_id: str):
    db = get_db()
    try:
        db.execute("DELETE FROM project_item_results WHERE project_item_id = ?", (item_id,))
        db.execute("DELETE FROM project_items WHERE id = ? AND project_id = ?",
                   (item_id, project_id))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/projects/{project_id}/items/copy-from/{source_project_id}")
def copy_project_items(project_id: str, source_project_id: str):
    """Copy all project_items from source project to target project."""
    db = get_db()
    try:
        source_items = db.execute(
            "SELECT * FROM project_items WHERE project_id = ? ORDER BY sort_order",
            (source_project_id,)).fetchall()
        count = 0
        # Build old→new ID mapping for source_item_id remapping
        id_map = {}
        for si in source_items:
            new_id = f"pi-{project_id}-{uuid.uuid4().hex[:8]}"
            id_map[si["id"]] = new_id
        for si in source_items:
            new_id = id_map[si["id"]]
            new_source = id_map.get(si["source_item_id"], si["source_item_id"]) if si["source_item_id"] else None
            db.execute(
                "INSERT INTO project_items (id, project_id, name, prompt, skill, "
                "output_mode, config_json, source_item_id, sort_order) "
                "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)",
                (new_id, project_id, si["name"], si["prompt"], si["skill"],
                 si["output_mode"], si["config_json"], new_source, si["sort_order"]))
            count += 1
        db.commit()
        return {"ok": True, "copied": count}
    finally:
        db.close()


# ── Project Item Results ──

@app.get("/api/projects/{project_id}/items/{item_id}/results")
def list_item_results(project_id: str, item_id: str):
    db = get_db()
    try:
        rows = db.execute(
            "SELECT * FROM project_item_results WHERE project_item_id = ? ORDER BY created_at DESC",
            (item_id,)).fetchall()
        return {"results": [dict(r) for r in rows]}
    finally:
        db.close()


@app.post("/api/projects/{project_id}/items/{item_id}/results")
def save_item_result(project_id: str, item_id: str, req: ProjectItemResultSave):
    db = get_db()
    try:
        existing = db.execute(
            "SELECT id FROM project_items WHERE id = ? AND project_id = ?",
            (item_id, project_id)).fetchone()
        if not existing:
            raise HTTPException(404, "Item not found")
        db.execute(
            "INSERT INTO project_item_results (project_item_id, content, content_type, "
            "file_path, quality_score) VALUES (?, ?, ?, ?, ?)",
            (item_id, req.content, req.content_type, req.file_path, req.quality_score))
        db.commit()
        rid = db.execute("SELECT last_insert_rowid()").fetchone()[0]
        row = db.execute("SELECT * FROM project_item_results WHERE id = ?", (rid,)).fetchone()
        return dict(row)
    finally:
        db.close()


# ── Project Copy ──

@app.post("/api/projects/{project_id}/copy")
def copy_project(project_id: str):
    """Copy a project and all its items (the project IS the template)."""
    db = get_db()
    try:
        src = db.execute("SELECT * FROM projects WHERE id = ?", (project_id,)).fetchone()
        if not src:
            raise HTTPException(404, "Source project not found")
        new_id = uuid.uuid4().hex[:12]
        from datetime import date
        today = date.today().strftime("%y%m%d")
        today_prefix = f"KH{today}-%"
        today_count = db.execute(
            "SELECT COUNT(*) FROM projects WHERE project_code LIKE ?", (today_prefix,)
        ).fetchone()[0]
        project_code = f"KH{today}-{today_count + 1:04d}"
        db.execute(
            "INSERT INTO projects (id, name, storage_path, copied_from_project_id, project_code) "
            "VALUES (?, ?, ?, ?, ?)",
            (new_id, f"{src['name']} (副本)", "", project_id, project_code))
        db.commit()
        # Copy project items
        copy_project_items(new_id, project_id)
        return {"ok": True, "project": {"id": new_id, "name": f"{src['name']} (副本)"}}
    finally:
        db.close()


# ── LLM Providers ──

from pydantic import BaseModel

class LLMProviderCreate(BaseModel):
    name: str
    api_key: str = ""
    base_url: str = "https://api.deepseek.com/v1"
    models: list[str] = []


@app.get("/api/llm/providers")
def list_llm_providers():
    db = get_db()
    try:
        rows = db.execute("SELECT * FROM llm_providers ORDER BY created_at").fetchall()
        providers = []
        for r in rows:
            p = dict(r)
            models_str = p.get("models", "[]")
            p["models"] = json.loads(models_str) if models_str else []
            if p.get("api_key"):
                key = p["api_key"]
                p["api_key"] = key[:8] + "***" if len(key) > 8 else "***"
            providers.append(p)
        return {"providers": providers}
    finally:
        db.close()


@app.post("/api/llm/providers")
def create_llm_provider(req: LLMProviderCreate):
    pid = uuid.uuid4().hex[:8]
    db = get_db()
    try:
        db.execute(
            "INSERT INTO llm_providers (id, name, api_key, base_url, models) VALUES (?, ?, ?, ?, ?)",
            (pid, req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False)))
        db.commit()
        return {"id": pid, "name": req.name}
    finally:
        db.close()


@app.put("/api/llm/providers/{provider_id}")
def update_llm_provider(provider_id: str, req: LLMProviderCreate):
    db = get_db()
    try:
        db.execute(
            "UPDATE llm_providers SET name=?, api_key=?, base_url=?, models=? WHERE id=?",
            (req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False), provider_id))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.delete("/api/llm/providers/{provider_id}")
def delete_llm_provider(provider_id: str):
    db = get_db()
    try:
        db.execute("DELETE FROM llm_providers WHERE id = ?", (provider_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/llm/providers/{provider_id}/test")
async def test_provider(provider_id: str):
    db = get_db()
    try:
        row = db.execute("SELECT * FROM llm_providers WHERE id = ?", (provider_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Provider not found")
        result = await test_connection(row["api_key"], row["base_url"])
        return result
    finally:
        db.close()


# ── TTS Providers ──

class TTSProviderCreate(BaseModel):
    name: str
    api_key: str = ""
    base_url: str = "https://dashscope.aliyuncs.com/api/v1"
    models: list[str] = []
    is_default: int = 0


class VoiceCreate(BaseModel):
    name: str
    provider_id: str
    voice_id: str
    description: str = ""
    is_default: int = 0


class VoiceUpdate(BaseModel):
    name: Optional[str] = None
    provider_id: Optional[str] = None
    voice_id: Optional[str] = None
    description: Optional[str] = None
    is_default: Optional[int] = None


@app.get("/api/tts/providers")
def list_tts_providers():
    db = get_db()
    try:
        rows = db.execute("SELECT * FROM tts_providers ORDER BY created_at").fetchall()
        providers = []
        for r in rows:
            p = dict(r)
            models_str = p.get("models", "[]")
            p["models"] = json.loads(models_str) if models_str else []
            if p.get("api_key"):
                key = p["api_key"]
                p["api_key"] = key[:8] + "***" if len(key) > 8 else "***"
            providers.append(p)
        return {"providers": providers}
    finally:
        db.close()


@app.post("/api/tts/providers")
def create_tts_provider(req: TTSProviderCreate):
    pid = uuid.uuid4().hex[:8]
    db = get_db()
    try:
        if req.is_default:
            db.execute("UPDATE tts_providers SET is_default = 0")
        db.execute(
            "INSERT INTO tts_providers (id, name, api_key, base_url, models, is_default) VALUES (?, ?, ?, ?, ?, ?)",
            (pid, req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False), req.is_default))
        db.commit()
        return {"id": pid, "name": req.name}
    finally:
        db.close()


@app.put("/api/tts/providers/{provider_id}")
def update_tts_provider(provider_id: str, req: TTSProviderCreate):
    db = get_db()
    try:
        if req.is_default:
            db.execute("UPDATE tts_providers SET is_default = 0")
        db.execute(
            "UPDATE tts_providers SET name=?, api_key=?, base_url=?, models=?, is_default=? WHERE id=?",
            (req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False), req.is_default, provider_id))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.delete("/api/tts/providers/{provider_id}")
def delete_tts_provider(provider_id: str):
    db = get_db()
    try:
        db.execute("DELETE FROM tts_providers WHERE id = ?", (provider_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/tts/providers/{provider_id}/test")
async def test_tts_provider(provider_id: str):
    import httpx
    db = get_db()
    try:
        row = db.execute("SELECT * FROM tts_providers WHERE id = ?", (provider_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Provider not found")
        # Test by calling the models list or a minimal synthesis check
        base_url = row["base_url"].rstrip("/")
        test_url = f"{base_url}/services/audio/tts/SpeechSynthesizer"
        payload = {
            "model": "cosyvoice-v3-flash",
            "input": {"text": "测试", "voice": "longanyang", "format": "mp3"},
        }
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                test_url,
                headers={"Authorization": f"Bearer {row['api_key']}", "Content-Type": "application/json"},
                json=payload,
            )
        if resp.status_code in (200, 201):
            return {"ok": True, "status": resp.status_code}
        return {"ok": False, "error": f"HTTP {resp.status_code}: {resp.text[:200]}"}
    except httpx.HTTPError as e:
        return {"ok": False, "error": str(e)}
    finally:
        db.close()


# ── ASR Providers ──

class ASRProviderCreate(BaseModel):
    name: str
    api_key: str = ""
    base_url: str = "https://dashscope.aliyuncs.com"
    models: list[str] = []
    is_default: int = 0


@app.get("/api/asr/providers")
def list_asr_providers():
    db = get_db()
    try:
        rows = db.execute("SELECT * FROM asr_providers ORDER BY created_at").fetchall()
        providers = []
        for r in rows:
            p = dict(r)
            models_str = p.get("models", "[]")
            p["models"] = json.loads(models_str) if models_str else []
            if p.get("api_key"):
                key = p["api_key"]
                p["api_key"] = key[:8] + "***" if len(key) > 8 else "***"
            providers.append(p)
        return {"providers": providers}
    finally:
        db.close()


@app.post("/api/asr/providers")
def create_asr_provider(req: ASRProviderCreate):
    pid = uuid.uuid4().hex[:8]
    db = get_db()
    try:
        if req.is_default:
            db.execute("UPDATE asr_providers SET is_default = 0")
        db.execute(
            "INSERT INTO asr_providers (id, name, api_key, base_url, models, is_default) VALUES (?, ?, ?, ?, ?, ?)",
            (pid, req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False), req.is_default))
        db.commit()
        return {"id": pid, "name": req.name}
    finally:
        db.close()


@app.put("/api/asr/providers/{provider_id}")
def update_asr_provider(provider_id: str, req: ASRProviderCreate):
    db = get_db()
    try:
        if req.is_default:
            db.execute("UPDATE asr_providers SET is_default = 0")
        db.execute(
            "UPDATE asr_providers SET name=?, api_key=?, base_url=?, models=?, is_default=? WHERE id=?",
            (req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False), req.is_default, provider_id))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.delete("/api/asr/providers/{provider_id}")
def delete_asr_provider(provider_id: str):
    db = get_db()
    try:
        db.execute("DELETE FROM asr_providers WHERE id = ?", (provider_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/asr/providers/{provider_id}/test")
async def test_asr_provider(provider_id: str):
    import httpx
    db = get_db()
    try:
        row = db.execute("SELECT * FROM asr_providers WHERE id = ?", (provider_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Provider not found")
        base_url = row["base_url"].rstrip("/")
        # Test by listing available models or hitting a get endpoint
        test_url = f"{base_url}/compatible-mode/v1/models"
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.get(
                test_url,
                headers={"Authorization": f"Bearer {row['api_key']}"},
            )
        if resp.status_code in (200, 201):
            return {"ok": True, "status": resp.status_code}
        # Try list files as lightweight connectivity check
        test_url2 = f"{base_url}/api/v1/files"
        async with httpx.AsyncClient(timeout=30) as client:
            resp2 = await client.get(
                test_url2,
                headers={"Authorization": f"Bearer {row['api_key']}"},
            )
        if resp2.status_code in (200, 401, 403):
            return {"ok": True, "status": resp2.status_code}
        return {"ok": False, "error": f"HTTP {resp2.status_code}: {resp2.text[:200]}"}
    except httpx.HTTPError as e:
        return {"ok": False, "error": str(e)}
    finally:
        db.close()


# ── Image Providers ──

class ImageProviderCreate(BaseModel):
    name: str
    api_key: str = ""
    base_url: str = "https://dashscope.aliyuncs.com/compatible-mode/v1"
    models: list[str] = []
    is_default: int = 0


@app.get("/api/image/providers")
def list_image_providers():
    db = get_db()
    try:
        rows = db.execute("SELECT * FROM image_providers ORDER BY created_at").fetchall()
        providers = []
        for r in rows:
            p = dict(r)
            models_str = p.get("models", "[]")
            p["models"] = json.loads(models_str) if models_str else []
            if p.get("api_key"):
                key = p["api_key"]
                p["api_key"] = key[:8] + "***" if len(key) > 8 else "***"
            providers.append(p)
        return {"providers": providers}
    finally:
        db.close()


@app.post("/api/image/providers")
def create_image_provider(req: ImageProviderCreate):
    pid = uuid.uuid4().hex[:8]
    db = get_db()
    try:
        if req.is_default:
            db.execute("UPDATE image_providers SET is_default = 0")
        db.execute(
            "INSERT INTO image_providers (id, name, api_key, base_url, models, is_default) VALUES (?, ?, ?, ?, ?, ?)",
            (pid, req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False), req.is_default))
        db.commit()
        return {"id": pid, "name": req.name}
    finally:
        db.close()


@app.put("/api/image/providers/{provider_id}")
def update_image_provider(provider_id: str, req: ImageProviderCreate):
    db = get_db()
    try:
        if req.is_default:
            db.execute("UPDATE image_providers SET is_default = 0")
        db.execute(
            "UPDATE image_providers SET name=?, api_key=?, base_url=?, models=?, is_default=? WHERE id=?",
            (req.name, req.api_key, req.base_url, json.dumps(req.models, ensure_ascii=False), req.is_default, provider_id))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.delete("/api/image/providers/{provider_id}")
def delete_image_provider(provider_id: str):
    db = get_db()
    try:
        db.execute("DELETE FROM image_providers WHERE id = ?", (provider_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/image/providers/{provider_id}/test")
async def test_image_provider(provider_id: str):
    import httpx
    db = get_db()
    try:
        row = db.execute("SELECT * FROM image_providers WHERE id = ?", (provider_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Provider not found")
        base_url = row["base_url"].rstrip("/")
        test_url = f"{base_url}/models"
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.get(
                test_url,
                headers={"Authorization": f"Bearer {row['api_key']}"},
            )
        if resp.status_code in (200, 201):
            return {"ok": True, "status": resp.status_code}
        return {"ok": False, "error": f"HTTP {resp.status_code}: {resp.text[:200]}"}
    except httpx.HTTPError as e:
        return {"ok": False, "error": str(e)}
    finally:
        db.close()


# ── Voice Library ──

@app.get("/api/voices")
def list_voices(provider_id: str = ""):
    db = get_db()
    try:
        if provider_id:
            rows = db.execute(
                "SELECT * FROM voices WHERE provider_id = ? ORDER BY is_default DESC, created_at",
                (provider_id,)).fetchall()
        else:
            rows = db.execute("SELECT * FROM voices ORDER BY is_default DESC, created_at").fetchall()
        return {"voices": [dict(r) for r in rows]}
    finally:
        db.close()


@app.post("/api/voices")
def create_voice(data: VoiceCreate):
    import uuid
    db = get_db()
    try:
        vid = uuid.uuid4().hex[:10]
        if data.is_default:
            db.execute("UPDATE voices SET is_default = 0")
        db.execute(
            "INSERT INTO voices (id, name, provider_id, voice_id, description, is_default) VALUES (?,?,?,?,?,?)",
            (vid, data.name, data.provider_id, data.voice_id, data.description, data.is_default),
        )
        db.commit()
        return {"id": vid, "ok": True}
    finally:
        db.close()


@app.put("/api/voices/{voice_id}")
def update_voice(voice_id: str, data: VoiceUpdate):
    db = get_db()
    try:
        existing = db.execute("SELECT * FROM voices WHERE id = ?", (voice_id,)).fetchone()
        if not existing:
            raise HTTPException(status_code=404, detail="Voice not found")
        updates = {}
        if data.name is not None:
            updates["name"] = data.name
        if data.provider_id is not None:
            updates["provider_id"] = data.provider_id
        if data.voice_id is not None:
            updates["voice_id"] = data.voice_id
        if data.description is not None:
            updates["description"] = data.description
        if data.is_default is not None:
            updates["is_default"] = data.is_default
            if data.is_default:
                db.execute("UPDATE voices SET is_default = 0")
        if updates:
            set_clause = ", ".join(f"{k} = ?" for k in updates)
            values = list(updates.values()) + [voice_id]
            db.execute(f"UPDATE voices SET {set_clause} WHERE id = ?", values)
            db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.delete("/api/voices/{voice_id}")
def delete_voice(voice_id: str):
    db = get_db()
    try:
        db.execute("DELETE FROM voices WHERE id = ?", (voice_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/voices/{voice_id}/preview")
async def preview_voice(voice_id: str):
    """Generate a short preview audio for a voice"""
    import httpx
    db = get_db()
    try:
        voice = db.execute(
            "SELECT v.*, p.api_key, p.base_url FROM voices v LEFT JOIN tts_providers p ON v.provider_id = p.id WHERE v.id = ?",
            (voice_id,)).fetchone()
        if not voice:
            raise HTTPException(status_code=404, detail="Voice not found")
        if not voice["api_key"]:
            raise HTTPException(status_code=400, detail="关联的 TTS 提供商未配置 API Key")

        base_url = voice["base_url"].rstrip("/")
        tts_url = f"{base_url}/services/audio/tts/SpeechSynthesizer"
        payload = {
            "model": "cosyvoice-v3-flash",
            "input": {"text": "你好，这是一条音色预览测试。", "voice": voice["voice_id"], "format": "mp3"},
        }
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                tts_url,
                headers={"Authorization": f"Bearer {voice['api_key']}", "Content-Type": "application/json"},
                json=payload,
            )
        if resp.status_code != 200:
            raise HTTPException(status_code=400, detail=f"预览生成失败: {resp.text[:200]}")
        result = resp.json()
        audio_url = result.get("output", {}).get("audio", {}).get("url", "")
        if not audio_url:
            raise HTTPException(status_code=400, detail="未获取到音频URL")
        return {"audio_url": audio_url}
    finally:
        db.close()


# ── LLM Calls ──

@app.post("/api/llm/generate")
async def llm_generate(req: LLMGenerateRequest):
    import time, sys
    t0 = time.time()
    print(f"[LLM-REQ] {time.strftime('%H:%M:%S')} provider={req.provider_id} model={req.model} "
          f"sys_len={len(req.system_prompt)} user_len={len(req.user_message)} temp={req.temperature}",
          flush=True)
    try:
        result = await generate(req.provider_id, req.model, req.system_prompt, req.user_message, req.temperature)
        dt = (time.time() - t0) * 1000
        print(f"[LLM-OK]  {time.strftime('%H:%M:%S')} model={req.model} dt={dt:.0f}ms "
              f"out_len={len(result) if result else 0}", flush=True)
        return {"content": result}
    except Exception as e:
        dt = (time.time() - t0) * 1000
        print(f"[LLM-ERR] {time.strftime('%H:%M:%S')} model={req.model} dt={dt:.0f}ms "
              f"err={e}", flush=True)
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/llm/generate-stream")
async def llm_generate_stream(req: LLMGenerateRequest):
    """Streaming LLM generate via SSE — provider-aware routing."""
    async def event_stream():
        try:
            async for text in generate_stream(
                req.provider_id, req.model,
                req.system_prompt, req.user_message,
                req.temperature,
            ):
                yield f"data: {json.dumps({'content': text}, ensure_ascii=False)}\n\n"
            yield "data: [DONE]\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@app.post("/api/llm/refine")
async def llm_refine(req: LLMRefineRequest):
    try:
        result = await refine(req.provider_id, req.model, req.instruction, req.selected_text, req.full_context)
        return {"content": result}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


# ── Image Generation ──

def _parse_host(base_url: str) -> str:
    """Extract scheme+host from base_url, dropping any path."""
    from urllib.parse import urlparse
    parsed = urlparse(base_url)
    return f"{parsed.scheme}://{parsed.netloc}"


# PPT layout → recommended image size
PPT_IMAGE_SIZES = {
    "full": "1280*720",       # 16:9 full-slide background
    "hero": "1344*576",       # 21:9 wide banner (image_hero)
    "content": "1280*800",    # 16:10 content illustration
    "square": "1024*1024",    # 1:1 card/icon
    "portrait": "960*1280",   # 3:4 vertical
}


@app.post("/api/image/generate")
async def image_generate(req: ImageGenerateRequest):
    import httpx
    db = get_db()
    try:
        if req.provider_id:
            row = db.execute(
                "SELECT * FROM image_providers WHERE id = ? AND is_enabled = 1",
                (req.provider_id,)).fetchone()
        else:
            row = db.execute(
                "SELECT * FROM image_providers WHERE is_default = 1 AND is_enabled = 1").fetchone()
            if not row:
                row = db.execute(
                    "SELECT * FROM image_providers WHERE is_enabled = 1 ORDER BY created_at").fetchone()
        if not row:
            raise HTTPException(status_code=400, detail="未配置图片生成提供商，请在模型设置中添加")
        provider = dict(row)

        # Resolve model: try exact match, then lowercase, then first configured
        saved_models = json.loads(provider["models"]) if provider["models"] else []
        model = req.model or (saved_models[0] if saved_models else "qwen-image-2.0-pro")
        # DashScope API expects lowercase model names
        model_lower = model.lower()

        host = _parse_host(provider["base_url"])
        url = f"{host}/api/v1/services/aigc/multimodal-generation/generation"

        # Build content array: reference images first, then prompt text
        content = []
        for img_url in req.reference_images:
            content.append({"image": img_url})
        content.append({"text": req.prompt})

        payload = {
            "model": model_lower,
            "input": {
                "messages": [
                    {"role": "user", "content": content}
                ]
            },
            "parameters": {
                "size": req.size,
                "n": req.n,
                "prompt_extend": req.prompt_extend,
                "watermark": req.watermark,
            },
        }
        if req.negative_prompt:
            payload["parameters"]["negative_prompt"] = req.negative_prompt
        if req.seed is not None:
            payload["parameters"]["seed"] = req.seed

        async with httpx.AsyncClient(timeout=120) as client:
            resp = await client.post(
                url,
                headers={
                    "Authorization": f"Bearer {provider['api_key']}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )

        if resp.status_code in (200, 201):
            data = resp.json()
            images = []
            choices = data.get("output", {}).get("choices", [])
            for choice in choices:
                for item in choice.get("message", {}).get("content", []):
                    if "image" in item:
                        images.append({"url": item["image"]})
            usage = data.get("usage", {})
            return {"ok": True, "images": images, "model": model_lower,
                    "usage": {"width": usage.get("width"), "height": usage.get("height"),
                              "count": usage.get("image_count")}}
        else:
            return {"ok": False, "error": f"HTTP {resp.status_code}: {resp.text[:500]}"}
    except httpx.HTTPError as e:
        raise HTTPException(status_code=400, detail=str(e))
    finally:
        db.close()


# ── Prompts ──

class PromptCreate(BaseModel):
    name: str
    category: str
    system_prompt: str = ""
    skill_template: str = ""


class PromptUpdate(BaseModel):
    name: str = None
    category: str = None
    system_prompt: str = None
    skill_template: str = None
    change_note: str = ""


class PromptImport(BaseModel):
    data: list


@app.get("/api/prompts/export")
def api_export_prompts():
    return {"prompts": export_prompts()}


@app.post("/api/prompts/import")
def api_import_prompts(req: PromptImport):
    return import_prompts(req.data)


@app.get("/api/prompts")
def api_list_prompts(category: str = None):
    return {"prompts": list_prompts(category)}


@app.post("/api/prompts")
def api_create_prompt(req: PromptCreate):
    prompt = create_prompt(req.name, req.category, req.system_prompt, req.skill_template)
    return prompt


@app.get("/api/prompts/{prompt_id}")
def api_get_prompt(prompt_id: str):
    prompt = get_prompt(prompt_id)
    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt not found")
    return prompt


@app.put("/api/prompts/{prompt_id}")
def api_update_prompt(prompt_id: str, req: PromptUpdate):
    prompt = update_prompt(
        prompt_id, req.name, req.category,
        req.system_prompt, req.skill_template, req.change_note)
    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt not found")
    return prompt


@app.delete("/api/prompts/{prompt_id}")
def api_delete_prompt(prompt_id: str):
    delete_prompt(prompt_id)
    return {"ok": True}


@app.get("/api/prompts/{prompt_id}/versions")
def api_list_versions(prompt_id: str):
    prompt = get_prompt(prompt_id)
    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt not found")
    return {"versions": prompt["versions"]}


@app.post("/api/prompts/{prompt_id}/rollback")
def api_rollback(prompt_id: str, req: dict):
    prompt = rollback_version(prompt_id, req["version"])
    if not prompt:
        raise HTTPException(status_code=404, detail="Prompt or version not found")
    return prompt


@app.post("/api/prompts/{prompt_id}/diff")
def api_diff(prompt_id: str, req: dict):
    result = diff_versions(prompt_id, req["version_a"], req["version_b"])
    if not result:
        raise HTTPException(status_code=404, detail="Versions not found")
    return result


@app.post("/api/prompts/{prompt_id}/set-default")
def api_set_default(prompt_id: str):
    set_default(prompt_id)
    return {"ok": True}


# ── Templates ──

class TemplateCreate(BaseModel):
    name: str
    type: str  # "ppt" or "sop"
    file_path: str = ""
    prompt: str = ""
    skill: str = ""
    rules: str = "{}"
    linked_skill_id: str = ""
    branding_config: str = "{}"


@app.get("/api/templates")
def list_templates(type: str = None):
    db = get_db()
    try:
        if type:
            rows = db.execute("SELECT * FROM templates WHERE type = ? ORDER BY created_at DESC", (type,)).fetchall()
        else:
            rows = db.execute("SELECT * FROM templates ORDER BY type, created_at DESC").fetchall()
        return {"templates": [dict(r) for r in rows]}
    finally:
        db.close()


@app.post("/api/templates")
def create_template(req: TemplateCreate):
    tid = uuid.uuid4().hex[:8]
    db = get_db()
    try:
        db.execute(
            "INSERT INTO templates (id, name, type, file_path, prompt, skill, rules, linked_skill_id, branding_config) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)",
            (tid, req.name, req.type, req.file_path, req.prompt, req.skill, req.rules, req.linked_skill_id, req.branding_config))
        db.commit()
        return {"id": tid, "name": req.name}
    finally:
        db.close()


@app.put("/api/templates/{template_id}/toggle-enabled")
def toggle_template_enabled(template_id: str):
    """Toggle the enabled state of a template.

    Only enabled templates appear in the Step 3 style selector.
    """
    db = get_db()
    try:
        row = db.execute("SELECT id, enabled FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not row:
            raise HTTPException(404, "模板不存在")
        new_state = 0 if row["enabled"] == 1 else 1
        db.execute("UPDATE templates SET enabled = ? WHERE id = ?", (new_state, template_id))
        db.commit()
        return {"ok": True, "enabled": new_state == 1}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
    finally:
        db.close()


@app.put("/api/templates/{template_id}")
def update_template(template_id: str, req: TemplateCreate):
    db = get_db()
    try:
        existing = db.execute("SELECT file_path, thumbnail_path FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not existing:
            raise HTTPException(404, "Template not found")
        # Guard: never overwrite a valid file_path with an empty one
        final_file_path = req.file_path if req.file_path else (existing["file_path"] or "")
        db.execute(
            "UPDATE templates SET name=?, type=?, file_path=?, prompt=?, skill=?, rules=?, linked_skill_id=?, branding_config=? WHERE id=?",
            (req.name, req.type, final_file_path, req.prompt, req.skill, req.rules, req.linked_skill_id, req.branding_config, template_id))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.delete("/api/templates/{template_id}")
def delete_template(template_id: str):
    db = get_db()
    try:
        row = db.execute("SELECT is_default FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not row:
            raise HTTPException(404, "Template not found")
        if row["is_default"]:
            raise HTTPException(400, "默认模板不可删除，请先将其他模板设为默认后再删除")
        db.execute("DELETE FROM templates WHERE id = ?", (template_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/templates/{template_id}/set-default")
def set_template_default(template_id: str):
    db = get_db()
    try:
        row = db.execute("SELECT * FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Template not found")
        db.execute("UPDATE templates SET is_default = 0 WHERE type = ?", (row["type"],))
        db.execute("UPDATE templates SET is_default = 1 WHERE id = ?", (template_id,))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


def _generate_pptx_thumbnail(pptx_path: str, template_id: str) -> str | None:
    """Export slide 1 as a full-slide thumbnail PNG. Falls back to extracting first embedded image."""
    # Preferred: full-slide PNG via PowerPoint COM
    thumb_path = _export_slide_thumbnail(pptx_path, template_id)
    if thumb_path:
        return thumb_path

    # Fallback: extract first embedded picture from slide 1
    try:
        from pptx import Presentation
        from pptx.shapes.picture import Picture
        prs = Presentation(pptx_path)
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if isinstance(shape, Picture):
                    image = shape.image
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    thumb_name = f"{template_id}_thumb.{ext}"
                    thumb_path = os.path.join(THUMBNAIL_DIR, thumb_name)
                    with open(thumb_path, "wb") as f:
                        f.write(image.blob)
                    return thumb_path
    except Exception as e:
        logging.getLogger("uvicorn").info(f"Thumbnail picture fallback failed: {e}")

    return None


def _export_slide_thumbnail(pptx_path: str, template_id: str) -> str | None:
    """Export slide 1 as a PNG thumbnail using PowerPoint COM."""
    try:
        import subprocess
        thumb_name = f"{template_id}_thumb.png"
        thumb_path = os.path.join(THUMBNAIL_DIR, thumb_name)
        ps_script = f'''
$ppt = New-Object -ComObject PowerPoint.Application
$ppt.Visible = 0
try {{
    $pres = $ppt.Presentations.Open("{pptx_path}", $true, $false, $false)
    $pres.Slides[1].Export("{thumb_path}", "PNG", 960, 540)
    $pres.Close()
}} finally {{
    $ppt.Quit()
    [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
}}
Write-Output "1"
'''
        result = subprocess.run(["powershell", "-Command", ps_script],
                              capture_output=True, text=True, timeout=30)
        if os.path.exists(thumb_path) and os.path.getsize(thumb_path) > 0:
            return thumb_path
    except Exception as e:
        logging.getLogger("uvicorn").info(f"Slide thumbnail fallback failed: {e}")
    return None


@app.post("/api/templates/{template_id}/upload")
async def upload_template_file(template_id: str, file: UploadFile = File(...)):
    db = get_db()
    try:
        existing = db.execute("SELECT id FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not existing:
            raise HTTPException(404, "Template not found")
        if not file.filename or not file.filename.endswith('.pptx'):
            raise HTTPException(400, "请上传 .pptx 格式的模板文件")
        tmpl_dir = os.path.join(BASE_DIR, "data", "templates")
        os.makedirs(tmpl_dir, exist_ok=True)
        safe_name = f"{template_id}_{file.filename}"
        file_path = os.path.join(tmpl_dir, safe_name)
        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)
        db.execute("UPDATE templates SET file_path = ? WHERE id = ?", (file_path, template_id))
        # Auto-generate thumbnail from first slide
        thumb_path = _generate_pptx_thumbnail(file_path, template_id)
        if thumb_path:
            db.execute("UPDATE templates SET thumbnail_path = ? WHERE id = ?", (thumb_path, template_id))
        db.commit()
        return {"ok": True, "file_path": file_path, "filename": file.filename,
                "thumbnail_path": thumb_path}
    finally:
        db.close()


@app.post("/api/templates/{template_id}/upload-thumbnail")
async def upload_template_thumbnail(template_id: str, file: UploadFile = File(...)):
    db = get_db()
    try:
        existing = db.execute("SELECT id FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not existing:
            raise HTTPException(404, "Template not found")
        if not file.filename:
            raise HTTPException(400, "No file provided")
        ext = os.path.splitext(file.filename)[1].lower()
        if ext not in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"):
            raise HTTPException(400, "请上传图片格式文件 (png/jpg/gif/webp/svg)")
        safe_name = f"{template_id}_thumb{ext}"
        file_path = os.path.join(THUMBNAIL_DIR, safe_name)
        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)
        db.execute("UPDATE templates SET thumbnail_path = ? WHERE id = ?", (file_path, template_id))
        db.commit()
        return {"ok": True, "thumbnail_path": file_path, "filename": file.filename}
    finally:
        db.close()


@app.post("/api/templates/{template_id}/reset-thumbnail")
def reset_template_thumbnail(template_id: str):
    db = get_db()
    try:
        row = db.execute("SELECT id, file_path FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not row:
            raise HTTPException(404, "Template not found")
        if not row["file_path"] or not os.path.exists(row["file_path"]):
            raise HTTPException(400, "请先上传 PPTX 模板文件")
        thumb_path = _generate_pptx_thumbnail(row["file_path"], template_id)
        if thumb_path:
            db.execute("UPDATE templates SET thumbnail_path = ? WHERE id = ?", (thumb_path, template_id))
        else:
            db.execute("UPDATE templates SET thumbnail_path = NULL WHERE id = ?", (template_id,))
        db.commit()
        return {"ok": True, "thumbnail_path": thumb_path}
    finally:
        db.close()


@app.get("/api/templates/{template_id}/file")
def serve_template_file(template_id: str):
    """Serve the template PPTX file — generated from prompt+SKILL+rules, or original upload as fallback."""
    db = get_db()
    try:
        row = db.execute(
            "SELECT file_path, name, prompt, skill, rules FROM templates WHERE id = ?",
            (template_id,)).fetchone()
        if not row:
            raise HTTPException(404, "Template not found")
        name = row["name"]
        prompt = row["prompt"] or ""
        skill = row["skill"] or ""
        rules_str = row["rules"] or "{}"
        original_path = row["file_path"]

        # Generated template cache path
        gen_dir = os.path.join(BASE_DIR, "data", "templates", "_generated")
        os.makedirs(gen_dir, exist_ok=True)
        gen_path = os.path.join(gen_dir, f"{template_id}.pptx")

        # Serve cached generated template if available
        if os.path.exists(gen_path):
            serve_path = gen_path
        elif skill and prompt:
            # Find original upload for layout source
            layout_source = original_path if (original_path and os.path.exists(original_path)) else None
            if not layout_source:
                raise HTTPException(404, "No layout source file available")

            # Try AI generation from prompt+SKILL+rules
            rules = {}
            try:
                rules = json.loads(rules_str)
            except Exception:
                pass

            prov = db.execute(
                "SELECT id, models FROM llm_providers WHERE is_enabled=1 LIMIT 1").fetchone()
            if prov:
                models = json.loads(prov["models"] or "[]")
                model = models[0] if models else ""
                if model:
                    from pptx import Presentation
                    from services.ppt_service import generate_template_pptx
                    try:
                        prs = Presentation(layout_source)
                        generate_template_pptx(prs, prompt, skill, rules, prov["id"], model, gen_path)
                        # DO NOT overwrite file_path — keep original as layout source for future regenerations
                        logging.getLogger("uvicorn").info(f"Template {template_id} generated from SKILL")
                        serve_path = gen_path
                    except Exception as e:
                        logging.getLogger("uvicorn").warning(
                            f"Template generation failed, serving original: {e}")
                        serve_path = layout_source
                else:
                    serve_path = layout_source
            else:
                serve_path = layout_source
        elif original_path and os.path.exists(original_path):
            serve_path = original_path
        else:
            raise HTTPException(404, "No template file available")

        if not os.path.exists(serve_path):
            raise HTTPException(404, "Template file not found on disk")

        # Also copy to global save path
        try:
            save_dir = _get_global_save_path()
            os.makedirs(save_dir, exist_ok=True)
            dest = os.path.join(save_dir, name + ".pptx")
            shutil.copy2(serve_path, dest)
            logging.getLogger("uvicorn").info(f"Template copied to {dest}")
        except Exception as e:
            logging.getLogger("uvicorn").warning(f"Failed to copy template to save path: {e}")

        from urllib.parse import quote
        safe_name = name + ".pptx"
        return FileResponse(serve_path,
                          media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                          filename=safe_name,
                          headers={"Content-Disposition": f"inline; filename*=UTF-8''{quote(safe_name)}"})
    finally:
        db.close()


SLIDES_DIR = os.path.join(BASE_DIR, "data", "slides")


def _export_pptx_slides(pptx_path: str, template_id: str) -> list[str]:
    """Export all slides of a PPTX as PNG images. Returns list of filenames."""
    import subprocess, uuid
    os.makedirs(SLIDES_DIR, exist_ok=True)
    out_dir = os.path.join(SLIDES_DIR, f"{template_id}_slides")
    os.makedirs(out_dir, exist_ok=True)
    # Check if already exported
    existing = [f for f in os.listdir(out_dir) if f.endswith(".png")]
    if existing:
        existing.sort()
        return [f"{template_id}_slides/{f}" for f in existing]
    # Use PowerShell + PowerPoint COM to export slides
    ps = f'''
$ppt = New-Object -ComObject PowerPoint.Application
$ppt.Visible = 0
$pres = $ppt.Presentations.Open("{pptx_path}")
$count = $pres.Slides.Count
for ($i = 1; $i -le $count; $i++) {{
    $pres.Slides[$i].Export("{out_dir}\\slide_$i.png", "PNG", 960, 540)
}}
$pres.Close()
$ppt.Quit()
[System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
Write-Output $count
'''
    try:
        result = subprocess.run(["powershell", "-Command", ps], capture_output=True, text=True, timeout=60)
        count = result.stdout.strip()
        logging.getLogger("uvicorn").info(f"Exported {count} slides for {template_id}")
    except Exception as e:
        logging.getLogger("uvicorn").info(f"Slide export failed: {e}")
    slides = sorted([f for f in os.listdir(out_dir) if f.endswith(".png")])
    return [f"{template_id}_slides/{f}" for f in slides]


@app.get("/api/templates/{template_id}/slide-thumb")
def get_template_slide_thumb(template_id: str):
    """Generate a simple PNG thumbnail of the first slide using python-pptx + Pillow.
    No COM required — reads shapes and text from the PPTX, renders a basic representation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from PIL import Image, ImageDraw, ImageFont
    import io

    db = get_db()
    try:
        row = db.execute("SELECT file_path FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not row or not row["file_path"]:
            raise HTTPException(404, "Template file not found")
        pptx_path = row["file_path"]
        # Prefer original upload for thumbnail
        gen_path = os.path.join(BASE_DIR, "data", "templates", "_generated", f"{template_id}.pptx")
        if os.path.exists(gen_path):
            pptx_path = gen_path
        if not os.path.exists(pptx_path):
            raise HTTPException(404, "PPTX file not found on disk")
    finally:
        db.close()

    prs = Presentation(pptx_path)
    if not prs.slides:
        raise HTTPException(404, "No slides in template")

    slide = prs.slides[0]
    sw = prs.slide_width or 12192000   # EMU
    sh = prs.slide_height or 6858000

    # Thumbnail size
    THUMB_W, THUMB_H = 640, 360
    scale_x = THUMB_W / sw
    scale_y = THUMB_H / sh

    img = Image.new("RGB", (THUMB_W, THUMB_H), "#F0EDE8")
    draw = ImageDraw.Draw(img)

    # Try to detect background from first shape that covers most of the slide
    bg_color = None
    for shape in slide.shapes:
        try:
            fill = shape.fill
            if fill and fill.type is not None:
                fc = fill.fore_color
                if fc and fc.type is not None:
                    try:
                        bg_color = fc.rgb
                        break
                    except Exception:
                        pass
        except Exception:
            pass

    if bg_color:
        try:
            hex_str = str(bg_color)
            if len(hex_str) == 6:
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
            else:
                r, g, b = 240, 237, 232
        except Exception:
            r, g, b = 240, 237, 232
        # Only use if not pure white/black
        if (r, g, b) not in ((255, 255, 255), (0, 0, 0)):
            img = Image.new("RGB", (THUMB_W, THUMB_H), (r, g, b))
            draw = ImageDraw.Draw(img)

    # Load a basic font
    font_paths = [
        "C:/Windows/Fonts/msyh.ttc",       # Microsoft YaHei
        "C:/Windows/Fonts/simsun.ttc",      # SimSun
        "C:/Windows/Fonts/arial.ttf",
    ]
    font_lg = None
    font_sm = None
    for fp in font_paths:
        if os.path.exists(fp):
            try:
                font_lg = ImageFont.truetype(fp, 24)
                font_sm = ImageFont.truetype(fp, 14)
                break
            except Exception:
                pass
    if font_lg is None:
        font_lg = ImageFont.load_default()
        font_sm = ImageFont.load_default()

    # Render text shapes
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = tf.text.strip()
        if not text:
            continue
        x = int(shape.left * scale_x) if shape.left else 0
        y = int(shape.top * scale_y) if shape.top else 0
        w = int(shape.width * scale_x) if shape.width else THUMB_W
        h = int(shape.height * scale_y) if shape.height else THUMB_H

        # Clip to image bounds
        x = max(0, min(x, THUMB_W - 10))
        y = max(0, min(y, THUMB_H - 10))
        w = min(w, THUMB_W - x)
        h = min(h, THUMB_H - y)

        # Determine text color from first run
        text_color = (40, 40, 40)
        try:
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.font.color and r.font.color.rgb:
                        cr = r.font.color.rgb
                        text_color = ((cr >> 16) & 0xFF, (cr >> 8) & 0xFF, cr & 0xFF)
                    break
                break
        except Exception:
            pass

        font = font_lg if any(r.font.size and r.font.size >= Pt(18) for p in tf.paragraphs for r in p.runs) else font_sm

        # Draw first 2 lines of text
        lines = text.split('\n')[:3]
        line_h = 16
        for li, line in enumerate(lines):
            if li * line_h >= h - 4:
                break
            # Truncate long lines
            max_chars = max(4, int(w / 10))
            if len(line) > max_chars:
                line = line[:max_chars - 2] + '..'
            draw.text((x + 4, y + 4 + li * line_h), line, fill=text_color, font=font_sm)

    # Save to in-memory PNG
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return StreamingResponse(buf, media_type="image/png",
                             headers={"Cache-Control": "public, max-age=3600"})


@app.get("/api/templates/{template_id}/slides-content")
def get_template_slides_content(template_id: str):
    """Return full slide shape data (position, color, font, text) for visual in-browser preview.
    Each shape is rendered as a positioned HTML element in the frontend."""
    db = get_db()
    try:
        row = db.execute("SELECT file_path FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not row or not row["file_path"]:
            raise HTTPException(404, "Template file not found")
        pptx_path = row["file_path"]
        gen_path = os.path.join(BASE_DIR, "data", "templates", "_generated", f"{template_id}.pptx")
        if os.path.exists(gen_path):
            pptx_path = gen_path
        if not os.path.exists(pptx_path):
            raise HTTPException(404, "PPTX file not found on disk")
    finally:
        db.close()

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    import copy

    prs = Presentation(pptx_path)

    def _rgb_str(fc):
        """Convert python-pptx color to '#RRGGBB' string or None."""
        try:
            if fc is None or fc.type is None:
                return None
            s = str(fc.rgb)
            if len(s) == 6:
                return '#' + s
        except Exception:
            pass
        return None

    def _extract_fill(shape):
        """Extract fill color from shape."""
        try:
            fill = shape.fill
            if fill and fill.type is not None:
                return _rgb_str(fill.fore_color)
        except Exception:
            pass
        return None

    def _extract_text_runs(shape):
        """Extract rich text runs from shape."""
        runs_data = []
        try:
            if not shape.has_text_frame:
                return []
            tf = shape.text_frame
            for p in tf.paragraphs:
                align = None
                try:
                    from pptx.enum.text import PP_ALIGN
                    align_map = {
                        PP_ALIGN.LEFT: 'left', PP_ALIGN.CENTER: 'center',
                        PP_ALIGN.RIGHT: 'right', PP_ALIGN.JUSTIFY: 'justify',
                    }
                    align = align_map.get(p.alignment, 'left')
                except Exception:
                    align = 'left'
                for r in p.runs:
                    size_pt = None
                    try:
                        if r.font.size:
                            size_pt = round(r.font.size / 12700.0, 1)
                    except Exception:
                        pass
                    color = None
                    try:
                        if r.font.color and r.font.color.rgb:
                            color = _rgb_str(r.font.color)
                    except Exception:
                        pass
                    bold = r.font.bold if r.font.bold is not None else False
                    italic = r.font.italic if r.font.italic is not None else False
                    runs_data.append({
                        "text": r.text,
                        "size": size_pt,
                        "color": color,
                        "bold": bold,
                        "italic": italic,
                        "align": align,
                    })
        except Exception:
            pass
        return runs_data

    def _shape_type_name(shape):
        """Get a simple shape type name."""
        try:
            name = shape.shape_type
            return str(name).split('.')[-1].split('(')[0].strip().upper() if name else 'UNKNOWN'
        except Exception:
            return 'UNKNOWN'

    slides = []
    for idx, slide in enumerate(prs.slides):
        shapes_data = []
        # Try to get slide background
        bg_color = None
        try:
            bg = slide.background
            if bg.fill and bg.fill.type is not None:
                bg_color = _rgb_str(bg.fill.fore_color)
        except Exception:
            pass

        for shape in slide.shapes:
            shape_info = {
                "left": shape.left if shape.left is not None else 0,
                "top": shape.top if shape.top is not None else 0,
                "width": shape.width if shape.width is not None else 0,
                "height": shape.height if shape.height is not None else 0,
                "rotation": shape.rotation if shape.rotation else 0,
                "fill": _extract_fill(shape),
                "name": shape.name or "",
                "s_type": _shape_type_name(shape),
                "runs": _extract_text_runs(shape),
            }
            shapes_data.append(shape_info)

        slides.append({
            "num": idx + 1,
            "bg_color": bg_color,
            "shapes": shapes_data,
        })

    return {
        "slides": slides,
        "total": len(slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
    }


@app.post("/api/templates/{template_id}/preview-slides")
def preview_template_slides(template_id: str):
    """Export all slides as images and return their URLs."""
    db = get_db()
    try:
        row = db.execute("SELECT file_path FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not row or not row["file_path"]:
            raise HTTPException(404, "Template file not found")
        pptx_path = row["file_path"]
        # Prefer cached generated template over original upload
        gen_path = os.path.join(BASE_DIR, "data", "templates", "_generated", f"{template_id}.pptx")
        if os.path.exists(gen_path):
            pptx_path = gen_path
        if not os.path.exists(pptx_path):
            raise HTTPException(404, "PPTX file not found on disk")
    finally:
        db.close()
    slides = _export_pptx_slides(pptx_path, template_id)
    return {"slides": slides}


@app.get("/api/slides/{path:path}")
def serve_slide_image(path: str):
    """Serve exported slide images."""
    filepath = os.path.join(SLIDES_DIR, path)
    if not os.path.exists(filepath):
        raise HTTPException(404, "Slide not found")
    return FileResponse(filepath, media_type="image/png")


@app.get("/api/thumbnails/{filename}")
def serve_thumbnail(filename: str):
    filepath = os.path.join(THUMBNAIL_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Thumbnail not found")
    media_map = {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".gif": "image/gif", ".svg": "image/svg+xml", ".webp": "image/webp",
    }
    ext = os.path.splitext(filename)[1].lower()
    return FileResponse(filepath, media_type=media_map.get(ext, "application/octet-stream"))


@app.get("/api/templates/for-stage/{stage_type}")
def list_templates_for_stage(stage_type: str):
    type_map = {"sop": "ppt", "daoPpt": "ppt", "yanxiPpt": "ppt"}
    db_type = type_map.get(stage_type, "ppt")
    db = get_db()
    try:
        if db_type == "ppt":
            rows = db.execute(
                "SELECT id, name, type, file_path, prompt, skill, branding_config, is_default FROM templates WHERE type = 'style' AND enabled = 1 ORDER BY created_at ASC").fetchall()
        else:
            rows = db.execute(
                "SELECT id, name, type, file_path, prompt, skill, branding_config, is_default FROM templates WHERE type = ? ORDER BY is_default DESC, created_at ASC",
                (db_type,)).fetchall()
        items = []
        for r in rows:
            rdict = dict(r)
            items.append({
                "id": rdict["id"],
                "name": rdict["name"],
                "type": rdict["type"],
                "prompt": rdict.get("prompt") or "",
                "skill": rdict.get("skill") or "",
                "isDefault": rdict.get("is_default") == 1,
                "hasFile": bool(rdict.get("file_path") and os.path.exists(rdict.get("file_path") or "")),
            })
        return {"templates": items}
    finally:
        db.close()


# ── Template Analysis ──

def extract_pptx_structure(file_path: str) -> dict:
    from pptx import Presentation
    from lxml import etree
    prs = Presentation(file_path)
    slides = []
    all_fonts = set()
    all_colors = set()
    all_theme_colors = set()
    all_theme_fonts = set()

    # Extract theme colors from slide master XML
    for master in prs.slide_masters:
        master_xml = master.element
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        # Theme colors (dk1, lt1, dk2, lt2, accent1-6, hlink, folHlink)
        for clr in master_xml.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr'):
            val = clr.get('val')
            if val:
                all_theme_colors.add(val)
                all_colors.add(val)
        for clr in master_xml.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}sysClr'):
            val = clr.get('lastClr')
            if val:
                all_theme_colors.add(val)
                all_colors.add(val)
        # Theme fonts (latin, ea, cs)
        for font_elem in master_xml.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}latin'):
            typeface = font_elem.get('typeface')
            if typeface:
                all_theme_fonts.add(typeface)
                all_fonts.add(typeface)
        for font_elem in master_xml.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}ea'):
            typeface = font_elem.get('typeface')
            if typeface:
                all_theme_fonts.add(typeface)
                all_fonts.add(typeface)

    for i, slide in enumerate(prs.slides):
        info = {"index": i + 1, "layout": slide.slide_layout.name, "placeholders": [], "shapes": []}
        for shape in slide.placeholders:
            font_info = {}
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        f = run.font
                        if f.name:
                            font_info["name"] = f.name
                            all_fonts.add(f.name)
                        if f.size:
                            font_info["size_pt"] = round(f.size / 12700, 1)
                        font_info["bold"] = f.bold
                        try:
                            if f.color and f.color.rgb:
                                all_colors.add(str(f.color.rgb))
                        except Exception:
                            pass
            info["placeholders"].append({
                "idx": shape.placeholder_format.idx,
                "type": str(shape.placeholder_format.type),
                "name": shape.name,
                "text_preview": (shape.text or "")[:100],
                "font": font_info
            })
        for shape in slide.shapes:
            if not shape.is_placeholder:
                shape_info = {
                    "type": str(shape.shape_type),
                    "name": shape.name,
                    "text_preview": (shape.text or "")[:100] if shape.has_text_frame else ""
                }
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            f = run.font
                            if f.name:
                                all_fonts.add(f.name)
                            try:
                                if f.color and f.color.rgb:
                                    all_colors.add(str(f.color.rgb))
                            except Exception:
                                pass
                info["shapes"].append(shape_info)
        slides.append(info)

    master_layouts = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            master_layouts.append(layout.name)

    return {
        "slide_count": len(slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_layouts": [lyt.name for lyt in prs.slide_layouts],
        "master_layouts": master_layouts,
        "fonts_used": sorted(list(all_fonts))[:30],
        "theme_fonts": sorted(list(all_theme_fonts))[:10],
        "colors_used": sorted(list(all_colors))[:30],
        "theme_colors": sorted(list(all_theme_colors))[:20],
        "typography_extracted": _extract_typography(prs),
        "slides": slides
    }


@app.post("/api/templates/{template_id}/analyze")
async def analyze_template(template_id: str, stage_type: str = "daoPpt",
                            provider_id: str = "", model: str = ""):
    """Mechanically extract visual data from PPTX and save to template.

    No AI call — colors, fonts, and typography are extracted directly from
    the PPTX file. Column config's rules structure is used as-is.
    """
    from services.ppt_designer import _extract_dominant_colors
    db = get_db()
    try:
        tmpl = db.execute("SELECT * FROM templates WHERE id = ?", (template_id,)).fetchone()
        if not tmpl:
            raise HTTPException(404, "模板不存在")
        if not tmpl["file_path"] or not os.path.exists(tmpl["file_path"]):
            raise HTTPException(400, "请先上传 .pptx 模板文件")

        # Use built-in default design_rules structure (column_configs rules removed)
        rules = {}

        # Mechanical extraction from PPTX
        structure = extract_pptx_structure(tmpl["file_path"])
        typography = structure.get("typography_extracted", {}) or {}
        dominant = _extract_dominant_colors(tmpl["file_path"])

        # Populate design_rules.colors from extracted data
        design_rules = rules.get("design_rules", {})
        if not isinstance(design_rules, dict):
            design_rules = {}

        design_rules["colors"] = {
            "primary": dominant.get("primary", "#C02E2E"),
            "accent": dominant.get("accent", "#FF6D01"),
            "background": dominant.get("background", "#FFFFFF"),
            "text": dominant.get("text", "#333333"),
            "light_text": dominant.get("light_text", "#FFFFFF"),
        }

        # Populate design_rules.fonts from extracted data
        fonts_used = structure.get("fonts_used", [])
        real_fonts = [f for f in fonts_used if not f.startswith('+')
                      and f not in ('Calibri', 'Arial')]  # theme refs & generic defaults
        font_name = real_fonts[0] if real_fonts else "Microsoft YaHei"
        # Fall back if likely a theme font reference
        if font_name.startswith('+'):
            font_name = "Microsoft YaHei"
        design_rules["fonts"] = {
            "font_name": font_name,
            "title_size": int(typography.get("title_font_size_pt") or 36),
            "body_size": int(typography.get("body_font_size_pt") or 18),
        }

        rules["design_rules"] = design_rules

        # Typography profile
        typography_profile = {
            "body_font_size_pt": typography.get("body_font_size_pt") or 18,
            "title_font_size_pt": typography.get("title_font_size_pt") or 36,
            "line_height_ratio": typography.get("line_height_ratio") or 1.2,
        }

        rules_json = json.dumps(rules, ensure_ascii=False)
        db.execute(
            "UPDATE templates SET typography_profile = ?, rules = ? WHERE id = ?",
            (json.dumps(typography_profile, ensure_ascii=False), rules_json, template_id))
        db.commit()

        return {
            "ok": True,
            "typography_profile": typography_profile,
            "rules": rules,
            "colors_used": structure.get("colors_used", []),
            "fonts_used": structure.get("fonts_used", []),
        }
    finally:
        db.close()












# ── Video ──

from services.video_service import download_video, get_progress


class VideoDownloadRequest(BaseModel):
    url: str
    cookies_path: Optional[str] = None
    project_id: Optional[str] = None
    asr_model: Optional[str] = "fun-asr"
    asr_provider_id: Optional[str] = None


COOKIES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "cookies")
os.makedirs(COOKIES_DIR, exist_ok=True)

@app.post("/api/video/upload-cookies")
async def api_upload_cookies(file: UploadFile = File(...)):
    if not file.filename or not file.filename.endswith('.txt'):
        raise HTTPException(400, "请上传 .txt 格式的 cookies 文件")
    file_id = uuid.uuid4().hex[:8]
    file_path = os.path.join(COOKIES_DIR, f"{file_id}.txt")
    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)
    return {"cookies_path": file_path, "filename": file.filename}


@app.post("/api/video/download")
def api_download_video(req: VideoDownloadRequest):
    result = download_video(req.url, req.cookies_path, req.project_id, req.asr_model or "fun-asr", req.asr_provider_id)
    return result


@app.get("/api/video/progress/{task_id}")
def api_video_progress(task_id: str):
    data = get_progress(task_id)
    svc_status = data.get("status", "")
    # Normalize status codes for frontend
    if svc_status == "done":
        fe_status = "completed"
    elif svc_status == "error":
        fe_status = "failed"
    else:
        fe_status = svc_status
    return {
        "task_id": task_id,
        "status": fe_status,
        "percent": data.get("progress", 0),
        "text": data.get("merged_text", "") or data.get("asr_text", "") or data.get("subtitle_text", "") or data.get("message", ""),
        "subtitle_text": data.get("subtitle_text", ""),
        "asr_text": data.get("asr_text", ""),
        "merged_text": data.get("merged_text", ""),
        "video_path": data.get("video_path", ""),
        "task_dir": data.get("task_dir", ""),
        "error": data.get("message", "") if svc_status == "error" else None,
    }


@app.get("/api/video/file")
def api_video_file(path: str = ""):
    """Serve a downloaded video file. Allows paths under VIDEO_DIR, data_dir, or global save_path."""
    import os as _os
    base = _os.path.dirname(_os.path.abspath(__file__))
    video_dir = _os.path.normcase(_os.path.normpath(_os.path.join(base, "data", "videos")))
    data_dir = _os.path.normcase(_os.path.normpath(_os.path.join(base, "data")))
    save_root = _os.path.normcase(_os.path.normpath(_get_global_save_path()))
    full = _os.path.normcase(_os.path.normpath(_os.path.abspath(path)))
    if not (full.startswith(video_dir) or full.startswith(data_dir + _os.sep) or full.startswith(save_root)):
        raise HTTPException(403, f"Access denied: {full}")
    if not _os.path.exists(full):
        raise HTTPException(404, f"File not found: {full}")
    return FileResponse(full)


@app.post("/api/video/extract-subtitles")
def api_extract_subtitles(req: dict):
    """Manually extract subtitles from a previously downloaded video task."""
    task_id = req["task_id"]
    project_id = req.get("project_id")
    progress = get_progress(task_id)
    if progress.get("status") != "done":
        return {"status": "error", "message": "视频尚未下载完成"}
    subtitle_text = progress.get("subtitle_text", "")
    # Auto-save as step1 if project_id provided
    if project_id and subtitle_text:
        db = get_db()
        try:
            existing = db.execute(
                "SELECT id FROM step_results WHERE project_id = ? AND step_name = ?",
                (project_id, "step1")).fetchone()
            if existing:
                db.execute(
                    "UPDATE step_results SET content = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?",
                    (subtitle_text, existing["id"]))
            else:
                db.execute(
                    "INSERT INTO step_results (project_id, step_name, content) VALUES (?, ?, ?)",
                    (project_id, "step1", subtitle_text))
            db.commit()
        finally:
            db.close()
    return {"subtitle_text": subtitle_text}


# ── PPT Generation ──

from services.ppt_service import generate_ppt, _extract_typography

@app.get("/api/projects/{project_id}/ppt-results")
def api_ppt_results(project_id: str):
    """Return all saved PPT generation results for a project (survives page reload)."""
    db = get_db()
    try:
        rows = db.execute(
            "SELECT step_name, content, updated_at FROM step_results "
            "WHERE project_id = ? AND (step_name LIKE '_ppt_result_%' OR step_name LIKE '_ppt_plan_%') "
            "ORDER BY updated_at DESC",
            (project_id,)).fetchall()
        results = []
        seen_run_ids = set()
        for row in rows:
            try:
                meta = json.loads(row["content"])
                rid = meta.get("run_id") or ""
                # Deduplicate: _ppt_plan_ and _ppt_result_ may share the same run_id
                if rid and rid in seen_run_ids:
                    continue
                if rid:
                    seen_run_ids.add(rid)
                meta["_saved_at"] = row["updated_at"]
                meta["_step_name"] = row["step_name"]
                results.append(meta)
            except Exception:
                pass
        return {"results": results}
    finally:
        db.close()

@app.get("/api/projects/{project_id}/ppt-status")
def api_ppt_status(project_id: str):
    """Polled by frontend every 10s during PPT generation.
    Returns current phase, slide counts, preview URL if ready."""
    from services.ppt_service import get_ppt_status
    status = get_ppt_status(project_id)
    if not status:
        return {"phase": "idle", "phase_label": "未在生成", "message": "没有正在进行的生成任务"}
    return status


@app.get("/api/projects/{project_id}/ppt-log")
def api_ppt_log(project_id: str):
    """Polled by frontend every 60s during PPT/outline generation.
    Returns timestamped log entries."""
    from services.ppt_service import get_ppt_log
    return {"logs": get_ppt_log(project_id)}


@app.post("/api/ppt/generate")
def api_generate_ppt(req: PPTGenerateRequest):
    import time, sys, datetime as _dt
    t0 = time.time()
    print(f"[PPT-REQ] {time.strftime('%H:%M:%S')} provider={req.provider_id} model={req.model} "
          f"content_len={len(req.content) if req.content else 0} template={req.template_id} "
          f"has_rules=?", flush=True)
    try:
        with open(os.path.join(BASE_DIR, "data", "ppt_req.log"), "a", encoding="utf-8") as _lf:
            _lf.write(f"[{_dt.datetime.now().isoformat()}] provider={req.provider_id} model={req.model} "
                      f"content_len={len(req.content) if req.content else 0} template={req.template_id}\n")
    except Exception: pass
    try:
        output_dir = None
        project_name = ""
        if req.project_id:
            output_dir = resolve_project_storage(req.project_id)
            proj = _get_project(req.project_id)
            if proj:
                project_name = proj["name"]

        # Use saved slide_plan from template if not provided in request.
        # But if caller sent content+model, let AI regenerate (content takes priority).
        slide_plan = req.slide_plan
        has_content_input = bool(req.content and req.content.strip() and req.provider_id and req.model)
        if slide_plan is None and req.template_id and not has_content_input:
            db = get_db()
            try:
                row = db.execute(
                    "SELECT slide_plan FROM templates WHERE id = ?",
                    (req.template_id,)).fetchone()
                if row and row["slide_plan"]:
                    try:
                        slide_plan = json.loads(row["slide_plan"])
                        print(f"[PPT-REQ] Using saved slide_plan ({len(slide_plan)} slides)", flush=True)
                    except Exception:
                        pass
            finally:
                db.close()
        if has_content_input:
            print(f"[PPT-REQ] AI regeneration with content ({len(req.content)} chars)", flush=True)

        # Validate prerequisites BEFORE calling expensive pipeline
        if not has_content_input and not slide_plan:
            missing: list[str] = []
            has_content = bool(req.content and req.content.strip())
            if not has_content:
                missing.append("没有文案内容，请先在 Stage 1 导入素材或在 Stage 2 生成文案")
            if not req.provider_id:
                missing.append("没有选择大模型提供商，请在左侧下拉框中选择")
            elif not req.model:
                missing.append("没有选择大模型，请在左侧下拉框中选择模型")
            if not req.template_id:
                missing.append("没有选择模板，请先在左侧选择模板")
            detail = "；".join(missing) if missing else "缺少必要参数，请检查后再试"
            raise HTTPException(status_code=400, detail=detail)

        filepath, generated_slides = generate_ppt(req.content, req.template_id, req.branding, output_dir, req.provider_id, req.model, slide_plan, project_name=project_name, column_id=req.column_id, color_scheme=req.color_scheme, temperature=req.temperature, project_id=req.project_id or "", temp_keyword=req.temp_keyword, temp_research=req.temp_research, temp_outline=req.temp_outline, temp_fill=req.temp_fill, temp_cards=req.temp_cards, temp_html=req.temp_html, temp_svg_batch=req.temp_svg_batch, temp_svg_single=req.temp_svg_single, temp_review=req.temp_review, temp_fix=req.temp_fix, temp_holistic=req.temp_holistic, temp_holistic_fix=req.temp_holistic_fix, temp_stage_outline=req.temp_stage_outline, temp_stage_generation=req.temp_stage_generation, temp_stage_review=req.temp_stage_review)
        if generated_slides is not None:
            slide_plan = generated_slides

        # Detect output type: SVG (returns HTML path) vs PPTX (returns .pptx path)
        is_svg = filepath and filepath.endswith(".html")
        is_pptx = filepath and filepath.endswith(".pptx")

        if is_svg:
            # SVG output — register run_id → actual directory for preview serving
            run_dir = os.path.dirname(filepath)
            run_id = os.path.basename(run_dir)
            _run_dirs[run_id] = run_dir
            _save_run_dirs(_run_dirs)
            # Save color_scheme metadata so recolor can deterministically know the source scheme
            cs = getattr(req, "color_scheme", None) or "deep-blue"
            with open(os.path.join(run_dir, "color_scheme.txt"), "w", encoding="utf-8") as _csf:
                _csf.write(cs)
            preview_url = f"/api/exports/{run_id}/index.html"
            zip_url = f"/api/ppt/export-zip/{run_id}"
            # Persist result metadata so page reload can restore without re-generating
            from services.ppt_service import _load_style_from_template
            result_meta = {
                "run_id": run_id,
                "preview_url": preview_url,
                "zip_url": zip_url,
                "slide_plan": slide_plan,
                "slide_count": len(generated_slides) if generated_slides else 0,
                "template_id": req.template_id,
                "style_id": _load_style_from_template(req.template_id),
                "format": "svg",
                "color_scheme": cs,
                "column_id": req.column_id,
                "generated_at": _dt.datetime.now().isoformat(),
            }
            try:
                with open(os.path.join(run_dir, "result.json"), "w", encoding="utf-8") as _rf:
                    json.dump(result_meta, _rf, ensure_ascii=False, indent=2)
            except Exception as _e:
                print(f"[PPT] Failed to persist result.json: {_e}", flush=True)
            # Persist to DB for frontend auto-load on page refresh (optional)
            if req.project_id:
                try:
                    save_step_meta(req.project_id, f"_ppt_result_{run_id}",
                                   json.dumps(result_meta, ensure_ascii=False))
                except Exception as _e:
                    print(f"[PPT] Failed to persist result meta: {_e}", flush=True)
            return {
                "format": "svg",
                "run_id": run_id,
                "preview_url": preview_url,
                "zip_url": zip_url,
                "slide_plan": slide_plan,
                "slide_count": len(generated_slides) if generated_slides else 0,
            }
        elif is_pptx:
            # PPTX output — legacy download URL
            filename = os.path.basename(filepath)
            params = []
            if req.project_id:
                params.append(f"project_id={req.project_id}")
            if project_name:
                safe_name = "".join(c for c in project_name if c.isalnum() or c in "._- ()（）").strip()
                params.append(f"name={safe_name}_PPT.pptx")
            download_url = f"/api/download/{filename}"
            if params:
                download_url += "?" + "&".join(params)
            return {"format": "pptx", "filename": filename, "download_url": download_url,
                    "slide_plan": slide_plan,
                    "slide_count": len(generated_slides) if generated_slides else 0}
        else:
            try:
                with open(os.path.join(BASE_DIR, "data", "ppt_req.log"), "a", encoding="utf-8") as _lf:
                    _lf.write(f"[{_dt.datetime.now().isoformat()}] 500: filepath={filepath!r} is_svg={is_svg} is_pptx={is_pptx}\n")
            except Exception: pass
            raise HTTPException(status_code=500, detail="PPT generation failed — no output produced")
    except Exception as e:
        import traceback as _tb
        try:
            with open(os.path.join(BASE_DIR, "data", "ppt_req.log"), "a", encoding="utf-8") as _lf:
                _lf.write(f"[{_dt.datetime.now().isoformat()}] EXCEPTION: {e}\n{_tb.format_exc()}\n")
        except Exception: pass
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/ppt/outline")
def api_ppt_outline(req: PPTPlanRequest):
    """Generate outline only (Phase 2 Research + Phase 4 Outline+Content).

    Returns outline_json (structured) and outline_text (natural-language,
    human-readable format for review and editing).
    """
    from services.ppt_service import _generate_outline_only
    db = get_db()
    try:
        column_id = req.column_id or ""
        cfg = None
        if column_id:
            cfg = db.execute(
                "SELECT prompt, skill FROM column_configs WHERE column_id = ?",
                (column_id,)).fetchone()

        column_prompt = cfg["prompt"] or "" if cfg else ""
        column_skill = cfg["skill"] or "" if cfg else ""
        rules = {}

        if req.template_id:
            row = db.execute(
                "SELECT rules FROM templates WHERE id = ?",
                (req.template_id,)).fetchone()
            if row and row["rules"]:
                try:
                    template_rules = json.loads(row["rules"])
                    for key in ("style_id", "layout_types", "page_rhythm", "design_principles"):
                        if key in template_rules:
                            rules[key] = template_rules[key]
                except Exception:
                    pass

        st = dict(
            keyword=req.temp_keyword or req.temperature or 0.3,
            research=req.temp_research or req.temperature or 0.7,
            outline=req.temp_outline or req.temperature or 1.0,
            fill=req.temp_fill or req.temperature or 1.0,
        )
        if req.temp_stage_outline > 0:
            st.update(keyword=req.temp_stage_outline, research=req.temp_stage_outline,
                      outline=req.temp_stage_outline, fill=req.temp_stage_outline)
        outline_json, outline_text = _generate_outline_only(
            req.provider_id, req.model, rules, req.content,
            column_prompt, column_skill, temperature=req.temperature,
            st=st, project_id=req.project_id or "")
        return {"outline_json": outline_json or [], "outline_text": outline_text or ""}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
    finally:
        db.close()


@app.post("/api/ppt/outline/convert")
def api_ppt_outline_convert(req: PPTPlanRequest):
    """Convert edited natural-language text back to structured JSON via LLM.

    Called when user saves edited outline content. Uses low-temperature LLM
    to extract structured JSON from free-form text.
    """
    from services.ppt_service import _human_text_to_json
    if not req.provider_id or not req.model:
        raise HTTPException(status_code=400, detail="需要选择大模型才能转换")
    try:
        result = _human_text_to_json(req.provider_id, req.model, req.content, req.slide_plan or [])
        return {"outline_json": result or []}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/api/ppt/plan")
def api_ppt_plan(req: PPTPlanRequest):
    """Generate slide plan only (no PPTX file). Returns JSON for user review.

    Uses column config's prompt + skill as AI system message and structure
    template. Template is only used for its rules (layout_types, page_rhythm).
    """
    from services.ppt_service import _generate_slides_staged
    db = get_db()
    try:
        # Load column config — the single source of truth for content logic
        column_id = req.column_id or ""
        cfg = None
        if column_id:
            cfg = db.execute(
                "SELECT prompt, skill FROM column_configs WHERE column_id = ?",
                (column_id,)).fetchone()

        column_prompt = cfg["prompt"] or "" if cfg else ""
        column_skill = cfg["skill"] or "" if cfg else ""
        rules = {}

        # Template provides layout_types via its rules (visual structure)
        if req.template_id:
            row = db.execute(
                "SELECT rules FROM templates WHERE id = ?",
                (req.template_id,)).fetchone()
            if row and row["rules"]:
                try:
                    template_rules = json.loads(row["rules"])
                    # Merge: template's layout_types + page_rhythm take priority
                    for key in ("style_id", "layout_types", "page_rhythm", "design_principles"):
                        if key in template_rules:
                            rules[key] = template_rules[key]
                except Exception:
                    pass

        st = dict(
            keyword=req.temp_keyword or req.temperature or 0.3,
            research=req.temp_research or req.temperature or 0.7,
            outline=req.temp_outline or req.temperature or 1.0,
            fill=req.temp_fill or req.temperature or 1.0,
            cards=req.temp_cards or req.temperature or 0.7,
            html=req.temp_html or req.temperature or 0.8,
        )
        if req.temp_stage_outline > 0:
            st.update(keyword=req.temp_stage_outline, research=req.temp_stage_outline,
                      outline=req.temp_stage_outline, fill=req.temp_stage_outline)
        if req.temp_stage_generation > 0:
            st.update(cards=req.temp_stage_generation, html=req.temp_stage_generation)
        slide_plan = _generate_slides_staged(
            req.provider_id, req.model, rules, req.content,
            column_prompt, column_skill, temperature=req.temperature,
            st=st
        )
        return {"slide_plan": slide_plan or []}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
    finally:
        db.close()


def _fix_per_slide_divs(html: str) -> str:
    """Repair each slide-wrapper's internal div balance so no slide
    can corrupt the DOM of adjacent slides.  Fixes the root cause of
    "pages 6-7 disappeared / page 5 layout wrong" bugs.

    Uses the interval between consecutive wrapper tags to define slide
    boundaries, then balances <div>/</div> counts per-interval by adding
    or removing closing tags at the tail of each interval."""
    tag = '<div class="slide-wrapper">'

    # Find all wrapper positions
    positions = []
    p = 0
    while True:
        idx = html.find(tag, p)
        if idx == -1:
            break
        positions.append(idx)
        p = idx + 1

    if not positions:
        return html

    result = []
    result.append(html[:positions[0]])  # preamble before first slide

    for i, start in enumerate(positions):
        end = positions[i + 1] if i + 1 < len(positions) else html.find('</body>', start)
        if end == -1:
            end = len(html)

        chunk = html[start:end]
        opens = chunk.count('<div')
        closes = chunk.count('</div>')
        diff = opens - closes

        if diff > 0:
            # Missing closing tags — append before the boundary
            chunk += '</div>' * diff
        elif diff < 0:
            # Excess closing tags — remove from the tail
            excess = -diff
            for _ in range(excess):
                p = chunk.rfind('</div>')
                if p >= 0:
                    chunk = chunk[:p] + chunk[p + 6:]

        result.append(chunk)

    # Append everything after the last wrapper's range
    last_end = positions[-1]
    body_end = html.find('</body>', last_end)
    if body_end == -1:
        body_end = len(html)
    # The last chunk already goes to </body> or end, so append the rest
    body_close_pos = html.find('</body>', positions[-1])
    if body_close_pos != -1:
        result.append(html[body_close_pos:])

    return ''.join(result)


def _ensure_backup(run_dir: str):
    """Create index_backup.html from index.html if no backup exists yet.
    Preserves the original generated version for restore."""
    index_path = os.path.join(run_dir, "index.html")
    backup_path = os.path.join(run_dir, "index_backup.html")
    if os.path.exists(index_path) and not os.path.exists(backup_path):
        shutil.copy2(index_path, backup_path)


@app.post("/api/ppt/recolor-slide")
def api_ppt_recolor_slide(
    run_id: str = Body(...),
    slide_seq: int = Body(...),
    style: str = Body("business"),
    color_scheme: str = Body("deep-blue"),
):
    """Recolor all slides to a new color scheme.

    New format (CSS variables): replace :root block — O(1) operation.
    Legacy format (hardcoded hex): per-slide hex replacement fallback.
    """
    from services.ppt_service import _recolor_slide_html, _build_root_vars, _load_scheme_data

    run_dir = _find_run_dir(run_id)
    if not run_dir:
        raise HTTPException(status_code=404, detail="Run not found")

    html_path = os.path.join(run_dir, "index.html")

    sid = style or "business"
    cs = color_scheme or "deep-blue"

    if not os.path.exists(html_path):
        raise HTTPException(status_code=404, detail="index.html not found")

    with open(html_path, encoding="utf-8") as f:
        full_html = f.read()

    # ── Detect format: CSS variables (has :root block) vs legacy hex ──
    import re
    if ":root {" in full_html:
        # ── New format: replace :root block ──
        new_scheme_data = _load_scheme_data(sid, cs)
        if not new_scheme_data:
            raise HTTPException(status_code=400, detail=f"Unknown color scheme: {cs}")
        new_root = _build_root_vars(new_scheme_data)

        # Replace :root { ... } block (first occurrence, up to matching })
        root_pattern = re.compile(r':root\s*\{[^}]*\}', re.DOTALL)
        new_html, count = root_pattern.subn(new_root, full_html, count=1)
        if count == 0:
            raise HTTPException(status_code=400, detail=":root block not found in HTML")

        full_html = new_html
        changed_count = 1
    else:
        # ── Legacy format: per-slide hex replacement ──
        slides_raw = re.findall(r'(<section[\s\S]*?</section>)', full_html, re.IGNORECASE)
        if not slides_raw:
            inner_pattern = re.compile(
                r'<div\s+class="slide-wrapper"\s*>\s*(<div\s[^>]*width\s*:\s*\d+px[^>]*>)',
                re.IGNORECASE
            )
            for m in inner_pattern.finditer(full_html):
                pos = m.end(1)
                depth = 1
                while pos < len(full_html) and depth > 0:
                    next_open = full_html.find('<div', pos)
                    next_close = full_html.find('</div>', pos)
                    if next_close == -1:
                        break
                    if next_open != -1 and next_open < next_close:
                        depth += 1
                        pos = next_open + 4
                    else:
                        depth -= 1
                        if depth == 0:
                            slides_raw.append(m.group(1) + full_html[m.end(1):next_close + 6])
                            break
                        pos = next_close + 6

        if not slides_raw:
            raise HTTPException(status_code=400, detail="No slides found in HTML")

        source_scheme = "deep-blue"
        rj_path = os.path.join(run_dir, "result.json")
        if os.path.exists(rj_path):
            try:
                rj_meta = json.loads(open(rj_path, "r", encoding="utf-8").read())
                source_scheme = rj_meta.get("color_scheme") or source_scheme
            except Exception:
                pass

        changed_count = 0
        for i, old_slide in enumerate(slides_raw):
            new_slide = _recolor_slide_html(sid, old_slide, cs, source_scheme)
            if new_slide != old_slide:
                full_html = full_html.replace(old_slide, new_slide)
                changed_count += 1

        if changed_count == 0:
            return {"ok": True, "changed": False, "message": "颜色无变化（可能已使用该色系）"}

    full_html = _fix_per_slide_divs(full_html)
    _ensure_backup(run_dir)
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(full_html)
    # Update color scheme tracking
    with open(os.path.join(run_dir, "color_scheme.txt"), "w", encoding="utf-8") as _csf:
        _csf.write(cs)
    # Sync color_scheme to result.json so regenerate picks it up
    rj_path = os.path.join(run_dir, "result.json")
    if os.path.exists(rj_path):
        with open(rj_path, "r+", encoding="utf-8") as _rjf:
            _rj_meta = json.loads(_rjf.read())
            _rj_meta["color_scheme"] = cs
            _rjf.seek(0)
            _rjf.truncate()
            json.dump(_rj_meta, _rjf, ensure_ascii=False, indent=2)

    return {"ok": True, "changed": True, "method": "css_vars" if ":root {" in full_html else "hex_replace"}


@app.post("/api/ppt/edit-slide")
def api_ppt_edit_slide(req: PPTEditSlideRequest):
    """Edit a single slide via natural language instruction.

    Flow:
      1. Read index.html from output directory
      2. Extract the target slide's HTML
      3. Send to LLM with editing instruction
      4. Run three code checks on the result
      5. If clean → write back → return ok
      6. If violation → return error with details
    """
    import re
    from services.ppt_service import (
        _safe_run_async, _build_edit_system_prompt, _run_edit_agent,
    )

    run_dir = _find_run_dir(req.run_id)
    if not run_dir:
        raise HTTPException(status_code=404, detail="Run not found")

    html_path = os.path.join(run_dir, "index.html")
    if not os.path.exists(html_path):
        raise HTTPException(status_code=404, detail="index.html not found")

    with open(html_path, encoding="utf-8") as f:
        full_html = f.read()

    seq = req.slide_seq

    # Count total slides (just for the prompt)
    slide_count = len(re.findall(r'<div\s+class="slide-wrapper"', full_html, re.IGNORECASE))
    if slide_count == 0:
        slide_count = len(re.findall(r'<section[\s>]', full_html, re.IGNORECASE))
    if seq < 1 or seq > max(slide_count, 1):
        raise HTTPException(status_code=400, detail=f"Slide {seq} out of range (1–{max(slide_count, 1)})")

    # Build the edit system prompt
    edit_system = _build_edit_system_prompt(req.style or "business", req.color_scheme or "deep-blue")

    edit_user = (
        f"以下是完整的 PPT HTML 文件（共 {slide_count} 页），每页是一个 slide-wrapper。\n\n"
        f"```html\n{full_html}\n```\n\n"
        f"修改要求：只修改第 {seq} 页，{req.instruction}\n\n"
        f"输出要求：返回修改后的完整 HTML 文件（所有页），用 ```html ``` 包裹。不要省略任何页。"
    )

    p_id = req.provider_id
    model = req.model
    if not p_id or not model:
        db = get_db()
        try:
            row = db.execute(
                "SELECT id, models FROM llm_providers WHERE is_enabled=1 LIMIT 1"
            ).fetchone()
            if row:
                p_id = row["id"]
                models = json.loads(row["models"]) if row["models"] else []
                model = models[0] if models else ""
        finally:
            db.close()

    if not p_id or not model:
        raise HTTPException(status_code=400, detail="没有可用的 LLM 提供商")

    try:
        raw = _safe_run_async(_run_edit_agent(p_id, model, edit_system, edit_user))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Agent 调用失败: {e}")

    # Extract the modified HTML from the response
    edited_html = ""
    m = re.search(r'```html\s*\n(.*?)\n```', raw, re.DOTALL)
    if m:
        edited_html = m.group(1).strip()
    else:
        # Fallback: try to find a complete HTML document
        m = re.search(r'(<!DOCTYPE html>[\s\S]*?</html>)', raw, re.IGNORECASE)
        if m:
            edited_html = m.group(1).strip()

    # Verify we got a valid HTML document
    if not edited_html or len(edited_html) < len(full_html) * 0.5:
        return {
            "ok": False,
            "slide_seq": seq,
            "error": "chat_reply",
            "detail": raw.strip() or "AI 未返回有效的完整 HTML，请换一种描述重试",
        }

    # Verify slide count is preserved
    new_slide_count = len(re.findall(r'<div\s+class="slide-wrapper"', edited_html, re.IGNORECASE))
    if new_slide_count == 0:
        new_slide_count = len(re.findall(r'<section[\s>]', edited_html, re.IGNORECASE))
    if new_slide_count != slide_count:
        return {
            "ok": False,
            "slide_seq": seq,
            "error": "chat_reply",
            "detail": f"AI 返回的页数不对（期望 {slide_count}，实际 {new_slide_count}），请重试",
        }

    # Quick div balance check
    divs_open = len(re.findall(r'<div\s', edited_html)) + len(re.findall(r'<div>', edited_html))
    divs_close = len(re.findall(r'</div>', edited_html))
    if abs(divs_open - divs_close) > 3:
        return {
            "ok": False,
            "slide_seq": seq,
            "error": "chat_reply",
            "detail": f"AI 返回的 HTML 标签不平衡（div: {divs_open}开/{divs_close}闭），请重试",
        }

    # Write to index.html directly (auto-backup preserves original)
    _ensure_backup(run_dir)
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(edited_html)

    return {"ok": True, "slide_seq": seq,
            "preview": True, "saved": True}


@app.post("/api/ppt/save-edit/{run_id}")
def api_save_edit(run_id: str):
    """No-op: edits now write directly to index.html."""
    return {"ok": True, "saved": True}


@app.post("/api/ppt/discard-edit/{run_id}")
def api_discard_edit(run_id: str):
    """Restore index.html from index_backup.html."""
    run_dir = _run_dirs.get(run_id)
    if not run_dir:
        candidate = os.path.join(EXPORT_DIR, run_id)
        if os.path.isdir(candidate):
            run_dir = candidate
    if not run_dir:
        raise HTTPException(status_code=404, detail="Export not found")

    backup_path = os.path.join(run_dir, "index_backup.html")
    if not os.path.exists(backup_path):
        raise HTTPException(status_code=400, detail="No backup to restore from")

    real_path = os.path.join(run_dir, "index.html")
    shutil.copy2(backup_path, real_path)
    return {"ok": True, "restored": True}


@app.put("/api/ppt/slide-source/{run_id}")
def api_slide_source(run_id: str, req: PPTSlideSourceRequest):
    """Write the full HTML document directly to index.html.

    Used by the source-code editor tab — the user edits the complete
    HTML document and applies it.
    """
    run_dir = _run_dirs.get(run_id)
    if not run_dir:
        candidate = os.path.join(EXPORT_DIR, run_id)
        if os.path.isdir(candidate):
            run_dir = candidate
    if not run_dir or not os.path.isdir(run_dir):
        raise HTTPException(status_code=404, detail="Run not found")

    _ensure_backup(run_dir)
    index_path = os.path.join(run_dir, "index.html")
    with open(index_path, "w", encoding="utf-8") as f:
        f.write(req.html)

    return {"ok": True, "saved": True}


@app.post("/api/ppt/regenerate-slide")
def api_ppt_regenerate_slide(req: PPTRegenerateSlideRequest):
    """Regenerate selected slides via the standard two-phase HTML pipeline.

    Reads existing slide structures from result.json, runs
    _stage2_html_per_slide() for the target slides, splices the
    results back into the full HTML deck, and writes to a separate
    index_regenerated.html file (never overwrites index.html).

    Writes progress to _regenerate_log.txt so the frontend can poll.
    """
    from services.ppt_service import (
        _stage2_html_per_slide, _assemble_html_deck,
        _resolve_color_vars, _load_scheme_data,
        _auto_fix_hardcoded_hex, _auto_fix_font_size,
        _get_canvas_dimensions,
    )
    import datetime as _dt

    run_dir = _find_run_dir(req.run_id)
    if not run_dir or not os.path.isdir(run_dir):
        raise HTTPException(status_code=404, detail="Run not found")

    log_path = os.path.join(run_dir, "_regenerate_log.txt")
    # Truncate old log so each regeneration starts fresh
    with open(log_path, "w", encoding="utf-8") as _lf:
        _lf.write("")
    def _log(msg: str):
        ts = _dt.datetime.now().strftime("%H:%M:%S")
        line = f"[{ts}] {msg}\n"
        with open(log_path, "a", encoding="utf-8") as _lf:
            _lf.write(line)
            _lf.flush()

    _log("开始重新生成...")

    import re as _re

    # ── Build slide_plan: prefer result.json, fallback to parsing index.html ──
    rj_path = os.path.join(run_dir, "result.json")
    style_id = "business"
    color_scheme = "deep-blue"
    column_id = ""
    if os.path.exists(rj_path):
        _log("读取 result.json...")
        with open(rj_path, encoding="utf-8") as _f:
            result_meta = json.loads(_f.read())
        slide_plan = result_meta.get("slide_plan", [])
        style_id = result_meta.get("style_id") or style_id
        color_scheme = result_meta.get("color_scheme") or color_scheme
        column_id = result_meta.get("column_id", "") or req.column_id
        _log(f"已加载 {len(slide_plan)} 页幻灯片结构，色系: {color_scheme}")
    else:
        _log("result.json 不存在，从 index.html 解析...")
        # Fallback: parse slides from index.html (old-format runs)
        column_id = req.column_id
        html_path = os.path.join(run_dir, "index.html")
        if not os.path.exists(html_path):
            _log("错误: index.html 不存在")
            raise HTTPException(status_code=404, detail="index.html not found")
        with open(html_path, encoding="utf-8") as _f:
            full_html = _f.read()
        sections = _re.findall(
            r'<section\s+class="slide"\s+data-seq="(\d+)"\s+data-type="([^"]*)"[^>]*>(.*?)</section>',
            full_html, _re.DOTALL
        )
        if not sections:
            _log("错误: index.html 中未找到幻灯片")
            raise HTTPException(status_code=400, detail="No slides found in index.html")
        slide_plan = []
        for seq_str, stype, raw in sections:
            seq = int(seq_str)
            h = _re.search(r'<h[12][^>]*>(.*?)</h[12]>', raw[:3000], _re.DOTALL)
            heading = _re.sub(r'<[^>]+>', '', h.group(1).strip()) if h else ""
            body_text = _re.sub(r'<[^>]+>', ' ', raw[:5000]).strip()
            body_text = _re.sub(r'\s+', ' ', body_text)[:1000]
            has_chart = bool(_re.search(r'<(svg|canvas|path|circle|rect|line)', raw[:3000]))
            chart_hint = ""
            if has_chart:
                chart_hint = _re.search(r'data-chart-hint="([^"]*)"', raw[:1000])
                if chart_hint:
                    chart_hint = chart_hint.group(1)
            slide_plan.append({
                "seq": seq, "type": stype, "layout": "hero_grid",
                "heading": heading, "body": body_text,
                "key_points": [], "kicker": "", "lead": "",
                "cards": [], "has_chart": has_chart, "chart_hint": chart_hint,
                "notes": "", "html": raw, "html_vars": raw,
            })
        slide_plan.sort(key=lambda s: s["seq"])
        # Fallback color_scheme for old runs: read from color_scheme.txt
        cs_path = os.path.join(run_dir, "color_scheme.txt")
        if os.path.exists(cs_path):
            with open(cs_path, encoding="utf-8") as _csf:
                color_scheme = _csf.read().strip() or color_scheme
        _log(f"已解析 {len(slide_plan)} 页幻灯片，色系(fallback): {color_scheme}")
    if not slide_plan:
        raise HTTPException(status_code=400, detail="No slides found")

    total = len(slide_plan)
    seqs = req.slide_seqs
    for seq in seqs:
        if seq < 1 or seq > total:
            raise HTTPException(status_code=400, detail=f"Slide {seq} out of range (1–{total})")

    _log(f"选中 {len(seqs)} 页待重新生成: {', '.join(str(s) for s in seqs)}")

    # Get provider/model (from request or fallback to default)
    p_id = req.provider_id
    model_str = req.model
    if not p_id or not model_str:
        db = get_db()
        try:
            row = db.execute(
                "SELECT id, models FROM llm_providers WHERE is_enabled=1 LIMIT 1"
            ).fetchone()
            if row:
                p_id = row["id"]
                models = json.loads(row["models"]) if row["models"] else []
                model_str = models[0] if models else ""
        finally:
            db.close()
    if not p_id or not model_str:
        raise HTTPException(status_code=400, detail="没有可用的 LLM 提供商")

    _log(f"使用模型: {model_str}")

    # Build structure list for selected slides
    redo_structure = []
    for seq in seqs:
        t = slide_plan[seq - 1]
        redo_structure.append({
            "seq": t.get("seq", seq),
            "type": t.get("type", "content"),
            "layout": t.get("layout", "hero_grid"),
            "heading": t.get("heading", ""),
            "body": t.get("body", ""),
            "key_points": t.get("key_points", []),
            "kicker": t.get("kicker", ""),
            "lead": t.get("lead", ""),
            "cards": t.get("cards", []),
            "has_chart": t.get("has_chart", False),
            "chart_hint": t.get("chart_hint", ""),
            "notes": t.get("notes", ""),
        })

    _log(f"正在调用 AI 生成 HTML（{len(redo_structure)} 页，并行 {min(len(redo_structure), 3)} 页）...")

    # Regenerate selected slides through the standard HTML pipeline
    html_slides = _stage2_html_per_slide(
        p_id, model_str, generate, redo_structure,
        style_id=style_id, color_scheme=color_scheme,
        parallel=min(len(redo_structure), 3), temperature=0.3,
        column_id=column_id,
    )

    if not html_slides:
        _log("错误: LLM 未返回 HTML")
        raise HTTPException(status_code=500, detail="Slide regeneration failed — LLM returned no HTML")

    _log(f"AI 生成完成，收到 {len(html_slides)} 页 HTML")

    # ── Apply auto-fixers (hex→var, font-size) BEFORE variable resolution ──
    scheme_data = _load_scheme_data(style_id, color_scheme)
    import re as _re_count
    for slide in html_slides:
        seq = slide.get("seq", 0)
        html_before = slide.get("html", "")
        hex_before = len(_re_count.findall(r'#[0-9a-fA-F]{3,6}\b', html_before))
        if scheme_data:
            slide["html"] = _auto_fix_hardcoded_hex(slide.get("html", ""), scheme_data, seq)
            slide["html_vars"] = _auto_fix_hardcoded_hex(slide.get("html_vars", slide.get("html", "")), scheme_data, seq)
        slide["html"] = _auto_fix_font_size(slide.get("html", ""), seq)
        slide["html_vars"] = _auto_fix_font_size(slide.get("html_vars", slide.get("html", "")), seq)
        hex_after = len(_re_count.findall(r'#[0-9a-fA-F]{3,6}\b', slide.get("html", "")))
        var_count = len(_re_count.findall(r'\{\{[a-z_0-9]+\}\}', slide.get("html", "")))
        _log(f"Slide {seq} 自动修正: {hex_before}→{hex_after} hex, {var_count} vars")

    # Splice regenerated slides back into the plan
    for slide in html_slides:
        seq = slide.get("seq", 0)
        idx = seq - 1
        if 0 <= idx < len(slide_plan):
            slide_plan[idx]["html"] = slide.get("html", "")
            slide_plan[idx]["html_vars"] = slide.get("html_vars", slide.get("html", ""))

    _log("正在解析颜色变量...")

    # Resolve color variables for resolved HTML deck
    resolved_slides = []
    for s in slide_plan:
        # Use html_vars (with {{primary}} placeholders) when available so
        # color re-resolution works. Fall back to plain html.
        html_v = s.get("html_vars", "") or s.get("html", "")
        if html_v and scheme_data:
            html_v = _resolve_color_vars(html_v, scheme_data, css_vars=True)
        resolved_slides.append({**s, "html": html_v})

    # Fix cover slide: if background is primary (dark), text must be light
    # otherwise title text is invisible (e.g. deep-blue: text=#1a202c on bg=#1a365d)
    if resolved_slides and scheme_data:
        cover = resolved_slides[0]
        primary = scheme_data.get("primary", "")
        secondary = scheme_data.get("secondary", "")
        text_color = scheme_data.get("text", "")
        if primary and text_color:
            cover_html = cover.get("html", "")
            # Search entire HTML for dark background (primary or secondary as bg)
            has_dark_bg = (f"background:{primary}" in cover_html
                          or f"background:{secondary}" in cover_html
                          or f"background: {primary}" in cover_html
                          or f"background: {secondary}" in cover_html)
            if has_dark_bg:
                # Replace text color with #ffffff (legal hardcoded white per VI spec)
                cover_html = cover_html.replace(f"color:{text_color}", "color:#ffffff")
                cover_html = cover_html.replace(f"color:{text_color};", "color:#ffffff;")
                cover_html = cover_html.replace(f"color: {text_color}", "color: #ffffff")
                cover_html = cover_html.replace(f"color: {text_color};", "color: #ffffff;")
                resolved_slides[0] = {**cover, "html": cover_html}
                _log("封面页修复：深色背景 → 白色文字")

    _log("正在组装完整 HTML...")

    # Get canvas dimensions for the column (defaults to 1280x720 if unknown)
    regen_canvas_w, regen_canvas_h = _get_canvas_dimensions(column_id)

    # Rebuild both decks
    title = slide_plan[0].get("heading", "") if slide_plan else "Presentation"
    deck_html = _assemble_html_deck(resolved_slides, title, style_id, scheme_data, canvas_w=regen_canvas_w, canvas_h=regen_canvas_h)
    if scheme_data:
        deck_html = _resolve_color_vars(deck_html, scheme_data, css_vars=True)
    deck_vars = _assemble_html_deck(
        [{**s, "html": s.get("html_vars", s.get("html", ""))} for s in slide_plan],
        title, style_id, scheme_data, canvas_w=regen_canvas_w, canvas_h=regen_canvas_h
    )

    _log(f"正在写入文件（{len(deck_html)} 字节）...")

    # Write to regenerated file (NEVER overwrite index.html)
    regen_html_path = os.path.join(run_dir, "index_regenerated.html")
    with open(regen_html_path, "w", encoding="utf-8") as f:
        f.write(deck_html)
    regen_vars_path = os.path.join(run_dir, "index_regenerated_vars.html")
    with open(regen_vars_path, "w", encoding="utf-8") as f:
        f.write(deck_vars)

    # Build partial mini-deck with only the regenerated slides
    regen_resolved = [s for s in resolved_slides if s.get("seq") in set(seqs)]
    partial_html = _assemble_html_deck(regen_resolved, title, style_id, scheme_data, total_slides=len(slide_plan), canvas_w=regen_canvas_w, canvas_h=regen_canvas_h)
    if scheme_data:
        partial_html = _resolve_color_vars(partial_html, scheme_data, css_vars=True)
    regen_partial_path = os.path.join(run_dir, "index_regenerated_partial.html")
    with open(regen_partial_path, "w", encoding="utf-8") as f:
        f.write(partial_html)

    preview_url = f"/api/exports/{req.run_id}/index_regenerated.html"
    partial_url = f"/api/exports/{req.run_id}/index_regenerated_partial.html"

    # Per-slide resolved HTML for targeted splice (not full-deck overwrite)
    regen_slides = [{"seq": s["seq"], "html": s["html"]} for s in regen_resolved]

    _log("DONE")

    return {"ok": True, "slide_seqs": seqs,
            "preview_url": preview_url, "partial_url": partial_url,
            "slides": regen_slides, "html": deck_html}


@app.put("/api/ppt/splice-slides/{run_id}")
def api_ppt_splice_slides(run_id: str, req: dict):
    """Replace specific slides in index.html with new HTML.

    Used by the regenerate flow: only the regenerated slides are spliced in,
    leaving other pages untouched.
    """
    import re as _re

    run_dir = _find_run_dir(run_id)
    if not run_dir or not os.path.isdir(run_dir):
        raise HTTPException(status_code=404, detail="Run not found")

    slides_data = req.get("slides", [])
    if not slides_data:
        raise HTTPException(status_code=400, detail="slides is required")

    log_path_splice = os.path.join(run_dir, "_splice_log.txt")
    def _log_splice(msg: str):
        import datetime as _dt_sp
        ts = _dt_sp.datetime.now().strftime("%H:%M:%S")
        line = f"[{ts}] {msg}\n"
        with open(log_path_splice, "a", encoding="utf-8") as _lf:
            _lf.write(line)
            _lf.flush()

    index_path = os.path.join(run_dir, "index.html")
    if not os.path.exists(index_path):
        raise HTTPException(status_code=404, detail="index.html not found")

    with open(index_path, "r", encoding="utf-8") as _f:
        html = _f.read()

    has_wrapper = '<div class="slide-wrapper"' in html[:5000]
    _log_splice(f"DEBUG has_wrapper={has_wrapper}, html_len={len(html)}, slides={len(slides_data)}, first_section={'<section' in html[:5000]}")

    replaced = 0
    for slide in slides_data:
        seq = slide["seq"]
        new_html = slide["html"]
        _log_splice(f"DEBUG slide seq={seq}, html_len={len(new_html)}")
        if has_wrapper:
            # Find the nth <div class="slide-wrapper" by counting wrappers,
            # then find its matching </div> by balancing div depth.
            wrapper_tag = f'<div class="slide-wrapper" data-seq="{seq}">'
            pos = html.find(wrapper_tag)
            if pos == -1:
                # Fallback: count nth wrapper by position (old format without data-seq)
                nth = seq - 1
                pos = 0
                for _ in range(nth + 1):
                    pos = html.find('<div class="slide-wrapper"', pos)
                    if pos == -1:
                        break
                    if _ < nth:
                        pos += len('<div class="slide-wrapper"')
            else:
                # Found by data-seq; use this position
                pass

            if pos == -1:
                # Last resort: try matching by seq within the slide content
                # (search for data-seq attribute anywhere inside a wrapper)
                tag_pat = _re.compile(
                    r'<div\s+class="slide-wrapper"[^>]*>',
                    _re.DOTALL
                )
                all_wrappers = list(tag_pat.finditer(html))
                if 0 <= seq - 1 < len(all_wrappers):
                    pos = all_wrappers[seq - 1].start()

            if pos != -1:
                # Find the wrapper opener end (>)
                tag_end = html.find('>', pos) + 1
                # Balance divs to find matching closer
                depth = 1
                i = tag_end
                while i < len(html) and depth > 0:
                    next_open = html.find('<div', i)
                    next_close = html.find('</div>', i)
                    if next_close == -1:
                        break
                    if next_open != -1 and next_open < next_close:
                        depth += 1
                        i = next_open + 4
                    else:
                        depth -= 1
                        if depth == 0:
                            close_end = next_close + 6  # after </div>
                            replacement = f'<div class="slide-wrapper" data-seq="{seq}">{new_html}</div>'
                            html = html[:pos] + replacement + html[close_end:]
                            replaced += 1
                            break
                        i = next_close + 6
        else:
            # Old format (<section> based) — primary path for section-based files
            pat = _re.compile(
                r'<section[^>]*\s+data-seq="' + str(seq)
                + r'"[^>]*>.*?</section>',
                _re.DOTALL
            )
            new_html_text, count = pat.subn(new_html, html, count=1)
            _log_splice(f"DEBUG else-branch seq={seq}, count={count}, pat_ok={pat.search(html) is not None}")
            if count > 0:
                html = new_html_text
                replaced += count

    # Write back to index.html (auto-backup preserves original)
    _ensure_backup(run_dir)

    # ── Post-splice auto-fix: clean font sizes & hex colors in ALL slides ──
    from services.ppt_service import _auto_fix_font_size, _auto_fix_hardcoded_hex, _load_scheme_data, _resolve_color_vars
    rj_path = os.path.join(run_dir, "result.json")
    style_id = "business"
    color_scheme = "deep-blue"
    if os.path.exists(rj_path):
        try:
            rj_meta = json.loads(open(rj_path, "r", encoding="utf-8").read())
            style_id = rj_meta.get("style_id") or style_id
            color_scheme = rj_meta.get("color_scheme") or color_scheme
        except Exception:
            pass
    _log_splice(f"色系(from result.json): {color_scheme}")
    scheme_data = _load_scheme_data(style_id, color_scheme)
    if scheme_data:
        hex_before = len(_re.findall(r'#[0-9a-fA-F]{3,6}\b', html))
        html = _auto_fix_hardcoded_hex(html, scheme_data, 0)
        html = _resolve_color_vars(html, scheme_data, css_vars=True)
        hex_after = len(_re.findall(r'#[0-9a-fA-F]{3,6}\b', html))
        _log_splice(f"Post-splice hex fix: {hex_before}→{hex_after}")
        # Transitional: if splice produced var(--*) refs but no :root block exists
        # (old-format index.html with new-format regenerated slides), inject :root
        if "var(--" in html and ":root {" not in html:
            from services.ppt_service import _build_root_vars
            root_block = _build_root_vars(scheme_data)
            style_close = html.find('</style>')
            if style_close > 0:
                html = html[:style_close] + '\n' + root_block + '\n' + html[style_close:]
                _log_splice("Transitional fix: injected :root block for CSS var compatibility")
        # Cover fix: dark bg → white text
        primary = scheme_data.get("primary", "")
        secondary = scheme_data.get("secondary", "")
        if primary:
            cover_start = html.find('<div class="slide-wrapper" data-seq="1">')
            if cover_start == -1:
                cover_start = html.find('<section class="slide" data-seq="1"')
            if cover_start != -1:
                cover_end = html.find('</div>', cover_start)
                if cover_end == -1:
                    cover_end = html.find('</section>', cover_start)
                if cover_end != -1:
                    cover_html = html[cover_start:cover_end]
                    has_dark_bg = (f"background:{primary}" in cover_html
                                  or f"background:{secondary}" in cover_html
                                  or f"background: {primary}" in cover_html
                                  or f"background: {secondary}" in cover_html)
                    if has_dark_bg:
                        text_color = scheme_data.get("text", "")
                        if text_color and text_color != "#ffffff":
                            cover_html_fixed = cover_html
                            for pat in [f"color:{text_color}", f"color: {text_color}", f"color:{text_color};", f"color: {text_color};"]:
                                cover_html_fixed = cover_html_fixed.replace(pat, "color:#ffffff" if ";" not in pat else "color:#ffffff;")
                            html = html[:cover_start] + cover_html_fixed + html[cover_end:]
                            _log_splice("Cover fix: dark background → white text")
    font_before_13 = len(_re.findall(r'font-size:\s*1[0-3]px', html))
    font_before_15 = len(_re.findall(r'font-size:\s*15px', html))
    html = _auto_fix_font_size(html, 0)
    font_after_13 = len(_re.findall(r'font-size:\s*1[0-3]px', html))
    font_after_15 = len(_re.findall(r'font-size:\s*15px', html))
    _log_splice(f"Post-splice font fix: 13px {font_before_13}→{font_after_13}, 15px {font_before_15}→{font_after_15}")

    with open(index_path, "w", encoding="utf-8") as _f:
        _f.write(html)

    return {"ok": True, "replaced": replaced, "total": len(slides_data)}


def _find_run_dir(run_id: str):
    """Resolve run directory with 3-tier fallback: memory → EXPORT_DIR → scan."""
    if '..' in run_id or '/' in run_id or '\\' in run_id:
        return None
    run_dir = _run_dirs.get(run_id)
    if run_dir and os.path.isdir(run_dir):
        return run_dir
    candidate = os.path.join(EXPORT_DIR, run_id)
    if os.path.isdir(candidate):
        _run_dirs[run_id] = candidate
        return candidate
    for base in [EXPORT_DIR] + _scan_output_bases():
        candidate = os.path.join(base, run_id)
        if os.path.isdir(candidate):
            _run_dirs[run_id] = candidate
            _save_run_dirs(_run_dirs)
            return candidate
    return None


@app.get("/api/ppt/regenerate-state/{run_id}")
def api_ppt_get_regenerate_state(run_id: str):
    """Load persisted regenerate-tab state so it survives modal close/reopen."""
    run_dir = _find_run_dir(run_id)
    if not run_dir:
        return {"state": None}
    state_path = os.path.join(run_dir, "_regenerate_state.json")
    if not os.path.exists(state_path):
        return {"state": None}
    try:
        with open(state_path, encoding="utf-8") as _f:
            return {"state": json.loads(_f.read())}
    except Exception:
        return {"state": None}


@app.post("/api/ppt/regenerate-state/{run_id}")
def api_ppt_save_regenerate_state(run_id: str, req: dict):
    """Persist regenerate-tab state to disk so it survives modal close/reopen."""
    run_dir = _find_run_dir(run_id)
    if not run_dir:
        raise HTTPException(status_code=404, detail="Run not found")
    state_path = os.path.join(run_dir, "_regenerate_state.json")
    with open(state_path, "w", encoding="utf-8") as _f:
        json.dump(req, _f, ensure_ascii=False)
    return {"ok": True}


@app.delete("/api/ppt/regenerate-state/{run_id}")
def api_ppt_clear_regenerate_state(run_id: str):
    """Clear persisted regenerate-tab state (user clicked discard)."""
    run_dir = _find_run_dir(run_id)
    if not run_dir:
        return {"ok": True}
    state_path = os.path.join(run_dir, "_regenerate_state.json")
    if os.path.exists(state_path):
        try:
            os.remove(state_path)
        except Exception:
            pass
    return {"ok": True}


@app.get("/api/ppt/regenerate-log/{run_id}")
def api_ppt_regenerate_log(run_id: str):
    """Poll the regenerate progress log. Returns the full log text.
    Frontend polls this every 10 seconds during regeneration."""
    run_dir = _run_dirs.get(run_id)
    if not run_dir:
        candidate = os.path.join(EXPORT_DIR, run_id)
        if os.path.isdir(candidate):
            run_dir = candidate
    if not run_dir or not os.path.isdir(run_dir):
        raise HTTPException(status_code=404, detail="Run not found")

    log_path = os.path.join(run_dir, "_regenerate_log.txt")
    if not os.path.exists(log_path):
        return {"log": "", "done": False}

    with open(log_path, encoding="utf-8") as f:
        text = f.read()

    done = text.rstrip().endswith("DONE")
    return {"log": text, "done": done}


# ── PPT Styles (17 PPT-Agent YAML styles) ──

@app.get("/api/ppt/styles")
def api_ppt_styles():
    """List all 17 PPT-Agent styles from the YAML style library.
    Returns a flat list of styles. Each style has a 'group' field
    (Professional, Creative, Tech / Dark, Thematic) and includes
    color palette, typography, mood, and use cases.
    """
    try:
        from services.svg_renderer import StyleLoader
        loader = StyleLoader()
        groups = loader.list_styles()
        flat = []
        for g in groups:
            for s in g.get("styles", []):
                flat.append(s)
        return {"styles": flat}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── VI & Prompt file editor (directory-aware) ──

VI_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                      "resources", "vi")

VI_SECTIONS = ["vi", "cover", "content", "data", "summary", "prompt"]


class VIFileUpdate(BaseModel):
    content: str


def _style_dir(style_id: str) -> str:
    return os.path.join(VI_DIR, style_id)


def _vi_section_path(style_id: str, section: str) -> str:
    """Get path for a VI sub-file. Supports subdirectories e.g. blocks/header, templates/homework_manual."""
    ext = ".yaml" if section == "tokens" else ".md"
    return os.path.join(_style_dir(style_id), f"{section}{ext}")


def _validate_vi_path(style_id: str, section: str, for_write: bool = False) -> str:
    """Resolve and validate a VI path stays within the style directory. Returns the safe path."""
    p = _vi_section_path(style_id, section)
    base = os.path.realpath(_style_dir(style_id))
    # Use normpath to collapse ../ sequences, then verify the result is under base
    normalized = os.path.normpath(os.path.abspath(p))
    if not normalized.startswith(os.path.normpath(base) + os.sep):
        raise HTTPException(status_code=403, detail="Path traversal denied")
    # For read paths that exist, additionally resolve symlinks via realpath
    if not for_write and os.path.exists(p):
        real_p = os.path.realpath(p)
        if not real_p.startswith(base + os.sep):
            raise HTTPException(status_code=403, detail="Path traversal denied")
    return p


@app.get("/api/ppt/styles/{style_id}/vi/files")
def api_list_style_vi_files(style_id: str):
    """List all VI sub-files in the style directory (recursive into subdirectories)."""
    d = _style_dir(style_id)
    if not os.path.isdir(d):
        return {"files": [], "exists": False}
    files = []
    from services.ppt_service import _page_type_sort_key
    for root, _dirs, filenames in os.walk(d):
        for f in filenames:
            if f.endswith((".md", ".yaml")):
                p = os.path.join(root, f)
                rel = os.path.relpath(p, d).replace("\\", "/")
                section = rel.rsplit(".", 1)[0]  # e.g. "blocks/header", "cover"
                files.append({
                    "name": rel,
                    "size": os.path.getsize(p),
                    "section": section,
                })
    files.sort(key=lambda x: _page_type_sort_key(x["section"]))
    return {"files": files, "exists": True, "dir": d}


@app.get("/api/ppt/styles/{style_id}/page-types")
def api_get_page_types(style_id: str):
    """Return available page_types for a style — scanned from VI directory.
    Each .md file (except vi.md, prompt.md) = one page_type.
    """
    from services.ppt_service import _scan_vi_page_types, _fallback_page_types
    try:
        types = _scan_vi_page_types(style_id)
        return {"page_types": types}
    except Exception:
        return {"page_types": _fallback_page_types()}


@app.get("/api/ppt/styles/{style_id}/color-schemes")
def api_list_color_schemes(style_id: str):
    """List available color schemes for a style from its tokens.yaml."""
    from services.ppt_service import _list_color_schemes
    schemes = _list_color_schemes(style_id)
    return {"color_schemes": schemes, "exists": len(schemes) > 0}


@app.get("/api/ppt/styles/{style_id}/vi/{section:path}")
def api_get_style_vi_section(style_id: str, section: str, color_scheme: str = ""):
    """Load a specific VI sub-file. Resolves {{color}} variables when color_scheme is provided."""
    from services.ppt_service import _load_scheme_data, _resolve_color_vars

    p = _validate_vi_path(style_id, section)
    if not os.path.exists(p):
        return {"content": "", "exists": False, "section": section}
    with open(p, "r", encoding="utf-8") as f:
        content = f.read()
    if color_scheme:
        scheme = _load_scheme_data(style_id, color_scheme)
        if scheme:
            content = _resolve_color_vars(content, scheme)
    return {"content": content, "exists": True, "section": section}


@app.put("/api/ppt/styles/{style_id}/vi/{section:path}")
def api_save_style_vi_section(style_id: str, section: str, body: VIFileUpdate):
    """Save a specific VI sub-file. Supports subdirectory sections."""
    d = _style_dir(style_id)
    os.makedirs(d, exist_ok=True)
    p = _validate_vi_path(style_id, section, for_write=True)
    os.makedirs(os.path.dirname(p), exist_ok=True)
    with open(p, "w", encoding="utf-8") as f:
        f.write(body.content)
    return {"ok": True, "path": p, "section": section}


# ── Scenario prompt files (per-column design rule overrides) ──

SCENARIOS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "resources", "scenarios")
SCENARIO_FILES = [
    "design-system.md",
    "outline-architect.md",
    "cognitive-design-principles.md",
    "reviewer.md",
    "svg-generator.md",
    "bento-grid-layout.md",
]


def _scenario_file_path(column_id: str, filename: str, for_write: bool = False) -> str:
    """Resolve a scenario file path. For read: custom → default → None.
    For write: always returns custom path (auto-creates dir)."""
    fname = os.path.basename(str(filename))
    if fname != filename or fname.startswith(".") or ".." in filename:
        raise HTTPException(status_code=400, detail="Invalid filename")
    if for_write:
        d = os.path.join(SCENARIOS_DIR, column_id)
        os.makedirs(d, exist_ok=True)
        return os.path.join(d, fname)
    # Read path: check custom first, then default
    p = os.path.join(SCENARIOS_DIR, column_id, fname)
    if os.path.exists(p):
        return p
    p = os.path.join(SCENARIOS_DIR, "_default", fname)
    if os.path.exists(p):
        return p
    return ""


@app.get("/api/scenarios/{column_id}/files")
def api_list_scenario_files(column_id: str):
    """List all 6 scenario files with their source (custom / default / missing)."""
    files = []
    for fname in SCENARIO_FILES:
        custom_path = os.path.join(SCENARIOS_DIR, column_id, fname)
        default_path = os.path.join(SCENARIOS_DIR, "_default", fname)
        if os.path.exists(custom_path):
            source = "custom"
            size = os.path.getsize(custom_path)
        elif os.path.exists(default_path):
            source = "default"
            size = os.path.getsize(default_path)
        else:
            source = "missing"
            size = 0
        files.append({"name": fname, "size": size, "source": source})
    return {"files": files, "column_id": column_id}


@app.get("/api/scenarios/{column_id}/files/{filename}")
def api_get_scenario_file(column_id: str, filename: str):
    """Get a scenario file content (custom → default → 404)."""
    p = _scenario_file_path(column_id, filename)
    if not p:
        return {"content": "", "exists": False, "filename": filename}
    with open(p, "r", encoding="utf-8") as f:
        return {"content": f.read(), "exists": True, "filename": filename}


@app.put("/api/scenarios/{column_id}/files/{filename}")
def api_save_scenario_file(column_id: str, filename: str, body: VIFileUpdate):
    """Save a scenario file to the per-column custom directory."""
    p = _scenario_file_path(column_id, filename, for_write=True)
    with open(p, "w", encoding="utf-8") as f:
        f.write(body.content)
    return {"ok": True, "path": p, "filename": filename}


# Backward-compatible: GET /vi returns the full combined VI for a style
@app.get("/api/ppt/styles/{style_id}/vi")
def api_get_style_vi(style_id: str, color_scheme: str = ""):
    """Load full VI — concatenates all sub-files if directory exists, else legacy file.
    Resolves {{color}} variables against color_scheme when provided."""
    from services.ppt_service import _load_scheme_data, _resolve_color_vars, _load_style_vi

    d = _style_dir(style_id)
    if os.path.isdir(d):
        parts = []
        vi_path = os.path.join(d, "vi.md")
        if os.path.exists(vi_path):
            with open(vi_path, "r", encoding="utf-8") as f:
                parts.append(f.read())
        for section in ["cover", "content", "data", "summary"]:
            sp = os.path.join(d, f"{section}.md")
            if os.path.exists(sp):
                with open(sp, "r", encoding="utf-8") as f:
                    parts.append(f.read())
        if parts:
            content = "\n\n---\n\n".join(parts)
            if color_scheme:
                scheme = _load_scheme_data(style_id, color_scheme)
                if scheme:
                    content = _resolve_color_vars(content, scheme)
            return {"content": content, "exists": True}
    # Legacy fallback
    content = _load_style_vi(style_id, color_scheme or "deep-blue")
    return {"content": content, "exists": bool(content)}


@app.put("/api/ppt/styles/{style_id}/vi")
def api_save_style_vi(style_id: str, body: VIFileUpdate):
    """Save VI — writes to vi.md in directory structure."""
    d = _style_dir(style_id)
    os.makedirs(d, exist_ok=True)
    p = os.path.join(d, "vi.md")
    with open(p, "w", encoding="utf-8") as f:
        f.write(body.content)
    return {"ok": True, "path": p}


@app.get("/api/ppt/styles/{style_id}/prompt")
def api_get_style_prompt(style_id: str):
    """Load a style's AI prompt markdown file."""
    # Directory structure first
    p = os.path.join(_style_dir(style_id), "prompt.md")
    if os.path.exists(p):
        with open(p, "r", encoding="utf-8") as f:
            return {"content": f.read(), "exists": True}
    # Legacy fallback
    from services.ppt_service import _load_style_prompt
    content = _load_style_prompt(style_id)
    return {"content": content, "exists": bool(content)}


@app.put("/api/ppt/styles/{style_id}/prompt")
def api_save_style_prompt(style_id: str, body: VIFileUpdate):
    """Save a style's AI prompt markdown file."""
    d = _style_dir(style_id)
    os.makedirs(d, exist_ok=True)
    p = os.path.join(d, "prompt.md")
    with open(p, "w", encoding="utf-8") as f:
        f.write(body.content)
    return {"ok": True, "path": p}


@app.get("/api/ppt/export-zip/{run_id}")
def api_export_svg_zip(run_id: str):
    """Download all SVG files + index.html for a generated deck as a ZIP archive."""
    import zipfile, io
    run_dir = _run_dirs.get(run_id)
    if not run_dir:
        candidate = os.path.join(EXPORT_DIR, run_id)
        if os.path.isdir(candidate):
            run_dir = candidate
    if not run_dir:
        for base in [EXPORT_DIR] + _scan_output_bases():
            candidate = os.path.join(base, run_id)
            if os.path.isdir(candidate):
                run_dir = candidate
                _run_dirs[run_id] = run_dir
                _save_run_dirs(_run_dirs)
                break
    if not run_dir:
        raise HTTPException(status_code=404, detail="Export not found")
    buf = io.BytesIO()
    try:
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_STORED) as zf:
            for fname in os.listdir(run_dir):
                fpath = os.path.join(run_dir, fname)
                if os.path.isfile(fpath):
                    zf.write(fpath, fname)
    except Exception as e:
        import traceback
        print(f"[ZIP ERROR] run_dir={run_dir!r} error={e}", flush=True)
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"ZIP creation failed: {e}")
    buf.seek(0)
    from urllib.parse import quote
    safe_name = quote(f"svg-deck-{run_id}.zip", safe="")
    return StreamingResponse(
        buf,
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{safe_name}"},
    )


@app.post("/api/ppt/save-images/{run_id}")
async def api_save_slide_images(run_id: str):
    """Render each slide as a 1280x720 PNG and save to the export directory."""
    import asyncio

    run_dir = _find_run_dir(run_id)
    if not run_dir:
        raise HTTPException(status_code=404, detail="Export not found")

    html_path = os.path.join(run_dir, "index.html")
    if not os.path.exists(html_path):
        raise HTTPException(status_code=404, detail="index.html not found")

    def _capture():
        try:
            from playwright.sync_api import sync_playwright
        except ImportError:
            raise HTTPException(status_code=500, detail="playwright not installed")

        saved = []
        with sync_playwright() as pw:
            browser = pw.chromium.launch()
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.goto("file:///" + html_path.replace("\\", "/"))
            page.wait_for_timeout(500)

            slides = page.query_selector_all(".slide-wrapper, section")
            if not slides:
                slides = page.query_selector_all('[style*="width"]')

            for i, slide in enumerate(slides):
                png_path = os.path.join(run_dir, f"slide_{i+1:02d}.png")
                slide.screenshot(path=png_path)
                saved.append(png_path)

            browser.close()
        return saved

    saved = await asyncio.to_thread(_capture)
    return {"ok": True, "saved": len(saved), "files": [os.path.basename(f) for f in saved],
            "dir": run_dir}


@app.post("/api/open-folder")
async def api_open_folder(req: dict):
    path = req.get("path", "")
    if not path or not os.path.isdir(path):
        raise HTTPException(status_code=400, detail="Path not found")
    os.startfile(path)
    return {"ok": True}


# ── SOP Export ──

from services.export_service import export_sop

class SOPExportRequest(BaseModel):
    content: str
    branding: dict = None
    project_id: Optional[str] = None


@app.post("/api/export/sop")
def api_export_sop(req: SOPExportRequest):
    try:
        output_dir = None
        project_name = ""
        if req.project_id:
            output_dir = resolve_project_storage(req.project_id)
            proj = _get_project(req.project_id)
            if proj:
                project_name = proj["name"]
        filepath = export_sop(req.content, req.branding, output_dir)
        filename = os.path.basename(filepath)
        params = []
        if req.project_id:
            params.append(f"project_id={req.project_id}")
        if project_name:
            safe_name = "".join(c for c in project_name if c.isalnum() or c in "._- ()（）").strip()
            params.append(f"name={safe_name}_SOP.docx")
        download_url = f"/api/download/{filename}"
        if params:
            download_url += "?" + "&".join(params)
        return {"filename": filename, "download_url": download_url}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


# ── File Download ──

@app.get("/api/download/{filename}")
def download_file(filename: str, project_id: str = None, name: str = None):
    download_name = name or filename
    if project_id:
        try:
            proj_dir = resolve_project_storage(project_id, auto_create=False)
            filepath = os.path.join(proj_dir, filename)
            if os.path.exists(filepath):
                return FileResponse(filepath, filename=download_name)
        except Exception:
            pass

    filepath = os.path.join(EXPORT_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(filepath, filename=download_name)


# ── TTS ──

@app.post("/api/tts/synthesize")
async def api_tts_synthesize(req: SynthesizeRequest):
    import httpx
    try:
        # Resolve TTS API key and base_url from provider or fallback to settings
        api_key = ""
        base_url = "https://dashscope.aliyuncs.com/api/v1"
        if req.provider_id:
            tts_db = get_db()
            try:
                provider = tts_db.execute(
                    "SELECT * FROM tts_providers WHERE id = ? AND is_enabled = 1",
                    (req.provider_id,)).fetchone()
                if provider:
                    api_key = provider["api_key"]
                    base_url = provider["base_url"]
            finally:
                tts_db.close()
        if not api_key:
            api_key = _get_setting("tts_api_key") or os.getenv("DASHSCOPE_API_KEY", "")
        if not api_key:
            raise HTTPException(status_code=400, detail="请先在项目配置中设置 TTS API Key")
        tts_url = base_url.rstrip("/") + "/services/audio/tts/SpeechSynthesizer"

        payload = {
            "model": req.model,
            "input": {
                "text": req.text,
                "voice": req.voice_id or "longanyang",
                "format": "mp3",
                "volume": req.volume,
                "rate": req.speed,
            },
        }
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                tts_url,
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json=payload,
            )
        if resp.status_code != 200:
            raise HTTPException(status_code=400, detail=f"TTS 失败: {resp.text[:200]}")
        result = resp.json()
        audio_url = result.get("output", {}).get("audio", {}).get("url", "")
        if not audio_url:
            raise HTTPException(status_code=400, detail="未获取到音频URL")

        # Download and save
        async with httpx.AsyncClient(timeout=60) as client:
            dl = await client.get(audio_url)

        output_dir = AUDIO_DIR
        if req.project_id:
            output_dir = resolve_project_storage(req.project_id)
        os.makedirs(output_dir, exist_ok=True)

        audio_name = f"tts_{os.urandom(4).hex()}.mp3"
        audio_path = os.path.join(output_dir, audio_name)
        with open(audio_path, "wb") as f:
            f.write(dl.content)

        serve_url = f"/api/audio/{audio_name}"
        params = []
        if req.project_id:
            params.append(f"project_id={req.project_id}")
            proj = _get_project(req.project_id)
            if proj:
                safe_name = "".join(c for c in proj["name"] if c.isalnum() or c in "._- ()（）").strip()
                params.append(f"name={safe_name}_音频.mp3")
        if params:
            serve_url += "?" + "&".join(params)
        return {"audio_url": serve_url, "filename": audio_name}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.get("/api/audio/{filename}")
def serve_audio(filename: str, project_id: str = None, name: str = None):
    download_name = name or filename
    if project_id:
        try:
            proj_dir = resolve_project_storage(project_id, auto_create=False)
            filepath = os.path.join(proj_dir, filename)
            if os.path.exists(filepath):
                return FileResponse(filepath, media_type="audio/mpeg", filename=download_name)
        except Exception:
            pass

    filepath = os.path.join(AUDIO_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Audio not found")
    return FileResponse(filepath, media_type="audio/mpeg")


# ── Logo Upload ──

@app.post("/api/upload/logo")
async def upload_logo(file: UploadFile = File(...)):
    """Upload a logo image file. Returns the filename for later retrieval."""
    import uuid as _uuid
    ext = os.path.splitext(file.filename or "logo.png")[1] or ".png"
    if ext.lower() not in (".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp", ".ico"):
        raise HTTPException(status_code=400, detail="不支持的图片格式，请上传 PNG/JPG/GIF/SVG/WebP/ICO")
    filename = f"logo_{_uuid.uuid4().hex[:8]}{ext}"
    filepath = os.path.join(LOGO_DIR, filename)
    content = await file.read()
    if len(content) > 5 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="文件大小不能超过 5MB")
    with open(filepath, "wb") as f:
        f.write(content)
    return {"filename": filename, "url": f"/api/logos/{filename}"}


@app.get("/api/logos/{filename}")
def serve_logo(filename: str):
    filepath = os.path.join(LOGO_DIR, filename)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Logo not found")
    media_map = {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".gif": "image/gif", ".svg": "image/svg+xml", ".webp": "image/webp",
        ".ico": "image/x-icon",
    }
    ext = os.path.splitext(filename)[1].lower()
    return FileResponse(filepath, media_type=media_map.get(ext, "application/octet-stream"))


# ── Settings ──

@app.get("/api/settings")
def get_settings():
    db = get_db()
    try:
        rows = db.execute("SELECT key, value FROM settings").fetchall()
        return {"settings": {r["key"]: r["value"] for r in rows}}
    finally:
        db.close()


@app.put("/api/settings")
def update_settings(req: dict):
    db = get_db()
    try:
        for key, value in req.items():
            existing = db.execute("SELECT key FROM settings WHERE key = ?", (key,)).fetchone()
            if existing:
                db.execute("UPDATE settings SET value = ? WHERE key = ?", (value, key))
            else:
                db.execute("INSERT INTO settings (key, value) VALUES (?, ?)", (key, value))
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.post("/api/verify-password")
def verify_password(req: dict):
    plain = req.get("password", "")
    if not plain:
        raise HTTPException(status_code=400, detail="密码不能为空")
    stored_hash = _get_setting("admin_password")
    if not stored_hash:
        raise HTTPException(status_code=400, detail="未设置密码")
    if _hash_password(plain) != stored_hash:
        raise HTTPException(status_code=403, detail="密码错误")
    return {"ok": True}


@app.get("/api/browse-folder")
def browse_folder():
    """Open native folder picker dialog and return selected path."""
    import tkinter.filedialog, tkinter
    try:
        root = tkinter.Tk()
        root.withdraw()
        root.attributes('-topmost', True)
        path = tkinter.filedialog.askdirectory(title="选择默认保存路径")
        root.destroy()
        return {"path": path if path else ""}
    except Exception:
        return {"path": ""}


VERSION = "1.0.0"
UPDATE_CHECK_URL = "https://raw.githubusercontent.com/qiliang166/yishao-agent/master/version.json"


@app.get("/api/version")
def api_version():
    return {"version": VERSION, "app": _get_site_name()}


@app.get("/api/check-update")
async def api_check_update():
    import httpx
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(UPDATE_CHECK_URL)
        if resp.status_code == 200:
            data = resp.json()
            latest = data.get("version", VERSION)
            return {
                "current": VERSION,
                "latest": latest,
                "has_update": _version_greater(latest, VERSION),
                "release_url": data.get("release_url", ""),
                "release_notes": data.get("release_notes", ""),
            }
    except Exception:
        pass
    return {
        "current": VERSION,
        "latest": VERSION,
        "has_update": False,
        "release_url": "",
        "error": "无法连接更新服务器",
    }


@app.post("/api/download-update")
async def api_download_update(req: dict):
    """Download the latest release .exe from GitHub."""
    import httpx
    import subprocess
    import tempfile

    download_url = req.get("download_url", "")
    if not download_url:
        # Fetch the latest release info from GitHub API to get the .exe URL
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                resp = await client.get(
                    "https://api.github.com/repos/qiliang166/yishao-agent/releases/latest",
                    headers={"Accept": "application/vnd.github+json", "X-GitHub-Api-Version": "2022-11-28"},
                )
            if resp.status_code != 200:
                raise HTTPException(status_code=400, detail="无法获取最新版本信息")
            release = resp.json()
            assets = release.get("assets", [])
            exe_asset = next((a for a in assets if a["name"].endswith(".exe")), None)
            if not exe_asset:
                raise HTTPException(status_code=400, detail="未找到安装包")
            download_url = exe_asset["browser_download_url"]
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"获取下载地址失败: {str(e)}")

    # Download to temp directory
    tmpdir = tempfile.gettempdir()
    local_path = os.path.join(tmpdir, "yishao-agent-setup.exe")
    try:
        async with httpx.AsyncClient(timeout=600, follow_redirects=True) as client:
            async with client.stream("GET", download_url) as resp:
                if resp.status_code != 200:
                    raise HTTPException(status_code=400, detail=f"下载失败 HTTP {resp.status_code}")
                total = int(resp.headers.get("content-length", 0))
                downloaded = 0
                with open(local_path, "wb") as f:
                    async for chunk in resp.aiter_bytes(1024 * 1024):
                        f.write(chunk)
                        downloaded += len(chunk)
        # Launch the installer
        subprocess.Popen([local_path], shell=True)
        return {"status": "installing", "path": local_path}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"下载失败: {str(e)}")


def _version_greater(a: str, b: str) -> bool:
    """Compare two semver strings. Returns True if a > b."""
    try:
        pa = [int(x) for x in a.split(".")]
        pb = [int(x) for x in b.split(".")]
        for i in range(max(len(pa), len(pb))):
            va = pa[i] if i < len(pa) else 0
            vb = pb[i] if i < len(pb) else 0
            if va > vb:
                return True
            if va < vb:
                return False
        return False
    except Exception:
        return False


from pydantic import BaseModel

class OpenFolderRequest(BaseModel):
    path: str

@app.post("/api/open-folder")
def open_folder(req: OpenFolderRequest):
    p = req.path.strip()
    if not p:
        raise HTTPException(400, "path required")
    if not os.path.exists(p):
        raise HTTPException(404, f"path not found: {p}")
    try:
        os.startfile(p)
        return {"ok": True}
    except Exception as e:
        raise HTTPException(500, str(e))

# ── Column Configs ──

@app.get("/api/column-configs")
def list_column_configs():
    db = get_db()
    try:
        rows = db.execute("SELECT * FROM column_configs ORDER BY sort_order").fetchall()
        return {"configs": [dict(r) for r in rows]}
    finally:
        db.close()


def _build_prompt_from_rules(rules: dict) -> str:
    """从 rules JSON 自动生成 prompt（约束条文版本）"""
    dr = rules.get("design_rules", {})
    lts = rules.get("layout_types", [])
    comp = rules.get("components", {})
    imgr = rules.get("image_rules", {})
    chk = rules.get("checklist", {})
    pr = rules.get("page_rhythm", {})
    dps = rules.get("design_principles", [])

    parts = []

    # Role — 从栏目配置的 prompt 提供，此处不硬编码。仅输出约束规则部分。
    role = rules.get("_role_override", "")
    if not role:
        # 兼容旧规则：若无 _role_override，使用栏目配置中的 prompt
        role = "根据栏目配置中定义的提示词角色要求执行。"
    parts.append(role)

    # Constraints section
    parts.append("\n## 约束规则")

    # Color discipline (from constraint text)
    cd = dr.get("color_discipline", "")
    if cd:
        parts.append(f"- 配色纪律：{cd}")

    # Font discipline
    fd = dr.get("font_discipline", "")
    if fd:
        parts.append(f"- 字体纪律：{fd}")

    # Layout discipline
    ld = dr.get("layout_discipline", "")
    if ld:
        parts.append(f"- 版式纪律：{ld}")

    # Layout sequence from page_rhythm
    if pr.get("sequence"):
        parts.append(f"- 版式顺序：{' → '.join(pr['sequence'])}。{pr.get('alternation_rule', '')}")

    # Image rules
    img_types = imgr.get("types", [])
    if img_types:
        parts.append(f"- 配图：{imgr.get('placement', '每页最多1张')}。类型：{'/'.join(t.get('name','') for t in img_types)}")

    # Components
    callouts = comp.get("callouts", [])
    if callouts:
        parts.append(f"- 标注组件：{'/'.join(c.get('name','') for c in callouts)}")
    stats = comp.get("stats", [])
    if stats:
        parts.append(f"- 数据组件：{'/'.join(s.get('name','') for s in stats)}")

    # P0 checklist
    p0 = chk.get("p0_must_pass", [])
    if p0:
        p0_items = "; ".join(item.get("item","") for item in p0[:4])
        parts.append(f"- 硬约束（P0）：{p0_items}")

    # Design principles
    if dps:
        dp_text = "；".join(f'{d.get("rule","")}' for d in dps[:4])
        parts.append(f"- 设计原则：{dp_text}")
        if len(dps) > 4:
            parts[-1] += "等"

    # Page count
    p2 = chk.get("p2_suggested", [])
    for item in p2:
        if "页数" in item.get("item", ""):
            parts.append(f"- {item['item']}")

    parts.append("\n直接输出PPT内容，严格按下方SKILL模板结构。")
    return "\n".join(parts)



@app.put("/api/column-configs/{config_id}")
def update_column_config(config_id: str, req: dict):
    db = get_db()
    try:
        existing = db.execute("SELECT id, column_id FROM column_configs WHERE id = ?", (config_id,)).fetchone()
        if not existing:
            raise HTTPException(404, "Config not found")
        if 'prompt' in req:
            db.execute("UPDATE column_configs SET prompt = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?", (req['prompt'], config_id))
        if 'skill' in req:
            db.execute("UPDATE column_configs SET skill = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?", (req['skill'], config_id))
        db.commit()
        row = db.execute("SELECT * FROM column_configs WHERE id = ?", (config_id,)).fetchone()
        return dict(row)
    finally:
        db.close()


@app.post("/api/column-configs/{config_id}/upload-template")
async def upload_column_template(config_id: str, file: UploadFile = File(...)):
    db = get_db()
    try:
        existing = db.execute("SELECT id, has_template FROM column_configs WHERE id = ?", (config_id,)).fetchone()
        if not existing:
            raise HTTPException(404, "Config not found")
        if not existing["has_template"]:
            raise HTTPException(400, "此栏目不支持模板上传")
        # Save template file
        tmpl_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data", "templates")
        os.makedirs(tmpl_dir, exist_ok=True)
        name = f"{config_id}_{file.filename}"
        path = os.path.join(tmpl_dir, name)
        content = await file.read()
        with open(path, "wb") as f:
            f.write(content)
        db.execute("UPDATE column_configs SET template_path = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?", (path, config_id))
        db.commit()
        # If it's a text file, read and return for AI analysis
        if file.filename.endswith(('.txt', '.md', '.docx')):
            try:
                text = content.decode('utf-8')
            except Exception:
                text = content.decode('gbk', errors='ignore')
            return {"ok": True, "path": path, "content": text}
        return {"ok": True, "path": path}
    finally:
        db.close()


if __name__ == "__main__":
    import sys, uvicorn
    port = int(sys.argv[1]) if len(sys.argv) > 1 and sys.argv[1].isdigit() else 8766
    uvicorn.run(app, host="0.0.0.0", port=port)
